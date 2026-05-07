#!/usr/bin/env python3
"""
Argos VSP — AFTER AMOSI Academic Deck Generator

Builds the academic-talk variant of the Argos VSP presentation. Differs
from the client-facing v9.1 deck (`generate_presentation.py`) in:

  * Different ordering: Opening → Problem → Evaluation → Proof →
    Confidence → Demo+Future → Appendix.
  * Hides Path B human-IS estimates (pre-study, not measurements).
  * Unhides the academic-content slides that the client deck hid
    (`slide_exec_summary`, `slide_wer_lies`, `slide_disagreement_blind`,
    `slide_disagreement_context`, `slide_metric_disagreement`,
    `slide_metric_disagreement_2`).
  * Unhides `slide_30b` (LLM Upgrade — Why It Matters) per academic plan.
  * Output filename: `Argos_VSP_AFTER_AMOSI_May2026.pptx`.

Plan reference: /home/ubuntu/.claude/plans/make-a-plan-to-idempotent-summit.md

Usage:
    python3 docs/_research-tools/generators/generate_after_amosi_presentation.py

Output:
    presentation_materials_20260224/Argos_VSP_AFTER_AMOSI_May2026.pptx
"""

from pathlib import Path

from pptx import Presentation

from presentation.config import SL_W, SL_H, _auto_num
from presentation.helpers import _fix_pptx_video_compat, _strip_orphan_animation_refs

# Override the default OUTPUT path defined in presentation/config.py.
# (config.py's OUTPUT points at the client deck filename — we want the
# academic variant to live alongside it, not overwrite it.)
OUTPUT = Path(
    "/home/ubuntu/presentation_materials_20260224/"
    "Argos_VSP_For_Orchard_May2026.pptx"
)

# ─── Slide imports ───────────────────────────────────────────────────────
# Each name is verified to exist (or flagged as a blocker — see footer).
# Blockers are imported in a tolerant try/except so the orchestrator script
# itself stays valid; failures are reported at runtime.

from presentation.slides_opening import (
    slide_01, slide_what_was_done_1, slide_what_was_done_2, slide_toc,
    slide_02, slide_visemes, slide_03, slide_data_flow,
    slide_04, slide_05, slide_06,
    slide_exec_summary, slide_wer_lies,
)
from presentation.slides_research import (
    slide_research_transition,
    slide_failure_anatomy_transition,
    slide_is_motivation, slide_is_intro,
    slide_is_weight_rationale, slide_is_calc_examples,
    slide_is_dimensions, slide_is_radar, slide_is_wer_scatter,
    slide_07, slide_metric_transition,
    slide_08,
    slide_failure_deep_1a, slide_failure_deep_1b, slide_failure_deep_2,
)
from presentation.slides_evaluation import (
    slide_llm_judge, slide_llm_judge_30,
    slide_judge_ex1, slide_judge_ex2, slide_judge_ex3,
    slide_judge_ex4, slide_judge_ex5, slide_judge_ex6,
    slide_disagreement_blind, slide_disagreement_context,
    slide_25d, slide_25e,
    slide_15,
    slide_metric_disagreement, slide_metric_disagreement_2,
)
from presentation.slides_engineering import (
    slide_17,           # programmatic (NOT slide_17_png)
    slide_dual_env,
)
from presentation.slides_future import (
    slide_future_transition,
    slide_24, slide_26, slide_26b,
    slide_30, slide_30b,
    slide_29,
    slide_arabic_roadmap, slide_arabic_avhubert, slide_arabic_changes,
    slide_31, slide_thank_you,
    slide_28, slide_confidence_scoring,
    slide_a1, slide_a8, slide_a11, slide_a11b, slide_a13,
    slide_a16, slide_a17,
    # New appendix functions added by the slides_future agent (this file):
    slide_human_is_path_b,
    slide_appendix_pca_loadings,
    slide_appendix_mcnemar_full,
)

# ─── Tolerant imports: functions still being authored by sibling agents ──
# Each block imports from the expected file; if the symbol does not yet
# exist, we substitute a placeholder that prints a clear warning at build
# time. The orchestrator script itself remains importable + runnable.

def _stub(name):
    """Create a placeholder slide-builder that prints a warning."""
    def _stub_builder(prs):
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(11), Inches(1.5)
        ).text_frame
        tx.text = f"[MISSING BUILDER: {name}]"
        for p in tx.paragraphs:
            for r in p.runs:
                r.font.size = Pt(28)
                r.font.color.rgb = RGBColor(0xE0, 0x6C, 0x75)
        print(f"    WARN: builder '{name}' not yet defined — placeholder used")
        return slide
    _stub_builder.__name__ = name
    return _stub_builder


# Each maybe-import: import-or-stub. Tracked as blockers; orchestrator
# still runs end-to-end, leaving placeholder slides where the sibling
# agents have not landed yet.
_BLOCKERS = []

def _try_import(module_path, name):
    import importlib
    try:
        mod = importlib.import_module(module_path)
        return getattr(mod, name)
    except (ImportError, AttributeError):
        _BLOCKERS.append(f"{module_path}.{name}")
        return _stub(name)

# Opening section additions
slide_diversity_of_inputs       = _try_import(
    "presentation.slides_opening", "slide_diversity_of_inputs")
# Evaluation section additions
slide_literature_metrics_problem = _try_import(
    "presentation.slides_evaluation", "slide_literature_metrics_problem")
# Confidence section (new — confidence-research slides, all in slides_evaluation per Agent E's scope)
slide_confidence_problem        = _try_import(
    "presentation.slides_evaluation", "slide_confidence_problem")
slide_two_layer_confidence_research = _try_import(
    "presentation.slides_evaluation", "slide_two_layer_confidence_research")
slide_per_word_confidence_distribution = _try_import(
    "presentation.slides_evaluation", "slide_per_word_confidence_distribution")
slide_band_reliability_overall  = _try_import(
    "presentation.slides_evaluation", "slide_band_reliability_overall")
slide_band_reliability_stratified = _try_import(
    "presentation.slides_evaluation", "slide_band_reliability_stratified")
slide_green_leakage_examples    = _try_import(
    "presentation.slides_evaluation", "slide_green_leakage_examples")
slide_three_thresholds          = _try_import(
    "presentation.slides_evaluation", "slide_three_thresholds")
slide_three_tier_policy_research = _try_import(
    "presentation.slides_evaluation", "slide_three_tier_policy_research")
slide_band_reliability_by_niv   = _try_import(
    "presentation.slides_evaluation", "slide_band_reliability_by_niv")
slide_agreement_aware_bands     = _try_import(
    "presentation.slides_evaluation", "slide_agreement_aware_bands")
slide_agreement_vs_conf_information = _try_import(
    "presentation.slides_evaluation", "slide_agreement_vs_conf_information")
slide_client_trust_calibration  = _try_import(
    "presentation.slides_evaluation", "slide_client_trust_calibration")
# N-best (judge) slides
slide_nbest_v3_judge_paired_tests = _try_import(
    "presentation.slides_evaluation", "slide_nbest_v3_judge_paired_tests")
slide_mbr_decision              = _try_import(
    "presentation.slides_evaluation", "slide_mbr_decision")
slide_v1_vs_v3_judge_lesson     = _try_import(
    "presentation.slides_evaluation", "slide_v1_vs_v3_judge_lesson")
# Demo (Obama trust/salvage/strip) and judge example videos
slide_demo_obama_trust          = _try_import(
    "presentation.slides_evaluation", "slide_demo_obama_trust")
slide_demo_obama_salvage        = _try_import(
    "presentation.slides_evaluation", "slide_demo_obama_salvage")
slide_demo_obama_strip          = _try_import(
    "presentation.slides_evaluation", "slide_demo_obama_strip")
slide_demo_judge_entity         = _try_import(
    "presentation.slides_evaluation", "slide_demo_judge_entity")
slide_demo_judge_vocab          = _try_import(
    "presentation.slides_evaluation", "slide_demo_judge_vocab")


def main():
    _auto_num[0] = 0  # Reset auto-numbering

    # NOTE: plot regeneration deliberately skipped — sibling agents have
    # already regenerated the MBR-IS plots; running the legacy generators
    # here would clobber them with stale versions.
    # (See generate_presentation.py:_regenerate_plots for the historical
    # pattern.)

    prs = Presentation()
    prs.slide_width = SL_W
    prs.slide_height = SL_H

    print("Generating AFTER AMOSI academic deck...")

    builders = [
        # ─── §0. Opening (4 slides) ──────────────────────────────────
        slide_01,                       # Title
        slide_what_was_done_1,          # What was done? (1/2)
        slide_what_was_done_2,          # What was done? (2/2)
        slide_toc,                      # TOC

        # ─── §1. The Problem (10 slides) ─────────────────────────────
        slide_02,                       # What is VSP?
        slide_visemes,                  # Visemes
        slide_03,                       # Three components
        slide_data_flow,                # Data flow
        slide_17,                       # 8-stage pipeline (programmatic)
        slide_04,                       # Benchmark
        slide_05,                       # Reality gap
        slide_06,                       # Same WER, different effects
        slide_diversity_of_inputs,      # Diversity of inputs
        slide_wer_lies,                 # WER: the metric that lies

        # ─── §2. Evaluation — How do you evaluate lip-reading? ──────
        # Cut: trimmed judge examples 6 → 4 (kept ex1, ex3, ex5, ex6 — the
        # most distinct failure modes). Examples 2, 4 collapse into ex5 in
        # speaker notes (same "core preserved despite WER" pattern).
        slide_research_transition,      # Section divider
        slide_literature_metrics_problem,
        slide_llm_judge,                # LLM-as-a-Judge intro
        slide_llm_judge_30,             # LLM judge deep dive
        slide_judge_ex1,                # Named entity swap
        slide_judge_ex3,                # Technical vocab drift
        slide_judge_ex5,                # Cooking domain confusion
        slide_judge_ex6,                # Topic hijack
        slide_disagreement_blind,
        slide_disagreement_context,
        slide_is_motivation,            # Why LLM Judge isn't enough
        slide_is_intro,                 # IS signals
        slide_is_weight_rationale,      # Do 6 signals measure 6 things?
        slide_is_calc_examples,         # IS in action
        slide_is_dimensions,            # IS dimensions
        slide_is_radar,                 # IS radar
        slide_is_wer_scatter,           # IS vs WER scatter
        slide_a16,                      # IS x judge tier (cross-tab)

        # ─── §3. Where It Works — and How It Fails ───────────────────
        # SECTION SPLIT v5 (per research-overview review):
        # §3a Capture (where it works) = slide_07 + slide_metric_transition
        # §3b Failure Anatomy = visible divider + slide_08 + deep_1a + 2 + salvage
        # Cut: slide_failure_deep_1b (overlapping content with 1a — taxonomy
        # categories now fit in one slide). slide_a13 cut (overlaps with
        # slide_failure_deep_2 — same "real failure examples" framing).
        # E1 (research-overview): visible section divider added between the
        # capture half and the failure-anatomy half so the audience sees the
        # split, not just the orchestrator comment.
        slide_07,                       # IS results: useful output (CAPTURE)
        slide_metric_transition,        # Three numbers transition (CAPTURE)
        slide_failure_anatomy_transition,  # E1: visible §3a → §3b divider
        slide_08,                       # Failure-mode taxonomy bar (FAILURE)
        slide_failure_deep_1a,          # Failure deep (consolidated)
        slide_failure_deep_2,           # Failure modes: examples
        slide_25d,                      # LLM salvage: three recoveries
        slide_25e,                      # LLM salvage: domain context

        # ─── §4. Confidence — Trustable output w/o GT (18 slides) ────
        # A1 (research-overview): math-first opener — slide_two_layer_
        # confidence_research moved BEFORE slide_confidence_problem so the
        # closed-form math lands as the §4 hook; the prose runtime-vs-
        # eval framing follows as a "what we do with it" callback.
        slide_two_layer_confidence_research,
        slide_confidence_problem,
        slide_confidence_scoring,
        slide_per_word_confidence_distribution,
        slide_band_reliability_overall,
        slide_band_reliability_stratified,
        slide_green_leakage_examples,
        slide_three_thresholds,
        slide_three_tier_policy_research,
        slide_band_reliability_by_niv,
        slide_agreement_aware_bands,
        slide_agreement_vs_conf_information,
        slide_client_trust_calibration,
        slide_28,                       # Phase 2: N-best (intro)
        slide_nbest_v3_judge_paired_tests,
        slide_mbr_decision,
        slide_v1_vs_v3_judge_lesson,

        # ─── §5. Demo + Future (12+ slides) ──────────────────────────
        slide_15,                       # Demo intro
        slide_demo_obama_trust,
        slide_demo_obama_salvage,
        slide_demo_obama_strip,
        slide_demo_judge_entity,
        slide_demo_judge_vocab,
        slide_future_transition,        # Section divider
        slide_24,                       # Reframing the starting point
        slide_26,                       # Five phases
        slide_26b,                      # IS improvement roadmap
        slide_30,                       # Stronger LLM + smart prompts
        slide_30b,                      # LLM upgrade: why it matters
        slide_29,                       # Fine-tuning
        slide_arabic_roadmap,
        slide_arabic_avhubert,
        slide_arabic_changes,
        slide_31,                       # Key takeaways
        slide_thank_you,                # Thank you

        # ─── Appendix (~9, hidden by default) ────────────────────────
        slide_a1,                       # Homophenes
        slide_a8,                       # IS component correlation
        slide_appendix_pca_loadings,    # PCA loadings (NEW)
        slide_human_is_path_b,          # Human-IS Path B (NEW, hidden)
        slide_a11,                      # LLM salvage: recoverable
        slide_a11b,                     # LLM salvage: examples
        slide_a17,                      # Context transition matrix
        slide_appendix_mcnemar_full,    # Full McNemar table (NEW)
        slide_dual_env,                 # Two environments
    ]
    total = len(builders)

    # ─── Hidden builders (academic deck) ─────────────────────────────
    # Per the academic plan, the appendix slides (and the HIDE'd Path B
    # estimates) are hidden by default. The deck-builder retains them
    # so a presenter can choose to reveal them on the fly.
    hidden_builders = {
        slide_a1,
        slide_a8,
        slide_appendix_pca_loadings,
        slide_human_is_path_b,
        slide_a11,
        slide_a11b,
        slide_a17,
        slide_appendix_mcnemar_full,
        slide_dual_env,
    }

    for i, builder in enumerate(builders, 1):
        print(f"  Slide {i:2d}/{total} {builder.__name__:42s} ...", end=" ")
        try:
            builder(prs)
            if builder in hidden_builders:
                prs.slides[-1]._element.set('show', '0')
                print("OK (hidden)")
            else:
                print("OK")
        except Exception as e:
            print(f"ERROR: {e}")

    # ─── Renumber slides (overwrite all bottom-left labels) ──────────
    # Main slides 1..N, appendix slides A1..AM. We use the actual
    # builders → slides correspondence (NOT the old per-slide "A"
    # prefix heuristic, which was broken when a slide moved between
    # main and appendix sections — e.g. slide_a16 promoted into the
    # main body kept its "A8" stamp and got mislabeled).
    #
    # Rule: any builder in `hidden_builders` is an appendix slide.
    # Otherwise it's a main-deck slide.
    from pptx.util import Inches, Pt

    # Build a parallel list of slide-objects in builder order. Some
    # builders (e.g. slide_is_intro) emit multiple sub-slides — we
    # account for that by walking prs.slides and matching on builder
    # boundaries by counting newly-added slides.
    slide_counter = 0
    appendix_num = 0
    slide_idx = 0  # index into prs.slides

    # Pre-compute the slide-count produced by each builder by re-running
    # the builders count-only via attribute reads (avoid running them
    # again — instead, infer from `prs.slides` length post-build).
    # We simply assume 1 builder = 1 slide in the typical case and let
    # any multi-slide builder consume extras. Multi-slide builders are
    # rare (slide_is_intro adds 2 extras based on observed render).
    is_appendix_builder = {b: (b in hidden_builders) for b in builders}

    # Walk prs.slides in order; for each slide, locate its builder by
    # advancing through the builders list as we pass section boundaries.
    builder_iter = iter(builders)
    current_builder = next(builder_iter, None)
    builder_slide_start_count = len(prs.slides)  # not the right snapshot; recompute below

    # Simpler approach: track the slide whose appendix flag we already
    # know from the builder list. The number of slides that each builder
    # adds is recorded in `_slide_to_builder` via builder identity from
    # `_element` tag — but python-pptx doesn't expose that, so we
    # approximate: walk prs.slides; use the bottom-left number-shape's
    # OLD text only as a "this slide came from an A-labelled builder"
    # hint AND combine with the hidden_builders set via slide.show.
    #
    # Concretely: a slide is appendix if EITHER its builder is in
    # hidden_builders (we know via prs.slides[-1]._element.get('show'))
    # OR its old text starts with "A".

    for slide in prs.slides:
        num_shape = None
        for shape in slide.shapes:
            if (shape.has_text_frame
                    and shape.width <= Inches(0.6)
                    and shape.left <= Inches(1.2)
                    and shape.top >= Inches(6.8)):
                num_shape = shape
                break
        if num_shape is None:
            continue

        old_text = num_shape.text_frame.text.strip()
        is_hidden = slide._element.get('show') == '0'

        # Appendix iff the slide is hidden. Every appendix slide in this
        # orchestrator is hidden; visible slides always get a sequential
        # main-deck number even if their source-time stamp started with
        # "A" (e.g. slide_a16 was promoted from appendix to main §2,
        # slide_a13 was promoted to main §3).
        if is_hidden:
            appendix_num += 1
            new_text = f"A{appendix_num}"
        else:
            slide_counter += 1
            new_text = str(slide_counter)

        if old_text != new_text:
            num_shape.text_frame.paragraphs[0].text = new_text

    print(f"  Renumbered: {slide_counter} main + {appendix_num} appendix")

    # ─── Save deck ───────────────────────────────────────────────────
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    _fix_pptx_video_compat(str(OUTPUT))
    _strip_orphan_animation_refs(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)} (hidden: {sum(1 for b in builders if b in hidden_builders)})")

    # ─── Number-verification step ────────────────────────────────────
    # Sanity-check that the canonical AFTER-AMOSI numbers actually appear
    # somewhere in the rendered text. Uses raw XML scan (cheap) rather
    # than re-parsing each slide.
    print("\nNumber verification:")
    import zipfile, re
    with zipfile.ZipFile(str(OUTPUT)) as zf:
        all_text = ""
        for name in zf.namelist():
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                all_text += zf.read(name).decode("utf-8", errors="ignore")
    # Strip XML tags so "2.547" inside <a:t>2.547</a:t> matches cleanly.
    plain = re.sub(r"<[^>]+>", "", all_text)
    checks = [
        ("2.547",  "MBR IS"),
        ("62%",    "NIV-Y+P (MBR)"),
        ("71%",    "judge MBR Y+P"),
        ("6%",     "trust gate FPR"),
    ]
    for needle, label in checks:
        status = "PASS" if needle in plain else "FAIL"
        print(f"  [{status}] {needle:8s} ({label})")

    # ─── Blocker report ──────────────────────────────────────────────
    if _BLOCKERS:
        print(f"\nBLOCKERS — {len(_BLOCKERS)} sibling-agent functions not yet defined:")
        for b in _BLOCKERS:
            print(f"  - {b}")
        print("Run the sibling slide-writer agents, then re-run this script.")


if __name__ == "__main__":
    main()
