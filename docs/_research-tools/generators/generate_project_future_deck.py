#!/usr/bin/env python3
"""Argos VSP - internal continuation deck for the Amosi talk (July 20 2026).

INTERNAL sponsor deck: scaffold for a ~45-minute talk by the departing lead
to Amosi (orchard manager, project sponsor) deciding continuation, scale and
staffing. Amosi missed the mid-July client meeting, so S3-S4 recap it.
Unlike client decks, N1-N10 does NOT apply here: WER/IS jargon, client
names, honest failure data and reference-vs-hypothesis material are allowed.
Slides carry anchors + numbers + prompts; the long form lives in speaker
notes (V5). 11 main slides + backup divider + 3 backup slides.

Output:
    presentation_materials_20260224/Argos_VSP_ProjectFuture_20260720.pptx

Run:
    cd /home/ubuntu/docs/_research-tools/generators && \
        /home/ubuntu/vsp-llm-yoad-venv/bin/python generate_project_future_deck.py

Numbers are sourced from:
    docs/evaluation/after_amosi_audit.md            (canonical metrics)
    docs/evaluation/egla_kafe/resolution_ablation.md (capture-chain story)
    docs/evaluation/egla_kafe/findings.md            (client eval + forensics)
    docs/evaluation/egla_kafe/guessing_game_answer_key.md (cold-read rates)
    docs/prompts/topic_label_experiment.md           (word-subset evidence)
    docs/paper/arabic-vsp-adaptation.md              (Arabic path)
    docs/guides/project-handover-july2026.md         (bets, briefs, succession)

Text geometry is sized to scripts/audit_pptx_text_render.py's estimator
(cpl = (w-0.1)/(0.55*pt/72), line height 1.2*pt): keep flagged strings
within the per-box line budget rather than shrinking fonts.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from presentation.config import (
    SL_W, SL_H, WHITE, TEAL, CORAL, LGRAY, MGRAY,
    GREEN, GOLD, NAVY2,
    MX, CT, CW, _auto_num,
)
from presentation.helpers import (
    new_slide, add_title, add_accent_line, add_text,
    add_rect, add_image, add_logo,
    set_notes, add_fade_transition, _strip_orphan_animation_refs, _fmt,
)

OUTPUT = Path("/home/ubuntu/presentation_materials_20260224/"
              "Argos_VSP_ProjectFuture_20260720.pptx")

POSTERS = Path("/home/ubuntu/presentation_materials_20260224/.poster_frames")
EK_PLOTS = Path("/home/ubuntu/docs/evaluation/egla_kafe/plots")


# ── Small local helpers (internal-deck conventions) ─────────────────────

def _slide_num12(slide, label):
    """Bottom-left slide number at the 12pt floor (helpers' version is 9pt).

    Box width kept <= 0.6in so the text-render audit's slide-num exemption
    (sw <= 0.6, sx <= 1.2, sy >= 6.8) recognizes it.
    """
    tb = slide.shapes.add_textbox(MX, SL_H - Inches(0.38),
                                  Inches(0.55), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.text = str(label)
    _fmt(p, size=Pt(12), color=MGRAY)
    return tb


def _foot(slide, text, y=Inches(6.55)):
    add_text(slide, text, MX + Inches(0.6), y, CW - Inches(1.2),
             Inches(0.35), size=Pt(16), color=LGRAY, italic=True,
             align=PP_ALIGN.CENTER)


def _finish(slide, notes="", num=None):
    """num=None -> next main number; num='B1' -> literal; num=False -> none."""
    if num is None:
        _auto_num[0] += 1
        num = _auto_num[0]
    add_logo(slide)
    if num is not False:
        _slide_num12(slide, num)
    if notes:
        set_notes(slide, notes)


def _card_row(slide, cards, cy, ch, gap=Inches(0.35),
              title_h=Inches(0.45), stat_h=Inches(0.45)):
    """Row of equal-width bordered cards: (title, body, stat_or_None, color).

    title_h: 0.45 for 1-line titles, 0.85 when a title may wrap to 2 lines.
    stat_h: 0.45 for 1-line stats, 0.6 when a stat may wrap to 2 lines.
    """
    n = len(cards)
    cw = (CW - (n - 1) * gap) / n
    body_y = Inches(0.12) + title_h + Inches(0.03)
    for i, (t, d, stat, col) in enumerate(cards):
        x = MX + i * (cw + gap)
        add_rect(slide, x, cy, cw, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, x + Inches(0.2), cy + Inches(0.12),
                 cw - Inches(0.4), title_h, size=Pt(24), bold=True,
                 color=col)
        add_text(slide, d, x + Inches(0.2), cy + body_y,
                 cw - Inches(0.4), ch - body_y - stat_h - Inches(0.1),
                 size=Pt(24), color=WHITE)
        if stat:
            add_text(slide, stat, x + Inches(0.2),
                     cy + ch - stat_h - Inches(0.05),
                     cw - Inches(0.4), stat_h, size=Pt(18),
                     color=LGRAY, italic=True)


def _card_stack(slide, cards, cy, ch, gap=Inches(0.2)):
    """Stack of full-width cards: (title, body, color)."""
    for i, (t, d, col) in enumerate(cards):
        y = cy + i * (ch + gap)
        add_rect(slide, MX, y, CW, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, MX + Inches(0.3), y,
                 Inches(4.6), ch, size=Pt(25), bold=True,
                 color=col, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, d, MX + Inches(5.0), y,
                 CW - Inches(5.3), ch, size=Pt(24), color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)


def _mini_stack(slide, cards, x, w, cy, ch, gap=Inches(0.15)):
    """Narrow stacked cards (title above body): (title, body, color)."""
    for i, (t, d, col) in enumerate(cards):
        y = cy + i * (ch + gap)
        add_rect(slide, x, y, w, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, x + Inches(0.2), y + Inches(0.08),
                 w - Inches(0.4), Inches(0.45), size=Pt(24), bold=True,
                 color=col)
        add_text(slide, d, x + Inches(0.2), y + Inches(0.55),
                 w - Inches(0.4), ch - Inches(0.62), size=Pt(24),
                 color=WHITE)


def _arrow(slide, x, y):
    add_text(slide, "→", x, y, Inches(0.5), Inches(0.5),
             size=Pt(28), bold=True, color=LGRAY, align=PP_ALIGN.CENTER)


def _para_lines(slide, lines, left, top, width, height, size=Pt(24),
                color=WHITE):
    """Multi-line body as separate <a:p> paragraphs (no bullet chars)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ln
        para.alignment = PP_ALIGN.LEFT
        _fmt(para, size=size, color=color)
    return tb


# ════════════════════════════════════════════════════════════════════════
# MAIN SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_f_thesis(prs):
    slide = new_slide(prs)
    add_rect(slide, 0, Inches(1.75), SL_W, Inches(0.03), fill_color=TEAL)
    add_text(slide,
             "A working product, deployed at a client, on a "
             "prototype-grade model - and the client has offered "
             "the path to continuation.",
             MX + Inches(0.4), Inches(2.0), CW - Inches(0.8), Inches(2.8),
             size=Pt(40), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide,
             "What’s proven   ·   What the client agreed to   "
             "·   What continuation requires",
             MX, Inches(5.05), CW, Inches(0.85),
             size=Pt(24), color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "Argos  ·  internal continuation review  ·  "
                    "July 2026",
             MX, Inches(6.4), CW, Inches(0.4),
             size=Pt(20), color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide,
            "Open with the thesis sentence verbatim, then the three "
            "questions of the talk. Context for Amosi: I hand over soon "
            "(3 workdays this week + 4 in early August); this session is "
            "the continuation decision, not a status meeting. The handover "
            "one-liner is deliberate: product works and is deployed; the "
            "model underneath is prototype-grade; the client meeting gave "
            "us a concrete path to 'continue'. Everything today hangs off "
            "those three clauses.")
    return slide


def slide_f_shipped(prs):
    slide = new_slide(prs)
    add_title(slide, "Proven and shipped")
    add_accent_line(slide)
    _card_row(slide, [
        ("Consensus decoding - default",
         "Judge-useful 71.1% vs 68.4%\nWER 63.8  (top-1: 64.1)",
         "1,497 wild segments · p = 0.0002", GREEN),
        ("Trust gate + word confidence",
         "≥30% green → 65.2% recall, 5.6% FPR",
         "word colors in every report", TEAL),
    ], CT + Inches(0.15), Inches(2.1), stat_h=Inches(0.6))
    _card_row(slide, [
        ("Phonetic substitution - GO",
         "2 corrections shipped\nfigured → forgot",
         "dual-engine agreement arm", GOLD),
        ("Guessing game · Jul 16",
         "7 videos, hypothesis-only, audio stripped",
         "zero leaks · awaiting their read", TEAL),
    ], CT + Inches(2.45), Inches(2.1), stat_h=Inches(0.6))
    _foot(slide, "All four live in production · IS 2.547 vs 2.532 · "
                 "sources in notes.", y=Inches(6.3))
    _finish(slide,
            "Prompt: 'all of this is live in production, not slideware.' "
            "Canonical source: docs/evaluation/after_amosi_audit.md. "
            "MBR consensus is the displayed output since May 2 (blind-judge "
            "useful rate 68.4->71.1% on 1,497 segments, paired McNemar "
            "p=0.0002; WER 64.05->63.84; IS 2.532->2.547). Trust gate at "
            ">=30% green words: 65.2% recall of useful segments at 5.6% "
            "false-positive rate (precision 95.6%, 630 segments trusted). "
            "Phonetic substitution ships only dual-engine-agreed word fixes "
            "(2 on the client set, both real script words). Guessing-game "
            "package: 412 MB zip shipped Jul 16, QA 14/14 streams, leak "
            "greps all zero.")
    return slide


def slide_f_meeting_shown(prs):
    slide = new_slide(prs)
    add_title(slide, "Mid-July client meeting - what we showed")
    add_accent_line(slide)
    _mini_stack(slide, [
        ("Their footage, end-to-end",
         "21 conversations · reports + subtitle videos", TEAL),
        ("Capability on good capture",
         "Best video: 3 of 4 turns understood", GREEN),
        ("It knows when it’s right",
         "Confident subset 70–92% useful", GOLD),
    ], MX, Inches(5.6), CT + Inches(0.25), Inches(1.45))
    add_image(slide, POSTERS / "ek_best_vs_worst.jpg",
              Inches(6.6), Inches(1.75), width=Inches(5.5))
    add_text(slide, "Their footage: best vs worst conditions",
             Inches(6.6), Inches(5.65), Inches(5.5), Inches(0.75),
             size=Pt(20), color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide,
            "Amosi missed this meeting - 90-second replay of what was "
            "on the wall (client side: Altman + Shier, ~1 hour). We showed "
            "the full evaluation of their own 21 conversations: per-video "
            "reports, subtitle videos, and the confidence story. The three "
            "proof points are the same trio the July 13 deck carried: best "
            "video (img_6825) 3 of 4 turns understood with the plot "
            "recoverable; self-flagged confident subset 70–92% useful "
            "(70% keeping top 10%, 92% keeping 3%); and capture conditions "
            "dominating everything downstream - expanded on the "
            "capture-chain slide. Vibe: supportive, not yet convinced; I "
            "stayed confident. Source: docs/evaluation/egla_kafe/"
            "findings.md.")
    return slide


def slide_f_meeting_landed(prs):
    slide = new_slide(prs)
    add_title(slide, "Where the client landed")
    add_accent_line(slide)
    _card_stack(slide, [
        ("Their posture",
         "Supportive, not yet convinced\nwants it to work", GREEN),
        ("Agreed in the room",
         "Capture protocol + re-shoot · guessing game now", TEAL),
        ("Delivered since",
         "Guessing-game package Jul 16\nno feedback yet", GOLD),
    ], CT + Inches(0.25), Inches(1.35), gap=Inches(0.18))
    _foot(slide, "In the room: “hard to tell it works unless you "
                 "know the topic.”", y=Inches(6.35))
    _finish(slide,
            "Their posture in one line: engaged and rooting for it, but "
            "not yet convinced by the raw output - their words: usable "
            "but requires sitting down and analyzing; scattered correct "
            "bits. They accepted the screen-recording explanation without "
            "pushback and worried their camera is too weak; I reassured "
            "them full HD with enough mouth pixels is fine. Agreed: they "
            "adopt the capture protocol, we run the guessing game (shipped "
            "Jul 16, hypothesis-only, no feedback yet), then a filming "
            "round 2. A follow-up client meeting still needs a date - "
            "that lands on whoever continues this.")
    return slide


def slide_f_capture_chain(prs):
    slide = new_slide(prs)
    add_title(slide, "Why their footage failed - the filming spec")
    add_accent_line(slide)
    add_text(slide,
             "Screen recordings - not their camera.   "
             "4K = 2K = 1080p  (175 pairs, n.s.)",
             MX, Inches(1.5), CW, Inches(0.85),
             size=Pt(24), bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    _mini_stack(slide, [
        ("1 · Framing first",
         "Mouth ≥100 px in frame", GREEN),
        ("2 · Original file exports",
         "Never screen-record", TEAL),
        ("3 · Resolution last",
         "Full HD suffices", GOLD),
    ], MX, Inches(5.6), Inches(2.5), Inches(1.18), gap=Inches(0.15))
    add_image(slide, POSTERS / "ek_iphone_vs_cam.jpg",
              Inches(7.0), Inches(2.5), width=Inches(5.2))
    add_text(slide,
             "iPhone vs screen-recorded camera - clear turns "
             "13.7% vs 0.8% (17×)",
             Inches(6.7), Inches(6.2), Inches(5.9), Inches(0.7),
             size=Pt(20), color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide,
            "Told to the client openly and accepted: their footage was "
            "ruined by HOW it was captured - screen recordings of the "
            "viewer app (odd sizes like 1258×696, 3–4 Mbps, UI "
            "chrome baked into pixels) - not by their camera. Same "
            "scene: iPhone IS 1.51 vs screen-rec 0.88 (p=2.3e-05); clearly "
            "conveyed 13.7% vs 0.8%. The ablation kills the resolution "
            "excuse: 4K/2K/1080p paired on 175 segments, IS "
            "1.554/1.505/1.564/1.586, every test n.s. - the pipeline "
            "warps every mouth to ~45 px in a 96×96 crop. Biggest "
            "measured effect was the second lossy encode. Frontality is "
            "the other real lever (IS 1.62 front vs 1.03 at 45°). "
            "This slide doubles as the round-2 filming spec; the viewer "
            "app has zoom, so original exports may even make a re-shoot "
            "unnecessary. Sources: docs/evaluation/egla_kafe/"
            "resolution_ablation.md, docs/evaluation/egla_kafe/"
            "findings.md.")
    return slide


def slide_f_envelope(prs):
    slide = new_slide(prs)
    add_title(slide, "What it actually delivers today")
    add_accent_line(slide)
    half = (CW - Inches(0.35)) / 2
    add_text(slide, "It delivers", MX, CT + Inches(0.1), half,
             Inches(0.45), size=Pt(24), bold=True, color=GREEN)
    add_text(slide, "It does not", MX + half + Inches(0.35),
             CT + Inches(0.1), half, Inches(0.45),
             size=Pt(24), bold=True, color=CORAL)
    _card_row(slide, [
        ("Gist-level comprehension",
         "Cold read: 33–56% of turns land",
         "guessing-game answer key, per video", GREEN),
        ("Names and places",
         "0% recovered - confident fabrications",
         "“Abu dhabi” at 0.997 confidence", CORAL),
    ], Inches(2.1), Inches(1.9))
    _card_row(slide, [
        ("Trustable green words",
         "Green common nouns ~82% correct",
         "~10% of words are green", GREEN),
        ("Unattended automation",
         "Analyst effort required - ~half of segments useful",
         "hallucination in 20.5% of segments", CORAL),
    ], Inches(4.15), Inches(1.9))
    _foot(slide, "“Useful” ≠ “automatic” - NIV useful 61.9% · "
                 "judge-useful 71.1%.", y=Inches(6.3))
    _finish(slide,
            "The honest envelope, stated before anyone else states it: "
            "gist-level comprehension with scattered correct bits, real "
            "analyst effort required. Cold-read rates from our own "
            "guessing-game answer key: 33–56% of turns land (per "
            "video); green common nouns ~82% correct but only ~10% "
            "coverage; names and places 0% recovered and dangerously "
            "confident when fabricated ('Abu dhabi' at 0.997). Roughly "
            "half the segments carry useful content (NIV-Y+P 61.9%, "
            "salvage-adjusted ~62%); hallucination still hits 20.5% of "
            "segments. Prompt: this honesty is exactly why the next slide "
            "- the persuasion plan - exists.")
    return slide


def slide_f_agreed_path(prs):
    slide = new_slide(prs)
    add_title(slide, "The path the client proposed")
    add_accent_line(slide)
    gap = Inches(0.5)
    cw = (CW - 3 * gap) / 4
    cards = [
        ("1 · Guessing game", "Shipped - awaiting their read",
         GOLD),
        ("2 · Filming round 2",
         "Long sentences\nFraming\nFile exports",
         GOLD),
        ("3 · Wow reel", "Cherry-picked - for their bosses",
         GOLD),
        ("4 · Their verdict",
         "“We need this - continue”", GREEN),
    ]
    cy, ch = Inches(1.75), Inches(3.35)
    for i, (t, d, col) in enumerate(cards):
        x = MX + i * (cw + gap)
        add_rect(slide, x, cy, cw, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, x + Inches(0.18), cy + Inches(0.15),
                 cw - Inches(0.36), Inches(0.95), size=Pt(24), bold=True,
                 color=col)
        _para_lines(slide, d.split("\n"), x + Inches(0.18),
                    cy + Inches(1.2), cw - Inches(0.36),
                    ch - Inches(1.35), size=Pt(24), color=WHITE)
        if i < 3:
            _arrow(slide, x + cw + Inches(0.02), Inches(3.2))
    add_text(slide,
             "Their offer: direct line to our bosses + ongoing tests, "
             "footage, feedback - possibly small budget.",
             MX, Inches(5.45), CW, Inches(0.95),
             size=Pt(26), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _finish(slide,
            "The heart of the decision case. Step 2 'realistic settings' "
            "as the client described: two people talking naturally, filmed "
            "from a distance, simple camera, longer sentences, direct "
            "exports - they defer to our judgment on specifics. Step "
            "3: emotional impact over literal usefulness; cherry-picking "
            "is acceptable in their culture. Step 4 is the payoff: our "
            "management wants a client saying 'it works well' before "
            "continuing - the client explicitly offered to talk to "
            "our management directly, to keep investing time in tests, "
            "filming and feedback, and possibly small money (not heavy "
            "- they avoid funding politics). Prompt: the client "
            "volunteered to be exactly the voice our bosses asked for.")
    return slide


def slide_f_new_asks(prs):
    slide = new_slide(prs)
    add_title(slide, "Two new technical asks from the client")
    add_accent_line(slide)
    half = (CW - Inches(0.35)) / 2
    x2 = MX + half + Inches(0.35)
    cy, ch = Inches(1.75), Inches(3.9)
    for x, col, title, lines, stat in [
        (MX, TEAL, "A · Predetermined word subsets",
         ["Anchor a wordlist",
          "Prompting tested: worse",
          "Route 1: constrained decode",
          "Route 2: vocab fine-tune"],
         "topic labels: WER 86.6→87.6 · 24% echo"),
        (x2, GOLD, "B · Arabic transition",
         ["~80% of stack carries over",
          "Retrain k-means + LLM",
          "Bottleneck: Arabic corpus",
          "2–3 months from green-light"],
         "ask: which Arabic - MSA / Levantine / Gulf?"),
    ]:
        add_rect(slide, x, cy, half, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, title, x + Inches(0.2), cy + Inches(0.15),
                 half - Inches(0.4), Inches(0.85), size=Pt(24), bold=True,
                 color=col)
        _para_lines(slide, lines, x + Inches(0.2), cy + Inches(1.05),
                    half - Inches(0.4), Inches(2.1), size=Pt(24),
                    color=WHITE)
        add_text(slide, stat, x + Inches(0.2), cy + ch - Inches(0.65),
                 half - Inches(0.4), Inches(0.6), size=Pt(18),
                 color=LGRAY, italic=True)
    _foot(slide, "Neither is promised - both are scoped · detail in "
                 "backup slides.", y=Inches(6.35))
    _finish(slide,
            "Both asks came from the client in the meeting. Word subsets: "
            "a genuinely good idea we have NOT built - be precise "
            "with Amosi. The only decode-time injection we measured "
            "(topic labels, 284 segments) made things worse: WER 86.6->"
            "87.6 and 24% of outputs echoed the instruction, because the "
            "model is QLoRA-locked to a fixed prompt. Two honest routes: "
            "constrained beam decoding / logit biasing (engineering, no "
            "code exists, testable), or vocab-conditioned fine-tuning, "
            "which rides the Llama-3.1 + 20K-segment bet. Arabic: "
            "AV-HuBERT is ~80% language-agnostic; retrain k-means (an "
            "Arabic k-means artifact already exists), swap in an Arabic-"
            "capable LLM (Jais/AceGPT/Llama-3 - with the caveat that "
            "multilingual pretraining can hurt VSR), retrain the Q-Former "
            "bridge. Real bottleneck: no Arabic LRS3-scale corpus - "
            "scraping is manageable but is the long pole. Framing "
            "precedent: 'roadmap mapped, work not started, 2–3 months "
            "from green-light.' Sources: docs/prompts/"
            "topic_label_experiment.md, docs/paper/"
            "arabic-vsp-adaptation.md.")
    return slide


def slide_f_closed_bets(prs):
    slide = new_slide(prs)
    add_title(slide, "What did not work - closed bets")
    add_accent_line(slide)
    _card_stack(slide, [
        ("Decode-param tuning",
         "13 experiments - baseline wins\nbest alt: +13% hallucinations", CORAL),
        ("Small-data LoRA",
         "1,273 segments overfit\nreal need: 20K–50K + stronger LLM", CORAL),
        ("xseg_merge",
         "Silent no-op; wired, it breaks 1 : 10.7 - retired", CORAL),
        ("Span substitution",
         "149 fixed / 303 broken on wild data - NO-GO", CORAL),
    ], CT + Inches(0.2), Inches(1.05), gap=Inches(0.15))
    _foot(slide, "Each closed with data and documented - the successor "
                 "never re-runs these.")
    _finish(slide,
            "Credibility slide: what we spent effort on and shut down. "
            "Decode-parameter tuning: 13 experiments (A–M); baseline "
            "beam=20/lenpen=0 most robust; the best alternative bought "
            "+0.08 IS at +13% hallucinations. LoRA fine-tuning on 1,273 "
            "segments: IS fell 2.487 → 2.312 (r=16) → 2.023 "
            "(r=64) with empty outputs 7.1% → 26.8% - a data-"
            "scale problem, not capacity; the honest requirement is "
            "20K–50K segments plus a stronger LLM. xseg_merge shipped "
            "as a silent no-op (0 neighbors on all 1,497); simulated "
            "fixed, it breaks ~10.7 segments for every one it fixes - "
            "retired. Span-level substitution broke twice what it fixed "
            "(149F/303B, p≈4e-13 in the wrong direction) - only "
            "the word-level dual-engine arm shipped.")
    return slide


def slide_f_continuation(prs):
    slide = new_slide(prs)
    add_title(slide, "What continuation requires")
    add_accent_line(slide)
    _card_row(slide, [
        ("People",
         "7 lead workdays remain\nfinish vs transfer",
         "draft split in notes · successor unstaffed", TEAL),
        ("Data",
         "Filming round 2 (client) · LRS3 sourcing",
         "LRS3 blocks the Llama-3.1 bet", TEAL),
    ], Inches(1.6), Inches(2.1), stat_h=Inches(0.6))
    _card_row(slide, [
        ("Compute",
         "Llama-3.1 ready - needs p4d quota (~$300, ~10h)",
         "M4.1 calibration: ~1 GPU-day", TEAL),
        ("Scale",
         "Continue at what level\nopen decision",
         "second group waiting: ~100h footage handed over", TEAL),
    ], Inches(3.9), Inches(2.1), stat_h=Inches(0.6))
    _foot(slide, "Decision owner: Amosi - discussion, not a pitch.",
          y=Inches(6.3))
    _finish(slide,
            "Deliberately neutral - no recommendation on the slide; "
            "mine is delivered verbally here. DRAFT finish-vs-transfer "
            "split for my 7 remaining workdays (edit live): FINISH - "
            "guessing-game feedback intake, scheduling the next client "
            "meeting, M4.1 confidence calibration if a GPU day is "
            "available, handover polish. TRANSFER - executing the "
            "Llama-3.1 migration (code ready, smoke-tested; blocked on "
            "LRS3 sourcing + p4d.24xlarge quota, ~$300 and ~10h when "
            "unblocked; swap alone ~1–2pp WER, the real unlock is "
            "context injection + domain-data scaling), the three 5-day "
            "briefs, filming round 2, Arabic scoping. Demand signal for "
            "the scale card: a second interested group already handed "
            "over ~100 hours of real footage (Q2 summary).")
    return slide


def slide_f_succession(prs):
    slide = new_slide(prs)
    add_title(slide, "Handover - where everything lives")
    add_accent_line(slide)
    _card_stack(slide, [
        ("Handover doc",
         "First-week checklist scripted - docs/guides/", GREEN),
        ("Ready-to-assign briefs",
         "3 × 5-day: quality pre-filter,\nmulti-speaker, additivity", TEAL),
        ("Working assets",
         "Deployed container · client eval re-runs in days", GOLD),
    ], CT + Inches(0.25), Inches(1.35), gap=Inches(0.18))
    _foot(slide, "Nothing lives only in my head.", y=Inches(6.35))
    _finish(slide,
            "Close: succession is prepared, not aspirational. "
            "docs/guides/project-handover-july2026.md carries system "
            "state, deployment doctrine, client commitments, open bets, "
            "traps, and a scripted first week (read the docs in order, "
            "run the 37-test suite, one video end-to-end, learn the "
            "overlay sync, own the client follow-ups). Three onboarding-"
            "sized 5-day project briefs are written and estimated: M15 "
            "video-quality pre-filter (high value, no dependencies - "
            "the intended ramp project), M12 multi-speaker (client-"
            "requested), M4.2 signal-additivity. The client-eval pipeline "
            "re-runs on any new footage in days, not weeks - whoever "
            "takes over can repeat the Egla-Kafe evaluation on round-2 "
            "footage immediately.")
    return slide


# ════════════════════════════════════════════════════════════════════════
# BACKUP SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_f_backup_divider(prs):
    slide = new_slide(prs)
    add_rect(slide, 0, Inches(3.0), SL_W, Inches(0.03), fill_color=TEAL)
    add_text(slide, "Backup", MX, Inches(3.25), CW, Inches(1.0),
             size=Pt(44), bold=True, color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide, num=False)
    return slide


def slide_b_resolution(prs):
    slide = new_slide(prs)
    add_title(slide, "Backup · resolution ablation")
    add_accent_line(slide)
    add_text(slide,
             "Every mouth is warped to ~45 px in a 96×96 crop",
             MX, Inches(1.5), CW, Inches(0.5),
             size=Pt(24), bold=True, align=PP_ALIGN.CENTER)
    add_image(slide, EK_PLOTS / "resolution_metrics_by_condition.png",
              MX, Inches(2.15), width=Inches(6.1))
    add_text(slide, "IS / WER per condition",
             MX, Inches(4.85), Inches(6.1), Inches(0.35),
             size=Pt(16), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)
    add_image(slide, EK_PLOTS / "resolution_mouthpx_vs_delta.png",
              Inches(7.1), Inches(2.15), width=Inches(5.3))
    add_text(slide, "Mouth px vs quality delta - no cliff",
             Inches(7.1), Inches(5.42), Inches(5.3), Inches(0.35),
             size=Pt(16), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)
    add_text(slide,
             "IS 1.554 / 1.505 / 1.564 / 1.586  (4K / re-encode / 2K / 1080p) - all n.s.\nBiggest measured effect: the second lossy encode",
             MX, Inches(6.0), CW, Inches(0.7),
             size=Pt(20), color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide,
            "For 'just buy a better camera' pushback. Paired sweep, same "
            "175 segments per arm, byte-identical segmentation, decode "
            "identical. Native 4K mouths are only 55–68 px (one "
            "subject 104 px); 1080p mouths 30–33 px get mildly "
            "upsampled toward the ~45 px canonical scale - which is "
            "why resolution washes out, and 4K's extra pixels are "
            "discarded at input. The model's own mean word probability "
            "is 0.40 at every resolution. Re-encoding once (control arm) "
            "cost as much IS as discarding 75% of pixels. Frontality is "
            "the lever that survives: front IS 1.62 vs 1.03 at 45° "
            "(p=2e-03). Source: docs/evaluation/egla_kafe/"
            "resolution_ablation.md.", num="B1")
    return slide


def slide_b_word_subsets(prs):
    slide = new_slide(prs)
    add_title(slide, "Backup · word subsets - evidence and routes")
    add_accent_line(slide)
    _card_stack(slide, [
        ("Motivation",
         "19% of segments (284)\ndomain-vocab confusion", TEAL),
        ("Evidence",
         "Topic labels: WER 86.6 → 87.6\n24% instruction echo",
         CORAL),
        ("Route 1 - decoding",
         "Constrained beam / logit bias - no code today", GREEN),
        ("Route 2 - training",
         "Vocab-conditioned fine-tune\nLlama 3.1 + 20K segments", GREEN),
    ], CT + Inches(0.2), Inches(1.05), gap=Inches(0.15))
    _foot(slide, "Prompting alone cannot fix domain vocabulary - the "
                 "model is prompt-locked.")
    _finish(slide,
            "The client's ask maps to Report 3 Strategy 3 (vocabulary "
            "constraints) - proposed in February, never implemented. "
            "The measured neighbor experiment: injecting topic labels on "
            "the 284 domain-confused segments made WER worse (86.6→"
            "87.6) and 24% of outputs echoed the instruction text instead "
            "of transcribing - the QLoRA fine-tune locked the model "
            "to one prompt format and extra tokens shift visual features "
            "into positions the mapping never saw. Mislabeled topics cost "
            "up to -10.3pp. So: route 1 modifies the beam search itself "
            "(penalize/boost wordlist tokens - clean engineering "
            "experiment, unbuilt); route 2 bakes vocab conditioning into "
            "fine-tuning, which needs the Llama-3.1 migration plus 20K+ "
            "segments. Sources: docs/prompts/report_3_prompt_engineering."
            "md §3.3, docs/prompts/topic_label_experiment.md.",
            num="B2")
    return slide


def slide_b_arabic(prs):
    slide = new_slide(prs)
    add_title(slide, "Backup · Arabic - what it takes")
    add_accent_line(slide)
    _card_stack(slide, [
        ("Transfers (~80%)",
         "AV-HuBERT visual stack is language-agnostic", GREEN),
        ("Must retrain",
         "k-means (artifact exists) · Arabic LLM · bridge",
         TEAL),
        ("Bottleneck",
         "No Arabic LRS3-scale corpus\ncollection required", CORAL),
        ("Estimate",
         "2–3 months from green-light\nask: which dialect?",
         GOLD),
    ], CT + Inches(0.2), Inches(1.05), gap=Inches(0.15))
    _foot(slide, "Same data-scale reality that closed the English "
                 "fine-tuning bets.")
    _finish(slide,
            "Lip/jaw geometry is largely universal, so the frozen "
            "AV-HuBERT encoder plus Arabic-retrained downstream is a "
            "viable first attempt (subtle gap: Arabic emphatics were "
            "never learned). Concretely: retrain k-means on Arabic audio "
            "(an arabic_flat_kmeans_200.bin artifact already exists), "
            "swap the LLM for an Arabic-capable model (Jais, AceGPT, "
            "multilingual Llama-3 - caveat: heavy multilingual "
            "pretraining measurably hurt VSR when we compared LLMs), "
            "retrain the Q-Former bridge + LoRA on Arabic pairs. The "
            "long pole is data: no Arabic LRS3 equivalent exists; "
            "broadcast/YouTube collection is manageable but real work. "
            "Phased plan and 2–3-month estimate: docs/paper/"
            "arabic-vsp-adaptation.md. Open question for the client: "
            "which Arabic (MSA / Levantine / Egyptian / Gulf).", num="B3")
    return slide


def main():
    _auto_num[0] = 0
    prs = Presentation()
    prs.slide_width = SL_W
    prs.slide_height = SL_H

    builders = [
        slide_f_thesis,
        slide_f_shipped,
        slide_f_meeting_shown,
        slide_f_meeting_landed,
        slide_f_capture_chain,
        slide_f_envelope,
        slide_f_agreed_path,
        slide_f_new_asks,
        slide_f_closed_bets,
        slide_f_continuation,
        slide_f_succession,
        slide_f_backup_divider,
        slide_b_resolution,
        slide_b_word_subsets,
        slide_b_arabic,
    ]
    total = len(builders)
    for i, b in enumerate(builders, 1):
        print(f"  Slide {i:2d}/{total}  {b.__name__} ...", end=" ")
        b(prs)
        print("OK")

    for s in prs.slides:
        add_fade_transition(s, speed="med")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    _strip_orphan_animation_refs(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
