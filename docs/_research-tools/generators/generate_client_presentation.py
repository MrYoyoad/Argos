#!/usr/bin/env python3
"""Argos VSP — Client Deck Generator (Round 5 — framing-v2 alignment).

Three-act 2-hour meeting structure per docs/CLIENT_MEETING_FRAMING_v2.md.
Calls existing builders from presentation/slides_*.py per the borrow-vs-
build policy in .claude/plans/i-need-to-create-proud-cupcake.md. Net-new
client builders live in presentation/slides_client.py.

Output:
    presentation_materials_20260224/Argos_VSP_Client_60slides_Apr2026.pptx
    (final count determined at build time)

Run:
    python3 docs/_research-tools/generators/generate_client_presentation.py
"""

from pathlib import Path

from pptx import Presentation

from presentation.config import SL_W, SL_H, _auto_num, MATERIALS

# ── Borrowed builders kept across rounds ────────────────────────────────
from presentation.slides_opening import (
    slide_data_flow,                    # pipeline overview (clean enough as-is)
)
from presentation.slides_engineering import (
    slide_web_ui,                       # UI overview
    slide_17_png,                       # 8-stage pipeline (full PNG) — appendix-only
)
from presentation.slides_future import (
    slide_thank_you,                    # close
)

# ── Round-5.1 academic borrows ──────────────────────────────────────────
# slide_judge_ex1-6 → replaced by slide_client_judge_ex1-6 (verdict tags
# instead of WER/WWER score strips per N9). The academic builders stay
# defined; we just don't import them here.
# slide_25e → DROPPED (jalapeno duplicates slide_client_judge_ex5).
# slide_failure_deep_1a + 1b → MERGED into slide_client_failure_taxonomy_full.
# slide_08 → DROPPED (duplicates the new failure_taxonomy slide).
# slide_arabic_roadmap → REPLACED by slide_client_arabic_high_level.
# Round 8 — slide_25d DROPPED (jesus phonetic-bridge example too sketchy
# for client deck; replacement would require editing academic-shared
# slides_evaluation.py). Whole slide dropped instead of card swap.
from presentation.slides_research import (
    slide_failure_deep_2,               # Taxonomy worked examples (3 ref/hyp cards)
)

# ── Net-new and reframed client builders ────────────────────────────────
from presentation.slides_client import (
    # Section 1 — Hello
    slide_client_title,
    slide_client_about_argos,
    slide_client_value_prop_headline,           # NEW 8.8 — outcomes-first dashboard

    # Section 2 — Background (deepened in Round 5.1)
    slide_client_what_is_vsr,
    slide_client_what_is_lipreading_not,    # NEW Round 5.1
    slide_client_canonical_scenario,        # NEW Round 5.1 (coffee shop)

    # Section 3 — Compared to today
    slide_client_compared_to_today,
    slide_client_human_ceiling,             # NEW Round 5.1 (why 45-52%)

    # Section 4 — Why hard
    slide_client_visemes,

    # Section 5 — Basic idea
    slide_client_pipeline_components,

    # Section 6 — What we built (multi-speaker moved here in Round 5.1)
    slide_client_what_we_built,
    slide_client_engineering_journey,
    slide_client_multi_speaker_today_vs_path,   # MOVED UP from §13

    # Section 7 — How to use it
    slide_client_demo_intro,
    slide_client_demo_video_embed,
    slide_client_demo_recap,

    # Section 8 — Real outputs
    slide_client_video_gallery,
    # slide_client_examples_intro DROPPED (Round 5.1 — overlaps with deep examples)
    slide_client_example_perfect,
    slide_client_example_simple_positive,   # NEW Round 7 — non-Obama clean win
    slide_client_realtalk_spectrum,         # NEW 8.8 — conversational two-people Trust/Salvage/Strip
    slide_client_clean_outputs_gallery,     # NEW Round 5.1 (6-tile gallery of clean outputs)
    slide_client_example_partial,
    # slide_client_example_flagged — DROPPED Round 5.10 from active deck
    # (builder kept defined in slides_client.py; can be re-imported if
    # the Obama-flagged example is wanted back). Same Obama segment 5
    # content lives on slide_client_halluc_real_example in the
    # hallucination trio.
    # 6 client-side judge-example wrappers (verdict-tag instead of WER score strip)
    slide_client_judge_ex1,
    slide_client_judge_ex2,
    slide_client_judge_ex3,
    slide_client_judge_ex4,
    slide_client_judge_ex5,
    slide_client_judge_ex6,

    # Section 9 — Headline
    slide_client_headline_numbers,

    # Section 10 — Why trust it
    slide_client_confidence_question,
    slide_client_two_layer_confidence,
    slide_client_word_color_coding,
    slide_client_three_tier_policy,        # NEW Round 5.6 — band-reliability + UI policy
    slide_client_how_to_read,              # NEW Round 5.9 — operational instruction (decision flow)
    slide_client_reader_example,           # NEW Round 5.9 — Salvage worked example
    slide_client_case_topic_shift,         # NEW Round 5.10 — topic-shift hallucination caught
    slide_client_case_strip_save,          # NEW Round 5.10 — Strip catches fluent fabrication
    slide_client_pitfalls,                 # NEW Round 5.9 — three rules every reviewer learns
    slide_client_halluc_problem,
    slide_client_halluc_real_example,
    slide_client_halluc_caught,
    slide_client_seq_confidence_correlation,
    slide_client_trust_dashboard,
    slide_client_failure_taxonomy_full,     # NEW Round 5.1 (merged 1a+1b) — Round 8: HIDDEN
    # slide_client_failure_worked_example dropped Round 8 — builder kept defined
    slide_client_hallucination_flag,
    slide_client_claims,                    # NEW Round 5.5 (credibility-hardening)
    slide_client_trust_without_ground_truth, # NEW Round 5.1 (deepest critique answer)
    slide_client_trust_operating_points,    # NEW Round 5.15 — permissive/moderate/strict thresholds
    slide_client_what_this_means,

    # Section 11 — Validation
    slide_client_validation_intro,
    slide_client_agreement_chart,
    slide_client_cross_config_stability,
    # slide_client_aggregation_safety DROPPED Round 8.3 — user "not needed"
    # slide_client_validation_summary DROPPED (Round 5.1 — restates 82% point)

    # Section 12 — Engineering (Round 5.2: bulked up)
    slide_client_pipeline_detailed,         # NEW Round 5.2 — 8-stage diagram from science deck
    slide_client_deployment_options,
    # slide_web_ui DROPPED (Round 5.1 — duplicate Live Demo placeholder)

    # Section 13 — What's next (Round 8 reorder: pilot first, then
    # multi-speaker, then options with Arabic first)
    slide_client_try_it_out,                # NEW Round 8 — pilot-first opener
    slide_client_quality_filter,
    slide_client_preprocessing_summary,
    slide_client_arabic_high_level,         # NEW Round 5.1 (replaces slide_arabic_roadmap)
    slide_client_next_milestone,            # NEW Round 5.8 (technical direction behind the ask)
    slide_client_feedback_loop_ask,         # NEW Round 5.13 — soft ask: end-user feedback
    slide_client_partnership_ask,           # NEW Round 5.1 (merged data_ask + investment_ask)

    # Section 14 — Close
    slide_client_recap,
    slide_client_integration_commitment,
    slide_client_next_steps,

    # Section transitions (factory)
    make_section_transition,
)

# Section transitions — only at major pivots, not between every section
_section_real_outputs = make_section_transition(
    "REAL OUTPUTS",
    "What the system actually produces, on real video",
)
_section_trust = make_section_transition(
    "WHY YOU CAN TRUST IT",
    "Per-word, per-segment, validated against an independent reviewer",
)
_section_validation = make_section_transition(
    "VALIDATION",
    "An independent reviewer agreed in 82% of cases",
)
_section_engineering = make_section_transition(
    "ENGINEERING UNDER THE HOOD",
    "What it actually takes to ship this",
)
_section_whats_next = make_section_transition(
    "WHAT'S NEXT",
    "",
)


OUTPUT = MATERIALS / "Argos_VSP_Client_v10_May2026.pptx"


def slide_pipeline_appendix(prs):
    """Round 7 — height-bound wrapper for the 8-stage academic pipeline PNG.

    `slide_17_png` uses `width=CW` which makes the image overflow the
    slide vertically (the source PNG is wide-aspect). Render a height-
    bound copy here for the appendix instead.
    """
    from presentation.helpers import new_slide, add_title, add_accent_line, add_image, add_logo, add_slide_num, set_notes
    from presentation.config import IMG, MX, MY, CT, CW, _auto_num
    from pptx.util import Inches
    slide = new_slide(prs)
    _auto_num[0] += 1
    add_title(slide, "8-Stage Automated Pipeline (appendix)")
    add_accent_line(slide)
    # Height-bound centered image
    img_h = Inches(5.3)
    img_w = Inches(11.5)
    img_x = (SL_W - img_w) // 2
    img_y = Inches(1.55)
    add_image(slide, "pipeline_detailed", img_x, img_y, width=img_w, height=img_h)
    add_logo(slide)
    add_slide_num(slide, _auto_num[0])
    set_notes(slide,
        "Appendix: full pipeline diagram showing 8 stages with ASR as "
        "evaluation-only side-branch. Data flow: .mp4 → normalize → mouth "
        "crop → LRS3 convert → manifests → K-means → LLM decode → reports. "
        "ASR (Whisper) provides reference text for WER/IS scoring only.")
    return slide


def main():
    _auto_num[0] = 0  # reset slide counter

    prs = Presentation()
    prs.slide_width = SL_W
    prs.slide_height = SL_H

    print("Generating client presentation (v10 — color pass + audit-locked fixes)...")

    builders = [
        # ═══════ ACT 1 — IT WORKS (0:00 – 0:30) ═══════
        # § Hello — v9.3: value_prop_headline (S3) DROPPED per user decision (a)
        # — S15 headline_numbers is the single headline slide.
        slide_client_title,
        slide_client_about_argos,
        # § Background — opens with WHAT and HOW so the audience has a
        # mental model before we go into use cases / comparison /
        # difficulty (Round 5.12: pipeline_components moved up from old
        # slide 9 to slide 4 per user "general overview at start" ask).
        slide_client_what_is_vsr,
        slide_client_pipeline_components,       # MOVED UP Round 5.12 — system-overview-first
        slide_client_what_is_lipreading_not,
        slide_client_canonical_scenario,
        # § Compared to today — Round-5 anchor
        slide_client_compared_to_today,
        slide_client_human_ceiling,             # NEW Round 5.1 (why 45-52%)
        # § Why hard — Round 8.8: visemes DROPPED (overlapped with human_ceiling)
        # slide_client_visemes,
        # § What we built — Round 8: multi-speaker moved BACK to §13
        # (What's next) per user feedback that the partnership ask
        # should lead with try-it-out → multi-speaker → options.
        slide_client_what_we_built,
        slide_client_deployment_options,        # MOVED 8.9 — deploy intent up front (cloud/on-prem)
        # § How to use it — Demo
        slide_client_demo_intro,
        slide_client_demo_video_embed,
        # slide_client_demo_recap moved to APPENDIX (Round 8.7 — was thin transition)
        # § Real outputs — Section transition + gallery + cleans + 3 Obama + 6 verdict-tagged
        _section_real_outputs,
        # Round 8.9 (#11+#12): structure is now examples → tiers → more examples
        slide_client_headline_numbers,          # headline numbers up front
        slide_client_example_perfect,           # E1 — Obama Trust
        slide_client_realtalk_spectrum,         # E2 — Conversational Trust/Salvage/Strip spectrum
        # ── tier explanation lands HERE (between early examples and later examples) ──
        _section_trust,
        slide_client_confidence_question,
        slide_client_two_layer_confidence,
        slide_client_word_color_coding,         # blue good / orange mid / purple bad
        slide_client_three_tier_policy,         # tier UI policy
        slide_client_how_to_read,               # decision flow
        # v9.3: slide_client_seq_confidence_correlation CUT per user lock
        # ── back to more examples now that audience knows how to read them ──
        slide_client_clean_outputs_gallery,     # E3 — gallery of six cleans
        slide_client_example_partial,           # E4 — Obama Salvage
        slide_client_judge_ex1,                 # E5 — named entity swap (judge_entity video, restored v9.1)
        slide_client_judge_ex3,                 # E6 — technical vocab drift (judge_router video)
        slide_client_judge_ex6,                 # E7 — topic hijack (judge_lights video)

        # ═══════ ACT 2 — Worked-example deep dives (post-tier-explanation) ═══════
        # Trust section moved up above (between early examples and later examples).
        # These are deeper case studies the audience can now read with the tier mental model.
        slide_client_reader_example,            # E8 — Salvage walk-through
        slide_client_case_topic_shift,          # E9 — Strip topic invented
        slide_client_case_strip_save,           # E10 — Strip fluent fabrication
        slide_client_halluc_caught,             # E11 — Strip hallucination flagged
        # slide_client_trust_dashboard moved to APPENDIX (duplicates headline)
        # slide_client_failure_taxonomy_full moved to APPENDIX (Round 8 dial-back)
        # Round 8: slide_client_failure_worked_example DROPPED (item 17 — not useful)
        # Round 7: slide_failure_deep_2 dropped from active deck.
        # Round 8: slide_25d DROPPED (item 16 — phonetic-bridge "jesus" example too sketchy)
        # slide_client_hallucination_flag moved to APPENDIX (duplicates headline)
        slide_client_claims,                        # NEW Round 5.5 — claims/non-claims credibility-harden
        slide_client_trust_without_ground_truth,    # NEW Round 5.1
        slide_client_trust_operating_points,        # NEW Round 5.15 — permissive/moderate/strict thresholds
        # slide_client_what_this_means moved to APPENDIX (thin transition)

        # § Validation (Round 5.1: validation_summary dropped)
        _section_validation,
        slide_client_validation_intro,
        slide_client_agreement_chart,
        slide_client_cross_config_stability,
        # slide_client_aggregation_safety dropped Round 8.3 (user "not needed")

        # ═══════ ACT 3 — PATH FORWARD (1:25 – 2:00) ═══════
        # § Engineering (Round 5.2: bulked up per user feedback — pipeline
        # diagram from science deck + engineering_journey moved here +
        # data_flow + deployment options. The engineering section now
        # carries real depth instead of a single architecture slide.)
        _section_engineering,
        # slide_client_pipeline_detailed moved to APPENDIX (Round 8 — bulk PNG, unclear)
        slide_data_flow,                        # Borrowed: 5-step model architecture (LoRA stripped)
        slide_client_engineering_journey,       # MOVED Round 5.2 from §What-we-built — deeper rewrite
        # Round 8.9: deployment_options moved to right after What We Built (before Real Outputs)

        # § What's next (Round 8 reorder per user: pilot-first, then
        # multi-speaker, then options with Arabic FIRST among options)
        _section_whats_next,
        slide_client_try_it_out,                # NEW Round 8 — pilot-first opener
        slide_client_multi_speaker_today_vs_path,   # MOVED Round 8 from §6 — engineering shipping
        # Then options:
        slide_client_arabic_high_level,         # FIRST option (Round 8 reorder)
        slide_client_quality_filter,
        slide_client_next_milestone,            # NEW Round 5.8 — what investment buys (stronger LLM + more data)
        slide_client_feedback_loop_ask,         # NEW Round 5.13 — soft client-feedback ask
        slide_client_partnership_ask,           # NEW Round 5.1 (merged ask) — partnership logistics

        # § Close — Round 8.9: next_steps DROPPED (user request — generic placeholder)
        slide_thank_you,                        # Borrowed: close

        # ═══════ APPENDIX (hidden in presentation, physically last in file) ═══════
        # Round 8.7 reorganization: every slide that was previously hidden
        # mid-deck is now collected here at the end of the file. This keeps
        # them available for Q&A reference while making sure they don't
        # appear mid-flow in PowerPoint Normal view.
        slide_client_demo_recap,                # was inline @ slide 13
        slide_client_judge_ex2,                 # was inline @ slide 20
        slide_client_judge_ex4,                 # was inline @ slide 22
        slide_client_judge_ex5,                 # was inline @ slide 23
        slide_client_halluc_problem,            # was inline @ slide 36
        slide_client_halluc_real_example,       # was inline @ slide 37
        slide_client_trust_dashboard,           # was inline @ slide 40
        slide_client_failure_taxonomy_full,     # was inline @ slide 41
        slide_client_hallucination_flag,        # was inline @ slide 42
        slide_client_what_this_means,           # was inline @ slide 46
        slide_client_pipeline_detailed,         # was inline @ slide 52
        slide_client_recap,                     # was inline @ slide 64
        slide_client_video_gallery,             # appendix since Round 7
        slide_pipeline_appendix,                # appendix since Round 7
    ]

    total = len(builders)
    for i, builder in enumerate(builders, 1):
        name = getattr(builder, "__name__", "section_transition")
        print(f"  Slide {i:2d}/{total}  {name} ...", end=" ")
        try:
            builder(prs)
            print("OK")
        except Exception as e:
            print(f"ERROR: {e}")

    # Round 5.12 — Apply slide-level fade transitions to every slide so
    # navigating to a slide gives the perception of motion. Native
    # PowerPoint transition; no per-shape entrance animation work needed.
    from presentation.helpers import add_fade_transition
    for s in prs.slides:
        add_fade_transition(s, speed='med')
    print(f"  FADE TRANSITIONS applied to all {len(prs.slides)} slides")

    # Round 5.11 — Mark redundant / now-redundant slides as hidden so they
    # stay in the .pptx file as backup but PowerPoint skips them during
    # presentation. python-pptx has no direct hide-slide API; set the
    # show="0" attribute on <p:sldId> directly. 1-based indices.
    # Round 8: indices recomputed after dropping slide_client_failure_worked_example
    # (was 43) and slide_25d (was 44). Everything from old slide 45 onward
    # shifts down by 2. Newly added: 42 (failure_taxonomy_full), 52
    # (aggregation_safety). Slide 41 (trust_dashboard) hidden because the
    # "62%" headline is being reframed to "71%" elsewhere — this slide is
    # redundant after that update.
    # Round 8.7 — all hidden slides moved to appendix block at file end.
    # Visible slides 1-55 (where 55 is thank_you); slides 56-69 are appendix.
    # Round 8.9: net -1 visible from 8.8 (dropped next_steps).
    HIDDEN_SLIDES = list(range(51, 65))  # v9.3: -2 from value_prop+seq_corr drops  # 53..66 inclusive — every appendix slide
    if HIDDEN_SLIDES:
        from pptx.oxml.ns import qn
        sldIdLst = prs.slides._sldIdLst
        sld_ids = list(sldIdLst.findall(qn("p:sldId")))
        for n in HIDDEN_SLIDES:
            if 1 <= n <= len(sld_ids):
                sld_ids[n - 1].set("show", "0")
                print(f"  HIDDEN  slide {n:2d} (preserved in file, skipped in presentation)")
            else:
                print(f"  WARNING: hidden-slide index {n} out of range")

    # v9.4 — Q&A backstops on thank_you (last visible slide). Notes only.
    # Audience never sees these; presenter uses them as a deflection map.
    from presentation.helpers import set_notes
    qa_notes = (
        "Q&A BACKSTOPS — what to deflect to:\n"
        "\n"
        "  • If asked about confidence stability across configs → slide 38 (cross-config).\n"
        "  • If asked about per-word calibration breakdown → slide 21 (three-tier policy).\n"
        "  • If asked about hallucination examples → slide 31 (E11), or hidden examples 55-58.\n"
        "  • If asked about failure mode counts → hidden slide 59 (failure modes).\n"
        "  • If asked about IS internals → IS is not on visible slides; deflect to per-segment "
        "confidence on slide 21.\n"
        "  • If asked about WER → 36% on the 1,497-segment baseline. WER is the wrong metric "
        "for lip reading — meaning preservation matters more. That's why we use the 65% useful number.\n"
        "  • If asked about pricing → don't quote a number. Pilot is fixed-fee, named scope. "
        "Production engagement scales with monthly volume. Set up a call with commercial team.\n"
        "  • If asked about other languages → architecture is language-flexible (AV-HuBERT visual "
        "encoder isn't English-locked). Adding a language is a funded engagement: source data + "
        "your reviewer time + 6-10 weeks.\n"
    )
    # Find thank_you slide (last visible). Iterate from end and find one with 'Thank' in title.
    for slide in list(prs.slides)[::-1]:
        for sh in slide.shapes:
            if sh.has_text_frame and "Thank" in sh.text:
                set_notes(slide, qa_notes)
                print("  Q&A BACKSTOPS notes set on thank_you slide")
                break
        else:
            continue
        break

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)} total ({len(prs.slides) - len(HIDDEN_SLIDES)} visible during presentation)")


if __name__ == "__main__":
    main()
