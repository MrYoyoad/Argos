"""
New / merged slide builders for Argos VSP v13 Amosi deck.

These slides have no direct equivalent in the existing slides_*.py modules.
Reverse-engineered from `presentation_materials_20260224/Argos_VSP_v13_Amosi_2.pptx`.
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .config import (
    IMG, VID,
    SL_W, SL_H, BG, WHITE, TEAL, CORAL, LGRAY, MGRAY, DGRAY,
    GREEN, YELLOW, GOLD, ORANGE, RED, DRED, NAVY2, NAVY3,
    BLUE, PURPLE,
    MX, MY, CT, CW, CH, SLW, SRG, SRL, SRW,
)
from .helpers import (
    new_slide, add_title, add_text, add_rect, add_image, add_bullets,
    add_accent_line, add_table, _finish,
)


# ─────────────────────────────────────────────────────────────────────────
# v13 #2 — Thesis
# ─────────────────────────────────────────────────────────────────────────
def slide_thesis(prs):
    """Thesis statement: Argos is a calibrated review tool, not automation."""
    slide = new_slide(prs)
    add_title(slide, "Thesis")
    add_accent_line(slide)

    headline = add_text(slide,
        "Argos is not automatic transcription.",
        MX + Inches(0.5), Inches(2.5), CW - Inches(1.0), Inches(0.6),
        size=Pt(28), color=CORAL, bold=True)

    body = add_text(slide,
        "It is a calibrated review tool for silent video — surfacing meaningful signal "
        "where it exists, marking uncertainty where it doesn't, and routing every segment "
        "through a trust layer so high-confidence output, partial signal, and unreliable "
        "output are visibly distinct.",
        MX + Inches(0.5), Inches(3.1), CW - Inches(1.0), Inches(1.8),
        size=Pt(24), color=WHITE)

    tagline = add_text(slide,
        "Reviewable visual-speech intelligence. Uncertainty attached.",
        MX, Inches(5.95), CW, Inches(0.5),
        size=Pt(22), color=TEAL, italic=True, align=PP_ALIGN.CENTER)

    _finish(slide, 0,
        "Thesis slide. Frames the deck's central claim: Argos is a calibrated review "
        "tool, NOT automatic transcription. The trust layer (per-word confidence "
        "bands + segment-level tiers) is the load-bearing feature that distinguishes "
        "this from an ASR-style system.",
        [[headline, body, tagline]], click_reveal=False)


# ─────────────────────────────────────────────────────────────────────────
# v13 #3 — What we claim — and what we do not claim
# ─────────────────────────────────────────────────────────────────────────
def slide_what_we_claim(prs):
    """Two-column claims / non-claims contrast."""
    slide = new_slide(prs)
    add_title(slide, "What we claim — and what we do not claim")
    add_accent_line(slide)

    sub = add_text(slide,
        "Where the line is. So you know what you're getting.",
        MX, CT, CW, Inches(0.4),
        size=Pt(20), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

    # Two columns
    col_w = (CW - Inches(0.3)) / 2
    cy_top = CT + Inches(0.5)
    row_h = Inches(0.7)
    row_gap = Inches(0.08)

    # Headers
    claim_header = add_rect(slide, MX, cy_top, col_w, Inches(0.5),
                            fill_color=NAVY2, border_color=GREEN, border_width=Pt(2),
                            corner_radius=True)
    add_text(slide, "WE CLAIM", MX, cy_top + Inches(0.08),
             col_w, Inches(0.35), size=Pt(22), color=GREEN, bold=True,
             align=PP_ALIGN.CENTER)

    nope_x = MX + col_w + Inches(0.3)
    nope_header = add_rect(slide, nope_x, cy_top, col_w, Inches(0.5),
                            fill_color=NAVY2, border_color=CORAL, border_width=Pt(2),
                            corner_radius=True)
    add_text(slide, "WE DO NOT CLAIM", nope_x, cy_top + Inches(0.08),
             col_w, Inches(0.35), size=Pt(22), color=CORAL, bold=True,
             align=PP_ALIGN.CENTER)

    claims = [
        "Recover useful speech candidates from video-only input.",
        "Per-word confidence flags numerics, names, and uncertain spans for review.",
        "Many dangerous failure modes flagged before review.",
        "Reduce reviewer workload by routing attention to suspicious spans.",
    ]
    non_claims = [
        "Perfect lip-reading.",
        "Reliably preserve negation, named entities, or numeric values without verification.",
        "Every hallucination is caught.",
        "Replace human judgment for time-critical decisions.",
    ]

    row_y = cy_top + Inches(0.6)
    shapes = []
    for claim, non_claim in zip(claims, non_claims):
        # Left card (claim)
        r1 = add_rect(slide, MX, row_y, col_w, row_h,
                      fill_color=NAVY2, border_color=GREEN, border_width=Pt(1),
                      corner_radius=True)
        add_text(slide, "✓  " + claim,
                 MX + Inches(0.2), row_y + Inches(0.10),
                 col_w - Inches(0.4), row_h - Inches(0.2),
                 size=Pt(18), color=WHITE)
        # Right card (non-claim)
        r2 = add_rect(slide, nope_x, row_y, col_w, row_h,
                      fill_color=NAVY2, border_color=CORAL, border_width=Pt(1),
                      corner_radius=True)
        add_text(slide, "✗  " + non_claim,
                 nope_x + Inches(0.2), row_y + Inches(0.10),
                 col_w - Inches(0.4), row_h - Inches(0.2),
                 size=Pt(18), color=WHITE)
        shapes.extend([r1, r2])
        row_y += row_h + row_gap

    # Bottom callout
    bot_y = row_y + Inches(0.05)
    bot_rect = add_rect(slide, MX, bot_y, CW, Inches(0.5),
                        fill_color=NAVY2, border_color=TEAL, border_width=Pt(2),
                        corner_radius=True)
    add_text(slide,
        "Not blind automation. Reviewable visual-speech intelligence with uncertainty attached.",
        MX, bot_y + Inches(0.08), CW, Inches(0.40),
        size=Pt(20), color=TEAL, italic=True, bold=True, align=PP_ALIGN.CENTER)

    _finish(slide, 0,
        "Claims/anti-claims contrast — sets honest scope for the rest of the deck.",
        [[sub] + shapes], click_reveal=False)


# ─────────────────────────────────────────────────────────────────────────
# v13 #33 — Intelligibility Score: 61.6% Useful Output
# ─────────────────────────────────────────────────────────────────────────
def slide_is_61pct_headline(prs):
    """Headline IS results: 6 signal pills, callout, tier bar chart."""
    slide = new_slide(prs)
    add_title(slide, "Intelligibility Score: 61.6% Useful Output")
    add_accent_line(slide)

    # Row of 6 IS signal pills at top
    signals = [
        ("Semantic\nSim", "(25%)", TEAL),
        ("Phonetic\nSim", "(15%)", TEAL),
        ("Inv.\nWER", "(15%)", TEAL),
        ("Inv.\nWWER", "(15%)", TEAL),
        ("NEA\nF1", "(15%)", CORAL),
        ("Length\nRatio", "(15%)", LGRAY),
    ]
    pill_w = Inches(1.75)
    pill_h = Inches(1.20)
    pill_gap = Inches(0.15)
    total_w = 6 * pill_w + 5 * pill_gap
    pill_start_x = (SL_W - total_w) // 2
    pill_y = CT - Inches(0.05)

    pill_shapes = []
    for i, (name, weight, color) in enumerate(signals):
        x = pill_start_x + i * (pill_w + pill_gap)
        r = add_rect(slide, x, pill_y, pill_w, pill_h,
                     fill_color=NAVY2, border_color=color, border_width=Pt(2),
                     corner_radius=True)
        add_text(slide, name, x, pill_y + Inches(0.18),
                 pill_w, Inches(0.60),
                 size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, weight, x, pill_y + Inches(0.80),
                 pill_w, Inches(0.30),
                 size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
        pill_shapes.append(r)

    # Callout text below pills
    callout_y = pill_y + pill_h + Inches(0.20)
    callout = add_text(slide,
        "IS ≥ 2.00 = Useful Output (Y+P): 61.6% — 2.4× what WER suggests (25.5%)\n"
        "Phonetic similarity: 41.5% mean, r=0.943 with IS (strongest single signal)",
        MX, callout_y, CW, Inches(0.85),
        size=Pt(15), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Horizontal bar chart of tiers
    tiers = [
        ("18.4% Excellent", 18.4, GREEN),
        ("21.4% Good",      21.4, TEAL),
        ("21.7% Fair",      21.7, GOLD),
        ("22.4% Poor",      22.4, ORANGE),
        ("16.0% Failed",    16.0, CORAL),
    ]
    bar_chart_y = callout_y + Inches(1.05)
    bar_h = Inches(0.38)
    bar_gap = Inches(0.12)
    bar_x = MX + Inches(2.50)
    max_bar_w_inches = 7.5
    max_pct = 25.0

    for i, (label, pct, color) in enumerate(tiers):
        y = bar_chart_y + i * (bar_h + bar_gap)
        add_text(slide, label, MX + Inches(0.20), y, Inches(2.10), bar_h,
                 size=Pt(15), color=WHITE, align=PP_ALIGN.RIGHT)
        bw = Inches(max_bar_w_inches * (pct / max_pct))
        add_rect(slide, bar_x + Inches(0.10), y, bw, bar_h,
                 fill_color=color, border_color=None, corner_radius=False)

    _finish(slide, 0,
        "Headline IS result: 61.6% of segments deliver useful output (NIV-Y+P, "
        "IS≥2.00 calibrated against Opus-as-Judge). 6 IS signal pills shown at top, "
        "tier distribution as horizontal bars. WER would call this 25.5% (failure) — "
        "IS reveals 61.6% are useful.",
        [pill_shapes, [callout]], click_reveal=False)


# ─────────────────────────────────────────────────────────────────────────
# v13 #38 — Four Numbers That Tell the Real Story
# ─────────────────────────────────────────────────────────────────────────
def slide_four_numbers(prs):
    """Stack of 4 headline numbers with multi-line descriptions and connecting arrows."""
    slide = new_slide(prs)
    add_title(slide, "Four Numbers That Tell the Real Story")
    add_accent_line(slide)

    rows = [
        ("25.5%", CORAL,
         "By WER (≤34%)\n"
         "381 of 1,497 videos appear useful\n"
         "WER says only 1 in 4 works"),
        ("61.6%", TEAL,
         "By IS (≥2.00): 922 of 1,497 videos\n"
         "deliver useful meaning\n"
         "IS reveals 3 in 5 carry meaning"),
        ("64.9%", GREEN,
         "By LLM Judge (Y+P): 971 of 1,497\n"
         "confirmed useful\n"
         "Expert judgment confirms 2 in 3 work"),
        ("~65%", PURPLE,
         "What users actually see (trust-gate output)\n"
         "65% recall of useful content at FPR 6%\n"
         "630 segments shown with full coloring  (TP 602 · FP 28)"),
    ]

    row_h = Inches(1.10)
    row_gap = Inches(0.20)
    start_y = CT + Inches(0.20)
    arrow_h = Inches(0.20)

    card_w = CW - Inches(1.5)
    card_x = MX + Inches(0.75)

    for i, (num, color, desc) in enumerate(rows):
        y = start_y + i * (row_h + row_gap)
        r = add_rect(slide, card_x, y, card_w, row_h,
                     fill_color=NAVY2, border_color=color, border_width=Pt(2),
                     corner_radius=True)
        add_text(slide, num,
                 card_x + Inches(0.25), y + Inches(0.15),
                 Inches(2.2), Inches(0.85),
                 size=Pt(44), color=color, bold=True, align=PP_ALIGN.LEFT)
        # Highlight the last row's first line in bold
        if i == 3:
            lines = desc.split("\n")
            add_text(slide, lines[0],
                     card_x + Inches(2.55), y + Inches(0.10),
                     card_w - Inches(2.8), Inches(0.35),
                     size=Pt(14), color=WHITE, bold=True)
            add_text(slide, "\n".join(lines[1:]),
                     card_x + Inches(2.55), y + Inches(0.45),
                     card_w - Inches(2.8), Inches(0.65),
                     size=Pt(13), color=LGRAY)
        else:
            add_text(slide, desc,
                     card_x + Inches(2.55), y + Inches(0.15),
                     card_w - Inches(2.8), Inches(0.95),
                     size=Pt(13), color=LGRAY)
        # Down arrow between rows
        if i < 3:
            arrow_y = y + row_h + Inches(0.005)
            arrow_color = rows[i + 1][1]  # arrow takes color of next card
            try:
                add_rect(slide, MX + CW/2 - Inches(0.075), arrow_y,
                         Inches(0.15), Inches(0.15),
                         fill_color=arrow_color, border_color=None,
                         corner_radius=False)
            except Exception:
                pass

    _finish(slide, 0,
        "Four numbers framing the story: 25.5% by WER → 61.6% by IS → "
        "64.9% by LLM Judge → ~65% trust-gate output. WER says fail, two "
        "independent metrics say usable, confidence makes it actionable.",
        [], click_reveal=False)


# ─────────────────────────────────────────────────────────────────────────
# v13 #47 — How IS and Confidence Correlate (merged)
# ─────────────────────────────────────────────────────────────────────────
def slide_is_conf_correlate_combined(prs):
    """v13: How IS and Confidence Correlate - two-column research layout.

    Left column 'Statistical Agreement' has two stacked boxes (scalar
    correlation + tier-level agreement). Right column 'Where They Land
    Together' has a 2x2 contingency table and a Why-This-Holds rationale.
    """
    slide = new_slide(prs)
    add_title(slide, "How IS and Confidence Correlate")
    add_accent_line(slide)

    col_w = Inches(5.85)
    gap_ = Inches(0.43)

    # ---- Left column ----
    L = []
    L.append(add_text(slide, "Statistical Agreement",
             MX, CT, col_w, Inches(0.40),
             size=Pt(22), color=CORAL, bold=True))

    # Box 1: Scalar correlation (teal border)
    b1_y = CT + Inches(0.55)
    b1_h = Inches(2.10)
    L.append(add_rect(slide, MX, b1_y, col_w, b1_h,
                      fill_color=NAVY2, border_color=TEAL,
                      border_width=Pt(2), corner_radius=True))
    L.append(add_text(slide, "Scalar Correlation  (per-segment, n = 1,427)",
             MX + Inches(0.25), b1_y + Inches(0.15),
             col_w - Inches(0.5), Inches(0.40),
             size=Pt(15), color=TEAL, bold=True))
    L.append(add_bullets(slide, [
        ("Pearson r ( IS , <log p> )  ~  0.74", {}),
        ("Spearman rho  ~  0.78  (rank correlation, robust to "
         "saturation at the ends)", {"bold": True}),
    ], MX + Inches(0.25), b1_y + Inches(0.60),
       col_w - Inches(0.5), Inches(1.40), size=Pt(13)))

    # Box 2: Tier-level agreement (green border)
    b2_y = b1_y + b1_h + Inches(0.20)
    b2_h = Inches(2.20)
    L.append(add_rect(slide, MX, b2_y, col_w, b2_h,
                      fill_color=NAVY2, border_color=GREEN,
                      border_width=Pt(2), corner_radius=True))
    L.append(add_text(slide, "Tier-Level Agreement",
             MX + Inches(0.25), b2_y + Inches(0.15),
             col_w - Inches(0.5), Inches(0.40),
             size=Pt(15), color=GREEN, bold=True))
    L.append(add_bullets(slide, [
        ("Cohen's k ( IS-tier , Confidence-tier )  ~  0.62  -  "
         "substantial agreement", {}),
        ("Raw agreement ~ 79% on diagonal.  Off-diagonal cases mostly "
         "adjacent  (Y <-> TRUST off-by-one).",
         {"bold": True}),
    ], MX + Inches(0.25), b2_y + Inches(0.60),
       col_w - Inches(0.5), Inches(1.55), size=Pt(13)))

    # ---- Right column ----
    rx = MX + col_w + gap_
    R = []
    R.append(add_text(slide, "Where They Land Together",
             rx, CT, col_w, Inches(0.40),
             size=Pt(22), color=CORAL, bold=True))
    R.append(add_text(slide,
        "2x2 contingency on 1,497 segments. Both metrics agree on the "
        "\"useful or not\" verdict 96% of the time.",
        rx, CT + Inches(0.45), col_w, Inches(0.75),
        size=Pt(13), color=WHITE, bold=True))

    headers = ["", "Conf trusted", "Conf stripped"]
    rows_data = [
        ["IS useful (Y+P)",   "898",  "24"],
        ["IS not useful (N)",  "30", "545"],
    ]
    row_colors = {
        0: {1: GREEN, 2: WHITE},
        1: {1: WHITE, 2: CORAL},
    }
    tbl = add_table(slide, headers, rows_data,
                    rx, CT + Inches(1.30), col_w,
                    row_height=Inches(0.50),
                    col_widths=[Inches(2.4), Inches(1.7), Inches(1.7)],
                    text_size=Pt(13), row_colors=row_colors)
    R.append(tbl)

    R.append(add_text(slide, "Why This Holds",
             rx, CT + Inches(3.10), col_w, Inches(0.35),
             size=Pt(15), color=BLUE, bold=True))
    R.append(add_text(slide,
        "IS uses the reference text. Confidence uses only the softmax. "
        "Different inputs, same answer.\n"
        "Both metrics independently calibrated to the same blind judge "
        "(k_IS-Opus = 0.818).",
        rx, CT + Inches(3.50), col_w, Inches(1.40),
        size=Pt(12), color=WHITE))
    R.append(add_text(slide,
        "That shared anchor forces them to track each other - "
        "mathematically, not just conceptually.",
        rx, CT + Inches(4.80), col_w, Inches(0.50),
        size=Pt(12), color=LGRAY, italic=True))

    _finish(slide, 0,
        "v13 merge of correlation + contingency. Left: statistical "
        "agreement (Spearman 0.78, kappa 0.62). Right: 2x2 contingency "
        "(96% on useful/not verdict). Source: intelligibility_scores.csv "
        "+ word_confidence.json on n=1,497.",
        [L, R], click_reveal=False)


# ─────────────────────────────────────────────────────────────────────────
# v13 #64 — 10-Stage Pipeline (merge of 17_a + 17_b, rebadged 8→10)
# ─────────────────────────────────────────────────────────────────────────
def slide_pipeline_10stage(prs):
    """10-stage pipeline flow diagram (matches v13 source layout).

    Top row: 1.Normalize -> 2.Mouth Crop ---> 4.LRS3 Convert (cyan)
    Middle: 3.ASR (cyan, eval-only, isolated)
    Bottom row: 5.Manifests -> 6.K-means -> 7.LLM Decode -> 10.Outputs
                (green green gold pink)
    Below LLM Decode: 8.Confidence + 9.Aggregate (purple, trust layer)
    Red dashed box around 6.K-means + 7.LLM Decode (academic-repo origin).
    """
    slide = new_slide(prs)
    add_title(slide, "10-Stage Pipeline: 8 Decode Stages + Trust Layer")
    add_accent_line(slide)

    dark_text = RGBColor(0x0D, 0x1B, 0x2A)

    # Box sizing
    box_w = Inches(1.85)
    box_h = Inches(1.05)
    gap_x = Inches(0.55)
    gap_y = Inches(0.50)

    y_top = CT + Inches(0.05)
    y_mid = y_top + box_h + Inches(0.20)
    y_bot = y_mid + box_h + Inches(0.65)
    y_trust = y_bot + box_h + Inches(0.10)

    start_x = MX + Inches(0.10)

    def stage_box(x, y, num_text, desc, color, w=None, h=None):
        ww = w or box_w
        hh = h or box_h
        r = add_rect(slide, x, y, ww, hh,
                     fill_color=color, border_color=None, corner_radius=True)
        add_text(slide, num_text,
                 x + Inches(0.05), y + Inches(0.10),
                 ww - Inches(0.10), Inches(0.40),
                 size=Pt(17), color=dark_text, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, desc,
                 x + Inches(0.05), y + Inches(0.50),
                 ww - Inches(0.10), hh - Inches(0.55),
                 size=Pt(13), color=dark_text,
                 align=PP_ALIGN.CENTER)
        return r

    # ─── Top row: 1, 2, gap, 4 ───
    x1 = start_x
    x2 = x1 + box_w + gap_x
    x4 = x2 + (box_w + gap_x) * 2
    stage_box(x1, y_top, "1. Normalize",    "HDR/10-bit\nconversion", TEAL)
    stage_box(x2, y_top, "2. Mouth Crop",   "Face detect\n& ROI",     TEAL)
    stage_box(x4, y_top, "4. LRS3 Convert", "Flat → LRS3\nformat",    TEAL)

    # Arrow 1 -> 2
    add_text(slide, "→",
             x1 + box_w - Inches(0.05), y_top + Inches(0.30),
             gap_x + Inches(0.10), Inches(0.45),
             size=Pt(28), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Long horizontal line 2 -> 4
    line_y = y_top + box_h // 2
    add_rect(slide, x2 + box_w, line_y - Inches(0.02),
             x4 - (x2 + box_w), Inches(0.04),
             fill_color=TEAL, border_color=None)

    # ─── ASR (between 2 and 4, lower row) ───
    x3 = x2 + (box_w + gap_x) - Inches(0.10)
    stage_box(x3, y_mid, "3. ASR", "Whisper\ntranscription", TEAL)
    add_text(slide, "evaluation only",
             x3, y_mid + box_h + Inches(0.02),
             box_w, Inches(0.28),
             size=Pt(12), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

    # Drop arrow from main flow (line) down to ASR
    add_text(slide, "↓",
             x3 + box_w // 2 - Inches(0.20), line_y + Inches(0.02),
             Inches(0.40), Inches(0.40),
             size=Pt(20), color=CORAL, bold=True, align=PP_ALIGN.CENTER)

    # Drop arrow from 1 down to 5 (left side)
    add_text(slide, "↓",
             x1 + box_w // 2 - Inches(0.20), y_top + box_h - Inches(0.05),
             Inches(0.40), Inches(0.50),
             size=Pt(20), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # ─── Bottom row: 5, 6, 7, 10 ───
    x5 = x1
    x6 = x5 + box_w + gap_x
    x7 = x6 + box_w + gap_x
    x10 = x7 + box_w + gap_x
    stage_box(x5,  y_bot, "5. Manifests",   "TSV + splits",            GREEN)
    stage_box(x6,  y_bot, "6. K-means",     "Feature\nclustering",     GREEN)
    stage_box(x7,  y_bot, "7. LLM Decode",  "AV-HuBERT +\nLLaMA-2",    GOLD)
    stage_box(x10, y_bot, "10. Outputs",    "Reports &\nburned video", CORAL)

    # Arrow 5 -> 6
    add_text(slide, "→",
             x5 + box_w - Inches(0.05), y_bot + Inches(0.30),
             gap_x + Inches(0.10), Inches(0.45),
             size=Pt(22), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Arrow 7 -> 10
    add_text(slide, "→",
             x7 + box_w - Inches(0.05), y_bot + Inches(0.30),
             gap_x + Inches(0.10), Inches(0.45),
             size=Pt(22), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Red outlined box around 6 + 7 ("Existed in academic repo")
    repo_pad = Inches(0.10)
    repo_x = x6 - repo_pad
    repo_y = y_bot - repo_pad
    repo_w = (x7 + box_w) - (x6 - repo_pad) + repo_pad
    repo_h = box_h + 2 * repo_pad
    add_rect(slide, repo_x, repo_y, repo_w, repo_h,
             fill_color=None, border_color=CORAL,
             border_width=Pt(2.0), corner_radius=True)
    add_text(slide, "Existed in academic repo",
             x6, y_bot - Inches(0.42),
             (x7 + box_w) - x6, Inches(0.28),
             size=Pt(12), color=CORAL, bold=True, align=PP_ALIGN.CENTER)

    # Drop arrow from ASR down to 7
    add_text(slide, "↓",
             x3 + box_w // 2 - Inches(0.20), y_mid + box_h + Inches(0.30),
             Inches(0.40), Inches(0.40),
             size=Pt(16), color=CORAL, bold=True, align=PP_ALIGN.CENTER)

    # ─── Trust layer (8 Confidence, 9 Aggregate) below 7+10 area ───
    trust_box_h = Inches(0.85)
    x8 = x7
    x9 = x10
    add_rect(slide, x8, y_trust, box_w, trust_box_h,
             fill_color=PURPLE, border_color=None, corner_radius=True)
    add_text(slide, "8. Confidence",
             x8 + Inches(0.05), y_trust + Inches(0.05),
             box_w - Inches(0.10), Inches(0.35),
             size=Pt(15), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Per-word + per-segment\ntrust scoring",
             x8 + Inches(0.05), y_trust + Inches(0.40),
             box_w - Inches(0.10), Inches(0.45),
             size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, x9, y_trust, box_w, trust_box_h,
             fill_color=PURPLE, border_color=None, corner_radius=True)
    add_text(slide, "9. Aggregate",
             x9 + Inches(0.05), y_trust + Inches(0.05),
             box_w - Inches(0.10), Inches(0.35),
             size=Pt(15), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Majority's favorite\nfrom 20 hypotheses",
             x9 + Inches(0.05), y_trust + Inches(0.40),
             box_w - Inches(0.10), Inches(0.45),
             size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)

    # Trust-layer caption
    add_text(slide,
             "Trust Layer  —  shipped, runs every output",
             x8, y_trust + trust_box_h + Inches(0.04),
             (x9 + box_w) - x8, Inches(0.30),
             size=Pt(13), color=PURPLE, italic=True, align=PP_ALIGN.CENTER)

    # Legend along the bottom of the slide, well below all stage boxes.
    # Trust caption is at y_trust + 0.85 + 0.04 = y_trust + 0.89; place legend
    # below that and to the LEFT half of slide so it doesn't clash with the
    # Trust Layer boxes on the right.
    legend_y = y_trust + trust_box_h + Inches(0.45)
    legend_x = MX
    legend_items = [
        ("Preprocessing",     TEAL),
        ("Feature Extraction", GREEN),
        ("LLM Inference",     GOLD),
        ("Output",            CORAL),
    ]
    lx = legend_x
    for label, color in legend_items:
        add_rect(slide, lx, legend_y, Inches(0.16), Inches(0.16),
                 fill_color=color, border_color=None)
        add_text(slide, label,
                 lx + Inches(0.22), legend_y - Inches(0.03),
                 Inches(1.50), Inches(0.25),
                 size=Pt(11), color=LGRAY)
        lx += Inches(1.65)

    _finish(slide, 0,
        "Pipeline as 10 stages flow diagram. Top row: preprocessing "
        "(Normalize -> Mouth Crop, with LRS3 conversion at the end). "
        "ASR runs independently as an evaluation reference. Bottom row: "
        "manifests -> k-means -> LLM decode -> outputs. The k-means and "
        "LLM decode stages are highlighted as existing in the academic "
        "repo. Below the decode, the Trust Layer (Confidence + "
        "Aggregate) was added April-May 2026 — this is what distinguishes "
        "Argos from raw research output.",
        [], click_reveal=False)


# ─────────────────────────────────────────────────────────────────────────
# v13 #66 — What it actually took (lift from slides_client)
# ─────────────────────────────────────────────────────────────────────────
def slide_engineering_journey_v13(prs):
    """Lift of slide_client_engineering_journey, reframed for academic audience."""
    # Lazy import to avoid circular issues in modules
    try:
        from .slides_client import slide_client_engineering_journey
        return slide_client_engineering_journey(prs)
    except (ImportError, AttributeError) as e:
        # Fallback stub
        slide = new_slide(prs)
        add_title(slide, "What it actually took — four passes, six months")
        add_accent_line(slide)
        add_text(slide, f"[lift target missing: {e}]",
                 MX, CT, CW, Inches(0.6),
                 size=Pt(20), color=CORAL, italic=True, align=PP_ALIGN.CENTER)
        return slide


# ─────────────────────────────────────────────────────────────────────────
# v13 #79 — Optional add-on: pre-filter low-quality clips (lift from slides_client)
# ─────────────────────────────────────────────────────────────────────────
def slide_quality_filter_v13(prs):
    """Lift of slide_client_quality_filter, reframed for academic audience.

    v13 source omits the two subtitle lines ("Available if you want it..."
    and "Each row = clips remaining...") that the client version uses.
    Strip them after the lift.
    """
    try:
        from .slides_client import slide_client_quality_filter
        result = slide_client_quality_filter(prs)
        # Remove the two subtitle text frames at y ~ 1.5 and 1.9 (16pt italic
        # and 11pt italic) so the layout matches the v13 source PNG.
        slide = prs.slides[-1]
        to_remove = []
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if "Available if you want it" in txt or "Each row = clips remaining" in txt:
                to_remove.append(shape)
        for shape in to_remove:
            shape._element.getparent().remove(shape._element)
        return result
    except (ImportError, AttributeError) as e:
        slide = new_slide(prs)
        add_title(slide, "Optional add-on — pre-filter low-quality clips before decode")
        add_accent_line(slide)
        add_text(slide, f"[lift target missing: {e}]",
                 MX, CT, CW, Inches(0.6),
                 size=Pt(20), color=CORAL, italic=True, align=PP_ALIGN.CENTER)
        return slide
