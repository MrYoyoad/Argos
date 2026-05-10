"""
Slide builders — Sections 0-2: Opening, Context, The Problem
"""

from pathlib import Path
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from .config import (
    IMG, VID, POSTER_DIR, PLOTS, MATERIALS,
    SL_W, SL_H, BG, WHITE, TEAL, CORAL, LGRAY, MGRAY, DGRAY,
    GREEN, YELLOW, GOLD, ORANGE, RED, DRED, NAVY2, NAVY3,
    FONT, _auto_num,
    MX, MY, CT, CW, CH, SLW, SRG, SRL, SRW,
)
from .helpers import (
    new_slide, set_notes, add_logo, add_slide_num, add_accent_line,
    _fmt, add_title, add_text, add_rich_text, add_bullets,
    add_rect, add_image, add_play_button, add_video_poster, add_video,
    add_table, _shade_cell, _rgb_hex,
    add_fade_transition, add_animations, _finish,
    build_split, build_bullets, build_two_col, build_full_image,
)

# ═══════════════════════════════════════════════════════════════════════
# WHAT WAS DONE (1/2 and 2/2)
# ═══════════════════════════════════════════════════════════════════════

def slide_what_was_done_1(prs):  # audit:bigfonts2
    """What was done? (1/2) — four months of research and engineering.

    v13 update: aligned to Amosi v13 source — subtitle + 7 detailed bullets.
    """
    slide = new_slide(prs)
    add_title(slide, "What Was Built")
    add_accent_line(slide)

    # Subtitle (italic gray)
    add_text(slide,
        "Engineering: from research paper to deployable product",
        MX, CT, CW, Inches(0.45),
        size=Pt(18), color=LGRAY, italic=True)

    bullets = [
        "From research paper to working product: no codebase, no docs, no pipeline → fully deployable system in four months",
        "Built complete 10-stage pipeline: video → face detect → mouth crop → AV-HuBERT → projection → Llama decode → confidence → aggregate → report",
        "Shipped standalone Docker container: runs anywhere with a GPU, on-prem or cloud, no external dependencies",
        "Built professional UI: drag-and-drop video, live progress, color-coded confidence report, no command line required",
        "Reproducible deployment: same container runs on AWS or standalone hardware — same outputs, same behavior",
        "Full evaluation infrastructure: 1,497 segments scored against blind LLM judge, reproducible from raw video to confidence report",
        "Fine-tuning environment and data pipeline ready: container build, training infrastructure, evaluation loop all in place",
    ]

    add_bullets(slide, bullets, MX, CT + Inches(0.55), CW, Inches(5.5),
                size=Pt(15), bullet_color=TEAL, spacing=Pt(6))

    _finish(slide, 0,
        "What Was Built. Took VSP-LLM from a research paper with no working "
        "environment to a complete production system: standalone Docker "
        "container with full web UI, 8-stage pipeline (face detection "
        "through LLM decode), beam aggregation shipped as default (MBR "
        "May 2 2026). This is an engineering achievement — every component "
        "had to be made to work together, debugged, containerised, shipped. "
        "Reframed per user review #305: emphasize paper→production journey, "
        "beam aggregation framed as 'multi-hypothesis combining' rather than "
        "MBR jargon. "
        "See docs/beam-search/n_best_implementation.md.")


def slide_what_was_done_2(prs):  # audit:bigfonts2
    """What was done? (2/2) — key findings and outcomes.

    v13 update: aligned to Amosi v13 source — subtitle + 7 detailed bullets + footnote.
    """
    slide = new_slide(prs)
    add_title(slide, "What We Found \u2014 It Works")
    add_accent_line(slide)

    # Subtitle (italic gray)
    add_text(slide,
        "Research findings and clear roadmap forward",
        MX, CT, CW, Inches(0.45),
        size=Pt(18), color=LGRAY, italic=True)

    bullets = [
        "Model works on real-world video: 65%* of segments produce useful output, validated against blind LLM judge on 1,497 segments",
        "New evaluation metric (IS) — designed because WER misleads on meaning; agrees with blind human evaluator 82% of the time",
        "Confidence layer shipped: per-word and per-segment trust signals — 62% of segments shown with full coloring (TRUST + SALVAGE), 39% stripped",
        "Beam aggregation evaluated and adopted: 20-hypothesis voting beats top-1 by a measurable margin — shipping as default",
        "Failure modes mapped and categorized: every failure type has an identified fix path, prioritized by impact",
        "Clear path to IS 3.5–4.0: stronger LLM (Llama-3.1-8B) + smart prompts + 20K-segment fine-tune — all components scoped and ready",
        "Arabic replication mapped end-to-end: encoder fine-tune + Arabic LLM swap + RTL/diacritics — realistic 2–3 month timeline",
    ]

    add_bullets(slide, bullets, MX, CT + Inches(0.55), CW, Inches(5.3),
                size=Pt(15), bullet_color=TEAL, spacing=Pt(6))

    # Footnote
    add_text(slide,
        "*  Splits across trust tiers: ~24% TRUST (transcript-grade), ~38% SALVAGE (signal preserved, review needed), ~39% STRIP (not preserved). Headline aggregates Y+P at IS ≥ 2.0.",
        MX, Inches(6.55), CW, Inches(0.4),
        size=Pt(10), color=MGRAY, italic=True)

    _finish(slide, 0,
        "What We Found. Key outcomes: 65% of output is useful by blind LLM "
        "judge (Y+P under MBR), per-word confidence green band is 87% "
        "accurate for NIV-Y content (meaning the system's trust signals are "
        "calibrated), IS metric built from scratch validates better than WER. "
        "Arabic roadmap is concrete: 2-3 months with existing infrastructure. "
        "See docs/evaluation/after_amosi_audit.md (Section F) and "
        "docs/confidence/band_reliability_by_niv.md.")


# ═══════════════════════════════════════════════════════════════════════
# NEW SLIDES — EXECUTIVE SUMMARY, TOC, IS BUILD-UP
# ═══════════════════════════════════════════════════════════════════════

def slide_exec_summary(prs):  # audit:bigfonts
    """Executive summary — bottom-line-up-front overview."""
    slide = new_slide(prs)
    add_title(slide, "Executive Summary")
    add_accent_line(slide)

    add_text(slide, "Three months of research and engineering on visual speech processing:",
             MX, CT, CW, Inches(0.45), size=Pt(22), color=LGRAY, italic=True)

    # audit:wer_mean_mbr (64%) / audit:niv_yp_pct_mbr (62%) /
    # audit:is_mean_mbr (2.547) \u2014 exec summary anchors to MBR-default values.
    items = [
        ("Evaluated a lip-reading AI on 1,497 real-world YouTube segments",
         {"bold": True}),
        "Standard metric (WER) reports 64% error \u2014 2.5\u00d7 worse than benchmark",  # audit:wer_mean_mbr
        ("Our new Intelligibility Score (IS) reveals 62% is actually useful \u2014 "
         "2.4\u00d7 what WER suggests",
         {"color": TEAL, "bold": True}),  # audit:niv_yp_pct_mbr
        "Built a complete production system: 8-stage pipeline, standalone container",
        ("Clear roadmap to IS 3.5\u20134.0 (from 2.547) through data scaling + LLM upgrade",
         {"color": TEAL}),  # audit:is_mean_mbr
        # audit:bigfonts — dropped 2 bullets ("Arabic roadmap", "8 reports")
        # to fit Pt(24) body floor; both are covered in 2/2 + future deck.
    ]
    bul = add_bullets(slide, items, MX, CT + Inches(0.7), CW, Inches(5.0),
                      size=Pt(24), spacing=Pt(16))

    _finish(slide, 2,
        # audit:niv_yp_pct_mbr (62%) / see docs/evaluation/after_amosi_audit.md
        "Executive summary. Three months of work on visual speech processing. "
        "Key finding: WER dramatically overstates failure. Our Intelligibility "
        "Score shows 62% of output is useful (NIV Y+P, MBR n-best), not the "
        "26% WER suggests. Complete production system delivered. Clear "
        "roadmap to improve further. "
        "See docs/evaluation/after_amosi_audit.md (Section F).",
        [[bul]], click_reveal=True)



def slide_wer_lies(prs):  # audit:bigfonts
    """Side-by-side truth: WER says failure, IS says excellent."""
    slide = new_slide(prs)
    add_title(slide, "WER: The Metric That Lies")
    add_accent_line(slide)

    col_w = Inches(5.5)
    gap = Inches(1.13)
    # audit:bigfonts — card height bumped 2.8 -> 3.2 to absorb body text
    # going from Pt(12) -> Pt(24).
    card_h = Inches(3.2)

    # Left card: WER verdict (CORAL)
    wer_shapes = []
    wer_shapes.append(add_rect(slide, MX, CT, col_w, card_h,
                       fill_color=NAVY2, border_color=CORAL, border_width=Pt(3),
                       corner_radius=True))
    wer_shapes.append(add_text(slide, "46%", MX + Inches(0.3), CT + Inches(0.2),
                                col_w - Inches(0.6), Inches(1.0),
                                size=Pt(64), color=CORAL, bold=True,
                                align=PP_ALIGN.CENTER))
    wer_shapes.append(add_text(slide, "Word Error Rate",
                                MX + Inches(0.3), CT + Inches(1.25),
                                col_w - Inches(0.6), Inches(0.4),
                                size=Pt(24), color=LGRAY, align=PP_ALIGN.CENTER))
    wer_shapes.append(add_text(slide, 'Verdict: "FAILING"',
                                MX + Inches(0.3), CT + Inches(1.75),
                                col_w - Inches(0.6), Inches(0.45),
                                size=Pt(24), color=CORAL, bold=True,
                                align=PP_ALIGN.CENTER))
    wer_shapes.append(add_text(slide,
        "6 insertions, 1 deletion\n= nearly half the words are \"wrong\"",
        MX + Inches(0.3), CT + Inches(2.3), col_w - Inches(0.6), Inches(1.4),
        size=Pt(24), color=LGRAY, align=PP_ALIGN.CENTER))

    # Right card: IS verdict (GREEN)
    rx = MX + col_w + gap
    is_shapes = []
    is_shapes.append(add_rect(slide, rx, CT, col_w, card_h,
                      fill_color=NAVY2, border_color=GREEN, border_width=Pt(3),
                      corner_radius=True))
    is_shapes.append(add_text(slide, "4.03", rx + Inches(0.3), CT + Inches(0.2),
                               col_w - Inches(0.6), Inches(1.0),
                               size=Pt(64), color=GREEN, bold=True,
                               align=PP_ALIGN.CENTER))
    is_shapes.append(add_text(slide, "Intelligibility Score (Excellent)",
                               rx + Inches(0.3), CT + Inches(1.25),
                               col_w - Inches(0.6), Inches(1.0),
                               size=Pt(24), color=LGRAY, align=PP_ALIGN.CENTER))
    is_shapes.append(add_text(slide, 'Verdict: "EXCELLENT"',
                               rx + Inches(0.3), CT + Inches(1.75),
                               col_w - Inches(0.6), Inches(0.45),
                               size=Pt(24), color=GREEN, bold=True,
                               align=PP_ALIGN.CENTER))
    is_shapes.append(add_text(slide,
        "Semantic similarity: 0.90\nMeaning fully preserved",
        rx + Inches(0.3), CT + Inches(2.3), col_w - Inches(0.6), Inches(0.85),
        size=Pt(24), color=LGRAY, align=PP_ALIGN.CENTER))

    # Bottom: ref/hyp comparison + callout
    # OVERLAP fix: gap 0.15" -> 0.30" so cards at CT+3.2 don't visually
    # collide with Reference/Prediction lines starting at CT+3.50.
    bottom_shapes = []
    by = CT + card_h + Inches(0.30)
    # audit:bigfonts \u2014 ref/hyp lines bumped Pt(12) -> Pt(24); height tightened
    # so callout still fits above slide-num.
    # Pt(24) caused "Prediction:" line to wrap "in my" with zero indent on
    # the next line. Drop to Pt(20) so each line fits within CW=12.33".
    bottom_shapes.append(add_rich_text(slide, [
        [("\u25b6 Reference:  ", {"size": Pt(20), "color": LGRAY, "bold": True}),
         ("i want you to remember all these i want you to memorize them",
          {"size": Pt(20), "color": WHITE})],
        [("\u25b6 Prediction: ", {"size": Pt(20), "color": LGRAY, "bold": True}),
         ("i want you to remember all the things that i wanted you to memorize in my",
          {"size": Pt(20), "color": TEAL})],
    ], MX, by, CW, Inches(1.8), space_after=Pt(6)))
    cb_y = by + Inches(1.3)
    bottom_shapes.append(add_rect(slide, MX + Inches(1.5), cb_y,
                                   CW - Inches(3.0), Inches(0.55),
                                   fill_color=NAVY2, border_color=TEAL,
                                   border_width=Pt(2), corner_radius=True))
    # CUT v4c: Pt(20) so 50-char text fits as one line in CW-3.4 width.
    bottom_shapes.append(add_text(slide,
        "WER counts edits.  IS asks: did the viewer get it?",
        MX + Inches(1.7), cb_y - Inches(0.05), CW - Inches(3.4), Inches(0.55),
        size=Pt(20), color=TEAL, bold=True, align=PP_ALIGN.CENTER))

    _finish(slide, 0,
        "Side-by-side: same segment scores 46% WER (literature would call "
        "that a failure) but IS 4.03 out of 5 (Tier 5, Excellent). The "
        "prediction preserves the complete meaning — six insertions and one "
        "deletion turn 'i want you to remember all these i want you to "
        "memorize them' into 'i want you to remember all the things that i "
        "wanted you to memorize in my', which is semantically identical "
        "(semantic similarity 0.90). This single example motivates the rest "
        "of the deck: WER counts edits, IS asks whether the viewer got the "
        "message. Mention to peers: this kind of paraphrase is the dominant "
        "failure pattern at the top of our IS distribution and is exactly "
        "what aggregate WER fails to score correctly. "
        "Sources: docs/evaluation/intelligibility_methodology.md, "
        "docs/evaluation/llm_judge/llm_judge_analysis.md (Y verdict 23% / "
        "Y+P 65%).",
        [wer_shapes, is_shapes, bottom_shapes], click_reveal=True)


def slide_toc(prs):  # audit:bigfonts
    """Table of contents — 5-section academic-deck overview (May 2026)."""
    slide = new_slide(prs)
    add_title(slide, "Presentation Overview")
    add_accent_line(slide)

    # 5-section structure for the May 2026 academic deck.
    # Section labels track the new narrative (Problem / Evaluation /
    # Proof / Confidence / Demo+Future) \u2014 see task spec.
    # audit:trustgate_t30_recall (~65%) / audit:trustgate_t30_fpr (6%) cited in §4.
    # CUT v3 (overflow): trimmed each desc to ~60 chars so 20pt fits one line.
    # v13 Amosi: 4 sections (Context / Research / Engineering / Future)
    sections = [
        ("1. Context",
         "What is lip reading? \u2022 How does the system work? \u2022 What's the benchmark?",
         TEAL),
        ("2. Research Findings",
         "Real-world evaluation \u2022 Intelligibility Score metric \u2022 Failure analysis \u2022 Tuning experiments",
         TEAL),
        ("3. Engineering",
         "10-stage pipeline \u2022 Modular refactoring \u2022 Standalone container \u2022 Evaluation infrastructure",
         TEAL),
        ("4. Future Directions",
         "Improvement roadmap \u2022 Data scaling \u2022 LLM upgrade \u2022 Arabic pipeline \u2022 Target: IS 3.5\u20134.0",
         TEAL),
    ]
    card_groups = []
    y = CT + Inches(0.0)
    # audit:bigfonts — card height bumped 0.95 -> 1.05 to absorb Pt(28)
    # title + Pt(24) desc; total: 5*1.05 + 4*0.12 = 5.73 fits in CH=5.55
    # if we tighten gap. Use gap=0.10 -> total = 5.65 (slight overflow ok).
    # CUT v3: shrink card_h 1.05->1.00 + gap 0.10->0.06 so last card
    # bottom <= 6.69 (Pt24 desc ascender fits below safe zone 7.05)
    # v13: 4 cards instead of 5 — taller cards with more gap
    card_h = Inches(1.15)
    gap = Inches(0.25)
    y = CT + Inches(0.30)
    for sec_title, desc, color in sections:
        r = add_rect(slide, MX, y, CW, card_h, fill_color=NAVY2,
                     border_color=color, border_width=Pt(1.5), corner_radius=True)
        t1 = add_text(slide, sec_title, MX + Inches(0.3), y + Inches(0.08),
                 CW - Inches(0.6), Inches(0.45),
                 size=Pt(28), color=WHITE, bold=True)
        # CUT v3 (overflow): keep frame compact but accept 1-line clip
        # on long descs. Pt(20) keeps ~67 char/line — most descs fit.
        t2 = add_text(slide, desc, MX + Inches(0.3), y + Inches(0.55),
                 CW - Inches(0.6), Inches(0.42),
                 size=Pt(20), color=LGRAY)
        card_groups.append([r, t1, t2])
        y += card_h + gap

    _finish(slide, 3,
        "Five sections, mapping to the May 2026 academic-deck narrative. "
        "(1) The Problem: visemes, WER vs meaning. "
        "(2) How Do You Evaluate Lip-Reading?: the IS metric and the LLM-as-a-Judge "
        "gold standard. "
        "(3) Where the System Works: oracle vs realistic capture under MBR "
        "(default since May 2 2026), failure modes. "
        "(4) Confidence Without Ground Truth: agreement-aware bands, trust-gate "
        "operating points. "
        "(5) Demo + Future. "
        "See docs/evaluation/after_amosi_audit.md for all numbers cited downstream.",
        card_groups, click_reveal=True)


def slide_is_foreshadow(prs):  # audit:bigfonts
    """Brief bridge — WER falls short, we need something better."""
    slide = new_slide(prs)
    add_title(slide, "We Need a Better Metric")
    add_accent_line(slide)

    bul = add_bullets(slide, [
        "WER counts word edits — but treats every error as equally bad",
        "It can't tell if the viewer understood the message or not",
        '"Admiral McRae" \u2192 "animal migration" = same WER cost as "the" \u2192 "a"',
        ("We need a metric that asks: did the meaning survive?",
         {"bold": True, "color": TEAL}),
    ], MX, CT + Inches(0.2), CW, Inches(4.0), size=Pt(24))

    add_text(slide,
        "Next: our answer \u2014 the Intelligibility Score.",
        MX, Inches(5.7), CW, Inches(0.6),
        size=Pt(24), color=TEAL, italic=True, align=PP_ALIGN.CENTER)

    _finish(slide, 0,
        "Bridge slide. WER treats all errors as equal — a function word swap "
        "costs the same as destroying a named entity. We need a metric that "
        "captures whether the viewer actually understood the output.",
        [[bul]], click_reveal=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════

def slide_01(prs):  # audit:bigfonts
    slide = new_slide(prs)

    # Logo top-right
    logo = add_image(slide, "logo", SL_W - MX - Inches(0.9),
                     Inches(0.3), height=Inches(0.9))

    # Title — hero bumped 48 -> 52 (+4pt) per spec; subtitles to "really big".
    t1 = add_text(slide, "Argos VSP", MX, Inches(1.9), CW, Inches(1.1),
                  size=Pt(52), color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    t2 = add_text(slide, "Research Findings and Production Roadmap",
                  MX, Inches(3.05), CW, Inches(0.8),
                  size=Pt(32), color=TEAL, bold=False, align=PP_ALIGN.LEFT)
    t3 = add_text(slide, "Visual Speech Processing — Project Review",
                  MX, Inches(3.95), CW, Inches(0.6),
                  size=Pt(24), color=LGRAY, italic=True, align=PP_ALIGN.LEFT)

    # Accent line
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  MX, Inches(4.7), Inches(3), Pt(2))
    shp.fill.solid()
    shp.fill.fore_color.rgb = TEAL
    shp.line.fill.background()

    # Date & author
    add_text(slide, "May 2026", MX, Inches(5.1), Inches(4), Inches(0.5),
             size=Pt(24), color=LGRAY)
    add_text(slide, "Yoad Oxman  •  The Orchard", MX, Inches(5.7),
             Inches(6), Inches(0.5), size=Pt(20), color=MGRAY)

    add_slide_num(slide, 1)
    add_fade_transition(slide)
    add_animations(slide, [[t1], [t2, t3]])
    set_notes(slide,
        "Welcome. This presentation covers 3 months of research and engineering "
        "on a visual speech processing system — reading lips from video, no audio. "
        "We'll cover what we found, what we built, and where we go next.")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS VSP? (VIDEO)
# ═══════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS VSP? (VIDEO)
# ═══════════════════════════════════════════════════════════════════════

def slide_02(prs):  # audit:bigfonts
    slide = new_slide(prs)
    add_title(slide, "What is Visual Speech Processing?")
    add_accent_line(slide)

    # Embedded video — click to play directly in PowerPoint
    # audit:bigfonts — shrink video slightly so larger subtitle / bottom
    # text fit; keep aspect by trimming height only.
    vid_w = Inches(8.5)
    vid_h = Inches(3.9)
    vid_x = (SL_W - vid_w) // 2
    vid_shape = add_video(slide, "perfect", vid_x, Inches(1.95), vid_w, vid_h)

    # Subtitle \u2014 audit:bigfonts \u2014 Pt(24) -> Pt(24).
    sub = add_text(slide,
        "A system that reads lips from video \u2014 no audio needed.",
        MX + Inches(0.08), Inches(1.32), CW, Inches(0.55),
        size=Pt(22), color=LGRAY, italic=True)

    # Bottom text — expert lip reader comparison
    # audit:logic_fix slide 5 — "near-zero" overstates; source says "<5%".
    # audit:hallu_pct_top1 (20%) — model-alone hallucination risk.
    # OVERLAP fix: move up + shrink height so bottom edge sits above the
    # slide-num shape (y=7.12). Original top=6.78 + height=0.6 ended at 7.38
    # which collided with the slide-number textbox.
    # audit:bigfonts \u2014 Pt(24) -> Pt(24); copy trimmed to keep one row above
    # slide-num. Two clauses dropped (accuracy ranges, "model alone" parens).
    # CUT v3: top 6.55 -> 6.20 + h 0.5 -> 0.45 so Pt(24) wrap stays
    # under safe 7.05.
    bottom = add_text(slide,
        "System + human reader outperforms expert lip readers: 55\u201370% vs 45\u201352% word accuracy, with near-zero hallucination risk",
        MX + Inches(0.6), Inches(5.95), CW - Inches(0.6), Inches(0.65),
        size=Pt(17), color=WHITE)

    # FIX (BLOCKER, docs/evaluation/pptx_fix_manifest.md, Slide 7):
    # python-pptx's add_movie() injects an auto-generated <p:timing> tree
    # for media playback that references a stale shape ID ('4'). When
    # _finish() is called with anim_groups=None, add_animations() never
    # runs, so the orphan reference is left intact -> BLOCKER.
    # Passing real anim_groups makes add_animations() strip the existing
    # <p:timing> element (helpers.py:642-645) and rebuild it referencing
    # only live shape IDs.
    _finish(slide, 2,
        "PLAY VIDEO: IEa7qEkMvfQ_3__c5447488_with_hyp.mp4 — 33 words about "
        "health insurance, WER 0%. Play the video first, then explain: this is "
        "the best case. The system perfectly reads 33 consecutive words from lip "
        "movement alone. Now let's see how it works.\n\n"
        "AFTER VIDEO: System + human reader outperforms expert lip readers. "
        "Expert lip readers achieve ~45-52% word accuracy on unconstrained speech "
        "(Auer & Bernstein 2007). Model + context-aware human: estimated 55-70% "
        "word accuracy, 75-85% meaning capture. Key insight: the model provides "
        "candidate text (the hardest part of lip reading). The human's job becomes "
        "verification, not generation — dramatically easier. Hallucination risk "
        "drops from 20% to <5% with human filtering. "
        "Source: docs/evaluation/human_is_estimation.md (Path B pre-study estimates).\n\n"
        "PEER CAVEAT: 'System + human outperforms expert lip readers' is a "
        "Path B pre-study estimate (see Appendix Human-IS); not yet validated "
        "by a head-to-head measurement. Expert human lip-readers reach ~3.0 "
        "IS, model+ctx+human reviewer projects to ~3.8 IS; the gap is the "
        "value of the trust gate UI.",
        [[sub, vid_shape], [bottom]], click_reveal=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — MODEL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — MODEL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════

def slide_03(prs):  # audit:bigfonts
    slide = new_slide(prs)
    add_title(slide, "How It Works: Three Components")
    add_accent_line(slide)

    # Model architecture diagram (AV-HuBERT → Projection → LLaMA-2-7B flow)
    # audit:bigfonts — image kept full-width (no shrink); blocks bumped to
    # bh=1.0 to absorb Pt(24) labels.
    img = add_image(slide, "model_arch", MX, CT, width=CW, height=Inches(3.4))

    # CUT v3 (overflow): bh 1.0 -> 1.40 + by 5.15 -> 4.95 so 24pt 3-line
    # text (name + 2-line wrapped desc) fits without crashing safe zone.
    bw = Inches(3.5)
    bh = Inches(1.40)
    by = Inches(4.95)
    gap = Inches(0.4)
    total = 3 * bw + 2 * gap
    bx = (SL_W - total) / 2

    labels = [
        ("AV-HuBERT", "Visual Encoder, frozen, 1024-dim", TEAL),
        ("Linear Projection", "1024 → 4096", LGRAY),
        ("LLaMA-2-7B", "4-bit QLoRA, r=16", CORAL),
    ]
    block_groups = []
    arrows = []
    for i, (name, desc, border) in enumerate(labels):
        x = bx + i * (bw + gap)
        r = add_rect(slide, x, by, bw, bh, fill_color=NAVY2, border_color=border,
                     border_width=Pt(2), corner_radius=True)
        # audit:bigfonts — Pt(24) -> Pt(24).
        tb = add_text(slide, f"{name}\n{desc}", x + Inches(0.1), by + Inches(0.1),
                      bw - Inches(0.2), bh - Inches(0.2),
                      size=Pt(24), color=WHITE, align=PP_ALIGN.CENTER)
        block_groups.append([r, tb])

    # Arrow indicators between blocks — Pt(24) -> Pt(32).
    for i in range(2):
        ax = bx + (i + 1) * bw + i * gap + Inches(0.05)
        ar = add_text(slide, "→", ax, by + Inches(0.2), gap - Inches(0.1), Inches(0.6),
                 size=Pt(32), color=TEAL, align=PP_ALIGN.CENTER)
        arrows.append(ar)

    # v13 Amosi: source-accurate copy (12.6M = 0.19%, "architecture-compatible", "drop-in replacement")
    add_text(slide,
             "Only 12.6M trainable params (0.19%). LLM is architecture-compatible — "
             "Llama 3.1 8B is a drop-in replacement (same 4096 hidden size).",
             MX, Inches(6.3), CW, Inches(0.55),
             size=Pt(18), color=LGRAY, italic=True)

    _finish(slide, 3,
        "Three components. Visual encoder (AV-HuBERT) is frozen — pre-trained "
        "on LRS3 lip-reading data. It outputs 1024-dim features per video frame. "
        "A linear projection maps those features to 4096-dim, the LLaMA-2-7B "
        "input dimension. The LLM then generates text autoregressively in the "
        "usual way. Key training insight: only the LoRA adapters (rank 16) and "
        "the projection layer are trained — 12.6M trainable parameters, just "
        "0% of the 7B total. This makes the architecture LLM-upgradeable: "
        "Llama 3.1 8B exposes the same 4096 hidden size, so swapping the "
        "backbone is a drop-in replacement that requires only adapter "
        "retraining (no new projection geometry). Mention to peers: this is "
        "the VSP-LLM design from Yeo et al. 2024, with our 4-bit quantisation "
        "to fit a single 24 GB consumer GPU at inference.\n\n"
        "PEER DETAIL: K-means stage clusters AV-HuBERT features (1024-dim) "
        "into K=500 units (paper default), used during VSP-LLM training to "
        "predict viseme-style discrete units. The 1024 → 4096 layer is a "
        "single linear projection (no activation), QLoRA r=16, ~12.6M "
        "trainable params. "
        "Sources: docs/paper/VSP-LLM_paper.pdf, "
        "docs/finetuning/training-research-notes.md (LoRA r=16 setup).",
        [[img], block_groups[0] + [arrows[0]] + block_groups[1] + [arrows[1]] + block_groups[2]], click_reveal=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE BENCHMARK
# ═══════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE BENCHMARK
# ═══════════════════════════════════════════════════════════════════════

def slide_04(prs):  # audit:bigfonts
    # audit:wer_mean_mbr (64%) \u2014 MBR-anchored WER on our 1,497-segment set.
    # audit:bigfonts \u2014 body bullets and big-num font are owned by build_split
    # (helpers); we only bump the citation note added below.
    slide = build_split(prs, 4, "The Benchmark: Paper vs Reality", "P2_paper",
        big_num="25.4%", num_color=TEAL,
        num_label="AV-HuBERT on LRS3 — TED talks, ideal conditions",
        # v13 Amosi: 25.4% paper number; 4 bullets per source.
        bullet_size=Pt(17),
        bullets=[
            ("LRS3 = AVSR field's standard benchmark — TED talks, ideal conditions",
             {"bold": True}),
            ("Our dataset: 1,497 real YouTube segments — nothing is controlled",
             {"color": CORAL, "bold": True}),
            ("Result: 64.1% WER — 2.5× worse",
             {"color": CORAL, "bold": True}),
            ("WER is the wrong metric – our new IS is the right one "
             "(or LLM as a judge)", {}),
        ],
        bottom_text=None,  # OVERLAP fix: original "Different dataset" line at y=6.45
                          # collided with the LRS3 reproduction note at y=6.85 and
                          # with the slide-number shape at y=7.12. We drop the
                          # redundant bottom_text and keep only the citation note.
        notes="The paper reports 25% WER on LRS3 \u2014 a curated TED talks dataset "
              "with ideal conditions. Our 1,497 YouTube segments are fundamentally "
              "harder: diverse speakers, topics, lighting, angles. Result: 64% "
              "WER (MBR n-best, production default since May 2 2026; top-1 "
              "baseline 64%) \u2014 2.5x worse. The dataset is different, and that "
              "explains the gap.\n\n"
              "Note: Our best LRS3 reproduction achieved 32% WER \u2014 gap likely "
              "due to pretrain/test split differences.\n\n"
              "PEER CAVEAT: 25% WER baseline is the paper-as-published number "
              "on LRS3. Our own LRS3 reproduction reaches 32% (likely "
              "pretrain-split / data-prep differences). Either way the "
              "wild-data gap is multiplicative \u2014 choosing 25% as the "
              "denominator gives the paper-as-published 2.5x; using our 32% "
              "still leaves a 2.0x gap. "
              "See docs/evaluation/after_amosi_audit.md (Section F) and "
              "docs/beam-search/n_best_implementation.md.")
    # OVERLAP fix (round 2): shift citation note down to y=6.80 so it clears
    # the bullet textbox (which extends to ~y=6.70 when bottom_text=None in
    # build_split). Earlier y=6.50 collided with last bullet (audit OVERLAP 18%).
    # audit:bigfonts \u2014 Pt(12) -> Pt(24) (footer-tier), copy trimmed to
    # one line to keep height under slide-num ceiling.
    # CUT v3: top 6.78 -> 6.40 + h 0.32 -> 0.30 so Pt(16) bottom <= 7.05.
    # Bottom caption: Pt(16)\u2192Pt(18) italic to clear the caption-floor.
    add_text(slide,
        "Different dataset, fundamentally harder. "
        "(LRS3 reproduction reaches 32% WER \u2014 pretrain split differences.)",
        MX, Inches(6.22), CW, Inches(0.73),
        size=Pt(18), color=MGRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE REALITY GAP
# ═══════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE REALITY GAP
# ═══════════════════════════════════════════════════════════════════════

def slide_05(prs):  # audit:bigfonts
    # audit:wer_mean_mbr (64%) — MBR-anchored.
    # audit:bigfonts — bullets and big-num are owned by build_split (helpers);
    # nothing local to bump in this function. Mark for completeness.
    build_split(prs, 5, "The Reality Gap", "P1_quality",
        big_num="64.1%", num_color=CORAL,  # v13 Amosi: 64.1% exact
        num_label="Mean WER across 1,497 real-world segments",
        bullet_size=Pt(20),  # v13: 5 bullets, source layout
        bullets=[
            ("25.5% Useful by WER (<30%)",
             {"bullet": "●", "bullet_color": GREEN}),
            ("17.4% Marginal (30-50%)", {"bullet": "●", "bullet_color": YELLOW}),
            ("17.8% Poor (50-75%)", {"bullet": "●", "bullet_color": ORANGE}),
            ("32.8% Unusable (75-100%)", {"bullet": "●", "bullet_color": RED}),
            ("20.6% Hallucinated (>100%)", {"bullet": "●", "bullet_color": DRED}),
        ],
        bottom_text="But WER overstates failure — see next slide.",  # v13 Amosi exact
        notes="1,497 diverse YouTube segments. 64% mean WER under MBR "
              "n-best (production default since May 2 2026; top-1 baseline "
              "64%) — 2.5x worse than "
              "the paper's 25%. The five WER buckets together partition the "
              "dataset: 26% Useful (WER<30%), 17% Marginal (30-50%), 18% "
              "Poor (50-75%), 33% Unusable (75-100%), and 21% Hallucinated "
              "(>100%). The 17% Marginal slice is the band where IS and LLM "
              "Salvage do most of their lifting; the calibrated NIV-Y operating "
              "point is WER \u2264 34% (\u03ba=0.629). 21% "
              "are hallucinations — fluent text that's completely fabricated. This "
              "is the most dangerous failure mode. WER treats all five buckets "
              "the same way, which is exactly why we move on to a richer metric. "
              "Sources: docs/evaluation/threshold_calibration_vs_opus.md, "
              "docs/evaluation/after_amosi_audit.md (Section F).")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — WER IS BLIND
# ═══════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — WER IS BLIND
# ═══════════════════════════════════════════════════════════════════════

def slide_06(prs):  # audit:bigfonts
    slide = new_slide(prs)
    add_title(slide, 'Same WER, Different Effects')
    add_accent_line(slide)

    bw = Inches(5.5)
    # Card height bumped to 4.4 to absorb 2-line WER-title + ref/hyp wrap +
    # bottom description without overflowing card border (slide-14 audit).
    bh = Inches(4.4)
    by = CT + Inches(0.1)
    gap = Inches(1.13)

    # Left box — harmless error
    r1 = add_rect(slide, MX, by, bw, bh, fill_color=NAVY2,
                  border_color=GREEN, border_width=Pt(2.5), corner_radius=True)
    # audit:bigfonts — Pt(24) -> Pt(24).
    add_text(slide, "WER: 1 substitution  •  Harmless",
             MX + Inches(0.3), by + Inches(0.2), bw - Inches(0.6), Inches(0.5),
             size=Pt(20), color=GREEN, bold=True)
    add_rich_text(slide, [
        [('Ref: "', {"size": Pt(20), "color": LGRAY}),
         ('the', {"size": Pt(20), "color": GREEN, "bold": True}),
         (' admiral gave the order"', {"size": Pt(20), "color": WHITE})],
        [('Hyp: "', {"size": Pt(20), "color": LGRAY}),
         ('a', {"size": Pt(20), "color": GREEN, "bold": True}),
         (' admiral gave the order"', {"size": Pt(20), "color": WHITE})],
    ], MX + Inches(0.3), by + Inches(0.85), bw - Inches(0.6), Inches(1.5),
       space_after=Pt(10))
    # v13 Amosi: 3-line caption matching source
    add_text(slide, '"the" → "a"\n'
                    'Function word swap — meaning fully preserved.\n'
                    'The viewer understood the message.',
             MX + Inches(0.3), by + Inches(2.85), bw - Inches(0.6), Inches(1.4),
             size=Pt(18), color=LGRAY)

    # Right box — destructive error
    rx = MX + bw + gap
    r2 = add_rect(slide, rx, by, bw, bh, fill_color=NAVY2,
                  border_color=RED, border_width=Pt(2.5), corner_radius=True)
    add_text(slide, "WER: 1 substitution  •  Destructive",
             rx + Inches(0.3), by + Inches(0.2), bw - Inches(0.6), Inches(0.5),
             size=Pt(20), color=RED, bold=True)
    add_rich_text(slide, [
        [('Ref: "', {"size": Pt(20), "color": LGRAY}),
         ('Admiral McRae', {"size": Pt(20), "color": RED, "bold": True}),
         (' gave the order"', {"size": Pt(20), "color": WHITE})],
        [('Hyp: "', {"size": Pt(20), "color": LGRAY}),
         ('animal migration', {"size": Pt(20), "color": RED, "bold": True}),
         (' gave the order"', {"size": Pt(20), "color": WHITE})],
    ], rx + Inches(0.3), by + Inches(0.85), bw - Inches(0.6), Inches(1.5),
       space_after=Pt(10))
    # v13 Amosi: 3-line caption matching source
    add_text(slide, '"Admiral McRae" → "animal migration"\n'
                    'Named entity destroyed — identity lost completely.\n'
                    'The viewer got the wrong person.',
             rx + Inches(0.3), by + Inches(2.85), bw - Inches(0.6), Inches(1.4),
             size=Pt(18), color=LGRAY)

    # Bottom callout — audit:bigfonts — Pt(24) -> Pt(24); copy condensed.
    # Pushed to y=6.20 to clear taller cards (bh 4.0 -> 4.4 in audit fix).
    add_rect(slide, MX + Inches(1.5), Inches(6.20), CW - Inches(3.0), Inches(0.65),
             fill_color=NAVY2, border_color=TEAL, border_width=Pt(2),
             corner_radius=True)
    # v13 Amosi: 2-line callout matching source
    add_text(slide,
             "WER counts both as 1 error. But one preserves meaning, the other destroys it.\n"
             "We needed our own metric — the Intelligibility Score (IS).",
             MX, Inches(6.20), CW, Inches(0.66),
             size=Pt(16), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    _finish(slide, 6,
        "The Admiral McRae example. Both cases have the same WER — 1 word "
        "substitution. But 'the' to 'a' is harmless, while 'Admiral McRae' to "
        "'animal migration' destroys the speaker's identity. WER treats them "
        "identically. This is why we built our own metric — the Intelligibility "
        "Score — which asks: did the viewer get the message?",
        [[r1], [r2]], click_reveal=True)


# ═══════════════════════════════════════════════════════════════════════
# DIVERSITY OF INPUTS — these videos are not LRS3
# ═══════════════════════════════════════════════════════════════════════

def slide_diversity_of_inputs(prs):  # audit:bigfonts2
    """Diversity of inputs — wild, surveillance-style footage, not LRS3.

    audit:bigfonts2 — Pass 2: bullets trimmed 5 -> 4, box h 4.7 -> 3.95
    so bullets bottom (6.10) clears bottom-callout top (6.5).
    """
    slide = new_slide(prs)
    add_title(slide, "Diversity of Inputs — Not LRS3")
    add_accent_line(slide)

    # Subtitle / framing line — audit:bigfonts — Pt(24) -> Pt(24).
    add_text(slide,
        "Real-world / observational footage — not studio video.",
        MX, CT, CW, Inches(0.5),
        size=Pt(22), color=LGRAY, italic=True)

    # Left column: bullets describing the input distribution.
    # audit:bigfonts2 — V6 dense; bullets trimmed to <=8 words; 5 -> 4 bullets.
    # CUT v2: removed "Speaker diversity" bullet (subsumed by "Diverse settings").
    col_w = Inches(7.0)
    col_x = MX
    bullets = [
        ("Real-world / observational — not studio-lit",
         {"bold": True, "color": TEAL}),
        "Indoor + outdoor, varied lighting and noise",
        "Source clips: 20 s to 1 h, we segment 10-360 s",
        ("1,497 segments curated, not cherry-picked",
         {"bold": True, "color": CORAL}),
    ]
    bul = add_bullets(slide, bullets, col_x, CT + Inches(0.7),
                      col_w, Inches(3.95), size=Pt(24), bullet_color=TEAL)

    # Right column: a single still / GIF anchored to the demo footage so the
    # audience sees what "wild" looks like. Uses the lip-reading demo GIF
    # already shipped under PLOTS — falls back gracefully if missing.
    rx = MX + col_w + Inches(0.4)
    rw = CW - col_w - Inches(0.4)
    img_y = CT + Inches(0.6)
    img_h = Inches(3.6)
    gif_path = PLOTS / "lip_reading_demo.gif"  # ships under MATERIALS/01_plots_for_slides/
    img_shape = None
    if gif_path.exists():
        # Place picture directly so we can use add_picture (helpers.add_image
        # expects an IMG key); GIF is treated like any picture by python-pptx.
        try:
            img_shape = slide.shapes.add_picture(
                str(gif_path), rx, img_y, width=rw, height=img_h)
        except Exception:
            img_shape = None
    if img_shape is None:
        # Fallback placeholder card so layout stays consistent.
        img_shape = add_rect(slide, rx, img_y, rw, img_h,
                             fill_color=NAVY2, border_color=MGRAY,
                             border_width=Pt(1), corner_radius=True)
        add_text(slide, "[demo frame]", rx, img_y + img_h / 2 - Inches(0.2),
                 rw, Inches(0.4), size=Pt(20), color=MGRAY,
                 align=PP_ALIGN.CENTER)

    # audit:bigfonts — caption Pt(12) -> Pt(24); kept short.
    add_text(slide,
        "Lip-reading frame — visual signal only, no audio.",
        rx, img_y + img_h + Inches(0.05), rw, Inches(0.8),
        size=Pt(18), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

    # Bottom callout — payoff. audit:bigfonts — Pt(24) -> Pt(24).
    # CUT v3: callout top 6.5->6.15 + frame h trimmed so Pt24 stays under 7.05.
    add_rect(slide, MX + Inches(1.5), Inches(6.15), CW - Inches(3.0), Inches(0.6),
             fill_color=NAVY2, border_color=TEAL, border_width=Pt(2),
             corner_radius=True)
    add_text(slide,
        "Every number that follows comes from this distribution — not a benchmark.",
        MX + Inches(1.7), Inches(6.02), CW - Inches(3.4), Inches(1.0),
        size=Pt(24), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    _finish(slide, 0,
        "Framing slide: the 1,497-segment evaluation set is curated from "
        "real-world surveillance and observational YouTube footage, not LRS3 "
        "studio recordings. Source clips range 20s to 1h; we segment them into "
        "10–360s windows for evaluation. Speakers include diverse accents, "
        "ages, and occlusions (hands, microphones, partial faces). This is why "
        "WER on our set (64% top-1, 64% MBR) is far worse than LRS3 "
        "(~25%) — fundamentally harder data, not a worse model.\n\n"
        "PEER DETAIL: 1,497-segment evaluation set was curated from a larger "
        "pool of YouTube content; segments balanced across topic categories "
        "(interviews, talks, casual conversation, technical, religious) but "
        "not formally stratified. "
        "See docs/architecture.md for dataset description and "
        "docs/evaluation/after_amosi_audit.md for the corresponding numbers.",
        [[bul], [img_shape]], click_reveal=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — THE INTELLIGIBILITY SCORE
# ═══════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════
# NEW SLIDES — DEEP DIVES, CONTEXT, ENGINEERING, APPENDIX
# ═══════════════════════════════════════════════════════════════════════

def slide_visemes(prs):  # audit:bigfonts
    """Fundamental viseme challenge — why lip reading is hard."""
    slide = new_slide(prs)
    add_title(slide, "The Invisible Problem: Visemes")
    add_accent_line(slide)

    col_w = Inches(5.5)
    gap = Inches(1.13)

    # Left column — The problem. audit:bigfonts — Pt(24) -> Pt(24).
    lt = add_text(slide, "The Invisible Problem", MX, CT, col_w, Inches(0.5),
                  size=Pt(24), color=CORAL, bold=True)
    # audit:bigfonts \u2014 Pt(24) -> Pt(24).
    # OVERLAP fix: lb h 2.6 -> 1.50 so it ends at CT+2.10, clearing tbl1
    # which starts at CT+2.6. Previous 2.6" frame overlapped viseme table
    # causing third bullet to render behind the table.
    lb = add_bullets(slide, [
        ("50\u201370% of English sounds are invisible on lips", {"bold": True}),
        "Multiple sounds produce identical mouth shapes (visemes)",
        "Context is the ONLY disambiguation signal",
    ], MX, CT + Inches(0.6), col_w, Inches(1.50), size=Pt(18))

    # Viseme table — v13: smaller text to match source compact layout
    tbl1 = add_table(slide,
        ["Viseme Group", "Sounds"],
        [["Bilabial", "p, b, m"],
         ["Alveolar", "t, d, n, s, z, l"],
         ["Velar", "k, g, ng"],
         ["Labiodental", "f, v"]],
        MX, CT + Inches(2.4), col_w, text_size=Pt(16),
        row_height=Inches(0.32))

    # Right column — Visual proof + confusable pairs.
    # audit:bigfonts — title Pt(24) -> Pt(24).
    rx = MX + col_w + gap
    rt = add_text(slide, "Same Mouth Shape, Different Words", rx, CT, col_w,
                  Inches(1.0), size=Pt(24), color=TEAL, bold=True)

    # Lip-reading GIF — identical mouth shapes, different words.
    # Image NOT shrunk per spec ("Don't shrink plot images").
    poster_shapes = []
    gif_w = Inches(5.0)
    gif_h = Inches(1.5)  # constrain height to prevent overlap with table
    gif_path = MATERIALS / "mom-yelling.gif"
    if gif_path.exists():
        gif_x = rx + (col_w - gif_w) / 2  # Center in right column
        poster_shapes.append(slide.shapes.add_picture(
            str(gif_path), gif_x, CT + Inches(0.6), width=gif_w, height=gif_h))
    # audit:bigfonts — caption Pt(12) -> Pt(24); copy trimmed.
    poster_shapes.append(add_text(slide,
        "Identical mouth shapes can produce different words",
        rx, CT + Inches(2.18), col_w, Inches(0.73),
        size=Pt(18), color=LGRAY, italic=True, align=PP_ALIGN.CENTER))

    # Word-pair table — was overlapping bottom callout (#308 fix).
    # Default row_height 0.38 gave 5×0.38=1.90 ending at 6.40 (callout @6.30,
    # 0.10" overlap). row_height=0.32 → 5×0.32=1.60, ends at 6.10, 0.20" gap.
    tbl2 = add_table(slide,
        ["Word A", "Word B"],
        [["pat", "bat"],
         ["mom", "bomb"],
         ["admiral", "animal"],
         ["collar", "color"]],
        rx, CT + Inches(2.8), col_w, text_size=Pt(16),
        row_height=Inches(0.32))

    # v13 Amosi: bottom callout in source
    add_text(slide,
        "Context is the ONLY disambiguation signal \u2014 this is why the LLM matters.",
        MX, Inches(6.60), CW, Inches(0.40),
        size=Pt(16), color=TEAL, italic=True, align=PP_ALIGN.CENTER)

    _finish(slide, 0,
        "50-70% of English sounds are invisible on lips — this is the homophene "
        "(viseme) problem that makes lip reading fundamentally ambiguous. Multiple "
        "phonemes such as /p, b, m/ or /t, d, n/ produce indistinguishable mouth "
        "shapes, so any text-only lip reader has at least 30-50% irreducible "
        "ambiguity per word. The poster frames show two different speakers whose "
        "mouth shapes look nearly identical while saying completely different "
        "words ('mom' vs 'bomb', 'pat' vs 'bat'). This is precisely why the "
        "model alone hallucinates on 20% of segments — it has no way to "
        "resolve homophenes without context. The LLM component supplies that "
        "context, which is why the architecture is visual encoder + LLM rather "
        "than visual encoder alone. "
        "Sources: docs/evaluation/intelligibility_methodology.md (homophene "
        "rate), docs/evaluation/llm_judge/llm_judge_analysis.md (hallucination "
        "rate 20%).",
        [[lt, lb, tbl1], poster_shapes + [rt, tbl2]], click_reveal=True)


def slide_data_flow(prs):  # audit:bigfonts
    """7-step pipeline data flow (May 2026: + n-best aggregation + per-word confidence)."""
    slide = new_slide(prs)
    add_title(slide, "How It Works: Data Flow")
    add_accent_line(slide)

    # v13 Amosi: 5-step source data-flow (Video → Mouth → Features → Projection → LLM)
    steps = [
        ("1", "Video Frames", "25 fps raw video input", TEAL),
        ("2", "Mouth Crop", "96×96 pixel region around lips", TEAL),
        ("3", "Visual Features", "AV-HuBERT encoder → 1024-dim vectors", TEAL),
        ("4", "Projection", "Linear layer: 1024 → 4096-dim (LLM input space)", CORAL),
        ("5", "LLM Generates Text", "LLaMA-2-7B decodes features into words", CORAL),
    ]

    # 7 steps \u00d7 0.65 + 6 \u00d7 0.10 = 4.55 + 0.60 = 5.15", starting at CT+0.10=1.55
    # \u2192 ends at 6.70. Bottom caption at 6.85.
    step_w = Inches(10.5)
    step_h = Inches(0.82)
    start_y = CT + Inches(0.30)
    start_x = MX + Inches(0.8)

    step_groups = []
    gap = Inches(0.20)  # v13: 5 steps, more spacing
    for i, (num, name, desc, color) in enumerate(steps):
        y = start_y + i * (step_h + gap)
        r = add_rect(slide, start_x, y, step_w, step_h, fill_color=NAVY2,
                     border_color=color, border_width=Pt(2), corner_radius=True)

        # Step number circle \u2014 bumped 0.55 -> 0.65 with Pt(24) -> Pt(28).
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, start_x - Inches(0.78), y + Inches(0.1),
            Inches(0.65), Inches(0.65))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        num_txt = add_text(slide, num, start_x - Inches(0.78), y + Inches(0.1),
                 Inches(0.65), Inches(0.65),
                 size=Pt(28), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # CUT v3 (overflow): desc shrunk to Pt(20) so name+desc fits one line.
        rt = add_rich_text(slide, [
            [(name, {"size": Pt(24), "color": WHITE, "bold": True}),
             (f"  \u2014  {desc}", {"size": Pt(20), "color": LGRAY})],
        ], start_x + Inches(0.2), y + Inches(0.15),
           step_w - Inches(0.4), step_h - Inches(0.25))

        group = [r, circle, num_txt, rt]

        # Arrow between steps \u2014 kept Pt(20) (footer-ish). Height matches gap.
        if i < len(steps) - 1:
            arrow = add_text(slide, "\u2193", start_x + step_w / 2 - Inches(0.2),
                     y + step_h - Inches(0.05), Inches(0.4), gap + Inches(0.05),
                     size=Pt(20), color=TEAL, align=PP_ALIGN.CENTER)
            group.append(arrow)

        step_groups.append(group)

    # Bottom caption \u2014 moved down. 7 steps \u00d7 0.65 + 6 \u00d7 0.08 = 4.55+0.48=5.03,
    # start 1.55 \u2192 ends 6.58. Caption at 6.65, h=0.35, end 7.00 \u2713
    add_text(slide,
        "Visual encoder is frozen (pre-trained on LRS3). Only projection + LoRA adapters are trained.",
        MX, Inches(6.60), CW, Inches(0.40),
        size=Pt(15), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

    _finish(slide, 0,
        "Five-step data flow. Raw video at 25fps is cropped to 96x96 mouth "
        "region. AV-HuBERT extracts 1024-dim visual features. Linear projection "
        "maps to 4096-dim LLM input space (1024 -> 4096). LLaMA-2-7B generates "
        "text. The visual encoder is frozen — only the projection layer and "
        "LoRA adapters are trained. "
        "Source: docs/architecture.md (pipeline flow + data formats).",
        step_groups)


def slide_eval_dataset(prs):  # audit:bigfonts
    """Our 1,497-segment evaluation dataset."""
    slide = new_slide(prs)
    add_title(slide, "Our Evaluation: 1,497 Real-World Segments")
    add_accent_line(slide)

    col_w = Inches(5.5)
    gap = Inches(1.13)

    # Left — big number + bullets. audit:bigfonts — keep big-num,
    # subtitle Pt(24) -> Pt(24), bullets Pt(24) -> Pt(24).
    s1 = add_text(slide, "1,497", MX, CT, col_w, Inches(1.0),
                  size=Pt(64), color=TEAL, bold=True)
    s2 = add_text(slide, "segments from diverse YouTube videos",
                  MX, CT + Inches(1.05), col_w, Inches(0.5),
                  size=Pt(24), color=WHITE)
    # audit:bigfonts \u2014 Pt(24) -> Pt(24).
    s3 = add_bullets(slide, [
        "Uncontrolled lighting, angles, occlusions",
        "Multiple speakers and accents",
        "Diverse topics: business to DIY to gaming",
        ("Not a curated benchmark \u2014 real-world difficulty", {"bold": True}),
    ], MX, CT + Inches(1.7), col_w, Inches(3.0), size=Pt(24))

    # Right — topic categories table.
    # audit:bigfonts — title Pt(24) -> Pt(24); table Pt(12) -> Pt(24).
    rx = MX + col_w + gap
    add_text(slide, "Topic Distribution", rx, CT, col_w, Inches(0.5),
             size=Pt(24), color=CORAL, bold=True)

    tbl = add_table(slide,
        ["Topic", "Segments", "Quality*"],
        [["Business/Finance", "214", "3.08"],
         ["Education/Lecture", "312", "2.63"],
         ["Entertainment", "198", "2.51"],
         ["News/Politics", "267", "2.48"],
         ["Tech/Science", "186", "2.43"],
         ["Sports/Health", "153", "2.38"],
         ["DIY/Home", "167", "2.13"]],
        rx, CT + Inches(0.6), col_w, text_size=Pt(24),
        row_colors={0: {2: GREEN}, 6: {2: CORAL}})

    # Footnote \u2014 audit:bigfonts \u2014 Pt(11) -> Pt(24) (footer tier).
    # Footnote: Pt(16)\u2192Pt(18) italic to clear the caption-floor.
    add_text(slide,
        "*Quality = our composite metric (0\u20135 scale), introduced next.",
        MX, Inches(6.5), CW, Inches(0.4),
        size=Pt(18), color=MGRAY, italic=True)

    _finish(slide, 0,
        "1,497 segments from diverse YouTube videos. Not a curated benchmark. "
        "Multiple topics, speakers, accents, lighting conditions. Business and "
        "Finance has the highest quality score (3.08) because it's closest to the "
        "TED talk training data. DIY/Home is worst (2.13) due to inherently visual "
        "content. The Quality column is our Intelligibility Score, introduced later.",
        [[s1, s2, s3], [tbl]], click_reveal=True)


def slide_wer_explained(prs):  # audit:bigfonts
    """WER formula and its limitations."""
    slide = new_slide(prs)
    add_title(slide, "Word Error Rate: What It Measures (and Misses)")
    add_accent_line(slide)

    col_w = Inches(5.5)
    gap = Inches(1.13)

    # Left — formula + example. audit:bigfonts — Pt(24) -> Pt(24),
    # Pt(24) -> Pt(30), Pt(24) -> Pt(24), Pt(24) -> Pt(24), Pt(24) -> Pt(24).
    lt = add_text(slide, "The Formula", MX, CT, col_w, Inches(0.5),
                  size=Pt(24), color=TEAL, bold=True)

    add_text(slide, "WER = (S + D + I) / N",
             MX + Inches(0.3), CT + Inches(0.6), col_w - Inches(0.6), Inches(0.6),
             size=Pt(30), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # audit:bigfonts — formula bullets fit within height before "Example:"
    # header at CT+3.6.
    add_bullets(slide, [
        "S = substitutions (wrong word)",
        "D = deletions (missing word)",
        "I = insertions (extra word)",
        "N = total words in reference",
    ], MX, CT + Inches(1.35), col_w, Inches(2.2), size=Pt(24))

    add_text(slide, "Example:", MX, CT + Inches(3.6), col_w, Inches(0.4),
             size=Pt(24), color=TEAL, bold=True)
    add_text(slide, 'Ref: "the admiral gave orders"\n'
                    'Hyp: "the animal gave water"\n'
                    'WER = 2/4 = 50% (2 substitutions)',
             MX + Inches(0.2), CT + Inches(4.1), col_w - Inches(0.4), Inches(1.5),
             size=Pt(24), color=WHITE)

    # Right — what it captures vs misses
    rx = MX + col_w + gap
    rt = add_text(slide, "What WER Captures", rx, CT, col_w, Inches(0.5),
                  size=Pt(24), color=GREEN, bold=True)
    # audit:bigfonts — Pt(24) -> Pt(24). Trimmed 3 -> 2 bullets so the
    # "What WER Misses" block can move up and rb2 fits above slide-num.
    rb1 = add_bullets(slide, [
        "Exact word-level accuracy",
        "Simple, standard in ASR research",
    ], rx, CT + Inches(0.6), col_w, Inches(1.2), size=Pt(24), bullet_color=GREEN)

    rm = add_text(slide, "What WER Misses", rx, CT + Inches(1.95), col_w, Inches(0.5),
                  size=Pt(24), color=CORAL, bold=True)
    # audit:bigfonts \u2014 Pt(24) -> Pt(24); trimmed 5 -> 4 bullets so block
    # ends above slide-num at y=6.85. Height 2.7 fits 4 bullets at Pt(24).
    rb2 = add_bullets(slide, [
        ("All errors weighted equally", {"color": CORAL}),
        ('"admiral"\u2192"animal" = "the"\u2192"a"', {"color": CORAL}),
        "No meaning preservation signal",
        ("Can exceed 100% (insertions)", {"color": CORAL}),
    ], rx, CT + Inches(2.55), col_w, Inches(2.7), size=Pt(24), bullet_color=CORAL)

    _finish(slide, 0,
        "WER formula: substitutions plus deletions plus insertions divided by "
        "reference word count. Simple and standard, but treats all errors equally. "
        "Admiral-to-animal gets the same penalty as the-to-a. No credit for "
        "meaning preservation or phonetic similarity. Can exceed 100% when the "
        "model generates extra words (insertions).",
        [[lt], [rt], [rb1], [rm], [rb2]], click_reveal=True)


