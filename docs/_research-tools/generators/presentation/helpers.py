"""
Presentation helpers — slide setup, text, shapes, images, tables,
animations, transitions, and reusable slide builders.
"""

import os
import subprocess
from pathlib import Path
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from .config import (
    IMG, VID, POSTER_DIR,
    SL_W, SL_H, BG, WHITE, TEAL, CORAL, LGRAY, MGRAY, DGRAY,
    GREEN, YELLOW, GOLD, ORANGE, RED, DRED, NAVY2, NAVY3,
    FONT, _auto_num,
    MX, MY, CT, CW, CH, SLW, SRG, SRL, SRW,
)

# ═══════════════════════════════════════════════════════════════════════
# HELPERS — SLIDE SETUP
# ═══════════════════════════════════════════════════════════════════════

def new_slide(prs):
    """Create blank slide with navy background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = BG
    return slide


def set_notes(slide, text):
    """Set speaker notes."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_logo(slide, size=Inches(0.35)):
    """Small logo in bottom-right."""
    p = IMG["logo"]
    if p.exists():
        return slide.shapes.add_picture(
            str(p), SL_W - MX - size, SL_H - Inches(0.12) - size, height=size)
    return None


def add_slide_num(slide, num):
    """Slide number in bottom-left."""
    tb = slide.shapes.add_textbox(MX, SL_H - Inches(0.38), Inches(0.5), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.text = str(num)
    _fmt(p, size=Pt(9), color=MGRAY)
    return tb


def add_accent_line(slide, top=Inches(1.2)):
    """Thin teal accent line below title area."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, top, CW, Pt(1.5))
    shp.fill.solid()
    shp.fill.fore_color.rgb = TEAL
    shp.line.fill.background()
    return shp

# ═══════════════════════════════════════════════════════════════════════
# HELPERS — TEXT
# ═══════════════════════════════════════════════════════════════════════

def _fmt(para, size=Pt(16), color=WHITE, bold=False, italic=False, font=FONT):
    """Format all runs in a paragraph, or create default run formatting."""
    for run in para.runs:
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    if not para.runs:
        para.font.size = size
        para.font.color.rgb = color
        para.font.bold = bold
        para.font.italic = italic
        para.font.name = font


def add_title(slide, text, top=MY, left=MX, width=None, size=Pt(32)):
    """Add title text box."""
    if width is None:
        width = CW
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _fmt(p, size=size, color=WHITE, bold=True)
    return tb


def add_text(slide, text, left, top, width, height, size=Pt(16),
             color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    """Add a simple text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        body_props = tf._txBody.bodyPr
        body_props.set('anchor', {
            MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b'
        }.get(anchor, 't'))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    _fmt(p, size=size, color=color, bold=bold, italic=italic)
    return tb


def add_rich_text(slide, runs_list, left, top, width, height,
                  align=PP_ALIGN.LEFT, space_after=Pt(6)):
    """
    Add text box with mixed formatting.
    runs_list: list of lists of (text, {props}) tuples — one inner list per paragraph.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for pi, para_runs in enumerate(runs_list):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        for text, props in para_runs:
            run = p.add_run()
            run.text = text
            run.font.name = props.get("font", FONT)
            run.font.size = props.get("size", Pt(16))
            run.font.color.rgb = props.get("color", WHITE)
            run.font.bold = props.get("bold", False)
            run.font.italic = props.get("italic", False)
    return tb


def add_bullets(slide, items, left, top, width, height, size=Pt(15),
                color=WHITE, bullet_color=TEAL, spacing=Pt(6)):
    """
    Add bullet list. items can be str or (str, {props}) for custom bullets.
    Returns the text box shape.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, props = item
        else:
            text, props = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT
        # Bullet character
        brun = p.add_run()
        bc = props.get("bullet", "\u2022")
        brun.text = f"{bc} "
        brun.font.size = size
        brun.font.color.rgb = props.get("bullet_color", bullet_color)
        brun.font.name = FONT
        brun.font.bold = True
        # Text
        trun = p.add_run()
        trun.text = text
        trun.font.size = size
        trun.font.color.rgb = props.get("color", color)
        trun.font.name = FONT
        trun.font.bold = props.get("bold", False)
    return tb

# ═══════════════════════════════════════════════════════════════════════
# HELPERS — SHAPES & IMAGES
# ═══════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill_color=NAVY2,
             border_color=None, border_width=Pt(1), corner_radius=None):
    """Add a filled rectangle (optionally rounded)."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = border_width
    else:
        shp.line.fill.background()
    return shp


def add_image(slide, key_or_path, left, top, width=None, height=None):
    """Add an image by key name or path. Returns shape or None."""
    path = key_or_path if isinstance(key_or_path, Path) else IMG.get(key_or_path)
    if path and path.exists():
        kwargs = {}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        if not kwargs:
            kwargs["width"] = Inches(4)
        return slide.shapes.add_picture(str(path), left, top, **kwargs)
    # Placeholder rectangle if image missing
    w = width or Inches(4)
    h = height or Inches(3)
    shp = add_rect(slide, left, top, w, h, fill_color=DGRAY, border_color=MGRAY)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    name = path.name if path else str(key_or_path)
    p.text = f"[{name}]"
    p.alignment = PP_ALIGN.CENTER
    _fmt(p, size=Pt(11), color=LGRAY)
    try:
        tf._txBody.bodyPr.set('anchor', 'ctr')
    except Exception:
        pass
    return shp


def add_play_button(slide, left, top, size=Inches(1.8)):
    """Create a play button icon (circle + triangle)."""
    # Outer circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = NAVY2
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(2.5)
    # Play triangle text
    inset = Inches(0.15)
    tb = slide.shapes.add_textbox(
        left + inset, top + inset, size - 2 * inset, size - 2 * inset)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u25b6"
    run.font.size = Pt(int(size / Inches(1) * 36))
    run.font.color.rgb = WHITE
    run.font.name = FONT
    try:
        tf._txBody.bodyPr.set('anchor', 'ctr')
    except Exception:
        pass
    return circle


def _extract_poster(vid_path, poster_path):
    """Extract first frame from video as a poster image for embedding."""
    poster_path.parent.mkdir(parents=True, exist_ok=True)
    if poster_path.exists():
        return poster_path
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-i", str(vid_path), "-vframes", "1",
             "-q:v", "2", str(poster_path)],
            capture_output=True, timeout=10, check=True)
        if poster_path.exists() and poster_path.stat().st_size > 0:
            return poster_path
    except Exception:
        pass
    return None


def add_video_poster(slide, vid_key, left, top, width, height):
    """Show a video poster frame (static thumbnail) with play button overlay.

    This avoids PowerPoint repair issues caused by python-pptx video embeds.
    Falls back to a dark placeholder with play icon if poster not available.
    """
    vid_path = VID.get(vid_key)
    poster = None
    if vid_path and vid_path.exists():
        poster = _extract_poster(vid_path, POSTER_DIR / f"{vid_key}.jpg")
    if poster and poster.exists():
        shape = slide.shapes.add_picture(str(poster), left, top, width, height)
    else:
        # Dark placeholder
        shape = add_rect(slide, left, top, width, height,
                         fill_color=NAVY2, border_color=MGRAY, border_width=Pt(1))
    # Play button overlay (centered, small)
    btn_size = min(width, height) * 0.3
    btn_x = left + (width - btn_size) // 2
    btn_y = top + (height - btn_size) // 2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, btn_x, btn_y, btn_size, btn_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0, 0, 0)
    # Semi-transparent effect via alpha on the XML element
    try:
        spPr = circle._element.find(qn('a:spPr'))
        if spPr is not None:
            sf = spPr.find(qn('a:solidFill'))
            if sf is not None:
                srgb = sf.find(qn('a:srgbClr'))
                if srgb is not None:
                    alpha = etree.SubElement(srgb, qn('a:alpha'))
                    alpha.set('val', '50000')  # 50% opacity
    except Exception:
        pass
    circle.line.fill.background()
    tb = slide.shapes.add_textbox(btn_x, btn_y, btn_size, btn_size)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u25b6"
    run.font.size = Pt(int(btn_size / Inches(1) * 28))
    run.font.color.rgb = WHITE
    run.font.name = FONT
    try:
        tf._txBody.bodyPr.set('anchor', 'ctr')
    except Exception:
        pass
    return shape


def add_video(slide, vid_key, left, top, width, height):
    """Embed an MP4 video with poster frame. Click to play in PowerPoint."""
    vid_path = VID.get(vid_key)
    if not vid_path or not vid_path.exists():
        return add_play_button(slide, left, top, size=min(width, height))
    poster = _extract_poster(vid_path, POSTER_DIR / f"{vid_key}.jpg")
    poster_arg = str(poster) if poster else None
    shape = slide.shapes.add_movie(
        str(vid_path), left, top, width, height,
        poster_frame_image=poster_arg, mime_type="video/mp4")
    # Fix python-pptx bug: add_movie creates <a:hlinkClick r:id="" action="ppaction://media"/>
    # with empty r:id, which triggers PowerPoint repair warning.  Remove empty hlinkClick.
    cNvPr = shape._element.find('.//' + qn('p:cNvPr'))
    if cNvPr is not None:
        for hlink in list(cNvPr):
            if hlink.tag.endswith('hlinkClick'):
                rid = hlink.get(qn('r:id'), '')
                if rid == '':
                    cNvPr.remove(hlink)
    return shape


# ═══════════════════════════════════════════════════════════════════════
# HELPERS — TABLES
# ═══════════════════════════════════════════════════════════════════════

def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.38),
              header_color=TEAL, text_size=Pt(11), col_widths=None,
              row_colors=None):
    """
    Add a styled table. row_colors: dict mapping row_index -> dict mapping
    col_index -> RGBColor for cell text coloring.
    """
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                        row_height * n_rows)
    tbl = tbl_shape.table
    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    # Header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        _shade_cell(cell, _rgb_hex(header_color))
        for p in cell.text_frame.paragraphs:
            _fmt(p, size=text_size, color=WHITE, bold=True)
            p.alignment = PP_ALIGN.CENTER
    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            # Alternating row shading
            if ri % 2 == 1:
                _shade_cell(cell, "152a40")
            else:
                _shade_cell(cell, "0d1b2a")
            txt_color = WHITE
            if row_colors and ri in row_colors and ci in row_colors[ri]:
                txt_color = row_colors[ri][ci]
            for p in cell.text_frame.paragraphs:
                _fmt(p, size=text_size, color=txt_color)
                p.alignment = PP_ALIGN.LEFT
    # Remove table borders for dark theme
    _style_table_borders(tbl, "1a3550")
    return tbl_shape


def _shade_cell(cell, hex_color):
    """Set cell background color."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing shading
    for child in list(tcPr):
        if child.tag.endswith('solidFill') or child.tag.endswith('shd'):
            tcPr.remove(child)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val', hex_color)


def _style_table_borders(tbl, hex_color):
    """Set subtle table borders."""
    tbl_xml = tbl._tbl
    tblPr = tbl_xml.tblPr
    if tblPr is None:
        tblPr = etree.SubElement(tbl_xml, qn('a:tblPr'))
    # Remove default style
    for attr in ['bandRow', 'firstRow', 'lastRow', 'bandCol']:
        tblPr.set(attr, '0')


def _rgb_hex(color):
    """Convert RGBColor to hex string."""
    return f"{color[0]:02x}{color[1]:02x}{color[2]:02x}"

# ═══════════════════════════════════════════════════════════════════════
# HELPERS — ANIMATIONS & TRANSITIONS
# ═══════════════════════════════════════════════════════════════════════

def add_fade_transition(slide, speed='med'):
    """Add fade transition to slide.

    ECMA-376 requires p:sld children in order:
        cSld, clrMapOvr, transition, timing, extLst
    So we must insert transition *before* any existing timing element,
    not just append it (python-pptx add_movie creates timing early).
    """
    sld = slide._element
    # Remove existing transition
    for child in list(sld):
        if child.tag.endswith('transition'):
            sld.remove(child)
    transition = etree.Element(qn('p:transition'))
    transition.set('spd', speed)
    transition.set('advClick', '1')
    etree.SubElement(transition, qn('p:fade'))
    # Insert before timing/extLst to maintain ECMA-376 order
    timing_el = sld.find(qn('p:timing'))
    if timing_el is not None:
        sld.insert(list(sld).index(timing_el), transition)
    else:
        sld.append(transition)


def _fix_pptx_video_compat(pptx_path):
    """Post-process saved PPTX to wrap video p:pic in mc:AlternateContent.

    python-pptx embeds videos as bare p:pic elements with p14:media extensions.
    PowerPoint expects these wrapped in mc:AlternateContent (Choice/Fallback)
    per ISO/IEC 29500.  Without the wrapper, PowerPoint shows a repair dialog.

    Uses lxml for proper XML manipulation instead of regex.
    """
    import zipfile, shutil, tempfile, copy as _copy

    MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
    P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

    tmp_fd, tmp_path = tempfile.mkstemp(suffix='.pptx')
    os.close(tmp_fd)
    changed = False

    with zipfile.ZipFile(pptx_path, 'r') as zin, \
         zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            raw = zin.read(item.filename)

            if (item.filename.startswith('ppt/slides/slide') and
                    item.filename.endswith('.xml') and
                    b'videoFile' in raw):
                tree = etree.fromstring(raw)
                root = tree  # <p:sld ...>

                # Register mc and p14 namespaces on root
                nsmap = dict(root.nsmap)
                need_update = False
                if 'mc' not in nsmap:
                    nsmap['mc'] = MC_NS
                    need_update = True
                if 'p14' not in nsmap:
                    nsmap['p14'] = P14_NS
                    need_update = True

                if need_update:
                    # Rebuild root with updated nsmap (lxml requires this)
                    new_root = etree.Element(root.tag, root.attrib, nsmap=nsmap)
                    for child in root:
                        new_root.append(child)
                    root = new_root

                # Set mc:Ignorable="p14" on root
                mc_ign = root.get(f'{{{MC_NS}}}Ignorable', '')
                if 'p14' not in mc_ign.split():
                    root.set(f'{{{MC_NS}}}Ignorable',
                             (mc_ign + ' p14').strip())

                # Find all p:pic elements containing a:videoFile
                video_pics = []
                for pic in root.iter(f'{{{P_NS}}}pic'):
                    if pic.find(f'.//{{{A_NS}}}videoFile') is not None:
                        video_pics.append(pic)

                for pic in video_pics:
                    parent = pic.getparent()
                    idx = list(parent).index(pic)

                    # Build mc:AlternateContent wrapper
                    ac = etree.Element(f'{{{MC_NS}}}AlternateContent')
                    choice = etree.SubElement(ac, f'{{{MC_NS}}}Choice')
                    choice.set('Requires', 'p14')
                    fallback = etree.SubElement(ac, f'{{{MC_NS}}}Fallback')

                    # Choice gets the original pic (with videoFile + p14:media)
                    choice.append(pic)  # moves pic out of parent

                    # Fallback gets a deep copy without videoFile and p:extLst
                    fb_pic = _copy.deepcopy(pic)
                    for vf in fb_pic.findall(f'.//{{{A_NS}}}videoFile'):
                        vf.getparent().remove(vf)
                    # Remove p:extLst containing p14:media from nvPr
                    for ext_lst in fb_pic.findall(f'.//{{{P_NS}}}extLst'):
                        ext_lst.getparent().remove(ext_lst)
                    fallback.append(fb_pic)

                    # Insert wrapper where pic was
                    parent.insert(idx, ac)
                    changed = True

                if video_pics:
                    raw = etree.tostring(root, xml_declaration=True,
                                         encoding='UTF-8', standalone=True)

            zout.writestr(item, raw)

    if changed:
        shutil.move(tmp_path, pptx_path)
    else:
        os.remove(tmp_path)


def _para_indices(shape):
    """Return indices of non-empty paragraphs in shape's text frame."""
    try:
        return [i for i, p in enumerate(shape.text_frame.paragraphs)
                if p.text.strip()]
    except (AttributeError, TypeError):
        return []


def add_animations(slide, groups, click_reveal=False, para_build=False):
    """
    Add click-to-advance entrance animations (Appear) with valid OOXML.

    Produces the three-level par nesting that PowerPoint 365 expects:
      timing > tnLst > par(tmRoot) > seq(mainSeq)
        per click:  par(indefinite) > par(delay=0) > par(presetID=1) > set

    groups       – list of lists of shape objects; each inner list is one
                   animation group (one "click step").
    click_reveal – True  → first group visible on slide entry;
                   False → all groups hidden until clicked.
    para_build   – True  → multi-paragraph text shapes (bullet lists)
                   reveal one paragraph per click instead of all at once.
    """
    if not groups:
        return
    groups = [g for g in groups if g]
    if not groups:
        return

    # ── Collect all animatable shapes ──────────────────────────────────
    all_anim_shapes = []
    for gi, group in enumerate(groups):
        for shape in group:
            if shape is None:
                continue
            try:
                _ = shape.shape_id
                all_anim_shapes.append((gi, shape))
            except AttributeError:
                continue
    if not all_anim_shapes:
        return

    # ── Expand groups into click steps ─────────────────────────────────
    # Each step is a list of (shape, para_idx | None) tuples.
    click_steps = []
    para_build_spids = set()   # shape IDs needing build="p"
    hidden_spids = set()       # shape IDs that go in bldLst

    for gi, group in enumerate(groups):
        is_entry = click_reveal and gi == 0
        step_items = []

        for shape in group:
            if shape is None:
                continue
            try:
                spid = shape.shape_id
            except AttributeError:
                continue

            pindices = _para_indices(shape) if para_build else []

            if len(pindices) > 1:
                # Multi-paragraph → per-paragraph click steps
                para_build_spids.add(spid)
                hidden_spids.add(spid)
                for pi in pindices:
                    click_steps.append([(shape, pi)])
            elif is_entry:
                pass   # visible on entry — no animation needed
            else:
                hidden_spids.add(spid)
                step_items.append((shape, None))

        if step_items:
            click_steps.append(step_items)

    if not click_steps:
        return

    # ── Remove existing timing ─────────────────────────────────────────
    sld = slide._element
    for child in list(sld):
        if child.tag.endswith('timing'):
            sld.remove(child)

    _id = [1]
    def nid():
        v = _id[0]; _id[0] += 1; return str(v)

    SE = etree.SubElement

    # ── Timing tree ────────────────────────────────────────────────────
    timing  = SE(sld, qn('p:timing'))
    tnLst   = SE(timing, qn('p:tnLst'))
    par0    = SE(tnLst, qn('p:par'))
    cTn0    = SE(par0, qn('p:cTn'))
    cTn0.set('id', nid()); cTn0.set('dur', 'indefinite')
    cTn0.set('restart', 'never'); cTn0.set('nodeType', 'tmRoot')
    root_cl = SE(cTn0, qn('p:childTnLst'))

    # Main click sequence
    seq     = SE(root_cl, qn('p:seq'))
    seq.set('concurrent', '1'); seq.set('nextAc', 'seek')
    cTn_seq = SE(seq, qn('p:cTn'))
    cTn_seq.set('id', nid()); cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')
    seq_cl  = SE(cTn_seq, qn('p:childTnLst'))

    for step in click_steps:
        # Level 1 – click step container (delay=indefinite → waits for click)
        p1 = SE(seq_cl, qn('p:par'))
        c1 = SE(p1, qn('p:cTn'))
        c1.set('id', nid()); c1.set('fill', 'hold')
        s1 = SE(c1, qn('p:stCondLst'))
        SE(s1, qn('p:cond')).set('delay', 'indefinite')
        ch1 = SE(c1, qn('p:childTnLst'))

        # Level 2 – group container (delay=0)
        p2 = SE(ch1, qn('p:par'))
        c2 = SE(p2, qn('p:cTn'))
        c2.set('id', nid()); c2.set('fill', 'hold')
        s2 = SE(c2, qn('p:stCondLst'))
        SE(s2, qn('p:cond')).set('delay', '0')
        ch2 = SE(c2, qn('p:childTnLst'))

        for si, (shape, para_idx) in enumerate(step):
            spid = str(shape.shape_id)
            delay_ms = si * 120    # subtle stagger within group

            # Level 3 – Appear effect (presetID=1)
            p3 = SE(ch2, qn('p:par'))
            c3 = SE(p3, qn('p:cTn'))
            c3.set('id', nid())
            c3.set('presetID', '1')         # Appear
            c3.set('presetClass', 'entr')
            c3.set('presetSubtype', '0')
            c3.set('fill', 'hold')
            c3.set('grpId', '0')
            c3.set('nodeType', 'clickEffect' if si == 0 else 'withEffect')
            s3 = SE(c3, qn('p:stCondLst'))
            SE(s3, qn('p:cond')).set('delay', str(delay_ms))
            ch3 = SE(c3, qn('p:childTnLst'))

            # <p:set> visibility → visible
            p_set = SE(ch3, qn('p:set'))
            cb = SE(p_set, qn('p:cBhvr'))
            ct = SE(cb, qn('p:cTn'))
            ct.set('id', nid()); ct.set('dur', '1'); ct.set('fill', 'hold')
            sc = SE(ct, qn('p:stCondLst'))
            SE(sc, qn('p:cond')).set('delay', '0')
            tgt = SE(cb, qn('p:tgtEl'))
            sp = SE(tgt, qn('p:spTgt'))
            sp.set('spid', spid)
            if para_idx is not None:
                tx = SE(sp, qn('p:txEl'))
                pr = SE(tx, qn('p:pRg'))
                pr.set('st', str(para_idx)); pr.set('end', str(para_idx))
            al = SE(cb, qn('p:attrNameLst'))
            SE(al, qn('p:attrName')).text = 'style.visibility'
            to = SE(p_set, qn('p:to'))
            SE(to, qn('p:strVal')).set('val', 'visible')

    # Prev / Next navigation
    for evt, tag in [('onPrev', 'p:prevCondLst'),
                     ('onNext', 'p:nextCondLst')]:
        cl = SE(seq, qn(tag))
        c  = SE(cl, qn('p:cond'))
        c.set('evt', evt); c.set('delay', '0')
        t  = SE(c, qn('p:tgtEl'))
        SE(t, qn('p:sldTgt'))

    # ── Build list — shapes that start HIDDEN ──────────────────────────
    if hidden_spids:
        bldLst = SE(timing, qn('p:bldLst'))
        seen = set()
        for _gi, shape in all_anim_shapes:
            spid = shape.shape_id
            if spid not in hidden_spids or spid in seen:
                continue
            seen.add(spid)
            bp = SE(bldLst, qn('p:bldP'))
            bp.set('spid', str(spid))
            bp.set('grpId', '0')
            bp.set('animBg', '1')
            if spid in para_build_spids:
                bp.set('build', 'p')

# ═══════════════════════════════════════════════════════════════════════
# REUSABLE SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════

def _finish(slide, num, notes, anim_groups=None, click_reveal=True,
            para_build=False):
    """Add logo, slide number, transition, animations, and notes.

    num: ignored for int values (auto-numbered); string values (e.g. "A1")
         are used directly for appendix slides.  Use None to skip numbering
         entirely (section dividers).
    click_reveal: if True, first animation group visible on entry, rest
        appear one-by-one on click.  Default True.
    para_build: if True, multi-paragraph text shapes (bullet lists) reveal
        one paragraph per click.  Default True.
    """
    if num is None:
        display_num = ""  # Section dividers — no number
    elif isinstance(num, int):
        _auto_num[0] += 1
        display_num = _auto_num[0]
    else:
        display_num = num  # Appendix labels like "A1"
    add_logo(slide)
    add_slide_num(slide, display_num)
    add_fade_transition(slide)
    if anim_groups:
        try:
            add_animations(slide, anim_groups, click_reveal=click_reveal,
                           para_build=para_build)
        except Exception:
            pass  # Graceful fallback — slides still render
    set_notes(slide, notes)
    return slide


def build_split(prs, num, title, image_key, notes, big_num=None,
                num_color=TEAL, num_label=None, bullets=None,
                bottom_text=None):
    """Split layout: text/numbers left, image right."""
    slide = new_slide(prs)
    t = add_title(slide, title)
    add_accent_line(slide)

    left_shapes = []
    top = CT

    if big_num:
        s = add_text(slide, big_num, MX, top, SLW, Inches(1.1),
                     size=Pt(72), color=num_color, bold=True)
        left_shapes.append(s)
        top += Inches(1.1)
        if num_label:
            s2 = add_text(slide, num_label, MX, top, SLW, Inches(0.5),
                          size=Pt(15), color=LGRAY)
            left_shapes.append(s2)
            top += Inches(0.55)

    if bullets:
        # Cap bullet height so it doesn't overlap bottom_text at y=6.4
        max_h = Inches(6.1) - top if bottom_text else CH - (top - CT) - Inches(0.3)
        s = add_bullets(slide, bullets, MX, top, SLW, max_h)
        left_shapes.append(s)

    if bottom_text:
        s = add_text(slide, bottom_text, MX, Inches(6.45), CW, Inches(0.5),
                     size=Pt(13), color=LGRAY, italic=True)
        left_shapes.append(s)

    img = add_image(slide, image_key, SRL, CT, width=SRW)
    _finish(slide, num, notes, [left_shapes, [img]], click_reveal=True)
    return slide


def build_bullets(prs, num, title, items, notes, subtitle=None):
    """Full-width bullet list slide."""
    slide = new_slide(prs)
    add_title(slide, title)
    add_accent_line(slide)

    top = CT
    if subtitle:
        s = add_text(slide, subtitle, MX, top, CW, Inches(0.4),
                     size=Pt(16), color=LGRAY, italic=True)
        top += Inches(0.5)

    bul = add_bullets(slide, items, MX, top, CW, CH - (top - CT))
    _finish(slide, num, notes, [[bul]], click_reveal=True)
    return slide


def build_two_col(prs, num, title, left_title, left_items, right_title,
                  right_items, notes, left_color=TEAL, right_color=CORAL):
    """Two-column text layout."""
    slide = new_slide(prs)
    add_title(slide, title)
    add_accent_line(slide)

    col_w = Inches(5.5)
    gap = Inches(1.13)

    # Left column
    lt = add_text(slide, left_title, MX, CT, col_w, Inches(0.4),
                  size=Pt(18), color=left_color, bold=True)
    lb = add_bullets(slide, left_items, MX, CT + Inches(0.5), col_w,
                     CH - Inches(0.5))

    # Right column
    rx = MX + col_w + gap
    rt = add_text(slide, right_title, rx, CT, col_w, Inches(0.4),
                  size=Pt(18), color=right_color, bold=True)
    rb = add_bullets(slide, right_items, rx, CT + Inches(0.5), col_w,
                     CH - Inches(0.5))

    _finish(slide, num, notes, [[lt, lb], [rt], [rb]], click_reveal=True)
    return slide


def build_full_image(prs, num, title, image_key, notes, subtitle=None,
                     bottom_text=None):
    """Full-width image with optional text."""
    slide = new_slide(prs)
    add_title(slide, title)
    add_accent_line(slide)

    top = CT
    if subtitle:
        add_text(slide, subtitle, MX, top, CW, Inches(0.35),
                 size=Pt(16), color=LGRAY, italic=True)
        top += Inches(0.4)

    # Cap image height to avoid overlapping bottom text or going off-screen
    if bottom_text:
        img_h = Inches(4.5)
    else:
        img_h = SL_H - top - Inches(0.2)  # stay within slide bounds
    img = add_image(slide, image_key, MX, top, width=CW, height=img_h)

    if bottom_text:
        add_text(slide, bottom_text, MX, Inches(6.3), CW, Inches(0.5),
                 size=Pt(13), color=LGRAY, italic=True)

    _finish(slide, num, notes, [[img]])
    return slide
