"""
Slide builders — Section 6 + Salvage + Demos
"""

from pathlib import Path
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from .config import (
    IMG, VID, POSTER_DIR,
    SL_W, SL_H, BG, WHITE, TEAL, CORAL, LGRAY, MGRAY, DGRAY,
    GREEN, YELLOW, GOLD, ORANGE, RED, DRED, NAVY2, NAVY3,
    BLUE, PURPLE,
    FONT, _auto_num,
    MX, MY, CT, CW, CH, SLW, SRG, SRL, SRW,
)
from .helpers import (
    new_slide, set_notes, add_logo, add_slide_num, add_accent_line,
    _fmt, add_title, add_text, add_rich_text, add_bullets,
    add_rect, add_image, add_play_button, add_video_poster, add_video,
    add_table, _shade_cell, _rgb_hex,
    add_fade_transition, add_animations, _finish,
    build_split, build_bullets, build_two_col, build_full_image,
)


# Demo slides + slide_15 demo intro (split from slides_evaluation.py May 2026)

def slide_15(prs):
    """Demo slide.

    BLOCKER fix (May 2026): pre-fix audit found orphan animation refs to
    spids ['4','7','10'] (pptx_visual_audit.json, slide 48). Cause: a
    prior version called add_animations directly with stale shape ids.
    Fixed by collecting actual shape handles into per-video click-reveal
    groups passed to _finish - the animation tree now references only
    the shape IDs that exist on this slide. Body font also bumped from
    11pt to 12pt to clear the readability floor.

    audit:bigfonts \u2014 desc text bumped to 18pt; description strings
    rewritten as 2-line summaries to fit. Originals (3-line, included
    \"meaning close, key verb flipped\" etc.) preserved in speaker notes.
    """
    slide = new_slide(prs)
    add_title(slide, "Demo: OK \u2192 Almost There \u2192 Hallucination")
    add_accent_line(slide)

    # Three embedded videos side by side - click each to play
    # VID dict mapping (confirmed correct):
    #   "smartphone"   -> ktMebjnZiSE_3  (consumers want bigger smartphone, IS 4.1)
    #   "street_photo" -> 2HddWQse8Mw_0  (street photography topic right, names lost, IS 2.9)
    #   "halluc"       -> 00MUdHQ7GGY_8  (carry strap -> holocaust denier)
    vid_w = Inches(3.6)
    vid_h = Inches(2.4)  # audit:bigfonts \u2014 was 2.7; shrunk 0.3" for caption room
    gap = Inches(0.4)
    total = 3 * vid_w + 2 * gap
    start_x = (SL_W - total) / 2
    vid_y = CT + Inches(0.1)

    # audit:bigfonts \u2014 descs trimmed to 2 lines so 18pt fits in h=1.05.
    # Leftmost OK tile uses obama_perfect (segment #14, WER 0%, IS 5.00,
    # 27/29 words at conf-high) \u2014 most convincing per-word coloring in the deck.
    vids = [
        ("obama_perfect", '"\u2026the tireless and heroic work of our military"\n\u2192 (perfect \u2014 27/29 words green)', "WER 0%  IS 5.00", TEAL),
        ("street_photo", '"james and will talk about street photography"\n\u2192 "i\'m here to talk about street photography"', "WER 56%  IS 2.9", CORAL),
        ("halluc", '"carry strap"\n\u2192 "holocaust denier"', "WER 100%  IS 0.8", RED),
    ]

    # audit:pptx_visual_audit_after_amosi.md slide 63 BLOCKER -
    # animation referenced movie shape ids ['7','10'] which are wrapped in
    # <mc:AlternateContent>. python-pptx's slide.shapes iterator and the
    # audit's iter_shapes_recursive both skip AlternateContent, so the
    # movies look "non-existent" even though they render fine. Fix: drop
    # video shapes from anim_groups; only animate the captions. Videos
    # still render (they're not animated, so they appear immediately).
    anim_groups = []
    for i, (key, desc, wer, color) in enumerate(vids):
        x = start_x + i * (vid_w + gap)
        add_video(slide, key, x, vid_y, vid_w, vid_h)
        wer_t = add_text(slide, wer, x, vid_y + vid_h + Inches(0.05), vid_w,
                 Inches(0.4), size=Pt(24), color=color, bold=True,
                 align=PP_ALIGN.CENTER)
        # audit:bigfonts — desc h bumped 0.7->1.30 for 18pt 2-line summary.
        desc_t = add_text(slide, desc, x, vid_y + vid_h + Inches(0.45), vid_w,
                 Inches(2.2), size=Pt(24), color=LGRAY,
                 align=PP_ALIGN.CENTER)
        anim_groups.append([wer_t, desc_t])

    # audit:bigfonts2 — foot moved 6.85 -> 6.55 to clear slide-num zone (7.12).
    foot = add_text(slide, "Click each video to play.",
             MX, Inches(6.55), CW, Inches(0.30),
             size=Pt(18), color=MGRAY, italic=True, align=PP_ALIGN.CENTER)
    anim_groups.append([foot])

    _finish(slide, 15,
        "Three demos side by side, picked to span the quality range. "
        "LEFT (Obama segment 14, 29 words, the most convincing per-word "
        "coloring in the deck): WER 0%, IS 5.00 — REF and HYP identical "
        "('…and our allies over the last 10 years thanks to the tireless "
        "and heroic work of our military and our counterterrorism "
        "professionals we've made great strides in that effort'). 27 of "
        "29 per-word bands GREEN at conf-high; the remaining 2 yellow. "
        "Sentence_conf 0.72 (Salvage band) — the *segment-level* "
        "confidence is conservative because of two yellow words, but the "
        "WER says zero, IS says perfect. This is exactly the kind of "
        "example that motivates the joint conf+agreement rule explained "
        "earlier in §4. "
        "CENTER ('james and will talk about street photography' → "
        "'i'm here to talk about street photography', IS 2.9): topic "
        "captured perfectly but speaker names lost — the near-miss zone. "
        "RIGHT ('carry strap' → 'holocaust denier', IS 0.8): "
        "hallucination, fluent but completely fabricated. Click each "
        "video to play. Mention to peers: this is the qualitative "
        "bridge from the confidence section into the demo segments that "
        "follow. "
        "Sources: docs/evaluation/intelligibility_methodology.md, "
        "docs/evaluation/llm_judge/llm_judge_analysis.md.",
        anim_groups, click_reveal=True)


def slide_15_perfect(prs):
    """Demo intro 1/2: the success case — Obama 29-word perfect transcription."""
    slide = new_slide(prs)
    add_title(slide, "Demo: When It Works (1/2)")
    add_accent_line(slide)

    sub = add_text(slide,
        "Clean speech, model nails 29/29 words — the upper-bound demo.",
        MX, CT, CW, Inches(0.45),
        size=Pt(22), color=LGRAY, italic=True)

    # Single video, much larger
    vid_w = Inches(8.0)
    vid_h = Inches(4.0)
    vid_x = (SL_W - vid_w) // 2
    vid_y = CT + Inches(0.6)
    add_video(slide, "obama_perfect", vid_x, vid_y, vid_w, vid_h)

    # Big metrics under video
    metrics = add_text(slide,
        "WER 0%   /   IS 5.00   /   27 of 29 words green",
        MX, vid_y + vid_h + Inches(0.15), CW, Inches(0.55),
        size=Pt(28), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Per-word breakdown caption — y trimmed 0.75→0.55 + h 0.50→0.40 to stay
    # under safe zone (was bot=7.30, now bot=6.95). vid_y+vid_h=6.05.
    caption = add_text(slide,
        '"…the tireless and heroic work of our military and our counterterrorism…"',
        MX, vid_y + vid_h + Inches(0.55), CW, Inches(0.40),
        size=Pt(20), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

    foot = add_text(slide,
        "Click video to play.",
        MX, Inches(6.65), CW, Inches(0.30),
        size=Pt(18), color=MGRAY, italic=True, align=PP_ALIGN.CENTER)

    _finish(slide, 0,
        "Demo intro 1/2 (#339 split): the success case. Obama segment 14, "
        "29 words, the most convincing per-word coloring in the deck. "
        "WER 0%, IS 5.00 — REF and HYP identical: '…and our allies over the "
        "last 10 years thanks to the tireless and heroic work of our "
        "military and our counterterrorism professionals we've made great "
        "strides in that effort'. 27 of 29 per-word bands GREEN at conf-high; "
        "the remaining 2 yellow. Sentence_conf 0.72 (Salvage band) — the "
        "*segment-level* confidence is conservative because of two yellow "
        "words, but WER says zero, IS says perfect. This motivates the joint "
        "conf+agreement rule explained in §4. "
        "Sources: docs/evaluation/intelligibility_methodology.md.",
        [[sub], [metrics, caption], [foot]], click_reveal=True)


def slide_15_failures(prs):
    """Demo intro 2/2: the failure spectrum — partial vs hallucination."""
    slide = new_slide(prs)
    add_title(slide, "Demo: When It Fails (2/2)")
    add_accent_line(slide)

    sub = add_text(slide,
        "Topic captured (CENTER) vs total fabrication (RIGHT) — the failure spectrum.",
        MX, CT, CW, Inches(0.45),
        size=Pt(22), color=LGRAY, italic=True)

    # Two videos side by side, much larger than the 3-video version
    vid_w = Inches(5.5)
    vid_h = Inches(3.5)
    gap = Inches(0.4)
    total = 2 * vid_w + gap
    start_x = (SL_W - total) // 2
    vid_y = CT + Inches(0.6)

    vids = [
        ("street_photo", '"james and will talk about street photography"\n→ "i\'m here to talk about street photography"',
         "WER 56%  /  IS 2.9  /  topic right, names lost", CORAL),
        ("halluc", '"carry strap"\n→ "holocaust denier and a computer hacker"',
         "WER 100%  /  IS 0.8  /  fluent fabrication", RED),
    ]

    anim_groups = []
    for i, (key, desc, wer, color) in enumerate(vids):
        x = start_x + i * (vid_w + gap)
        add_video(slide, key, x, vid_y, vid_w, vid_h)
        wer_t = add_text(slide, wer, x, vid_y + vid_h + Inches(0.10),
                         vid_w, Inches(0.45),
                         size=Pt(20), color=color, bold=True,
                         align=PP_ALIGN.CENTER)
        # desc h 1.40→0.85: was overflowing slide bottom (bot=7.55 > 7.5).
        # Now bot = 5.55+0.60+0.85 = 7.00, within safe zone.
        desc_t = add_text(slide, desc, x, vid_y + vid_h + Inches(0.60),
                          vid_w, Inches(0.85),
                          size=Pt(18), color=LGRAY, italic=True,
                          align=PP_ALIGN.CENTER)
        anim_groups.append([wer_t, desc_t])

    foot = add_text(slide,
        "Click each video to play.",
        MX, Inches(6.65), CW, Inches(0.30),
        size=Pt(18), color=MGRAY, italic=True, align=PP_ALIGN.CENTER)
    anim_groups.append([foot])

    _finish(slide, 0,
        "Demo intro 2/2 (#339 split): the failure spectrum. "
        "LEFT ('james and will talk about street photography' → "
        "'i'm here to talk about street photography', IS 2.9): topic captured "
        "perfectly but speaker names lost — the near-miss zone. "
        "RIGHT ('carry strap' → 'holocaust denier', IS 0.8): hallucination, "
        "fluent but completely fabricated — the worst-case failure mode. "
        "Together the two clips bracket the IS scale around the production "
        "Salvage tier (1.0–2.99). Click each video to play. "
        "Sources: docs/evaluation/intelligibility_methodology.md, "
        "docs/evaluation/llm_judge/llm_judge_analysis.md.",
        [[sub]] + anim_groups, click_reveal=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — IS VALIDATION: CLAUDE-AS-JUDGE
# ═══════════════════════════════════════════════════════════════════════

def _demo_research_slide(prs, *, title, video_key, ref, hyp_runs,
                         metrics_line, badge_text, badge_color, body, notes):
    """Shared layout for the five research-flavored demo slides.

    Centered hero video + REF (label/text) + colour-coded HYP + metrics line +
    research observation body. Per-word band cited as observation, not pitch.
    audit:bigfonts — REF/HYP bumped to 18pt, body observation 16pt; video
    shrunk slightly (6.0->5.4 wide, 3.4->2.9 high) to free room for big-font
    REF + HYP + body all fitting between y=5.0 and y=7.30. Body observation
    intentionally trimmed to <=2 lines so 16pt fits in a single textbox.
    """
    slide = new_slide(prs)
    add_title(slide, title)
    add_accent_line(slide)

    badge_w = Inches(2.4)
    # Metrics-line caption: Pt(16)→Pt(18) to clear the italic-caption floor.
    sub = add_text(slide, metrics_line,
             MX, Inches(1.5), CW - badge_w - Inches(0.2), Inches(0.4),
             size=Pt(18), color=LGRAY, italic=True)
    badge_x = MX + CW - badge_w
    badge_box = add_rect(slide, badge_x, Inches(1.5), badge_w, Inches(0.4),
             fill_color=NAVY3, border_color=badge_color, border_width=Pt(1.0))
    # CUT v3 (overflow): 0.3 -> 0.45 + Pt(24) -> Pt(20) so badge text fits.
    badge_t = add_text(slide, badge_text,
             badge_x, Inches(1.50), badge_w, Inches(0.45),
             size=Pt(20), bold=True, color=badge_color, align=PP_ALIGN.CENTER)

    # Shrunk video slightly (5.4x2.6) to free more vertical space for REF/HYP.
    vid_w = Inches(5.4)
    vid_h = Inches(2.6)
    vid_x = (SL_W - vid_w) // 2
    vid_y = Inches(2.00)
    vid = add_video(slide, video_key, vid_x, vid_y, vid_w, vid_h)

    # REF block — starts just below video (4.65").
    add_text(slide, "REFERENCE",
             MX, Inches(4.68), CW, Inches(0.28),
             size=Pt(22), bold=True, color=LGRAY)
    ref_t = add_text(slide, ref,
             MX, Inches(4.96), CW, Inches(0.48),
             size=Pt(18), color=LGRAY, italic=True)

    # HYPOTHESIS block.
    add_text(slide, "HYPOTHESIS  (per-word confidence)",
             MX, Inches(5.52), CW, Inches(0.28),
             size=Pt(20), bold=True, color=WHITE)
    # remark #342 fix: shrunk hyp_t h from 1.20 to 0.95 so bottom (5.80+0.95=6.75)
    # no longer overlaps body_t (was at 6.85). body_t moved to 6.78 with h=0.27
    # so its bottom (7.05) lands exactly on the safe-zone floor.
    hyp_t = add_rich_text(slide, [hyp_runs],
             MX, Inches(5.80), CW, Inches(0.95))


    # Body observation (caption, sits below HYP at 6.75; ends at 7.05 safe-zone).
    body_t = add_text(slide, body,
             MX, Inches(6.78), CW, Inches(0.27),
             size=Pt(18), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)
    _finish(slide, 0, notes,
            [[sub, badge_box, badge_t], [ref_t], [hyp_t], [body_t]],
            click_reveal=True)


def _demo_card_slide_with_ref_runs(prs, *, title, ref_runs, hyp_runs,
                                   wer, is_score, mean_prob, min_word_conf,
                                   metric_color, sub_caption,
                                   ref_label_extra, hyp_label_extra,
                                   badge_text, badge_color,
                                   body_lead, body_lead_color, body_tail,
                                   notes=""):
    """Variant of _demo_card_slide accepting rich-text reference + per-label suffixes.

    Used for INSPECT/SALVAGE slides where the reference highlights domain
    vocabulary in bold.
    """
    slide = new_slide(prs)
    add_title(slide, title)
    add_accent_line(slide)

    # Metrics row + tier badge (same as _demo_card_slide)
    metrics_runs = [
        ("WER ", {"size": Pt(20), "color": LGRAY}),
        (f"{wer}", {"size": Pt(20), "color": metric_color, "bold": True}),
        ("   /   IS ", {"size": Pt(20), "color": LGRAY}),
        (f"{is_score}", {"size": Pt(20), "color": metric_color, "bold": True}),
        ("   /   mean_prob ", {"size": Pt(20), "color": LGRAY}),
        (f"{mean_prob}", {"size": Pt(20), "color": metric_color, "bold": True}),
        ("   /   min word conf ", {"size": Pt(20), "color": LGRAY}),
        (f"{min_word_conf}", {"size": Pt(20), "color": metric_color, "bold": True}),
    ]
    badge_w = Inches(2.6)
    metrics_t = add_rich_text(slide, [metrics_runs],
             MX, Inches(1.40), CW - badge_w - Inches(0.2), Inches(0.45))

    badge_x = MX + CW - badge_w
    badge_box = add_rect(slide, badge_x, Inches(1.40), badge_w, Inches(0.55),
             fill_color=badge_color, border_color=badge_color,
             border_width=Pt(1.0), corner_radius=True)
    badge_t = add_text(slide, badge_text,
             badge_x, Inches(1.42), badge_w, Inches(0.50),
             size=Pt(20), bold=True, color=NAVY3, align=PP_ALIGN.CENTER)

    add_text(slide, sub_caption,
             MX, Inches(1.95), CW, Inches(0.30),
             size=Pt(13), color=PURPLE, italic=True)

    # REFERENCE
    add_text(slide, f"REFERENCE{ref_label_extra}",
             MX, Inches(2.40), CW, Inches(0.30),
             size=Pt(16), bold=True, color=LGRAY)
    ref_box = add_rect(slide, MX, Inches(2.72), CW, Inches(1.55),
             fill_color=NAVY2, border_color=NAVY3,
             border_width=Pt(0.5), corner_radius=True)
    ref_t = add_rich_text(slide, [ref_runs],
             MX + Inches(0.25), Inches(2.84), CW - Inches(0.5), Inches(1.35))

    # HYPOTHESIS
    add_text(slide, f"HYPOTHESIS{hyp_label_extra}",
             MX, Inches(4.45), CW, Inches(0.30),
             size=Pt(16), bold=True, color=WHITE)
    hyp_box = add_rect(slide, MX, Inches(4.77), CW, Inches(1.50),
             fill_color=NAVY2, border_color=badge_color,
             border_width=Pt(1.5), corner_radius=True)
    hyp_t = add_rich_text(slide, [hyp_runs],
             MX + Inches(0.30), Inches(4.90), CW - Inches(0.6), Inches(1.30))

    # Body annotation
    body_runs = [
        (body_lead, {"size": Pt(16), "color": body_lead_color, "italic": True}),
        (body_tail, {"size": Pt(16), "color": PURPLE, "italic": True}),
    ]
    body_t = add_rich_text(slide, [body_runs],
             MX, Inches(6.40), CW, Inches(0.60))

    if notes:
        set_notes(slide, notes)


def _demo_card_slide(prs, *, title, ref, hyp_runs,
                     wer, is_score, mean_prob, min_word_conf,
                     metric_color, sub_caption,
                     badge_text, badge_color,
                     body_lead, body_lead_color, body_tail, notes):
    """Card-based demo layout (no video).

    Used for the TRUST and INSPECT demo slides where the source deck shows
    big rounded REF / HYP boxes and a solid-fill TIER badge instead of a
    video hero. Layout mirrors source slides 59 and 60 in
    Argos_VSP_v13_Amosi_2.pptx.
    """
    slide = new_slide(prs)
    add_title(slide, title)
    add_accent_line(slide)

    # ── Metrics row (top, color-highlighted values) ──────────────────
    # Build a single rich-text line with values colored by metric_color.
    metrics_runs = [
        ("WER ", {"size": Pt(20), "color": LGRAY}),
        (f"{wer}", {"size": Pt(20), "color": metric_color, "bold": True}),
        ("   /   IS ", {"size": Pt(20), "color": LGRAY}),
        (f"{is_score}", {"size": Pt(20), "color": metric_color, "bold": True}),
        ("   /   mean_prob ", {"size": Pt(20), "color": LGRAY}),
        (f"{mean_prob}", {"size": Pt(20), "color": metric_color, "bold": True}),
        ("   /   min word conf ", {"size": Pt(20), "color": LGRAY}),
        (f"{min_word_conf}", {"size": Pt(20), "color": metric_color, "bold": True}),
    ]
    badge_w = Inches(2.4)
    metrics_t = add_rich_text(slide, [metrics_runs],
             MX, Inches(1.40), CW - badge_w - Inches(0.2), Inches(0.45))

    # ── TIER badge (solid fill on the right) ────────────────────────
    badge_x = MX + CW - badge_w
    badge_box = add_rect(slide, badge_x, Inches(1.40), badge_w, Inches(0.55),
             fill_color=badge_color, border_color=badge_color,
             border_width=Pt(1.0), corner_radius=True)
    badge_t = add_text(slide, badge_text,
             badge_x, Inches(1.42), badge_w, Inches(0.50),
             size=Pt(20), bold=True, color=NAVY3, align=PP_ALIGN.CENTER)

    # ── Sub-caption (joint band rule etc.) ───────────────────────────
    sub_t = add_text(slide, sub_caption,
             MX, Inches(1.95), CW, Inches(0.30),
             size=Pt(14), color=MGRAY, italic=True)

    # ── REFERENCE label + rounded box ────────────────────────────────
    add_text(slide, "REFERENCE",
             MX, Inches(2.45), CW, Inches(0.30),
             size=Pt(16), bold=True, color=LGRAY)

    ref_box = add_rect(slide, MX, Inches(2.78), CW, Inches(1.10),
             fill_color=NAVY2, border_color=NAVY3,
             border_width=Pt(0.5), corner_radius=True)
    ref_t = add_text(slide, ref,
             MX + Inches(0.25), Inches(2.88), CW - Inches(0.5), Inches(0.95),
             size=Pt(20), color=WHITE)

    # ── HYPOTHESIS label + rounded box (color-coded text) ───────────
    add_text(slide, "HYPOTHESIS",
             MX, Inches(4.10), Inches(2.0), Inches(0.30),
             size=Pt(16), bold=True, color=WHITE)
    add_text(slide, " (per-word confidence — joint rule)",
             MX + Inches(1.55), Inches(4.10), CW - Inches(1.6), Inches(0.30),
             size=Pt(14), color=MGRAY, italic=True)

    hyp_box = add_rect(slide, MX, Inches(4.43), CW, Inches(1.40),
             fill_color=NAVY2, border_color=badge_color,
             border_width=Pt(1.5), corner_radius=True)
    hyp_t = add_rich_text(slide, [hyp_runs],
             MX + Inches(0.30), Inches(4.60), CW - Inches(0.6), Inches(1.15))

    # ── Body annotation (two-part: colored lead + grey tail) ────────
    body_runs = [
        (body_lead, {"size": Pt(18), "color": body_lead_color, "bold": True}),
        ("  " + body_tail, {"size": Pt(18), "color": LGRAY, "italic": True}),
    ]
    body_t = add_rich_text(slide, [body_runs],
             MX, Inches(6.10), CW, Inches(0.50))

    _finish(slide, 0, notes,
            [[metrics_t, badge_box, badge_t, sub_t],
             [ref_box, ref_t], [hyp_box, hyp_t], [body_t]],
            click_reveal=True)


def slide_live_example_intro(prs):
    """Early example — Obama 14 (29 words, all green) — sets the bar.

    Placed near the start of the deck (after 'What is VSP?') to give the
    audience a concrete, visceral taste of what the system delivers
    BEFORE diving into pipeline + benchmarks + metrics. Same Obama clip
    used later in slide 60's demo intro for callback.
    """
    # Full transcript with per-word coloring (WER 0% — all 29 words match).
    runs = [
        ("and our allies over the last 10 years thanks to the ", {"size": Pt(18), "color": BLUE}),
        ("tireless and heroic ", {"size": Pt(18), "color": BLUE, "bold": True}),
        ("work of our military and our ", {"size": Pt(18), "color": BLUE}),
        ("counterterrorism", {"size": Pt(18), "color": GOLD, "italic": True}),
        (" professionals we've made great strides in that effort", {"size": Pt(18), "color": BLUE}),
    ]
    _demo_research_slide(prs,
        title="Live example — clean speech, perfect transcription",
        video_key="obama_perfect",
        ref="and our allies over the last 10 years thanks to the tireless "
            "and heroic work of our military and our counterterrorism "
            "professionals we've made great strides in that effort",
        hyp_runs=runs,
        metrics_line="WER 0.00%   /   IS 5.00 (Excellent)   /   29 words   "
                     "/   27 of 29 GREEN  (per-word conf)",
        badge_text="TIER: TRUST",
        badge_color=BLUE,
        body="Reference and prediction are identical (WER 0%). 27 of 29 "
             "per-word confidence bands GREEN — the model is sure, and "
             "it's right.",
        notes="Early-deck taste: this is the upper bound. Obama bin "
              "Laden announcement segment 14 (41.95-45.55 s, 29 words). "
              "REF and HYP are character-for-character identical: "
              "'and our allies over the last 10 years thanks to the "
              "tireless and heroic work of our military and our "
              "counterterrorism professionals we've made great strides "
              "in that effort'. WER = 0%, IS = 5.00 (Tier 5 Excellent), "
              "27 of 29 per-word confidence bands GREEN at high "
              "conf-only threshold (this Obama decode predates "
              "VSP_NBEST=1 so the joint conf+agreement rule is not "
              "applied — but the visual is the same, and at 27/29 green "
              "the upgrade would not change the picture). Why this "
              "slide opens the deck: research peers see the punchline "
              "first — 'yes, the system can deliver perfect output on "
              "real-world video' — before the metric debate. This "
              "primes the rest of §1 ('but is it always like this?' → "
              "no, average is 62% NIV-Y+P, hence the metric work in "
              "§2). Same clip recurs as the leftmost tile in slide 60's "
              "demo intro, anchoring a callback. "
              "Source: flat_runs_archive/20260430_234843/"
              "client_outputs/report/report.csv "
              "(050111_OsamaBinLadenStatement_HD_14_004195_004555).")


def slide_failure_live_demo(prs):
    """Two live failure mode examples side by side: Partial Failure + Hallucination.

    LEFT: street_photo — topic captured, names lost (WER 56%, IS 2.91, NIV-Y+P).
    RIGHT: halluc — total hallucination, fluent fabrication (WER 100%, IS 0.81).
    Chosen to span the failure spectrum: "not that bad" vs "worst case".
    Column width 5.85" triggers the narrow-column audit exemption (18pt floor).
    """
    slide = new_slide(prs)
    add_title(slide, "Demo: Where Word Accuracy and Meaning Diverge")
    add_accent_line(slide)

    # Three columns side by side: smartphone (semantic flip), street_photo
    # (names lost), halluc (full hallucination).  Mirrors source slide 61.
    col_w = Inches(3.95)
    gap = Inches(0.20)
    total = 3 * col_w + 2 * gap
    start_x = (SL_W - total) / 2

    vid_h = Inches(2.40)
    vid_y = CT + Inches(0.05)

    tiles = [
        {
            "key": "smartphone",
            "label": "WER 28%   ·   IS 4.1",
            "caption_runs": [
                ("“consumers want a bigger smartphone”\n",
                 {"size": Pt(15), "color": WHITE, "italic": True}),
                ("→ misleading polarity flip:\n",
                 {"size": Pt(14), "color": ORANGE, "italic": True}),
                ("IS 4.1 says polarity flips won’t upgrade\n",
                 {"size": Pt(13), "color": LGRAY, "italic": True}),
                ("their smartphone\n\n",
                 {"size": Pt(13), "color": LGRAY, "italic": True}),
                ("Polarity-sensitive content requires human\n",
                 {"size": Pt(12), "color": MGRAY, "italic": True}),
                ("review regardless of IS.",
                 {"size": Pt(12), "color": MGRAY, "italic": True}),
            ],
        },
        {
            "key": "street_photo",
            "label": "WER 56%   ·   IS 2.9",
            "caption_runs": [
                ("“james and will talk about street\n",
                 {"size": Pt(15), "color": WHITE, "italic": True}),
                ("photography”\n",
                 {"size": Pt(15), "color": WHITE, "italic": True}),
                ("→ “i’m here to talk about street\n",
                 {"size": Pt(15), "color": ORANGE, "italic": True}),
                ("photography”\n\n",
                 {"size": Pt(15), "color": ORANGE, "italic": True}),
                ("Topic right, speaker names lost (IS 2.9).",
                 {"size": Pt(13), "color": MGRAY, "italic": True}),
            ],
        },
        {
            "key": "halluc",
            "label": "WER 100%   ·   IS 0.1",
            "caption_runs": [
                ("“carry strap”\n→ “holoc"
                 "aust denier”",
                 {"size": Pt(15), "color": RED, "italic": True}),
            ],
        },
    ]

    anim_groups = []
    for i, t in enumerate(tiles):
        x = start_x + i * (col_w + gap)

        add_video(slide, t["key"], x, vid_y, col_w, vid_h)

        badge = add_text(slide, t["label"],
                x, vid_y + vid_h + Inches(0.05), col_w, Inches(0.30),
                size=Pt(14), color=MGRAY, bold=True, align=PP_ALIGN.CENTER)

        cap_t = add_rich_text(slide, [t["caption_runs"]],
                x, vid_y + vid_h + Inches(0.40), col_w, Inches(2.65))

        anim_groups.append([badge, cap_t])

    _finish(slide, 0,
        "Two live failure examples chosen to span the spectrum: 'not that bad' "
        "vs worst case. "
        "LEFT — Right Topic, Names Lost (street photography): segment "
        "2HddWQse8Mw_0__8ecb0409_00_000000_000072. "
        "Reference: 'james and will talk about street photography and other'. "
        "Hypothesis: 'i'm here to talk about street photography and all the'. "
        "WER 55.6%, IS 2.91 (Tier 3 Fair, NIV-Y+P), mean_prob 0.602 (Strip band). "
        "Teaching point: IS 2.91 says this is USEFUL — the topic 'street "
        "photography' survived, only the speaker names (james, will) were lost. "
        "A viewer who needs to know the subject gets something real from this. "
        "The confidence (Strip, mean_prob 0.602) flagged uncertainty correctly — "
        "the model didn't know whose talk it was, and it showed that. "
        "RIGHT — Hallucination (carry strap): segment "
        "00MUdHQ7GGY_8__b1480c7a_00_000000_000194. "
        "Reference: 'it doesn't have a carry strap but you can put it on your "
        "shoulder pretty easily or just carry it with your hand'. "
        "Hypothesis: 'this is david irving he's a holocaust denier and a "
        "computer hacker who broke into the nuremberg trials'. "
        "WER 100%, IS 0.81 (Tier 1 Failed), mean_prob 0.468 (Strip band). "
        "This is the worst failure mode: the LLM generated a completely unrelated, "
        "fluent, harmful-sounding sentence. Length ratio >> 1 (hallucination "
        "pattern). Strip auto-flags it. "
        "Narrative: use LEFT to show 'sometimes it fails gracefully — useful "
        "content comes through, only specifics are lost'; use RIGHT to show "
        "'sometimes it fails catastrophically — the system flags these.' "
        "Sources: english_full_nbest_eval/report_v2/report.csv.",
        anim_groups, click_reveal=True)


def slide_demo_obama_trust(prs):
    """TRUST exemplar — non-trivial speaker, AI talk, full agreement-aware rule.

    Replaced Obama segment 14 (prepared political statement, too easy a TRUST
    example) with a non-Obama IS=5.00 / WER=0% segment: South-Asian speaker
    with Indian English accent talking about "this wave of artificial
    intelligence". Same TRUST tier (mean_prob 0.988), but a harder visual
    target — proves the model handles diverse speakers, not just rehearsed
    political video. VSP_NBEST=1 sidecar available, so the joint
    conf+agreement rule applied.

    Card layout (no video) to match the source deck's TRUST tier slide.
    """
    # Per-word colors all GREEN (TRUST band) — every word high-conf + high-agreement.
    runs = [
        ("to this wave of ", {"size": Pt(22), "color": GREEN, "bold": True}),
        ("artificial intelligence ", {"size": Pt(22), "color": GREEN, "bold": True}),
        ("that is slowly taking place", {"size": Pt(22), "color": GREEN, "bold": True}),
    ]
    _demo_card_slide(prs,
        title="Demo — TRUST tier: AI talk, Indian-English speaker",
        ref="to this wave of artificial intelligence that is slowly taking place",
        hyp_runs=runs,
        wer="0%", is_score="5.00", mean_prob="0.988", min_word_conf="0.94",
        metric_color=GREEN,
        sub_caption="VSP_NBEST = 1   ·   joint band rule  (p₁ ≥ 0.95  ∧  α ≥ 0.80)",
        badge_text="TIER:  TRUST",
        badge_color=GREEN,
        body_lead="Every word GREEN under the joint rule.",
        body_lead_color=GREEN,
        body_tail="Diverse speaker, Indian-English accent, perfect transcription — system at its best.",
        notes="Non-Obama TRUST exemplar (segment "
              "K0h33Ps7vz4_11__e66d3063_00_000000_000103, ~4s, IS=5.00 / "
              "WER=0% / mean_word_prob=0.988). South-Asian speaker with "
              "Indian English accent, talking about 'this wave of "
              "artificial intelligence'. Reference + hypothesis identical: "
              "'to this wave of artificial intelligence that is slowly "
              "taking place'. Why this is the right TRUST opener: a "
              "rehearsed political statement (Obama) is the easiest "
              "possible visual target — clear lighting, frontal face, "
              "minimal head movement. This clip carries higher visemic "
              "difficulty (different accent, different mouth shapes, "
              "casual gesture) and the model still hits IS=5. Joint "
              "conf+agreement rule applied (VSP_NBEST=1 sidecar available). "
              "NIV-Y green-band reliability for this population is 94% "
              "per the band-reliability stratification earlier in this "
              "section. Anchors the Trust → Salvage → Strip tour that "
              "follows on the next two slides. "
              "Sources: english_full_nbest_eval/report_v2/report.csv, "
              "docs/confidence/band_reliability_by_niv.md.")


def slide_demo_obama_salvage(prs):
    """Obama segment 31 - partial recovery (TRUST badge under conf-only)."""
    # audit:after_amosi_asset_fixes.md fix #9 - title previously claimed
    # "Salvage Tier (partial recovery)" but the rendered video shows the
    # TRUST badge: Obama seg 31 mean_prob = 0.920 which is above T_safe
    # (0.82), so the production tier is TRUST and not Salvage. The Obama
    # decode predates VSP_NBEST=1, so the joint conf+agreement rule is
    # not applied (no agreement sidecar) - the slide silently falls back
    # to conf-only band painting. Title now discloses the badge mismatch
    # and the metrics_line reports the correct mean_prob.
    runs = [
        ("...as president bush ", {"size": Pt(22), "color": BLUE}),
        ("said", {"size": Pt(22), "color": ORANGE, "bold": True}),
        (" america will never seek permission... (conf-only fallback)",
         {"size": Pt(22), "color": BLUE}),
    ]
    _demo_research_slide(prs,
        title="Demo — Obama: Trust Tier (conf-only fallback)",
        video_key="obama_partial",
        ref="(see speaker notes; Obama bin Laden announcement, segment #31, 92.90-96.50 s)",
        hyp_runs=runs,
        # OVERLAP fix (May 2026): metrics shortened to fit one 18pt line.
        metrics_line="WER ~22%  /  IS ~3.5  /  mean_prob 0.920  "
                     "(TRUST under T_safe=0.82)",
        badge_text="TIER: TRUST",
        badge_color=BLUE,
        # audit:bigfonts — body trimmed from 5 sentences to 2 to fit 16pt.
        # Cut: "The TRUST badge (mean_prob>=0.82) reflects the conf-only
        # fallback because this Obama decode predates VSP_NBEST=1; the joint
        # rule would likely demote the segment."  (kept in speaker notes.)
        body="Failed: 'president bush did' → 'said' (orange). Saved: LLM judge Y+P — "
             "meaning preserved. Note: LLM-as-Judge confirms what metrics miss.",
        notes="Obama bin Laden announcement, segment #31 (92.90-96.50 s). "
              "WER ~22%, IS ~3.5, mean_prob = 0.920 — above T_safe (0.82), "
              "so the production tier under the conf-only fallback is "
              "TRUST not Salvage. The original slide title called this "
              "'Salvage' for narrative continuity, but the badge mismatch "
              "is real: this Obama decode predates VSP_NBEST=1 (shipped "
              "April 30 2026) so no agreement sidecar exists, and the "
              "per-word band rule silently falls back to top1_conf only. "
              "Re-running stage 8 (outputs.sh::run_outputs) on a "
              "VSP_NBEST=1 decode would likely push this segment to "
              "Salvage under the joint rule because the substituted "
              "'said' almost certainly has low beam_agreement. The "
              "reviewer still triages off the orange word; this slide "
              "demonstrates the narrative-vs-production-tier gap that "
              "the joint rule closes. Mention to peers: this is the "
              "clearest worked example of why the agreement axis pulls "
              "high-conf-but-disagreed words out of green into yellow. "
              "Sources: docs/confidence/band_reliability_by_niv.md, "
              "docs/beam-search/n_best_implementation.md.")


def slide_demo_obama_strip(prs):
    """INSPECT-tier example — circadian/hormone segment, structure preserved.

    Replaced Obama segment 5 with a non-Obama INSPECT exemplar: the model
    keeps the "tells us when to X" repeating-clause structure but loses
    the domain vocabulary (cortisol/testosterone/light cycles → eat/turns/
    stops). mean_prob = 0.739 sits squarely in the Salvage band (0.65 ≤
    mean_prob < 0.82), and the clip was rendered with the full
    agreement-aware joint rule (VSP_NBEST=1 sidecar available).
    """
    # audit:after_amosi_narrative_actions.md fix #13 - first INSPECT
    # mention; gloss it inline so the audience knows it is the
    # production label for what the research literature calls 'Salvage'.
    # Per-word color sample: structure-preserved (BLUE) "tells us when to make"
    # vs domain-vocab replacement (ORANGE/PURPLE).
    # Per-word colors illustrating structure preserved (green) but
    # vocabulary substituted (purple/orange) under the joint rule.
    runs = [
        ("the job prescription takes into account ", {"size": Pt(18), "color": PURPLE, "bold": True}),
        ("our ", {"size": Pt(18), "color": GREEN, "bold": True}),
        ("environment tells us ", {"size": Pt(18), "color": GREEN, "bold": True}),
        ("what ", {"size": Pt(18), "color": ORANGE, "bold": True}),
        ("to ", {"size": Pt(18), "color": GREEN, "bold": True}),
        ("eat ", {"size": Pt(18), "color": PURPLE, "bold": True}),
        ("tells ", {"size": Pt(18), "color": GREEN, "bold": True}),
        ("us ", {"size": Pt(18), "color": ORANGE, "bold": True}),
        ("where ", {"size": Pt(18), "color": PURPLE, "bold": True}),
        ("to ", {"size": Pt(18), "color": ORANGE, "bold": True}),
        ("make ", {"size": Pt(18), "color": GREEN, "bold": True}),
        ("turns ", {"size": Pt(18), "color": PURPLE, "bold": True}),
        ("tells ", {"size": Pt(18), "color": GREEN, "bold": True}),
        ("us ", {"size": Pt(18), "color": ORANGE, "bold": True}),
        ("when ", {"size": Pt(18), "color": GREEN, "bold": True}),
        ("to make ", {"size": Pt(18), "color": GREEN, "bold": True}),
        ("stops ", {"size": Pt(18), "color": PURPLE, "bold": True}),
        ("basically switches on", {"size": Pt(18), "color": GREEN, "bold": True}),
    ]
    # Reference with domain-vocab words bolded.
    ref_runs = [
        ("us and couples us to light cycles in our environment tells us when to ",
         {"size": Pt(18), "color": WHITE}),
        ("sleep", {"size": Pt(18), "color": WHITE, "bold": True}),
        (" tells us when to make ", {"size": Pt(18), "color": WHITE}),
        ("cortisol", {"size": Pt(18), "color": WHITE, "bold": True}),
        (" tells us when to make\n", {"size": Pt(18), "color": WHITE}),
        ("testosterone", {"size": Pt(18), "color": WHITE, "bold": True}),
        (" basically switches on", {"size": Pt(18), "color": WHITE}),
    ]
    _demo_card_slide_with_ref_runs(prs,
        title="Demo — SALVAGE: structure preserved, vocabulary lost",
        ref_runs=ref_runs,
        hyp_runs=runs,
        wer="43%", is_score="2.66", mean_prob="0.739", min_word_conf="0.12",
        metric_color=GOLD,
        sub_caption="VSP_NBEST = 1   ·   joint band rule   ·   "
                    "purple flags substituted nouns: prescription (0.12), "
                    "turns (0.18), stops (0.23)",
        ref_label_extra="  (domain vocabulary highlighted)",
        hyp_label_extra="  (structure preserved — vocabulary substituted)",
        badge_text="TIER:  SALVAGE",
        badge_color=GOLD,
        body_lead="Repeating \"tells us when to X\" structure preserved (green/yellow); "
                  "domain nouns lost (",
        body_lead_color=WHITE,
        body_tail="sleep → eat, cortisol → turns, testosterone → stops). "
                  "Per-word confidence flagged the substituted nouns at p₁ ≤ 0.23 — "
                  "a reader is told exactly where to verify.",
        notes="Non-Obama INSPECT exemplar (segment "
              "9HanJOCw2Sc_11__19c7ec4e_00_000000_000261, ~13s). "
              "Reference: '… couples us to light cycles in our "
              "environment tells us when to sleep tells us when to make "
              "cortisol tells us when to make testosterone basically "
              "switches on'. Hypothesis: 'the job prescription takes "
              "into account our environment tells us what to eat tells "
              "us where to make turns tells us when to make stops "
              "basically switches on'. WER 43%, NEA F1 43% (lost "
              "entities: cortisol, couples, cycles, light, sleep, "
              "testosterone), IS 2.66 (Tier 3 Fair, NIV-Y+P). mean_prob "
              "= 0.739, sequence_conf in the Salvage band (0.65–0.82). "
              "min word conf = 0.12 on entity-replacement tokens "
              "(stops). Why this is a clean INSPECT: the "
              "circadian/hormones → behaviour-script semantic drift is "
              "a recoverable failure mode for a viewer who knows the "
              "topic, but the model still flagged the entity replacements "
              "via the joint conf+agreement rule. Per-word colors render "
              "from the agreement-aware sidecar (this decode had "
              "VSP_NBEST=1, unlike Obama segments). Sources: "
              "english_full_nbest_eval/report_v2/report.csv, "
              "docs/confidence/band_reliability_by_niv.md.")


def slide_demo_judge_entity(prs):
    """Judge entity slide - now shows STRIP badge for rogers / PV / will (research).

    Per render-log finding, judge_entity now shows STRIP badge under joint
    rule (rogers / pv / will all flagged red).
    """
    runs = [
        ("market ",       {"size": Pt(24), "color": BLUE}),
        ("research ",     {"size": Pt(24), "color": BLUE}),
        ("firm ",         {"size": Pt(24), "color": BLUE}),
        ("rogers ",       {"size": Pt(24), "color": PURPLE, "bold": True}),
        ("research ",     {"size": Pt(24), "color": BLUE}),
        ("is ",           {"size": Pt(24), "color": BLUE}),
        ("forecasting ",  {"size": Pt(24), "color": BLUE}),
        ("pv ",           {"size": Pt(24), "color": PURPLE, "bold": True}),
        ("installations ",{"size": Pt(24), "color": BLUE}),
        ("will ",         {"size": Pt(24), "color": PURPLE, "bold": True}),
        ("reach",         {"size": Pt(24), "color": BLUE}),
    ]
    # audit:after_amosi_asset_fixes.md fix #10 - mean_prob updated from
    # ~0.71 to 0.624 to match
    # english_full_nbest_eval/report_v2/report.csv (segment
    # 4D634qUi2BI_0__93a9f2b4_00_000000_000122, sentence_confidence=0.624,
    # is_label=Strip). Also fix #11 - explicit VSP_NBEST=1 disclosure
    # added (slide pulls per-word colours from the agreement-aware
    # sidecar, unlike the Obama slides).
    _demo_research_slide(prs,
        title="Demo - Strip: entity swap auto-flagged",
        video_key="judge_entity",
        ref="market research firm bernreuter research is forecasting "
            "pv installations could reach",
        hyp_runs=runs,
        # OVERLAP fix (May 2026): metrics shortened to fit one 18pt line.
        metrics_line="WER 18%  /  IS 4.55  /  mean_prob 0.624  (Strip)",
        badge_text="TIER: STRIP",
        badge_color=PURPLE,
        # audit:bigfonts — body trimmed (was: "Research observation: the
        # entity-swap tokens 'rogers', 'pv', 'will' are auto-flagged red
        # under the joint rule. Strengthens the entity-swap narrative.")
        body="Entity-swap tokens 'rogers', 'pv', 'will' auto-flagged red "
             "under the joint rule.",
        notes="Source: slides_client.py::slide_client_judge_ex1 (reframed "
              "for research). bernreuter -> rogers entity swap; PV / "
              "will also flagged red under the joint conf+agreement rule "
              "(per render-log inspection). WER 18%, IS 4.55 (Excellent), "
              "LLM judge Y. mean_prob = 0.624 per "
              "english_full_nbest_eval/report_v2/report.csv "
              "(prior copy said ~0.71, corrected per "
              "after_amosi_asset_fixes.md fix #10). The badge is STRIP "
              "not TRUST because the joint rule pushes the entity-swap "
              "tokens to red (per-word agreement low even though top-1 "
              "conf may be high). DECODE MODE: this clip was decoded "
              "with VSP_NBEST=1 so the agreement-aware sidecar is "
              "present and the full joint rule (top1_conf>=0.95 AND "
              "beam_agreement>=0.80) is applied (unlike the Obama clips "
              "earlier in this section, which fall back to conf-only "
              "band painting). audit:judge_v3_y_count_baseline picks "
              "this segment up; the per-word band overlay separates "
              "the firm-name swap from the rest.")


def slide_demo_judge_vocab(prs):
    """Judge router slide - technical-vocab drift (research framing)."""
    # audit:bigfonts — trimmed from 18 tokens to 11 to fit one line at 18pt
    # in the demo helper's HYP textbox. Cut the leading filler clause
    # ("we need a radically different approach we") — the band-painted swap
    # ("design"/"roads") is what the slide demonstrates.
    runs = [
        ("must ",       {"size": Pt(24), "color": ORANGE}),
        ("indeed ",     {"size": Pt(24), "color": ORANGE}),
        ("find ",       {"size": Pt(24), "color": BLUE}),
        ("a ",          {"size": Pt(24), "color": BLUE}),
        ("way ",        {"size": Pt(24), "color": BLUE}),
        ("we ",         {"size": Pt(24), "color": BLUE}),
        ("can ",        {"size": Pt(24), "color": BLUE}),
        ("design ",     {"size": Pt(24), "color": PURPLE}),
        ("existing ",   {"size": Pt(24), "color": BLUE}),
        ("roads ",      {"size": Pt(24), "color": PURPLE, "bold": True}),
        ("...",         {"size": Pt(24), "color": LGRAY}),
    ]
    # audit:after_amosi_asset_fixes.md fix #11 - explicit VSP_NBEST=1
    # disclosure added so the audience knows the per-word reds came from
    # the full joint conf+agreement rule (not the conf-only fallback
    # used by the Obama clips earlier in this section).
    _demo_research_slide(prs,
        title="Demo - Salvage: technical vocabulary drift",
        video_key="judge_router",
        ref="we need a radically different approach we basically need "
            "to find a way how we can take existing routers existing "
            "switches existing links and enable them for research",
        hyp_runs=runs,
        # OVERLAP fix (May 2026): metrics shortened to fit one 18pt line.
        metrics_line="WER 52%  /  IS 3.02  /  mean_prob ~0.78  (Salvage)",
        badge_text="TIER: SALVAGE",
        badge_color=ORANGE,
        # audit:bigfonts — body trimmed from 4 clauses to 2. Cut: "per-word
        # reds isolate the swaps; reviewer can patch the vocab" (in notes).
        body="Argument structure preserved (green spine). Domain terms "
             "drift: 'routers/switches/links' -> 'roads/structures/reuse'.",
        notes="Networking-research segment where the model preserved the "
              "argument structure but swapped the domain vocabulary "
              "(routers -> roads, switches -> structures, links -> reuse). "
              "WER 52%, IS 3.02 (Good tier), LLM judge P, mean_prob "
              "~0.78 (Salvage tier). DECODE MODE: this clip was decoded "
              "with VSP_NBEST=1 enabled so the agreement-aware sidecar is "
              "present and the full joint rule (top1_conf >= 0.95 AND "
              "beam_agreement >= 0.80) is applied — unlike the Obama "
              "clips earlier in this section, which fall back to "
              "conf-only band painting. Per-word colours show a green "
              "argument spine and red domain-vocabulary swaps. The "
              "per-word band isolates the exact tokens that need a "
              "domain-aware re-decode pass — exactly the case Mission 8 "
              "(topic-aware prompting) is designed to address. Mention "
              "to peers: this is the cleanest demonstration of "
              "per-tier reliability paying off in production. "
              "Sources: docs/confidence/band_reliability_by_niv.md, "
              "docs/evaluation/llm_judge/llm_judge_analysis.md.")
