#!/usr/bin/env python3
"""
Argos VSP — PPTX Presentation Generator

Creates the complete 80-slide "Argos VSP: Research Findings and Production
Roadmap" presentation with professional styling, real images, entrance
animations, fade transitions, and speaker notes.

Usage:
    python3 docs/_research-tools/generators/generate_presentation.py

Output:
    presentation_materials_20260224/Argos_VSP_Project_Review.pptx

Architecture:
    This file is a slim orchestrator. All logic lives in the presentation/ package:
        presentation/config.py           — paths, colors, layout constants
        presentation/helpers.py          — slide setup, text, shapes, animations
        presentation/slides_opening.py   — Sections 0-2: Opening, Context, Problem
        presentation/slides_research.py  — Sections 3-5: Research, Understanding, Tuning
        presentation/slides_evaluation.py— Section 6 + Salvage + Demos
        presentation/slides_engineering.py — Section 7: Engineering
        presentation/slides_future.py    — Section 8: Future + Appendix
"""

from pptx import Presentation

from presentation.config import SL_W, SL_H, OUTPUT, _auto_num
from presentation.helpers import _fix_pptx_video_compat

from presentation.slides_opening import (
    slide_01, slide_what_was_done_1, slide_what_was_done_2,
    slide_exec_summary, slide_wer_lies, slide_toc,
    slide_02, slide_visemes, slide_03, slide_data_flow,
    slide_04, slide_05, slide_06,
)
from presentation.slides_research import (
    slide_research_transition, slide_is_motivation, slide_is_intro,
    slide_is_weight_rationale,
    slide_is_calc_examples, slide_is_radar, slide_is_wer_scatter,
    slide_07, slide_metric_transition,
    slide_failure_deep_1a, slide_failure_deep_1b, slide_failure_deep_2, slide_08,
    slide_failure_deep_3,
)
from presentation.slides_evaluation import (
    slide_is_deep_dive, slide_metric_disagreement, slide_metric_disagreement_2,
    slide_two_eval_systems, slide_llm_judge, slide_context_eval,
    slide_llm_context_engine,
    slide_25, slide_25d, slide_25e,
    slide_14b, slide_15,
    slide_llm_judge_30, slide_judge_ex1, slide_judge_ex2,
    slide_judge_ex3, slide_judge_ex4, slide_judge_ex5, slide_judge_ex6,
    slide_disagreement_blind, slide_disagreement_context,
)
from presentation.slides_engineering import (
    slide_eng_transition, slide_three_repos,
    slide_17, slide_17_png, slide_18, slide_20,
    slide_web_ui, slide_21, slide_dual_env,
)
from presentation.slides_future import (
    slide_future_transition, slide_24,
    slide_26, slide_26b, slide_confidence_scoring,
    slide_27, slide_28, slide_data_scaling, slide_price_tag,
    slide_29, slide_30, slide_30b,
    slide_arabic_roadmap, slide_arabic_avhubert, slide_arabic_changes,
    slide_31, slide_thank_you,
    slide_a1, slide_a8, slide_a11, slide_a11b,
    slide_a13, slide_a15, slide_a16, slide_a17,
)


def _regenerate_plots():
    """Re-generate derived plot PNGs before building the PPTX."""
    import subprocess, sys
    scripts = [
        "docs/_research-tools/generators/generate_is_wer_scatter.py",
    ]
    for script in scripts:
        print(f"  Pre-generating: {script}")
        subprocess.run([sys.executable, script], check=True)


def main():
    _auto_num[0] = 0  # Reset auto-numbering

    _regenerate_plots()

    prs = Presentation()
    prs.slide_width = SL_W
    prs.slide_height = SL_H

    print("Generating presentation...")

    builders = [
        # --- Section 0: Opening (slides 1-6) ---
        slide_01,                    # 1: Title
        slide_what_was_done_1,       # 2: What was done? (1/2)
        slide_what_was_done_2,       # 3: What was done? (2/2)
        slide_exec_summary,          # 4: Executive Summary
        slide_wer_lies,              # 5: WER: The Metric That Lies
        slide_toc,                   # 6: Presentation Overview
        # --- Section 1: Context (slides 7-13) ---
        slide_02,                    # 7: What is VSP?
        slide_visemes,               # 8: Visemes
        slide_03,                    # 9: Three Components
        slide_data_flow,             # 10: Data Flow
        slide_04,                    # 11: Benchmark
        slide_05,                    # 12: Reality Gap
        slide_06,                    # 13: Same WER, Different Effects
        # --- Section 2: Research Findings (slides 14-34) ---
        slide_research_transition,   # 14: RESEARCH FINDINGS
        slide_llm_judge,             # 15: LLM-as-a-Judge (moved here)
        slide_llm_judge_30,          # 16: LLM Judge Deep Dive
        slide_judge_ex1,             # 17: Judge Example 1
        slide_judge_ex2,             # 18: Judge Example 2
        slide_judge_ex3,             # 19: Judge Example 3
        slide_judge_ex4,             # 20: Judge Example 4
        slide_judge_ex5,             # 21: Judge Example 5
        slide_judge_ex6,             # 22: Judge Example 6
        slide_is_motivation,         # 23: Why LLM Judge Not Enough
        slide_is_intro,              # 24-26: IS Signals (3 sub-slides)
        slide_is_weight_rationale,   # 27: Do 6 Signals Measure 6 Things?
        slide_is_calc_examples,      # 28: IS in Action
        slide_is_radar,              # 29: Model Comparison: IS Profiles
        slide_is_wer_scatter,        # 30: The Gap: Where WER Lies Most
        slide_07,                    # 31: IS Results: 61.6% Useful
        slide_two_eval_systems,      # 32: Two Evaluation Systems (moved here)
        slide_context_eval,          # 33: IS: Calibrated Surrogate Metric
        slide_disagreement_blind,    # 34: Where IS and Judge Disagree (blind)
        slide_disagreement_context,  # 35: Context Exposes Hidden Failures
        slide_metric_transition,     # 36: Three Numbers That Tell the Real Story
        # --- Section 3: Failure Analysis (slides 37-40) ---
        slide_08,                    # 37: Failure Mode Taxonomy (bar chart)
        slide_failure_deep_1a,       # 38: Failure Taxonomy 1/2
        slide_failure_deep_1b,       # 39: Failure Taxonomy 2/2
        slide_failure_deep_2,        # 40: Failure Modes: Real Examples
        # --- Section 4: Validation & Salvage (slides 41-48) ---
        slide_is_deep_dive,          # 41: IS Validation
        slide_metric_disagreement,   # 42: When Metrics Disagree pt 1 (unhidden)
        slide_metric_disagreement_2, # 43: When Metrics Disagree pt 2
        slide_25,                    # 44: IS Calibrated Surrogate for LLM
        slide_25d,                   # 43: LLM Salvage: Three Recoveries
        slide_25e,                   # 44: LLM Salvage: Domain Context
        slide_14b,                   # 45: Video Gallery
        slide_15,                    # 46: Demo
        # --- Section 5: Engineering (slides 47-55) ---
        slide_eng_transition,        # 47: ENGINEERING
        slide_three_repos,           # 48: Starting Point
        slide_17,                    # 49: 8-Stage Pipeline (animated)
        slide_17_png,                # 50: 8-Stage Pipeline (PNG version)
        slide_18,                    # 51: Engineering Journey
        slide_20,                    # 52: Standalone Container
        slide_web_ui,                # 53: Live Demo
        slide_21,                    # 54: Pipeline Intelligence
        slide_dual_env,              # 55: Two Environments
        # --- Section 6: Future Directions (slides 56-72) ---
        slide_future_transition,     # 56: FUTURE DIRECTIONS
        slide_24,                    # 57: Starting from 61.6%
        slide_26,                    # 58: Five Phases
        slide_26b,                   # 59: IS Improvement Roadmap
        slide_confidence_scoring,    # 60: Phase 1: Confidence Scoring
        slide_27,                    # 61: Confidence Summary
        slide_28,                    # 62: Phase 2: N-Best
        slide_data_scaling,          # 63: Data Scaling
        slide_price_tag,             # 64: Price Tag
        slide_29,                    # 65: Fine-Tuning
        slide_30,                    # 66: Stronger LLM + Smart Prompts (unhidden)
        slide_llm_context_engine,    # 67: LLM Is a Context Engine (moved here)
        slide_30b,                   # 68: LLM Upgrade: Why It Matters (unhidden)
        slide_failure_deep_3,        # 69: Failure Modes: Impact & Fixes
        slide_arabic_roadmap,        # 70: Arabic Roadmap
        slide_arabic_avhubert,       # 71: AV-HuBERT
        slide_arabic_changes,        # 72: Arabic Adaptation
        # --- Section 7: Closing (slides 73-74) ---
        slide_31,                    # 73: Key Takeaways
        slide_thank_you,             # 74: Thank You
        # --- Appendix (slides 75-82) ---
        slide_a1,                    # 75: A1: Homophenes
        slide_a8,                    # 76: A3: IS Component Correlation
        slide_a11,                   # 77: A4: LLM Salvage Recoverable
        slide_a11b,                  # 78: A5: LLM Salvage Examples
        slide_a13,                   # 79: A6: Failure Mode Examples
        slide_a15,                   # 80: A7: Video Gallery Map
        slide_a16,                   # 81: A8: LLM Judge x IS Tier
        slide_a17,                   # 82: A9: Context transition matrix
    ]
    total = len(builders)

    # Slides to hide (matching FINAL_PRESENTATION.pptx user edits)
    hidden_builders = {
        slide_exec_summary,          # Executive Summary
        slide_wer_lies,              # WER: The Metric That Lies
        slide_is_deep_dive,          # IS Validation
        slide_metric_disagreement,   # When Metrics Disagree pt 1
        slide_metric_disagreement_2, # When Metrics Disagree pt 2
        slide_disagreement_blind,    # Where IS and Judge Disagree (new)
        slide_disagreement_context,  # Context Exposes Hidden Failures (new)
        slide_17,                    # 8-Stage Pipeline (animated)
        slide_30b,                   # LLM Upgrade: Why It Matters
        slide_a11b,                  # A5: LLM Salvage Examples (not relevant)
    }

    for i, builder in enumerate(builders, 1):
        print(f"  Slide {i:2d}/{total} ...", end=" ")
        try:
            builder(prs)
            if builder in hidden_builders:
                prs.slides[-1]._element.set('show', '0')
                print("OK (hidden)")
            else:
                print("OK")
        except Exception as e:
            print(f"ERROR: {e}")

    # ── Fix slide numbering (overwrite all bottom-left labels) ──
    from pptx.util import Inches, Pt
    appendix_num = 0
    in_appendix = False
    slide_counter = 0
    for slide in prs.slides:
        # Find the slide-number text box: bottom-left, narrow, at y > 6.5"
        num_shape = None
        for shape in slide.shapes:
            if (shape.has_text_frame
                    and shape.width <= Inches(0.6)
                    and shape.left <= Inches(1.2)
                    and shape.top >= Inches(6.8)):
                num_shape = shape
                break
        if num_shape is None:
            continue

        old_text = num_shape.text_frame.text.strip()
        # Detect appendix slides by "A" prefix
        if old_text.startswith("A"):
            in_appendix = True
            appendix_num += 1
            new_text = f"A{appendix_num}"
        else:
            slide_counter += 1
            new_text = str(slide_counter)

        if old_text != new_text:
            num_shape.text_frame.paragraphs[0].text = new_text

    print(f"  Renumbered: {slide_counter} main + {appendix_num} appendix slides")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    _fix_pptx_video_compat(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
