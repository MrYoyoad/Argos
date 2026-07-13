#!/usr/bin/env python3
"""Argos VSP — Egla-Kafe roadmap & partnership deck (companion to TOOL D).

Deck 2 of the July 2026 client meeting: Deck 1
(generate_egla_kafe_client_deck.py) tells the data story on the client's
own footage; this deck carries the roadmap and the continuation ask.
Borrows the client-approved builders from the May 2026 client deck
(presentation/slides_client.py) per the borrow-don't-rewrite rule; net-new
slides follow the Egla-Kafe deck's visual conventions.

Output:
    presentation_materials_20260224/Argos_VSP_EglaKafe_Roadmap_20260713.pptx

Run:
    /home/ubuntu/vsp-llm-yoad-venv/bin/python \
        docs/_research-tools/generators/generate_egla_kafe_roadmap_deck.py

Numbers are sourced from:
    docs/evaluation/egla_kafe/findings.md          (client-footage results)
    docs/CLIENT_MEETING_FRAMING.md                 (canonical client numbers)
    MEMORY.md Key Project Numbers                  (68->71% judge-useful)
N1-N10 number rules apply: no WER / kappa / MBR jargon on slides; "1 in N"
framing; every % paired with its denominator.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from presentation.config import (
    SL_W, SL_H, WHITE, TEAL, CORAL, LGRAY, MGRAY,
    GREEN, YELLOW, GOLD, ORANGE, RED, NAVY2, NAVY3,
    MX, CT, CW, _auto_num,
)
from presentation.helpers import (
    new_slide, add_title, add_accent_line, add_text,
    add_bullets, add_rect, add_logo, add_slide_num,
    set_notes, add_fade_transition, _strip_orphan_animation_refs,
)

# Borrowed client-approved builders (shown in the May 2026 meeting).
from presentation.slides_client import (
    slide_client_compared_to_today,
    slide_client_multi_speaker_today_vs_path,
    slide_client_quality_filter,
    slide_client_data_ask,
    slide_client_investment_ask,
    slide_client_next_steps,
)

OUTPUT = Path("/home/ubuntu/presentation_materials_20260224/"
              "Argos_VSP_EglaKafe_Roadmap_20260713.pptx")


# ── Small local helpers (Egla-Kafe deck conventions) ────────────────────

def _foot(slide, text):
    add_text(slide, text, MX + Inches(0.6), Inches(6.05), CW - Inches(1.2),
             Inches(0.4), size=Pt(16), color=LGRAY, italic=True,
             align=PP_ALIGN.CENTER)


def _finish(slide, notes=""):
    _auto_num[0] += 1
    add_logo(slide)
    add_slide_num(slide, _auto_num[0])
    if notes:
        set_notes(slide, notes)


def _card_row(slide, cards, cy, ch, gap=Inches(0.35)):
    """Row of equal-width bordered cards: (title, body, stat_or_None, color)."""
    n = len(cards)
    cw = (CW - (n - 1) * gap) / n
    for i, (t, d, stat, col) in enumerate(cards):
        x = MX + i * (cw + gap)
        add_rect(slide, x, cy, cw, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, x + Inches(0.2), cy + Inches(0.15),
                 cw - Inches(0.4), Inches(0.85), size=Pt(24), bold=True,
                 color=col)
        add_text(slide, d, x + Inches(0.2), cy + Inches(1.05),
                 cw - Inches(0.4), ch - Inches(1.6), size=Pt(20),
                 color=WHITE)
        if stat:
            add_text(slide, stat, x + Inches(0.2), cy + ch - Inches(0.55),
                     cw - Inches(0.4), Inches(0.45), size=Pt(18),
                     color=LGRAY, italic=True)


def _card_stack(slide, cards, cy, ch, gap=Inches(0.2)):
    """Stack of full-width cards: (title, body, color)."""
    for i, (t, d, col) in enumerate(cards):
        y = cy + i * (ch + gap)
        add_rect(slide, MX, y, CW, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, MX + Inches(0.3), y,
                 Inches(4.6), ch, size=Pt(25), bold=True,
                 color=col, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, d, MX + Inches(5.0), y,
                 CW - Inches(5.3), ch, size=Pt(20), color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
# SLIDES (net-new)
# ════════════════════════════════════════════════════════════════════════

def slide_r_title(prs):
    slide = new_slide(prs)
    add_rect(slide, 0, Inches(2.55), SL_W, Inches(0.03), fill_color=TEAL)
    add_text(slide, "From evaluation to production",
             MX, Inches(1.55), CW, Inches(1.0),
             size=Pt(44), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Roadmap & partnership",
             MX, Inches(2.75), CW, Inches(0.8),
             size=Pt(30), color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide,
             "What your footage proved — and what we build next",
             MX, Inches(3.7), CW, Inches(0.6),
             size=Pt(24), color=LGRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Argos  —  The Orchard      ·      July 2026",
             MX, Inches(5.6), CW, Inches(0.5),
             size=Pt(20), color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide,
            "Deck 2 of 2. Deck 1 told the story of the Egla-Kafe footage; "
            "this deck turns it into decisions: adopt the capture protocol, "
            "share data, green-light the next phase as a partnership. "
            "Target 13 minutes, then open discussion.")
    return slide


def slide_r_since_we_met(prs):
    slide = new_slide(prs)
    add_title(slide, "Since we last met")
    add_accent_line(slide)
    _card_stack(slide, [
        ("Your footage, evaluated",
         "21 conversations analyzed end-to-end; per-video report "
         "and subtitle videos delivered", TEAL),
        ("Consensus reading shipped",
         "Each segment now combines 20 candidate readings — "
         "useful outputs up from 68% to 71%", GREEN),
        ("Trust bands in every report",
         "A green word is right about 9 times in 10 in "
         "trusted segments", GOLD),
    ], CT + Inches(0.15), Inches(1.35), gap=Inches(0.17))
    _foot(slide, "Useful-output rate: independent automated reviewer, "
                 "1,497 real-world segments.")
    _finish(slide,
            "Credibility beat: the product moved since the May meeting. "
            "(1) The Egla-Kafe evaluation itself — deliverables in their "
            "hands. (2) N-best consensus decoding (say 'combining twenty "
            "candidate readings', never 'MBR') shipped as the production "
            "default — blind judge useful rate 68.4% -> 71.1% on 1,497 "
            "segments, statistically significant. (3) Per-word trust bands "
            "with the agreement-aware rule; green ~95% correct within "
            "Trust-tier segments, ~90% overall.")
    return slide


def slide_r_what_eval_proved(prs):
    slide = new_slide(prs)
    add_title(slide, "What the evaluation proved")
    add_accent_line(slide)
    _card_row(slide, [
        ("It works on\ngood footage",
         "Best video: 3 of 4 turns understood, full plot recovered",
         "iPhone 4K, frontal", GREEN),
        ("It knows when\nit's right",
         "The confident subset is 70–92% useful",
         "self-flagged, no reference needed", TEAL),
        ("Capture decides\nthe outcome",
         "Resolution and frontality dominate everything downstream",
         "the one robust lever", GOLD),
    ], CT + Inches(0.2), Inches(3.6))
    add_text(slide,
             "Capability is proven — coverage is the frontier.",
             MX, CT + Inches(4.1), CW, Inches(0.6),
             size=Pt(26), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _finish(slide,
            "The bridge from Deck 1 to the roadmap. Three proof points: "
            "(1) best clip img_6825 — 72.7% of turns understood in context, "
            "plot fully recoverable; (2) gating by the model's own "
            "confidence reproduces English-benchmark quality (70% useful "
            "keeping the top 10%, 92% keeping 3%); (3) iPhone-vs-camera and "
            "frontal-vs-45 are the only statistically robust levers "
            "(p=2e-5 / p=2e-3, speaker notes only). Everything on the "
            "roadmap attacks coverage: better input, more speakers, "
            "stronger model.")
    return slide


def slide_r_capture_protocol(prs):
    slide = new_slide(prs)
    add_title(slide, "The free win: film for the model")
    add_accent_line(slide)
    add_text(slide,
             "No new software, no waiting — the biggest lift available today",
             MX, CT, CW, Inches(0.5), size=Pt(24), color=LGRAY,
             align=PP_ALIGN.CENTER)
    _card_row(slide, [
        ("Native, high-res",
         "Record on the device — never a screen-recording",
         "clear turns: 1 in 7 vs 1 in 125", GREEN),
        ("Frontal faces",
         "A face at 45° loses nearly everything",
         "useful turns: 27% front vs 3% at 45°", TEAL),
        ("Fill the frame",
         "Closer camera, more pixels on the lips",
         "resolution is signal", GOLD),
    ], CT + Inches(0.65), Inches(3.1))
    add_text(slide,
             "Offer: a re-shoot pilot — same scenes, captured right.",
             MX, CT + Inches(4.05), CW, Inches(0.55),
             size=Pt(24), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _finish(slide,
            "The 380px screen-recording was the confound — their camera "
            "hardware is likely fine at native resolution ('why was our "
            "camera so bad?' -> it was a screen-rec of a composite, with UI "
            "chrome). Clearly-conveyed turns: iPhone 13.7% (~1 in 7) vs "
            "screen-rec 0.8% (~1 in 125); useful turns front 27.4% vs 45-deg "
            "3.2%. Ask them to adopt the protocol AND to run a re-shoot "
            "pilot so the next evaluation isolates their true ceiling.")
    return slide


def slide_r_continuity(prs):
    slide = new_slide(prs)
    add_title(slide, "Built to continue")
    add_accent_line(slide)
    _card_stack(slide, [
        ("In your hands today",
         "Standalone system on your hardware, user guide, "
         "this evaluation's full report pack", GREEN),
        ("Repeatable evaluation",
         "The same pipeline re-runs on any new footage "
         "you bring — days, not months", TEAL),
        ("Scoped next projects",
         "Multi-speaker and quality pre-filter are written up, "
         "estimated, and ready to staff", GOLD),
    ], CT + Inches(0.15), Inches(1.35), gap=Inches(0.17))
    _foot(slide, "The next phase is staffed as part of the partnership — "
                 "people, data, and training together.")
    _finish(slide,
            "Continuity slide — succession is deliberately part of the ask. "
            "Everything is documented and packaged: the deployed container, "
            "the client user guide, the repeatable client-eval pipeline "
            "(docs/guides/client-lipread-eval.md), and ready-to-staff "
            "project briefs (multi-speaker, video-quality pre-filter). "
            "Spoken, not on slide: the current lead is handing over in the "
            "coming weeks; committing the next phase now lets the incoming "
            "team start from these packaged assets rather than from zero.")
    return slide


def slide_r_thanks(prs):
    slide = new_slide(prs)
    add_title(slide, "Three decisions we're asking for")
    add_accent_line(slide)
    _card_stack(slide, [
        ("1 · Film for the model",
         "Adopt the capture protocol; run the re-shoot "
         "pilot with us", GREEN),
        ("2 · Share footage",
         "Your recordings become the training data for "
         "your production model", TEAL),
        ("3 · Green-light the next phase",
         "Multi-speaker + a stronger model, built "
         "together as a partnership", GOLD),
    ], CT + Inches(0.15), Inches(1.3))
    add_text(slide,
             "Thank you — Argos · The Orchard",
             MX, Inches(6.1), CW, Inches(0.5),
             size=Pt(22), color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide,
            "Close on the three concrete decisions: capture protocol + "
            "re-shoot pilot (free, immediate), data contribution, and the "
            "partnership green-light. Specifics of budget and staffing in "
            "follow-up, per the investment-ask doctrine.")
    return slide


def main():
    _auto_num[0] = 0
    prs = Presentation()
    prs.slide_width = SL_W
    prs.slide_height = SL_H

    builders = [
        slide_r_title,
        slide_r_since_we_met,
        slide_r_what_eval_proved,
        slide_client_compared_to_today,          # borrowed — recap anchor
        slide_r_capture_protocol,
        slide_client_multi_speaker_today_vs_path,  # borrowed — their canonical case
        slide_client_quality_filter,             # borrowed — optional add-on
        slide_client_data_ask,                   # borrowed — bridge to the ask
        slide_client_investment_ask,             # borrowed — partnership ask
        slide_r_continuity,
        slide_client_next_steps,                 # borrowed — presenter customizes
        slide_r_thanks,
    ]
    total = len(builders)
    for i, b in enumerate(builders, 1):
        print(f"  Slide {i:2d}/{total}  {b.__name__} ...", end=" ")
        b(prs)
        print("OK")

    for s in prs.slides:
        add_fade_transition(s, speed="med")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    _strip_orphan_animation_refs(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
