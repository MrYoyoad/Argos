#!/usr/bin/env python3
"""Argos VSP — Egla-Kafe client deck (TOOL D).

A short (~12-slide) honest-positive client deck summarizing the lip-reading
evaluation of the Egla-Kafe footage. Reuses the brand palette, layout
constants, and battle-tested helpers (video embed + poster extraction +
PowerPoint repair-dialog fixes) from the existing `presentation/` package.

Output:
    presentation_materials_20260224/Argos_VSP_EglaKafe_20260625.pptx

Run:
    /home/ubuntu/vsp-llm-yoad-venv/bin/python \
        docs/_research-tools/generators/generate_egla_kafe_client_deck.py

Numbers are sourced from:
    docs/evaluation/egla_kafe/findings.md
    docs/evaluation/egla_kafe/per_video_understanding.md
STYLE_GUIDE.md: body floor 24pt, titles 32-40pt, <=4 bullets / <=8 words,
visual-first, tier colors green=trust / yellow=salvage / red(purple)=strip.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette + layout + helpers from the existing deck package.
from presentation.config import (
    VID,
    SL_W, SL_H, BG, WHITE, TEAL, CORAL, LGRAY, MGRAY,
    GREEN, YELLOW, GOLD, ORANGE, RED, NAVY2, NAVY3, BLUE, PURPLE,
    FONT, MX, MY, CT, CW,
)
from presentation.helpers import (
    new_slide, add_title, add_accent_line, add_text, add_rich_text,
    add_bullets, add_rect, add_image, add_video, add_logo, add_slide_num,
    set_notes, add_fade_transition, _fix_pptx_video_compat,
    _strip_orphan_animation_refs,
)

# ── Egla-Kafe deliverable assets ────────────────────────────────────────
DELIV = Path("/home/ubuntu/datasets/clients/egla_kafe/deliverables")
PLOTS = DELIV / "plots"
CLIPS = DELIV / "clips"

PLOT = {
    "recovery_ladder": PLOTS / "recovery_ladder.png",
    "levers":          PLOTS / "levers.png",
    "category_trust":  PLOTS / "category_trust.png",
    "confidence_gate": PLOTS / "confidence_gate.png",
    "calibration":     PLOTS / "calibration.png",
}

# Register clips into the shared VID dict so add_video() can resolve them
# (it keys off VID and handles poster extraction + the empty-hlink repair fix).
_CLIPS = {
    "ek_active_speaker": CLIPS / "active_speaker_overlay.mp4",
    "ek_iphone_vs_cam":  CLIPS / "iphone_vs_camera.mp4",
    "ek_best_vs_worst":  CLIPS / "best_vs_worst.mp4",
    "ek_conf_colored":   CLIPS / "confidence_colored.mp4",
}
VID.update(_CLIPS)

OUTPUT = Path("/home/ubuntu/presentation_materials_20260224/"
              "Argos_VSP_EglaKafe_20260713.pptx")

# ── Small local helpers ─────────────────────────────────────────────────

_num = [0]


def _foot(slide, text):
    """Small footnote/source line in the footer safe zone (16pt, ends <=6.45)."""
    add_text(slide, text, MX + Inches(0.6), Inches(6.05), CW - Inches(1.2),
             Inches(0.4), size=Pt(16), color=LGRAY, italic=True,
             align=PP_ALIGN.CENTER)


def _finish(slide, notes=""):
    _num[0] += 1
    add_logo(slide)
    add_slide_num(slide, _num[0])
    if notes:
        set_notes(slide, notes)


def _img_fit(path, max_w, max_h):
    """Return (w, h) Inches preserving aspect ratio inside the box."""
    from PIL import Image
    iw, ih = Image.open(str(path)).size
    ar = iw / ih
    box_ar = (max_w / max_h)
    if ar > box_ar:
        w = max_w
        h = Inches(max_w / Inches(1) / ar)
    else:
        h = max_h
        w = Inches(max_h / Inches(1) * ar)
    return w, h


def _centered_image(slide, path, top, max_w, max_h):
    w, h = _img_fit(path, max_w, max_h)
    x = (SL_W - w) // 2
    add_image(slide, path, x, top, width=w, height=h)
    return x, w, h


def _video_centered(slide, key, top, max_w, max_h):
    """Embed a video centered, preserving aspect ratio of the source."""
    import subprocess, json
    p = VID[key]
    out = subprocess.run(
        ["ffprobe", "-v", "error", "-select_streams", "v:0",
         "-show_entries", "stream=width,height", "-of", "json", str(p)],
        capture_output=True, text=True)
    st = json.loads(out.stdout)["streams"][0]
    ar = st["width"] / st["height"]
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = Inches(max_w / Inches(1) / ar)
    else:
        h = max_h
        w = Inches(max_h / Inches(1) * ar)
    x = (SL_W - w) // 2
    add_video(slide, key, x, top, w, h)
    return x, w, h


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = new_slide(prs)
    # Centered hero title block.
    add_rect(slide, 0, Inches(2.55), SL_W, Inches(0.03), fill_color=TEAL)
    add_text(slide, "Visual Speech Recognition",
             MX, Inches(1.55), CW, Inches(1.0),
             size=Pt(44), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Egla-Kafe footage evaluation",
             MX, Inches(2.75), CW, Inches(0.8),
             size=Pt(30), color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Reading speech from the lips alone — no audio",
             MX, Inches(3.7), CW, Inches(0.6),
             size=Pt(24), color=LGRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Argos  —  The Orchard      ·      July 2026",
             MX, Inches(5.6), CW, Inches(0.5),
             size=Pt(20), color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide,
            "Honest-positive framing. We evaluated our visual speech "
            "recognition model on your real conversation footage (no audio "
            "track used). This deck: what we built, what it recovers, what "
            "you can trust, the honest limits, and recommendations.")
    return slide


def slide_what_we_built(prs):
    slide = new_slide(prs)
    add_title(slide, "What we built")
    add_accent_line(slide)
    # 3-step pipeline cards across the top.
    steps = [
        ("1", "Find the speaker", "Who is talking, from mouth\nmotion — no audio", TEAL),
        ("2", "Read the lips", "The model transcribes the\nactive speaker's lips", GOLD),
        ("3", "Evaluate", "Score each turn; flag what\nto trust vs verify", GREEN),
    ]
    n = len(steps)
    gap = Inches(0.35)
    cw = (CW - (n - 1) * gap) / n
    cy = CT
    ch = Inches(2.05)
    for i, (numlbl, t, d, col) in enumerate(steps):
        x = MX + i * (cw + gap)
        add_rect(slide, x, cy, cw, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, numlbl, x + Inches(0.2), cy + Inches(0.12),
                 Inches(0.7), Inches(0.6), size=Pt(30), bold=True, color=col)
        add_text(slide, t, x + Inches(0.2), cy + Inches(0.7),
                 cw - Inches(0.4), Inches(0.5), size=Pt(24), bold=True)
        add_text(slide, d, x + Inches(0.2), cy + Inches(1.2),
                 cw - Inches(0.4), Inches(0.8), size=Pt(20), color=LGRAY)
    # Active-speaker overlay clip below.
    cap_y = cy + ch + Inches(0.15)
    add_text(slide,
             "Active-speaker stream: green box marks who is talking",
             MX, cap_y, CW, Inches(0.4), size=Pt(20), color=TEAL,
             align=PP_ALIGN.CENTER)
    _video_centered(slide, "ek_active_speaker", cap_y + Inches(0.45),
                    Inches(9.5), Inches(2.35))
    _finish(slide,
            "Footage had no usable audio, so we detect the speaker visually "
            "from mouth-openness variance (speech oscillates; a smile is a "
            "sustained stretch), build a single active-speaker stream, then "
            "lip-read each turn and score it. All 11 scene-1/2 streams: "
            "speaker alternation 1.0.")
    return slide


def slide_recovery_ladder(prs):
    slide = new_slide(prs)
    add_title(slide, "On good footage, the conversation comes through")
    add_accent_line(slide)
    # Lead positive headline.
    add_text(slide,
             "Up to 73% of turns understood on the best-captured video",
             MX, CT, CW, Inches(0.55), size=Pt(26), bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)
    _centered_image(slide, PLOT["recovery_ladder"], CT + Inches(0.65),
                    Inches(8.6), Inches(3.8))
    _foot(slide, "Context-aware recovery (Y+P): share of turns a viewer "
                 "grasps with conversation context.")
    _finish(slide,
            "The pessimistic average hides a high ceiling. On the best video "
            "(iPhone 4K, frontal, Military script) a context-aware viewer "
            "grasps 73% of turns and the full plot is recoverable. The model "
            "isn't incapable on this footage — it locks on intermittently; "
            "capture quality and confidence-gating decide how much you get.")
    return slide


def slide_levers(prs):
    slide = new_slide(prs)
    add_title(slide, "Three controllable levers")
    add_accent_line(slide)
    # Left: 3 lever cards. Right: iphone-vs-camera clip.
    left_w = Inches(5.4)
    cards = [
        ("Capture quality", "iPhone 4K beats screen-rec",
         "IS 1.51 vs 0.88", GREEN),
        ("Frontality", "Face the camera, not 45°",
         "IS 1.62 vs 1.03", TEAL),
        ("Content", "Distinctive words survive",
         "Military > small-talk", GOLD),
    ]
    cy = CT
    ch = Inches(1.45)
    gap = Inches(0.2)
    for i, (t, d, stat, col) in enumerate(cards):
        y = cy + i * (ch + gap)
        add_rect(slide, MX, y, left_w, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, MX + Inches(0.2), y + Inches(0.12),
                 left_w - Inches(0.4), Inches(0.5), size=Pt(26), bold=True,
                 color=col)
        add_text(slide, d, MX + Inches(0.2), y + Inches(0.65),
                 left_w - Inches(0.4), Inches(0.45), size=Pt(22),
                 color=WHITE)
        add_text(slide, stat, MX + Inches(0.2), y + Inches(1.05),
                 left_w - Inches(0.4), Inches(0.4), size=Pt(20),
                 color=LGRAY, italic=True)
    # Right column: clip on top, plot below.
    rx = MX + left_w + Inches(0.4)
    rw = CW - left_w - Inches(0.4)
    add_text(slide, "iPhone 4K  vs  client camera",
             rx, CT - Inches(0.02), rw, Inches(0.4), size=Pt(20),
             color=TEAL, align=PP_ALIGN.CENTER)
    # Center the clip within the right column.
    import subprocess, json
    p = VID["ek_iphone_vs_cam"]
    st = json.loads(subprocess.run(
        ["ffprobe", "-v", "error", "-select_streams", "v:0",
         "-show_entries", "stream=width,height", "-of", "json", str(p)],
        capture_output=True, text=True).stdout)["streams"][0]
    ar = st["width"] / st["height"]
    vh = Inches(2.5)
    vw = Inches(2.5 * ar)
    vx = rx + (rw - vw) // 2
    add_video(slide, "ek_iphone_vs_cam", vx, CT + Inches(0.45), vw, vh)
    _centered_image_in_box(slide, PLOT["levers"], rx, CT + Inches(3.15),
                           rw, Inches(2.4))
    _finish(slide,
            "Capture quality and frontality are the only statistically "
            "robust levers (Mann-Whitney: iPhone-vs-camera p=2e-5, "
            "front-vs-45° p=2e-3). Content helps via distinctive words "
            "surviving + cross-turn redundancy, not a per-segment IS gap. "
            "All three are within the client's control at capture time.")
    return slide


def _centered_image_in_box(slide, path, bx, by, bw, bh):
    from PIL import Image
    iw, ih = Image.open(str(path)).size
    ar = iw / ih
    if (bw / bh) > ar:
        h = bh
        w = Inches(bh / Inches(1) * ar)
    else:
        w = bw
        h = Inches(bw / Inches(1) / ar)
    x = bx + (bw - w) // 2
    add_image(slide, path, x, by, width=w, height=h)


def slide_what_you_can_trust(prs):
    slide = new_slide(prs)
    add_title(slide, "What you can trust")
    add_accent_line(slide)
    add_text(slide, "Trust the gist and topic. Verify names and numbers.",
             MX, CT, CW, Inches(0.5), size=Pt(26), bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)
    # Plot left, trust/verify cards right.
    lw = Inches(6.6)
    _centered_image_in_box(slide, PLOT["category_trust"], MX, CT + Inches(0.65),
                           lw, Inches(3.8))
    rx = MX + lw + Inches(0.3)
    rw = CW - lw - Inches(0.3)
    # Trust card (green).
    add_rect(slide, rx, CT + Inches(0.65), rw, Inches(1.85),
             fill_color=NAVY2, border_color=GREEN, border_width=Pt(2),
             corner_radius=True)
    add_text(slide, "TRUST", rx + Inches(0.2), CT + Inches(0.72),
             rw - Inches(0.4), Inches(0.5), size=Pt(26), bold=True,
             color=GREEN)
    add_bullets(slide, [
        ("Common nouns — 82% correct", {"color": WHITE}),
        ("Topic / gist of the turn", {"color": WHITE}),
    ], rx + Inches(0.2), CT + Inches(1.25), rw - Inches(0.4), Inches(1.2),
        size=Pt(22), bullet_color=GREEN)
    # Verify card (red/purple).
    add_rect(slide, rx, CT + Inches(2.7), rw, Inches(1.95),
             fill_color=NAVY2, border_color=RED, border_width=Pt(2),
             corner_radius=True)
    add_text(slide, "VERIFY", rx + Inches(0.2), CT + Inches(2.77),
             rw - Inches(0.4), Inches(0.5), size=Pt(26), bold=True,
             color=RED)
    add_bullets(slide, [
        ("Names / places — hallucinated", {"color": WHITE}),
        ("Numbers / dates — confident, unsafe", {"color": WHITE}),
    ], rx + Inches(0.2), CT + Inches(3.3), rw - Inches(0.4), Inches(1.3),
        size=Pt(22), bullet_color=RED)
    _foot(slide, "P(correct) for high-confidence words, by category. "
                 "All 778 scored segments.")
    _finish(slide,
            "A gated output gives a topical skeleton, not full sentences. "
            "Green common-nouns are ~82% reliable but only ~10% of nouns are "
            "recovered — a scatter of trustworthy content words to fuse "
            "with context. Names/places are 0% recovered AND the model emits "
            "confident fake entities (Abu Dhabi 0.997, Wikipedia). Numbers "
            "~73% green but still need verification.")
    return slide


def slide_confidence_flag(prs):
    slide = new_slide(prs)
    add_title(slide, "The model flags its own good output")
    add_accent_line(slide)
    add_text(slide,
             "On the confident subset: 70–92% of turns are useful",
             MX, CT, CW, Inches(0.5), size=Pt(26), bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)
    # Plot left, colored clip right.
    lw = Inches(7.4)
    _centered_image_in_box(slide, PLOT["confidence_gate"], MX, CT + Inches(0.6),
                           lw, Inches(3.85))
    rx = MX + lw + Inches(0.3)
    rw = CW - lw - Inches(0.3)
    add_text(slide, "Per-word confidence", rx, CT + Inches(0.55), rw,
             Inches(0.4), size=Pt(20), color=TEAL, align=PP_ALIGN.CENTER)
    # Confidence clip is portrait (256x312) — fit by height.
    import subprocess, json
    p = VID["ek_conf_colored"]
    st = json.loads(subprocess.run(
        ["ffprobe", "-v", "error", "-select_streams", "v:0",
         "-show_entries", "stream=width,height", "-of", "json", str(p)],
        capture_output=True, text=True).stdout)["streams"][0]
    ar = st["width"] / st["height"]
    vh = Inches(3.6)
    vw = Inches(3.6 * ar)
    vx = rx + (rw - vw) // 2
    add_video(slide, "ek_conf_colored", vx, CT + Inches(1.0), vw, vh)
    _foot(slide, "Gate by the model's own confidence — no reference text "
                 "needed. conf≥0.7 keeps top 10% at 70% useful.")
    _finish(slide,
            "The model knows when it's right. Gating by its own word "
            "confidence reproduces the ceiling: conf≥0.6 keeps a quarter "
            "at 55% useful; conf≥0.7 keeps the top 10% at 70% useful, "
            "WER 65% (English-benchmark quality); conf≥0.8 keeps 3% at "
            "92% useful. Deliver the sure parts; flag the rest.")
    return slide


def slide_best_vs_worst(prs):
    slide = new_slide(prs)
    add_title(slide, "Best vs worst conditions")
    add_accent_line(slide)
    # Two condition cards flanking, clip centered.
    add_text(slide,
             "Same model, same pipeline — capture decides the outcome",
             MX, CT, CW, Inches(0.5), size=Pt(24), color=LGRAY,
             align=PP_ALIGN.CENTER)
    _video_centered(slide, "ek_best_vs_worst", CT + Inches(0.6),
                    Inches(8.2), Inches(3.9))
    # Small badges left/right under headline.
    add_text(slide, "BEST: frontal 4K → plot recoverable",
             MX, Inches(6.15), Inches(6.0), Inches(0.4), size=Pt(20),
             color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "WORST: 45° / low-res → near-total loss",
             SL_W - MX - Inches(6.0), Inches(6.15), Inches(6.0), Inches(0.4),
             size=Pt(20), color=RED, bold=True, align=PP_ALIGN.CENTER)
    _finish(slide,
            "Left: best-conditions turn (frontal, high-res) — the model "
            "locks on and the line is readable. Right: worst-conditions turn "
            "(45° profile / 380px screen-rec) — near-total loss, "
            "fluent hallucination. The model is frontal-trained; lip-reading "
            "collapses as the face turns to profile.")
    return slide


def slide_honest_limits(prs):
    slide = new_slide(prs)
    add_title(slide, "Honest limits")
    add_accent_line(slide)
    # Limits as plain cards (left), claim/non-claim box (right).
    lw = Inches(6.4)
    limits = [
        ("Names & numbers", "Hallucinated — never trust unverified", RED),
        ("Profile / low-res", "45° or screen-rec not usable", ORANGE),
        ("Word coverage", "~10% of words recovered overall", YELLOW),
    ]
    cy = CT
    ch = Inches(1.4)
    gap = Inches(0.2)
    for i, (t, d, col) in enumerate(limits):
        y = cy + i * (ch + gap)
        add_rect(slide, MX, y, lw, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, MX + Inches(0.2), y + Inches(0.12),
                 lw - Inches(0.4), Inches(0.5), size=Pt(26), bold=True,
                 color=col)
        add_text(slide, d, MX + Inches(0.2), y + Inches(0.7),
                 lw - Inches(0.4), Inches(0.55), size=Pt(22), color=WHITE)
    # What we claim / do NOT claim box.
    rx = MX + lw + Inches(0.35)
    rw = CW - lw - Inches(0.35)
    bh = Inches(2.45)
    add_rect(slide, rx, CT, rw, bh, fill_color=NAVY3,
             border_color=GREEN, border_width=Pt(2), corner_radius=True)
    add_text(slide, "What we CLAIM", rx + Inches(0.2), CT + Inches(0.1),
             rw - Inches(0.4), Inches(0.4), size=Pt(22), bold=True,
             color=GREEN)
    add_bullets(slide, [
        ("Gist & topic on good footage", {"color": WHITE}),
        ("A self-flagged confident subset", {"color": WHITE}),
    ], rx + Inches(0.2), CT + Inches(0.6), rw - Inches(0.4), Inches(1.7),
        size=Pt(21), bullet_color=GREEN)
    by = CT + bh + Inches(0.25)
    bh2 = Inches(2.45)
    add_rect(slide, rx, by, rw, bh2, fill_color=NAVY3,
             border_color=RED, border_width=Pt(2), corner_radius=True)
    add_text(slide, "What we do NOT claim", rx + Inches(0.2), by + Inches(0.1),
             rw - Inches(0.4), Inches(0.4), size=Pt(22), bold=True,
             color=RED)
    add_bullets(slide, [
        ("Verbatim transcripts", {"color": WHITE}),
        ("Reliable names, numbers, dates", {"color": WHITE}),
    ], rx + Inches(0.2), by + Inches(0.6), rw - Inches(0.4), Inches(1.7),
        size=Pt(21), bullet_color=RED)
    _finish(slide,
            "Stated plainly and constructively. The 122% per-segment average "
            "is a coverage problem, not a capability ceiling — but on "
            "this footage word coverage is ~10%, profile/low-res clips are "
            "not usable, and names/numbers are hallucinated. We claim gist + "
            "a trustworthy confident subset; we do not claim verbatim text or "
            "reliable specifics.")
    return slide


def slide_recommendations(prs):
    slide = new_slide(prs)
    add_title(slide, "Recommendations")
    add_accent_line(slide)
    recs = [
        ("\U0001F4F9", "Capture native, high-res",
         "Frontal, 1080p+, no screen-recording", GREEN),
        ("\U0001F9ED", "Keep faces frontal",
         "Avoid profile / 45° angles", TEAL),
        ("\U0001F6E1️", "Use confidence-gating",
         "Deliver sure turns; flag the rest", GOLD),
        ("\U0001F4DD", "Expect gist, not verbatim",
         "Read for meaning; verify specifics", ORANGE),
    ]
    n = len(recs)
    gap = Inches(0.3)
    cw = (CW - (n - 1) * gap) / n
    cy = CT + Inches(0.4)
    ch = Inches(3.4)
    for i, (icon, t, d, col) in enumerate(recs):
        x = MX + i * (cw + gap)
        add_rect(slide, x, cy, cw, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, icon, x, cy + Inches(0.3), cw, Inches(0.8),
                 size=Pt(40), align=PP_ALIGN.CENTER)
        add_text(slide, t, x + Inches(0.15), cy + Inches(1.25),
                 cw - Inches(0.3), Inches(0.9), size=Pt(24), bold=True,
                 color=col, align=PP_ALIGN.CENTER)
        add_text(slide, d, x + Inches(0.15), cy + Inches(2.2),
                 cw - Inches(0.3), Inches(1.0), size=Pt(20), color=LGRAY,
                 align=PP_ALIGN.CENTER)
    _finish(slide,
            "Four practical levers, all within the client's control. The "
            "biggest wins are upstream of the model: native high-res frontal "
            "capture, plus confidence-gating downstream and reading for gist "
            "rather than verbatim text.")
    return slide


def slide_whats_next(prs):
    slide = new_slide(prs)
    add_title(slide, "What's next")
    add_accent_line(slide)
    nexts = [
        ("Dynamic active-speaker", "Per-frame face tracking; use the "
         "masters' audio for audio-visual speaker detection", TEAL),
        ("More & better data", "Frontal, high-res captures to lift "
         "coverage beyond the confident subset", GREEN),
        ("Calibrated confidence", "Tune the trust gate to your "
         "accept/verify tradeoff", GOLD),
    ]
    cy = CT + Inches(0.15)
    ch = Inches(1.25)
    gap = Inches(0.22)
    for i, (t, d, col) in enumerate(nexts):
        y = cy + i * (ch + gap)
        add_rect(slide, MX, y, CW, ch, fill_color=NAVY2,
                 border_color=col, border_width=Pt(1.5), corner_radius=True)
        add_text(slide, t, MX + Inches(0.3), y,
                 Inches(5.0), ch, size=Pt(27), bold=True,
                 color=col, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, d, MX + Inches(5.4), y,
                 CW - Inches(5.7), ch, size=Pt(22), color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide,
             "The model is capable on well-captured footage — the path "
             "forward is more of the right input and smarter gating.",
             MX, cy + 3 * (ch + gap) + Inches(0.08), CW, Inches(0.6),
             size=Pt(22), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)
    _finish(slide,
            "Phase B: dynamic per-frame cropping + audio-visual active-speaker "
            "detection (the masters have audio). More frontal high-res data "
            "to raise coverage. Calibrate the confidence gate to the client's "
            "accept/verify tradeoff.")
    return slide


def slide_thank_you(prs):
    slide = new_slide(prs)
    add_rect(slide, 0, Inches(3.0), SL_W, Inches(0.03), fill_color=TEAL)
    add_text(slide, "Thank you", MX, Inches(2.0), CW, Inches(1.0),
             size=Pt(48), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Trust the gist · verify specifics · capture frontal & "
             "high-res",
             MX, Inches(3.3), CW, Inches(0.6), size=Pt(24), color=TEAL,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Argos — The Orchard", MX, Inches(5.4), CW,
             Inches(0.5), size=Pt(20), color=LGRAY, align=PP_ALIGN.CENTER)
    _finish(slide, "Close. One-line takeaway repeats the operating doctrine.")
    return slide


def main():
    _num[0] = 0
    prs = Presentation()
    prs.slide_width = SL_W
    prs.slide_height = SL_H

    builders = [
        slide_title,
        slide_what_we_built,
        slide_recovery_ladder,
        slide_levers,
        slide_what_you_can_trust,
        slide_confidence_flag,
        slide_best_vs_worst,
        slide_honest_limits,
        slide_recommendations,
        # slide_whats_next dropped for the July 2026 meeting — the roadmap
        # now lives in the companion deck (generate_egla_kafe_roadmap_deck.py)
        # so it isn't presented twice.
        slide_thank_you,
    ]
    total = len(builders)
    for i, b in enumerate(builders, 1):
        print(f"  Slide {i:2d}/{total}  {b.__name__} ...", end=" ")
        b(prs)
        print("OK")

    for s in prs.slides:
        add_fade_transition(s, speed="med")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    # Post-process: wrap embedded videos for PowerPoint (no repair dialog),
    # and scrub any orphan animation refs.
    _fix_pptx_video_compat(str(OUTPUT))
    _strip_orphan_animation_refs(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
