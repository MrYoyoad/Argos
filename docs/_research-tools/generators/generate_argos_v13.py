#!/usr/bin/env python3
"""
Argos VSP — v13 Amosi Deck Generator

Reverse-engineered from `presentation_materials_20260224/Argos_VSP_v13_Amosi_2.pptx`
(99 slides, 23 hidden). Reuses existing builders from `presentation/slides_*.py`,
patches titles for renamed slides via `with_title()`, and adds 4 new builders
+ 2 merge builders in `slides_v13_new.py`.

Source: see /tmp/v13_mapping.tsv (slide-by-slide inventory).

Usage:
    python3 docs/_research-tools/generators/generate_argos_v13.py

Output:
    presentation_materials_20260224/Argos_VSP_v13_Amosi_2_generated.pptx
"""
from pathlib import Path
from pptx import Presentation

from presentation.config import SL_W, SL_H, _auto_num
from presentation.helpers import _fix_pptx_video_compat, _strip_orphan_animation_refs

OUTPUT = Path(
    "/home/ubuntu/presentation_materials_20260224/"
    "Argos_VSP_v13_Amosi_2_generated.pptx"
)

# ─── Existing builders ─────────────────────────────────────────────────
from presentation.slides_opening import (
    slide_01, slide_what_was_done_1, slide_what_was_done_2, slide_toc,
    slide_02, slide_visemes, slide_03, slide_data_flow,
    slide_04, slide_05, slide_06,
    slide_exec_summary, slide_wer_lies,
)
from presentation.slides_research import (
    slide_research_transition,
    slide_failure_anatomy_transition,
    slide_is_motivation, slide_is_intro_a, slide_is_intro_b, slide_is_intro_c,
    slide_is_weight_rationale, slide_is_calc_examples,
    slide_is_dimensions, slide_is_radar, slide_is_wer_scatter,
    slide_07, slide_metric_transition,
    slide_08,
    slide_failure_deep_1a, slide_failure_deep_1b, slide_failure_deep_2,
    slide_failure_deep_3,
)
from presentation.slides_evaluation import (
    slide_llm_judge, slide_llm_judge_30,
    slide_judge_ex1, slide_judge_ex2, slide_judge_ex3,
    slide_judge_ex4, slide_judge_ex5, slide_judge_ex6,
    slide_disagreement_blind, slide_disagreement_context,
    slide_25d, slide_25e,
    slide_metric_disagreement, slide_metric_disagreement_2,
    slide_two_eval_systems, slide_context_eval,
    slide_is_deep_dive,
    slide_llm_context_engine,
    slide_25, slide_25b, slide_25c,
)
from presentation.slides_demo import (
    slide_15,
    slide_live_example_intro,
    slide_failure_live_demo,
)
from presentation.slides_engineering import (
    slide_17_a, slide_17_b, slide_17_png,
    slide_dual_env,
    slide_18, slide_19, slide_20, slide_21,
    slide_three_repos, slide_web_ui,
)
from presentation.slides_future import (
    slide_future_transition,
    slide_24, slide_26, slide_26b,
    slide_30, slide_30b,
    slide_29_phases, slide_29_lift,
    slide_arabic_roadmap, slide_arabic_avhubert, slide_arabic_changes,
    slide_31, slide_thank_you,
    slide_28, slide_confidence_scoring,
    slide_a1, slide_a8, slide_a11, slide_a11b, slide_a13, slide_a15,
    slide_a16, slide_a17,
    slide_human_is_path_b,
    slide_appendix_pca_loadings,
    slide_appendix_mcnemar_full,
    slide_data_scaling, slide_price_tag,
)
from presentation.slides_demo import (
    slide_demo_obama_trust,
    slide_demo_obama_salvage,
    slide_demo_obama_strip,
    slide_demo_judge_entity,
    slide_demo_judge_vocab,
)
from presentation.slides_confidence import (
    slide_confidence_problem,
    slide_two_layer_confidence_research,
    slide_per_word_confidence_distribution,
    slide_band_reliability_overall,
    slide_band_reliability_stratified,
    slide_green_leakage_examples,
    slide_three_thresholds,
    slide_three_tier_policy_research,
    slide_band_reliability_by_niv,
    slide_agreement_aware_bands,
    slide_agreement_vs_conf_information,
    slide_client_trust_calibration,
    slide_is_conf_correlation,
    slide_is_conf_contingency,
    slide_nbest_v3_judge_paired_tests,
    slide_mbr_decision,
    slide_v1_vs_v3_judge_lesson,
)

# ─── New v13-only builders ─────────────────────────────────────────────
from presentation.slides_v13_new import (
    slide_thesis,                # v13 #2
    slide_what_we_claim,         # v13 #3
    slide_is_61pct_headline,     # v13 #33 — "Intelligibility Score: 61.6% Useful Output"
    slide_four_numbers,          # v13 #38
    slide_is_conf_correlate_combined,  # v13 #47 — merge of conf_correlation + contingency
    slide_pipeline_10stage,      # v13 #64 — merge of 17_a + 17_b as 10-stage
    slide_engineering_journey_v13,   # v13 #66 — lift from slides_client
    slide_quality_filter_v13,    # v13 #79 — lift from slides_client
)


def with_title(builder, new_title):
    """Wrap a builder so its slide(s) get their title rewritten.

    Patches the FIRST text-frame shape near the top of each new slide added.
    Preserves all formatting (font, color, size) by only mutating run text.
    Works for builders that add exactly 1 slide. For multi-slide builders,
    only the first slide's title is patched.
    """
    def wrapped(prs):
        n_before = len(prs.slides)
        result = builder(prs)
        if len(prs.slides) <= n_before:
            return result
        slide = prs.slides[n_before]
        # Title heuristic: textframe near top with non-empty text
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if not tf.text.strip():
                continue
            # Title sits in the top ~1.2" of the slide
            if shape.top is None or shape.top > 1097280:  # 1.2 inches in EMU
                continue
            # Heuristic: title shapes are wide (>5") and have bold text
            if shape.width and shape.width < 4572000:  # 5 inches
                continue
            # Replace run text preserving format of first run
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = new_title
                for r in tf.paragraphs[0].runs[1:]:
                    r.text = ''
            else:
                tf.text = new_title
            break
        return result
    wrapped.__name__ = f"{builder.__name__}__as__{new_title[:30].replace(' ', '_')}"
    return wrapped


def hide(builder):
    """Mark builder's slide as hidden (show=0 on the sldId element)."""
    def wrapped(prs):
        n_before = len(prs.slides)
        result = builder(prs)
        if len(prs.slides) > n_before:
            slide = prs.slides[n_before]
            # python-pptx exposes slide.element.set('show', '0') via XML
            from lxml import etree
            sldId_elem = slide.element
            sldId_elem.set('show', '0')
        return result
    wrapped.__name__ = builder.__name__ + "__hidden"
    return wrapped


def main():
    _auto_num[0] = 0

    builders = [
        # ── §0: Opening (8 slides) ────────────────────────────────────
        slide_01,                                               # 1: Argos VSP
        slide_thesis,                                           # 2: Thesis [NEW]
        slide_what_we_claim,                                    # 3: What we claim — and what we do not claim [NEW]
        with_title(slide_what_was_done_1, "What was done? (1/2)"),  # 4
        with_title(slide_what_was_done_2, "What was done? (2/2)"),  # 5
        hide(slide_exec_summary),                               # 6: HIDDEN
        hide(slide_wer_lies),                                   # 7: HIDDEN
        slide_toc,                                              # 8: Presentation Overview

        # ── §1: The Problem (7 slides) ────────────────────────────────
        slide_02,                                               # 9
        slide_visemes,                                          # 10
        slide_data_flow,                                        # 11
        slide_03,                                               # 12
        slide_04,                                               # 13
        slide_05,                                               # 14
        slide_06,                                               # 15

        # ── §2: Research Findings — Evaluation (20 slides) ────────────
        slide_research_transition,                              # 16
        slide_llm_judge,                                        # 17
        hide(slide_llm_judge_30),                               # 18: HIDDEN
        with_title(slide_judge_ex1, "Judge Example 1 — Verdict Y: Names Swapped, Meaning Kept"),  # 19
        hide(slide_judge_ex2),                                  # 20: HIDDEN
        with_title(slide_judge_ex3, "Judge Example 2 — Verdict P: Domain Drifted, Structure Kept"),  # 21
        hide(slide_judge_ex4),                                  # 22: HIDDEN
        hide(slide_judge_ex5),                                  # 23: HIDDEN
        with_title(slide_judge_ex6, "Judge Example 3 — Verdict P at low IS: Fluent but Wrong Topic"),  # 24
        slide_is_motivation,                                    # 25
        slide_is_intro_a,                                       # 26
        slide_is_intro_b,                                       # 27
        slide_is_intro_c,                                       # 28
        slide_is_weight_rationale,                              # 29
        slide_is_calc_examples,                                 # 30
        slide_is_radar,                                         # 31
        slide_is_wer_scatter,                                   # 32
        slide_is_61pct_headline,                                # 33: NEW
        slide_two_eval_systems,                                 # 34
        slide_context_eval,                                     # 35: "IS: A Calibrated Surrogate Metric"
        hide(slide_disagreement_blind),                         # 36: HIDDEN
        hide(slide_disagreement_context),                       # 37: HIDDEN
        slide_four_numbers,                                     # 38: NEW

        # ── §3: Confidence (10 slides) ────────────────────────────────
        with_title(slide_confidence_problem, "How Do You Know When to Trust an Output?"),  # 39
        slide_two_layer_confidence_research,                    # 40
        with_title(slide_confidence_scoring, "How the Report Handles Uncertainty"),  # 41
        with_title(slide_three_tier_policy_research, "Per-Tier Reliability — What the Numbers Say"),  # 42
        with_title(slide_band_reliability_stratified, "Green Reliability vs Segment Quality — the 0.65 Cliff"),  # 43
        with_title(slide_per_word_confidence_distribution, "Joint Rule vs Legacy: How Band Counts Shifted"),  # 44
        slide_agreement_vs_conf_information,                    # 45
        slide_green_leakage_examples,                           # 46
        slide_is_conf_correlate_combined,                       # 47: merged

        # ── §4: Failure modes + Salvage (11 slides) ───────────────────
        hide(slide_08),                                         # 48: HIDDEN
        slide_failure_deep_1a,                                  # 49
        slide_failure_deep_1b,                                  # 50
        with_title(slide_failure_live_demo, "Failure Modes: Real Examples"),  # 51
        hide(slide_is_deep_dive),                               # 52: HIDDEN
        hide(slide_metric_disagreement),                        # 53: HIDDEN
        hide(slide_metric_disagreement_2),                      # 54: HIDDEN
        hide(with_title(slide_25, "IS: A Calibrated Surrogate for LLM Judgment")),  # 55: HIDDEN
        slide_25d,                                              # 56
        hide(slide_25e),                                        # 57: HIDDEN
        with_title(slide_a15, "Curated Examples — Video Gallery"),  # 58

        # ── §5: Demo (3 slides) ───────────────────────────────────────
        slide_demo_obama_trust,                                 # 59
        with_title(slide_demo_obama_strip, "Demo — SALVAGE: structure preserved, vocabulary lost"),  # 60
        with_title(slide_failure_live_demo, "Demo: Where Word Accuracy and Meaning Diverge"),  # 61

        # ── §6: Engineering (10 slides) ───────────────────────────────
        # Note: this section is titled "ENGINEERING" in v13. Using a divider style transition.
        slide_research_transition,                              # 62: ENGINEERING divider (override label)
        slide_three_repos,                                      # 63
        slide_pipeline_10stage,                                 # 64: merged 17_a/b NEW
        hide(slide_17_png),                                     # 65: HIDDEN
        slide_engineering_journey_v13,                          # 66: NEW (lifted from slides_client)
        hide(slide_18),                                         # 67: HIDDEN
        slide_20,                                               # 68: Standalone Container
        hide(slide_web_ui),                                     # 69: HIDDEN (Live Demo)
        slide_21,                                               # 70: Pipeline Intelligence
        slide_dual_env,                                         # 71

        # ── §7: Future (18 slides) ────────────────────────────────────
        slide_future_transition,                                # 72
        with_title(slide_24, "Starting from 61.6%, Not 25%"),   # 73
        slide_26,                                               # 74
        slide_failure_deep_3,                                   # 75: Failure Modes: Impact & What Fixes Them
        hide(slide_26b),                                        # 76: HIDDEN
        with_title(slide_28, "Aggregating 20 Hypotheses: Majority's Favorite Sentence"),  # 77
        hide(with_title(slide_nbest_v3_judge_paired_tests, "N-best Aggregation: McNemar Paired Tests")),  # 78: HIDDEN
        slide_quality_filter_v13,                               # 79: NEW (lifted from slides_client)
        slide_data_scaling,                                     # 80
        hide(slide_price_tag),                                  # 81: HIDDEN
        hide(with_title(slide_29_phases, "Fine-Tuning: Limited Data, Limited Gains")),  # 82: HIDDEN
        slide_30,                                               # 83
        slide_llm_context_engine,                               # 84
        hide(slide_30b),                                        # 85: HIDDEN
        slide_arabic_roadmap,                                   # 86
        slide_arabic_avhubert,                                  # 87
        slide_arabic_changes,                                   # 88
        slide_31,                                               # 89
        slide_thank_you,                                        # 90

        # ── Appendix (9 slides, renumbered) ──────────────────────────
        slide_a1,                                               # 91: A1
        with_title(slide_a8, "A2: IS Component Correlation"),   # 92
        with_title(slide_a11, "A3: LLM Salvage — Recoverable Segments"),  # 93
        hide(with_title(slide_a11b, "A5: LLM Salvage — Curated Examples")),  # 94: HIDDEN
        with_title(slide_a13, "A4: Failure Mode Examples"),     # 95
        with_title(slide_a15, "A5: Video Gallery — All Example Segments"),  # 96
        with_title(slide_a16, "A6: LLM Judge × IS Tier Cross-Tabulation"),  # 97
        with_title(slide_a17, "A7: Context Evaluation — Transition Details"),  # 98
        with_title(slide_v1_vs_v3_judge_lesson,
                   "A8: Why a Dual-Conf Judge Prompt — a Methodology Lesson"),  # 99
    ]

    print(f"Target: 99 slides (23 hidden)")
    print(f"Builders: {len(builders)}")

    prs = Presentation()
    prs.slide_width = SL_W
    prs.slide_height = SL_H

    for i, builder in enumerate(builders, 1):
        try:
            builder(prs)
        except Exception as e:
            import traceback
            print(f"  Slide {i}/{len(builders)} {builder.__name__}: ERROR — {e}")
            traceback.print_exc()
            raise

    prs.save(str(OUTPUT))
    _fix_pptx_video_compat(str(OUTPUT))
    _strip_orphan_animation_refs(str(OUTPUT))

    # Re-count hidden by inspecting saved file
    saved = Presentation(str(OUTPUT))
    hidden_count = sum(1 for s in saved.slides if s.element.get('show', '1') == '0')
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(saved.slides)} (hidden: {hidden_count})")


if __name__ == "__main__":
    import sys
    sys.path.insert(0, str(Path(__file__).parent))
    main()
