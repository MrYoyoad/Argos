#!/usr/bin/env python3
"""
Comprehensive audit of Argos_VSP_Project_Review.pptx (30-slide generated deck).

Checks:
  1. Physical overlaps (>0.5 sq.in, excluding intentional text-in-container)
  2. Off-screen elements (beyond 13.333" x 7.5")
  3. Too-small text (<9pt)
  4. Empty/missing images (placeholder shapes with [filename] text)
  5. Shape count per slide (flag >25 shapes)
  6. Content verification (title text per slide)
  7. Video/media embedding
  8. Animation groups (click-triggered vs auto-play)
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

PPTX_PATH = "/home/ubuntu/presentation_materials_20260224/Argos_VSP_Project_Review.pptx"

# Slide dimensions (standard 16:9 widescreen)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Thresholds
OVERLAP_THRESHOLD_SQIN = 0.5
MIN_FONT_PT = 9
MAX_SHAPES = 25

# Namespaces
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def emu_to_in(emu):
    return emu / 914400.0


def rect_from_shape(shape):
    """Return (left, top, right, bottom) in inches."""
    try:
        l = emu_to_in(shape.left)
        t = emu_to_in(shape.top)
        r = l + emu_to_in(shape.width)
        b = t + emu_to_in(shape.height)
        return (l, t, r, b)
    except Exception:
        return None


def overlap_area(r1, r2):
    """Compute overlap area in sq.in between two rectangles."""
    x_overlap = max(0, min(r1[2], r2[2]) - max(r1[0], r2[0]))
    y_overlap = max(0, min(r1[3], r2[3]) - max(r1[1], r2[1]))
    return x_overlap * y_overlap


def shape_type_name(shape):
    try:
        return str(shape.shape_type).split('.')[-1].split('(')[0].strip()
    except Exception:
        return "UNKNOWN"


def shape_auto_type(shape):
    """Get the auto-shape sub-type (e.g. ROUNDED_RECTANGLE)."""
    try:
        el = shape._element
        prstGeom = el.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prstGeom is not None:
            return prstGeom.get('prst', '')
        return ''
    except Exception:
        return ''


def is_intentional_overlap(s1, s2):
    """
    Determine if an overlap is intentional (text-inside-container).
    Exclude: Rounded Rect + TextBox, Oval + TextBox, any shape fully contained
    in another that has text (label-on-shape pattern).
    """
    t1 = shape_type_name(s1)
    t2 = shape_type_name(s2)

    r1 = rect_from_shape(s1)
    r2 = rect_from_shape(s2)
    if not r1 or not r2:
        return False

    # Check if one shape fully contains the other
    def contains(outer, inner):
        return (outer[0] <= inner[0] + 0.05 and outer[1] <= inner[1] + 0.05 and
                outer[2] >= inner[2] - 0.05 and outer[3] >= inner[3] - 0.05)

    # Text box inside a container shape
    text_types = {'TEXT_BOX', 'PLACEHOLDER'}
    container_types = {'AUTO_SHAPE', 'FREEFORM', 'ROUNDED_RECTANGLE'}

    if t1 in text_types and contains(r2, r1):
        return True
    if t2 in text_types and contains(r1, r2):
        return True

    # Rounded rectangle + text box is almost always intentional
    auto1 = shape_auto_type(s1)
    auto2 = shape_auto_type(s2)
    rounded = {'roundRect', 'oval', 'ellipse', 'rect'}
    if auto1 in rounded and t2 in text_types:
        return True
    if auto2 in rounded and t1 in text_types:
        return True

    # If one fully contains the other and the inner has text, it's a label
    if contains(r1, r2):
        if hasattr(s2, 'has_text_frame') and s2.has_text_frame and s2.text_frame.text.strip():
            return True
    if contains(r2, r1):
        if hasattr(s1, 'has_text_frame') and s1.has_text_frame and s1.text_frame.text.strip():
            return True

    return False


def get_title(slide):
    """Extract the title text from a slide."""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()

    for shape in slide.shapes:
        if shape.has_text_frame:
            try:
                ph = shape.placeholder_format
                if ph and ph.idx in (0, 15):
                    return shape.text_frame.text.strip()
            except Exception:
                pass

    # Fallback: largest text at the top
    best = None
    best_size = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt > best_size:
                        txt = para.text.strip()
                        if txt and len(txt) < 120:
                            best = txt
                            best_size = run.font.size.pt

    if best:
        return best

    # Last resort: first text found
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            txt = shape.text_frame.text.strip()
            if len(txt) < 120:
                return txt

    return "(no title found)"


def check_small_text(shape):
    """Return list of (run_text, font_size_pt) for runs below MIN_FONT_PT."""
    issues = []
    if not shape.has_text_frame:
        return issues
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            sz = run.font.size
            if sz is not None:
                pt = sz.pt
                if pt < MIN_FONT_PT:
                    snippet = run.text[:60].replace('\n', ' ')
                    if snippet.strip():
                        issues.append((snippet, pt, shape.name))
    return issues


def check_empty_image_placeholder(shape):
    """Check if a shape is an empty image placeholder."""
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip()
        if txt.startswith('[') and txt.endswith(']') and len(txt) < 100:
            return txt
    return None


def get_animation_info(slide):
    """Parse slide XML for animation information."""
    slide_xml = slide._element

    # Main sequence: p:timing/p:tnLst/p:par/p:cTn/p:childTnLst/p:seq
    seq_nodes = slide_xml.findall('.//p:timing//p:seq', NSMAP)
    has_timing = len(slide_xml.findall('.//p:timing', NSMAP)) > 0

    # Count groups under the main sequence
    main_groups = []
    for seq in seq_nodes:
        ctn = seq.find('p:cTn', NSMAP)
        if ctn is not None:
            child_list = ctn.find('p:childTnLst', NSMAP)
            if child_list is not None:
                pars = child_list.findall('p:par', NSMAP)
                main_groups.extend(pars)

    num_groups = len(main_groups)

    # Analyze each group for trigger type
    click_triggered = 0
    auto_play = 0

    for group in main_groups:
        # Look at the start condition
        ctn = group.find('.//p:cTn', NSMAP)
        if ctn is not None:
            stCondLst = ctn.find('p:stCondLst', NSMAP)
            if stCondLst is not None:
                conds = stCondLst.findall('p:cond', NSMAP)
                for cond in conds:
                    evt = cond.get('evt', '')
                    delay = cond.get('delay', '')
                    if evt == 'onClick':
                        click_triggered += 1
                        break
                    elif delay == '0' or evt == 'onBegin':
                        auto_play += 1
                        break
                else:
                    # Default: if no specific event, it's click-triggered
                    click_triggered += 1
            else:
                click_triggered += 1

    # Count individual animation effects
    anim_effects = (
        len(slide_xml.findall('.//p:timing//p:anim', NSMAP)) +
        len(slide_xml.findall('.//p:timing//p:animEffect', NSMAP)) +
        len(slide_xml.findall('.//p:timing//p:set', NSMAP)) +
        len(slide_xml.findall('.//p:timing//p:animMotion', NSMAP)) +
        len(slide_xml.findall('.//p:timing//p:animRot', NSMAP)) +
        len(slide_xml.findall('.//p:timing//p:animScale', NSMAP))
    )

    return {
        'num_groups': num_groups,
        'click_triggered': click_triggered,
        'auto_play': auto_play,
        'total_effects': anim_effects,
        'has_timing': has_timing,
    }


def check_embedded_media(slide):
    """Check for embedded media (video/audio) in slide relationships."""
    media = []
    try:
        for rel in slide.part.rels.values():
            reltype = str(rel.reltype).lower()
            if 'video' in reltype or 'audio' in reltype:
                target = str(rel.target_ref) if hasattr(rel, 'target_ref') else str(rel.target_partname)
                media.append(f"{rel.reltype.split('/')[-1]} -> {target}")
    except Exception:
        pass

    # Also check for movie/video shapes
    for shape in slide.shapes:
        try:
            el = shape._element
            # Look for video elements in the shape XML
            videos = el.findall('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}videoFile')
            if videos:
                media.append(f"videoFile in '{shape.name}'")
            # Check for hlinkClick with action=ppaction://media
            hlinks = el.findall('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick')
            for hl in hlinks:
                action = hl.get('action', '')
                if 'media' in action.lower() or 'ppaction' in action.lower():
                    media.append(f"media link in '{shape.name}': {action}")
        except Exception:
            pass

    return media


def run_audit():
    print("=" * 90)
    print("COMPREHENSIVE PRESENTATION AUDIT")
    print(f"File: {PPTX_PATH}")
    print("=" * 90)

    prs = Presentation(PPTX_PATH)

    sw = emu_to_in(prs.slide_width)
    sh = emu_to_in(prs.slide_height)
    print(f"\nSlide dimensions: {sw:.3f}\" x {sh:.3f}\"  (expected {SLIDE_W_IN}\" x {SLIDE_H_IN}\")")
    if abs(sw - SLIDE_W_IN) > 0.01 or abs(sh - SLIDE_H_IN) > 0.01:
        print("  *** WARNING: Slide dimensions differ from expected! ***")
    else:
        print("  OK")

    print(f"Total slides: {len(prs.slides)}")

    # Counters
    total_overlaps = 0
    total_offscreen = 0
    total_smalltext = 0
    total_emptyimg = 0
    total_cluttered = 0
    slides_with_media = []
    slides_with_anims = []
    all_issues = []

    for idx, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        title = get_title(slide)
        slide_issues = []

        print(f"\n{'─' * 90}")
        print(f"SLIDE {idx:2d} | {len(shapes):2d} shapes | \"{title[:70]}\"")
        print(f"{'─' * 90}")

        # ── 5. Shape count ──
        if len(shapes) > MAX_SHAPES:
            msg = f"CLUTTERED: {len(shapes)} shapes (max {MAX_SHAPES})"
            print(f"  !! {msg}")
            slide_issues.append(msg)
            total_cluttered += 1

        # ── 2. Off-screen elements ──
        offscreen = []
        tol = 0.05
        for shape in shapes:
            r = rect_from_shape(shape)
            if r is None:
                continue
            l, t, ri, b = r
            reasons = []
            if l < -tol:
                reasons.append(f"left={l:.2f}\"")
            if t < -tol:
                reasons.append(f"top={t:.2f}\"")
            if ri > sw + tol:
                reasons.append(f"right={ri:.2f}\" (>{sw:.1f}\")")
            if b > sh + tol:
                reasons.append(f"bottom={b:.2f}\" (>{sh:.1f}\")")
            if reasons:
                offscreen.append((shape.name, shape_type_name(shape),
                                  f"({l:.2f}\", {t:.2f}\", {ri:.2f}\", {b:.2f}\")",
                                  ', '.join(reasons)))

        if offscreen:
            print(f"  OFF-SCREEN ({len(offscreen)}):")
            for name, stype, bounds, reason in offscreen:
                print(f"    - '{name}' [{stype}] {bounds}: {reason}")
            total_offscreen += len(offscreen)
            for item in offscreen:
                slide_issues.append(f"Off-screen: '{item[0]}' {item[3]}")

        # ── 1. Physical overlaps ──
        rects = []
        for shape in shapes:
            r = rect_from_shape(shape)
            if r:
                rects.append((shape, r))

        overlaps = []
        for i in range(len(rects)):
            for j in range(i + 1, len(rects)):
                s1, r1 = rects[i]
                s2, r2 = rects[j]
                area = overlap_area(r1, r2)
                if area > OVERLAP_THRESHOLD_SQIN:
                    if is_intentional_overlap(s1, s2):
                        continue
                    overlaps.append((s1.name, shape_type_name(s1),
                                     s2.name, shape_type_name(s2), area))

        if overlaps:
            print(f"  OVERLAPS > {OVERLAP_THRESHOLD_SQIN} sq.in ({len(overlaps)}):")
            for n1, t1, n2, t2, area in overlaps:
                print(f"    - '{n1}' [{t1}] x '{n2}' [{t2}]: {area:.2f} sq.in")
            total_overlaps += len(overlaps)
            for item in overlaps:
                slide_issues.append(f"Overlap: '{item[0]}' x '{item[2]}': {item[4]:.2f} sq.in")

        # ── 3. Too-small text ──
        small_all = []
        for shape in shapes:
            small = check_small_text(shape)
            small_all.extend(small)

        if small_all:
            print(f"  SMALL TEXT < {MIN_FONT_PT}pt ({len(small_all)}):")
            for snippet, pt, sname in small_all:
                print(f"    - \"{snippet}\" @ {pt:.1f}pt (in '{sname}')")
            total_smalltext += len(small_all)
            for item in small_all:
                slide_issues.append(f"Small text: \"{item[0]}\" @ {item[1]:.1f}pt")

        # ── 4. Empty/missing images ──
        empty_imgs = []
        for shape in shapes:
            result = check_empty_image_placeholder(shape)
            if result:
                empty_imgs.append((shape.name, result))

        if empty_imgs:
            print(f"  EMPTY/MISSING IMAGES ({len(empty_imgs)}):")
            for name, txt in empty_imgs:
                print(f"    - '{name}': {txt}")
            total_emptyimg += len(empty_imgs)
            for item in empty_imgs:
                slide_issues.append(f"Empty image: '{item[0]}': {item[1]}")

        # ── 7. Embedded media ──
        media = check_embedded_media(slide)
        if media:
            print(f"  EMBEDDED MEDIA ({len(media)}):")
            for m in media:
                print(f"    - {m}")
            slides_with_media.append((idx, title, media))

        # ── 8. Animations ──
        anim = get_animation_info(slide)
        if anim['num_groups'] > 0 or anim['total_effects'] > 0:
            trigger_desc = []
            if anim['click_triggered'] > 0:
                trigger_desc.append(f"{anim['click_triggered']} click-triggered")
            if anim['auto_play'] > 0:
                trigger_desc.append(f"{anim['auto_play']} auto-play")
            trigger_str = ', '.join(trigger_desc) if trigger_desc else "unknown trigger"
            print(f"  ANIMATIONS: {anim['num_groups']} group(s), {anim['total_effects']} effect(s) ({trigger_str})")
            slides_with_anims.append((idx, title, anim))
        else:
            print(f"  Animations: none")

        # Status
        if not slide_issues:
            print(f"  >> CLEAN")

        if slide_issues:
            all_issues.append((idx, title, slide_issues))

    # ═══════════════════════════════════════════════════════════════════
    # SUMMARY
    # ═══════════════════════════════════════════════════════════════════
    print(f"\n{'=' * 90}")
    print("AUDIT SUMMARY")
    print(f"{'=' * 90}")
    print(f"Total slides:           {len(prs.slides)}")
    print(f"Off-screen elements:    {total_offscreen}")
    print(f"Unintentional overlaps: {total_overlaps}")
    print(f"Small text (<{MIN_FONT_PT}pt):     {total_smalltext}")
    print(f"Empty/missing images:   {total_emptyimg}")
    print(f"Cluttered slides (>{MAX_SHAPES}):  {total_cluttered}")
    print(f"Slides with media:      {len(slides_with_media)}")
    print(f"Slides with animations: {len(slides_with_anims)}")

    total_issue_count = total_overlaps + total_offscreen + total_smalltext + total_emptyimg + total_cluttered
    print(f"\nTOTAL ISSUES: {total_issue_count}")

    # ── Slides with issues ──
    if all_issues:
        print(f"\n{'─' * 90}")
        print(f"SLIDES WITH ISSUES ({len(all_issues)} slides)")
        print(f"{'─' * 90}")
        for slide_num, title, issues in all_issues:
            print(f"  Slide {slide_num:2d} — \"{title[:60]}\" ({len(issues)} issue(s)):")
            for iss in issues:
                print(f"    - {iss}")
    else:
        print("\n  All slides are clean!")

    # ── Title index ──
    print(f"\n{'─' * 90}")
    print("SLIDE TITLE INDEX")
    print(f"{'─' * 90}")
    for idx, slide in enumerate(prs.slides, 1):
        title = get_title(slide)
        shapes = list(slide.shapes)
        flags = []
        # Check if this slide has issues
        for sn, st, si in all_issues:
            if sn == idx:
                flags.append(f"{len(si)} issues")
        for sn, st, sm in slides_with_media:
            if sn == idx:
                flags.append("MEDIA")
        for sn, st, sa in slides_with_anims:
            if sn == idx:
                flags.append(f"ANIM({sa['num_groups']}g)")
        flag_str = f"  [{', '.join(flags)}]" if flags else ""
        print(f"  {idx:2d}. ({len(shapes):2d} shapes) {title[:68]}{flag_str}")

    # ── Media summary ──
    print(f"\n{'─' * 90}")
    print(f"MEDIA EMBEDDING SUMMARY")
    print(f"{'─' * 90}")
    if slides_with_media:
        for sn, st, sm in slides_with_media:
            print(f"  Slide {sn:2d}: \"{st[:55]}\"")
            for m in sm:
                print(f"           {m}")
    else:
        print("  No embedded media found in any slide.")

    # ── Animation summary ──
    print(f"\n{'─' * 90}")
    print(f"ANIMATION SUMMARY")
    print(f"{'─' * 90}")
    if slides_with_anims:
        print(f"  {'Slide':>5s}  {'Groups':>6s}  {'Click':>5s}  {'Auto':>4s}  {'Effects':>7s}  Title")
        print(f"  {'─'*5}  {'─'*6}  {'─'*5}  {'─'*4}  {'─'*7}  {'─'*40}")
        for sn, st, sa in slides_with_anims:
            print(f"  {sn:5d}  {sa['num_groups']:6d}  {sa['click_triggered']:5d}  {sa['auto_play']:4d}  {sa['total_effects']:7d}  {st[:45]}")
        total_groups = sum(sa['num_groups'] for _, _, sa in slides_with_anims)
        total_effects = sum(sa['total_effects'] for _, _, sa in slides_with_anims)
        print(f"  {'─'*5}  {'─'*6}  {'─'*5}  {'─'*4}  {'─'*7}")
        print(f"  TOTAL  {total_groups:6d}  {sum(sa['click_triggered'] for _,_,sa in slides_with_anims):5d}  {sum(sa['auto_play'] for _,_,sa in slides_with_anims):4d}  {total_effects:7d}")
    else:
        print("  No animations found in any slide.")

    print(f"\n{'=' * 90}")
    print("AUDIT COMPLETE")
    print(f"{'=' * 90}")

    return total_issue_count


if __name__ == "__main__":
    issues = run_audit()
    sys.exit(0 if issues == 0 else 1)
