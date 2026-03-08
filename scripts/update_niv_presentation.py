#!/usr/bin/env python3
"""
Comprehensive PPTX Update Script — 23 Items
Reads PRESNETATION_AFTER_NIV.pptx, applies all changes, saves as Argos_VSP_Final.pptx.
"""

import copy
import csv
import os
import re
import sys
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Constants ────────────────────────────────────────────────────────────────

INPUT_PATH = "presentation_materials_20260224/PRESNETATION_AFTER_NIV.pptx"
OUTPUT_PATH = "presentation_materials_20260224/Argos_VSP_Final.pptx"

# Colors
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
NAVY2 = RGBColor(0x14, 0x25, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
CORAL = RGBColor(0xE0, 0x6C, 0x75)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
GOLD = RGBColor(0xFF, 0xD5, 0x4F)
LGRAY = RGBColor(0xAA, 0xAA, 0xAA)
MGRAY = RGBColor(0x99, 0x99, 0x99)
DGRAY = RGBColor(0x66, 0x66, 0x66)

# Slide dimensions (widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MX = Inches(0.60)  # margin x
CW = Inches(12.13)  # content width
CT = Inches(1.45)  # content top

# OOXML namespaces
nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


# ─── Helpers ──────────────────────────────────────────────────────────────────

def find_shape_by_text(slide, text_fragment, case_sensitive=True):
    """Find shape containing text fragment."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = shape.text_frame.text
            if case_sensitive:
                if text_fragment in full:
                    return shape
            else:
                if text_fragment.lower() in full.lower():
                    return shape
    return None


def find_shapes_by_text(slide, text_fragment, case_sensitive=True):
    """Find ALL shapes containing text fragment."""
    results = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = shape.text_frame.text
            if case_sensitive:
                if text_fragment in full:
                    results.append(shape)
            else:
                if text_fragment.lower() in full.lower():
                    results.append(shape)
    return results


def replace_text_in_shape(shape, old, new):
    """Replace text in a shape, preserving formatting."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                changed = True
    return changed


def replace_text_in_slide(slide, old, new):
    """Replace text across all shapes in a slide."""
    changed = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            if replace_text_in_shape(shape, old, new):
                changed = True
    return changed


def replace_in_notes(slide, old, new):
    """Replace text in slide notes."""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame
        for para in notes.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)


def _ensure_notes_placeholder(notes_slide):
    """Ensure notes slide has a notes body placeholder. Creates one if missing."""
    if notes_slide.notes_text_frame is not None:
        return notes_slide.notes_text_frame
    # Create a notes body placeholder (ph type=body, idx=1)
    spTree = notes_slide._element.find('.//' + '{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    if spTree is None:
        return None
    nsmap_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp = etree.SubElement(spTree, f'{{{nsmap_p}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{nsmap_p}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{nsmap_p}}}cNvPr', attrib={'id': '3', 'name': 'Notes Placeholder 2'})
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{nsmap_p}}}cNvSpPr')
    spLocks = etree.SubElement(cNvSpPr, f'{{{nsmap_a}}}spLocks', attrib={'noGrp': '1'})
    nvPr = etree.SubElement(nvSpPr, f'{{{nsmap_p}}}nvPr')
    ph = etree.SubElement(nvPr, f'{{{nsmap_p}}}ph', attrib={'type': 'body', 'idx': '1'})
    spPr = etree.SubElement(sp, f'{{{nsmap_p}}}spPr')
    txBody = etree.SubElement(sp, f'{{{nsmap_a}}}txBody')
    bodyPr = etree.SubElement(txBody, f'{{{nsmap_a}}}bodyPr')
    lstStyle = etree.SubElement(txBody, f'{{{nsmap_a}}}lstStyle')
    p = etree.SubElement(txBody, f'{{{nsmap_a}}}p')
    endParaRPr = etree.SubElement(p, f'{{{nsmap_a}}}endParaRPr', attrib={'lang': 'en-US'})
    return notes_slide.notes_text_frame


def set_notes(slide, text):
    """Set slide notes (replacing all content)."""
    try:
        notes_slide = slide.notes_slide  # creates if needed
    except Exception:
        return  # Can't create notes for this slide type
    notes_tf = _ensure_notes_placeholder(notes_slide)
    if notes_tf is None:
        return
    # Clear existing
    for i in range(len(notes_tf.paragraphs) - 1, 0, -1):
        p = notes_tf.paragraphs[i]._p
        p.getparent().remove(p)
    notes_tf.paragraphs[0].text = text


def append_to_notes(slide, text):
    """Append text to existing slide notes."""
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return
    notes_tf = notes_slide.notes_text_frame
    if notes_tf is None:
        return
    existing = notes_tf.text.strip()
    if existing:
        set_notes(slide, existing + "\n\n" + text)
    else:
        set_notes(slide, text)


def add_textbox(slide, text, left, top, width, height,
                size=Pt(14), color=WHITE, bold=False, italic=False,
                align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_textbox(slide, bullets, left, top, width, height,
                       size=Pt(14), color=WHITE, font_name="Calibri"):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet_text}"
        p.font.size = size
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox


def add_animated_bullets(slide, bullets, left, top, width, line_height=Inches(0.55),
                         size=Pt(15), text_color=WHITE, bullet_color=None,
                         font_name="Calibri"):
    """Add bullet points as individual text boxes with click-to-appear animations.

    Each bullet appears on click. bullet_color controls the bullet marker color
    (defaults to TEAL). Returns list of created text boxes.
    """
    if bullet_color is None:
        bullet_color = RGBColor(0x00, 0xB4, 0xD8)  # TEAL
    boxes = []
    for i, bullet_text in enumerate(bullets):
        y = top + line_height * i
        txBox = slide.shapes.add_textbox(left, y, width, Inches(0.50))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # Colored bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = size
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = font_name
        run_bullet.font.bold = True
        # Text content
        run_text = p.add_run()
        run_text.text = bullet_text
        run_text.font.size = size
        run_text.font.color.rgb = text_color
        run_text.font.name = font_name
        add_click_animation(slide, txBox)
        boxes.append(txBox)
    return boxes


def set_shape_text(shape, text, preserve_format=True):
    """Set shape text, optionally preserving first run's formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if preserve_format and tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        # Remove extra paragraphs
        for i in range(len(tf.paragraphs) - 1, 0, -1):
            p = tf.paragraphs[i]._p
            p.getparent().remove(p)
        # Remove extra runs
        first_p = tf.paragraphs[0]
        for i in range(len(first_p.runs) - 1, 0, -1):
            r = first_p.runs[i]._r
            r.getparent().remove(r)
        first_p.runs[0].text = text
    else:
        tf.paragraphs[0].text = text


def set_bullets_text(shape, bullets, color=None):
    """Replace bullet list text in a shape, preserving formatting of first run per para."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Preserve format from first paragraph's first run
    fmt = {}
    if tf.paragraphs and tf.paragraphs[0].runs:
        r0 = tf.paragraphs[0].runs[0]
        fmt['size'] = r0.font.size
        fmt['bold'] = r0.font.bold
        fmt['name'] = r0.font.name
        if color:
            fmt['color'] = color
        elif r0.font.color and r0.font.color.rgb:
            fmt['color'] = r0.font.color.rgb

    # Clear all paragraphs except first
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)

    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            # Clear extra runs
            for j in range(len(p.runs) - 1, 0, -1):
                r = p.runs[j]._r
                r.getparent().remove(r)
            if p.runs:
                p.runs[0].text = text
            else:
                p.text = text
        else:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = text

        # Apply formatting
        for run in p.runs:
            if 'size' in fmt:
                run.font.size = fmt['size']
            if 'bold' in fmt:
                run.font.bold = fmt['bold']
            if 'name' in fmt:
                run.font.name = fmt['name']
            if 'color' in fmt:
                run.font.color.rgb = fmt['color']


def is_hidden(slide):
    """Check if a slide is hidden."""
    sld = slide._element
    show = sld.get('show')
    return show == '0'


def add_click_animation(slide, shape):
    """Add a click-to-appear animation to a shape using OOXML."""
    # Get the shape's spTree and the slide's timing element
    sld = slide._element
    sp = shape._element

    # Get shape ID
    sp_id = sp.find('.//p:cNvPr', nsmap)
    if sp_id is None:
        sp_id = sp.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}cNvPr')
    if sp_id is None:
        # Try direct child nvSpPr/cNvPr
        nvSpPr = sp.find('{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
        if nvSpPr is None:
            nvSpPr = sp.find('{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}nvSpPr')
        if nvSpPr is not None:
            sp_id = nvSpPr.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            if sp_id is None:
                sp_id = nvSpPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')

    # Fallback: search more broadly
    if sp_id is None:
        for elem in sp.iter():
            if elem.tag.endswith('}cNvPr') or elem.tag == 'cNvPr':
                sp_id = elem
                break

    if sp_id is None:
        print(f"  Warning: Could not find shape ID for animation, skipping")
        return

    shape_id = sp_id.get('id')
    shape_name = sp_id.get('name', 'Shape')

    # Check if timing already exists
    timing = sld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
    if timing is None:
        timing = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}timing')

    tnLst = timing.find('{http://schemas.openxmlformats.org/presentationml/2006/main}tnLst')
    if tnLst is None:
        tnLst = etree.SubElement(timing, '{http://schemas.openxmlformats.org/presentationml/2006/main}tnLst')

    # Get or create the main sequence
    par = tnLst.find('{http://schemas.openxmlformats.org/presentationml/2006/main}par')
    if par is None:
        par = etree.SubElement(tnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}par')

    cTn = par.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cTn')
    if cTn is None:
        cTn = etree.SubElement(par, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn',
                               attrib={'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'})

    childTnLst = cTn.find('{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
    if childTnLst is None:
        childTnLst = etree.SubElement(cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')

    # Find or create the main sequence node
    seq = childTnLst.find('{http://schemas.openxmlformats.org/presentationml/2006/main}seq')
    if seq is None:
        seq = etree.SubElement(childTnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}seq',
                               attrib={'concurrent': '1', 'nextAc': 'seek'})
        seq_cTn = etree.SubElement(seq, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn',
                                   attrib={'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'})
        etree.SubElement(seq_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
        # Add prevCondLst and nextCondLst
        prevCond = etree.SubElement(seq, '{http://schemas.openxmlformats.org/presentationml/2006/main}prevCondLst')
        cond_prev = etree.SubElement(prevCond, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond',
                                     attrib={'evt': 'onPrev', 'delay': '0'})
        etree.SubElement(cond_prev, '{http://schemas.openxmlformats.org/presentationml/2006/main}tgtEl').append(
            etree.SubElement(etree.Element('dummy'), '{http://schemas.openxmlformats.org/presentationml/2006/main}sldTgt'))
        nextCond = etree.SubElement(seq, '{http://schemas.openxmlformats.org/presentationml/2006/main}nextCondLst')
        cond_next = etree.SubElement(nextCond, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond',
                                     attrib={'evt': 'onNext', 'delay': '0'})
        etree.SubElement(cond_next, '{http://schemas.openxmlformats.org/presentationml/2006/main}tgtEl').append(
            etree.SubElement(etree.Element('dummy'), '{http://schemas.openxmlformats.org/presentationml/2006/main}sldTgt'))
    else:
        seq_cTn = seq.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cTn')

    seq_childTnLst = seq_cTn.find('{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
    if seq_childTnLst is None:
        seq_childTnLst = etree.SubElement(seq_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')

    # Count existing animations to get next ID
    existing_count = len(seq_childTnLst)
    base_id = 3 + existing_count * 5

    # Build the animation XML for appear-on-click
    # Level 1 par (click trigger)
    par1 = etree.SubElement(seq_childTnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}par')
    cTn1 = etree.SubElement(par1, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn',
                            attrib={'id': str(base_id), 'fill': 'hold'})
    stCondLst1 = etree.SubElement(cTn1, '{http://schemas.openxmlformats.org/presentationml/2006/main}stCondLst')
    etree.SubElement(stCondLst1, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond',
                     attrib={'delay': 'indefinite'})
    childTnLst1 = etree.SubElement(cTn1, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')

    # Level 2 par
    par2 = etree.SubElement(childTnLst1, '{http://schemas.openxmlformats.org/presentationml/2006/main}par')
    cTn2 = etree.SubElement(par2, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn',
                            attrib={'id': str(base_id + 1), 'fill': 'hold'})
    stCondLst2 = etree.SubElement(cTn2, '{http://schemas.openxmlformats.org/presentationml/2006/main}stCondLst')
    etree.SubElement(stCondLst2, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond',
                     attrib={'delay': '0'})
    childTnLst2 = etree.SubElement(cTn2, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')

    # Level 3 par (the actual animation)
    par3 = etree.SubElement(childTnLst2, '{http://schemas.openxmlformats.org/presentationml/2006/main}par')
    cTn3 = etree.SubElement(par3, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn',
                            attrib={'id': str(base_id + 2), 'presetID': '1', 'presetClass': 'entr',
                                    'presetSubtype': '0', 'fill': 'hold',
                                    'nodeType': 'clickEffect'})
    stCondLst3 = etree.SubElement(cTn3, '{http://schemas.openxmlformats.org/presentationml/2006/main}stCondLst')
    etree.SubElement(stCondLst3, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond',
                     attrib={'delay': '0'})
    childTnLst3 = etree.SubElement(cTn3, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')

    # Set effect (appear)
    aset = etree.SubElement(childTnLst3, '{http://schemas.openxmlformats.org/presentationml/2006/main}set')
    cBhvr = etree.SubElement(aset, '{http://schemas.openxmlformats.org/presentationml/2006/main}cBhvr')
    cTn_set = etree.SubElement(cBhvr, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn',
                                attrib={'id': str(base_id + 3), 'dur': '1', 'fill': 'hold'})
    stCondLst_set = etree.SubElement(cTn_set, '{http://schemas.openxmlformats.org/presentationml/2006/main}stCondLst')
    etree.SubElement(stCondLst_set, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond',
                     attrib={'delay': '0'})

    tgtEl = etree.SubElement(cBhvr, '{http://schemas.openxmlformats.org/presentationml/2006/main}tgtEl')
    spTgt = etree.SubElement(tgtEl, '{http://schemas.openxmlformats.org/presentationml/2006/main}spTgt',
                             attrib={'spid': shape_id})
    attrNameLst = etree.SubElement(cBhvr, '{http://schemas.openxmlformats.org/presentationml/2006/main}attrNameLst')
    attrName = etree.SubElement(attrNameLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}attrName')
    attrName.text = 'style.visibility'

    to = etree.SubElement(aset, '{http://schemas.openxmlformats.org/presentationml/2006/main}to')
    val = etree.SubElement(to, '{http://schemas.openxmlformats.org/presentationml/2006/main}val')
    # Use the 'val' attribute instead of text
    # Actually for <p:to><p:val val="visible"/></p:to>
    # The correct structure is <p:set><p:to><p:strVal val="visible"/></p:to></p:set>
    # Let me fix: remove <p:val> and use proper structure
    to.getparent().remove(to)
    to2 = etree.SubElement(aset, '{http://schemas.openxmlformats.org/presentationml/2006/main}to')
    strVal = etree.SubElement(to2, '{http://schemas.openxmlformats.org/presentationml/2006/main}strVal',
                              attrib={'val': 'visible'})


def hide_slide(slide):
    """Mark a slide as hidden."""
    slide._element.set('show', '0')


def add_rounded_rect(slide, left, top, width, height, fill_color=NAVY2,
                     border_color=TEAL, border_width=Pt(2)):
    """Add a rounded rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    return shape


# ─── Item 1: Professional Language Rewrite (Slides 2-3) ──────────────────────

def rewrite_what_was_done(prs):
    """Item 1: Clean up English on slides 2-3, fix coloring."""
    print("Item 1: Rewriting slides 2-3...")

    slide2 = prs.slides[1]  # 0-indexed
    slide3 = prs.slides[2]

    # Slide 2 bullet replacements
    s2_replacements = [
        ("Started from a conceptual paper, with little to no documentation, no working environments and no pipeline",
         "Started from a research paper with no working environment, no documentation, and no pipeline"),
        ("Standard metric (WER) was the basic metric – which is unsatisfactory to actually measure what is needed",
         "Standard metric (WER) proved insufficient for measuring actual output quality"),
        ("Built a complete end to end pipeline, including serious pre processing (STT, Face crop and post processing",
         "Built a complete end-to-end pipeline, including preprocessing (STT, face crop) and post-processing"),
        ("Created a working environment and migrated into a container on a standalone container including complete and professional UI handling",
         "Created a working environment and deployed to a standalone container with a professional UI"),
        ("Evaluated extensively the performance of the model, including the ideation of a new metric that describes if meaning is preserved.",
         "Evaluated the model extensively, including designing a new metric that measures whether meaning is preserved"),
        ("Started fine tuning the model including environment and data",
         "Started fine-tuning the model, including environment setup and data preparation"),
        ("Created clear future clear plan on how to improve performance + generalize from English to Arabic",
         "Created a clear plan to improve performance and generalize from English to Arabic"),
    ]

    for old, new in s2_replacements:
        # Try exact match first, then partial
        if not replace_text_in_slide(slide2, old, new):
            # Try with leading bullet
            replace_text_in_slide(slide2, old.lstrip("• "), new)

    # Slide 3 bullet replacements
    s3_replacements = [
        ("Proved that the model performs well, and makes about 65% of videos to be useful by LLM as a judge",
         "Proved the model performs well \u2014 about 65% of videos are useful by LLM judge assessment"),
        ("the IS metric provided shows high agreement with LLM as a judge, and can be computed on the standalone computer",
         "The IS metric shows high agreement with the LLM judge and can be computed on the standalone computer"),
        ("Semantic meaning, phonetic meaning and NEA are the critical factors in understanding when the model performs well.",
         "Semantic meaning, phonetic similarity, and named entity accuracy are the critical factors in understanding model performance"),
        ("Full failure analysis was made including suggested improvements",
         "Full failure analysis completed with suggested improvements"),
        ("Fully reproduceable container building + model deployment strategy between AWS and standalone computer",
         "Fully reproducible container build and model deployment between AWS and standalone computer"),
        ("Weeks away from vastly improving the base model with confidence scoring, aggregating outputs and stronger LLM",
         "Close to improving the base model through confidence scoring, output aggregation, and a stronger LLM"),
        ("Full plan to repeat the cookbook to produce an Arabic model in 2-3 months",
         "Full plan to replicate the approach for an Arabic model in 2\u20133 months"),
    ]

    for old, new in s3_replacements:
        if not replace_text_in_slide(slide3, old, new):
            replace_text_in_slide(slide3, old.lstrip("• "), new)

    # Fix coloring: change all-teal bold bullets to white text
    for slide in [slide2, slide3]:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color.rgb == TEAL and run.font.bold:
                            run.font.color.rgb = WHITE


# ─── Item 4: LLM-as-a-Judge Explanation (Slide 16) ───────────────────────────

def expand_llm_judge_slide(prs):
    """Item 4: Fill in LLM judge methodology on slide 16."""
    print("Item 4: Expanding LLM judge slide 16...")

    slide = prs.slides[15]  # 0-indexed

    # Remove placeholder text "More explanation here! For claude"
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "More explanation here" in shape.text_frame.text or "For claude" in shape.text_frame.text:
                # Clear this shape and replace with methodology
                tf = shape.text_frame
                tf.clear()
                lines = [
                    "Methodology:",
                    "• Claude Opus received each ref+hyp pair blind (no metrics visible)",
                    "• 3-level holistic judgment: Y (fully conveyed), P (partial), N (lost)",
                    "• 30 duplicate pairs for reliability: 86.7% exact agreement",
                    "• \u03ba = 0.690 (Y threshold) and \u03ba = 0.818 (Y+P threshold)",
                    "• Used as gold standard to calibrate IS thresholds",
                ]
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(13) if i > 0 else Pt(15)
                    p.font.color.rgb = WHITE if i == 0 else LGRAY
                    p.font.bold = (i == 0)
                    p.font.name = "Calibri"
                    p.space_after = Pt(4)
                break

    # Update speaker notes
    append_to_notes(slide,
        "LLM-as-a-Judge methodology: Claude Opus 4.6 evaluated all 1,497 ref+hyp pairs in a blind setting "
        "(no metrics, no IS scores visible). Each pair received a 3-level holistic judgment: Y (meaning clearly conveyed), "
        "P (partial \u2014 some meaning preserved), N (meaning lost or fabricated). "
        "Intra-rater reliability tested with 30 duplicate pairs: 86.7% exact agreement. "
        "Results: Y=345 (23.0%), P=626 (41.8%), N=526 (35.1%). Y+P=64.9%. "
        "NIV threshold calibration: IS \u2265 3.80 maps to Y (\u03ba=0.690), IS \u2265 2.00 maps to Y+P (\u03ba=0.818). "
        "IS beats WER at both operating points: \u03ba 0.690 vs 0.629 (Y), 0.818 vs 0.777 (Y+P).")


# ─── Item 5: IS Motivation Slide (Slide 18) ──────────────────────────────────

def fill_is_motivation_slide(prs):
    """Item 5: Fill the empty 'Why LLM as a judge is not enough?' slide."""
    print("Item 5: Filling IS motivation slide 18...")

    slide = prs.slides[17]  # 0-indexed

    # Clear everything and rebuild cleanly
    clear_slide_content(slide)

    # Title
    add_textbox(slide, "Why LLM as a Judge Is Not Enough",
                MX, Inches(0.35), CW, Inches(0.65),
                size=Pt(32), color=WHITE, bold=True)

    # Subtitle
    add_textbox(slide, "Five reasons we built the Intelligibility Score (IS) framework",
                MX, Inches(0.95), CW, Inches(0.35),
                size=Pt(15), color=MGRAY, italic=True)

    # Accent line
    line = slide.shapes.add_shape(1, MX, Inches(1.32), CW, Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # 5 reason cards — two columns
    reasons = [
        ("\u2460 Deployment Constraint",
         "IS runs offline — no API, no GPU, no internet.\nEssential for air-gapped container deployment."),
        ("\u2461 Determinism",
         "IS produces identical scores every time.\nLLM judges vary ~13% on re-evaluation."),
        ("\u2462 Continuous Signal",
         "IS is 0.0–5.0 continuous from 6 signals.\nLLM judge outputs coarse Y/P/N categories."),
        ("\u2463 Known Biases",
         "12+ documented systematic biases in LLM judges.\nIS is a fixed formula with none."),
        ("\u2464 Decomposability",
         "IS breaks into 6 named signals → diagnose\nexactly what failed. LLM returns verdict only."),
    ]

    # Layout: 2 columns, 3 rows (last row has 1 centered card)
    col_w = Inches(5.85)
    card_h = Inches(1.25)
    row_gap = Inches(0.15)
    col_gap = Inches(0.43)
    start_y = Inches(1.55)

    for i, (title, body) in enumerate(reasons):
        row = i // 2
        col = i % 2
        if i == 4:
            # Center the 5th card
            x = MX + (CW - col_w) / 2
        else:
            x = MX + col * (col_w + col_gap)
        y = start_y + row * (card_h + row_gap)

        # Card background
        card = slide.shapes.add_shape(1, x, y, col_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY2
        card.line.fill.background()

        # Card title
        add_textbox(slide, title,
                    x + Inches(0.20), y + Inches(0.12), col_w - Inches(0.40), Inches(0.30),
                    size=Pt(16), color=TEAL, bold=True)

        # Card body
        add_textbox(slide, body,
                    x + Inches(0.20), y + Inches(0.45), col_w - Inches(0.40), Inches(0.75),
                    size=Pt(13), color=LGRAY)

    # Bottom takeaway
    add_textbox(slide,
        "IS runs in production; LLM-as-a-Judge audits the IS framework. You need both.",
        MX, Inches(5.85), CW, Inches(0.40),
        size=Pt(15), color=GOLD, bold=True, italic=True, align=PP_ALIGN.CENTER)

    # Speaker notes
    set_notes(slide,
        "Why IS, not just LLM judge? Five reasons from our analysis doc:\n"
        "1. Deployment: IS runs on bare Python, no GPU/API/internet — essential for standalone container.\n"
        "2. Determinism: IS gives exact same score every time. LLM judges have ~13% inconsistency.\n"
        "3. Continuous signal: IS is 0-5 continuous vs coarse Y/P/N. Enables statistical testing.\n"
        "4. Known biases: 12+ documented LLM judge biases (position, verbosity, self-enhancement). IS has none.\n"
        "5. Decomposability: IS breaks into 6 components mapping to specific failure modes. "
        "When IS drops, you can query exactly which signal failed.\n\n"
        "Bottom line: IS is the operational metric (deployed, deterministic, free). "
        "LLM judge is the validation oracle (catches edge cases the formula misses).")


# ─── Item 11: LRS3 Reproduction Comment (Slide 11) ───────────────────────────

def add_lrs3_comment(prs):
    """Item 11: Add LRS3 reproduction note to slide 11."""
    print("Item 11: Adding LRS3 reproduction comment...")

    slide = prs.slides[10]  # 0-indexed

    # Add a note bullet near the bottom
    add_textbox(slide,
        "Note: Our best LRS3 reproduction achieved 32% WER \u2014 gap likely due to pretrain/test split differences.",
        MX, Inches(6.0), CW, Inches(0.4),
        size=Pt(11), color=DGRAY, italic=True)

    # Update speaker notes
    append_to_notes(slide,
        "LRS3 reproduction: We tested 4 decode variants on 197 LRS3 segments. "
        "Best result was 32% WER (V1 with all non-empty outputs). "
        "Gap from paper's 25.4% likely due to pretrain vs test split overlap and small sample size. "
        "Details in docs/evaluation/lrs3_decode_experiment.md.")


# ─── Item 13: Reframe Three Numbers (Slide 28) ───────────────────────────────

def reframe_three_numbers(prs):
    """Item 13: Change three numbers from error rates to 'how many videos useful'."""
    print("Item 13: Reframing three numbers slide 28...")

    slide = prs.slides[27]  # 0-indexed

    # Replace the three number values and descriptions
    # Box 1: 64.1% → 25.5%
    replace_text_in_slide(slide, "64.1%", "25.5%")
    replace_text_in_slide(slide, "What WER reports", "By WER (\u226434%)")
    replace_text_in_slide(slide, '"Two-thirds of words are wrong"',
                         "381 of 1,497 videos appear useful\nWER says only 1 in 4 works")
    replace_text_in_slide(slide, "Two-thirds of words are wrong",
                         "381 of 1,497 videos appear useful\nWER says only 1 in 4 works")

    # Box 2: 39.9% → 61.6%
    replace_text_in_slide(slide, "39.9%", "61.6%")
    replace_text_in_slide(slide, "What IS reveals: 597 of 1,497 segments",
                         "By IS (\u22652.00): 922 of 1,497 videos")
    replace_text_in_slide(slide, "deliver genuinely useful output",
                         "deliver useful meaning\nIS reveals 3 in 5 carry meaning")

    # Box 3: keep 64.9%
    replace_text_in_slide(slide, "Opus-as-a-Judge confirms: Y+P = 971/1,497",
                         "By LLM Judge (Y+P): 971 of 1,497")
    replace_text_in_slide(slide, "2 in 3 segments deliver useful output",
                         "confirmed useful\nExpert judgment confirms 2 in 3 work")

    # Update speaker notes
    set_notes(slide,
        "Three numbers answering the same question: 'How many videos are useful?' "
        "WER (\u226434%): 25.5% (381/1,497) \u2014 most pessimistic, uses NIV-calibrated threshold. "
        "IS (\u22652.00): 61.6% (922/1,497) \u2014 captures useful meaning. "
        "LLM Judge (Y+P): 64.9% (971/1,497) \u2014 expert confirms. "
        "Progressive revelation: WER dramatically understates quality.")


# ─── Item 16: LLM Swap Requires Training ─────────────────────────────────────

def clarify_llm_swap_training(prs):
    """Item 16: Emphasize LLM upgrade requires training across slides 58,60-62."""
    print("Item 16: Clarifying LLM swap requires training...")

    # We need to find these slides by content since numbering may shift
    for i, slide in enumerate(prs.slides):
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                title_text += shape.text_frame.text + " "

        # Slide with "LLM Swap" or "drop-in" text
        if "drop-in" in title_text.lower() or "LLM Swap" in title_text:
            replace_text_in_slide(slide, "drop-in swap", "architecture-compatible upgrade")
            replace_text_in_slide(slide, "Drop-in swap", "Architecture-compatible upgrade")
            replace_text_in_slide(slide, "1–2 hours setup", "requires adapter retraining")
            replace_text_in_slide(slide, "1-2 hours setup", "requires adapter retraining")
            replace_text_in_slide(slide, "1\u20132 hours", "2\u20134 weeks + retraining")

        if "swappable" in title_text.lower() or "is swappable" in title_text.lower():
            replace_text_in_slide(slide, "is swappable", "is architecture-compatible")
            replace_text_in_slide(slide, "LLM is swappable", "LLM is upgradeable")
            replace_text_in_slide(slide, "trivial 1-2 hour swap", "requires LoRA adapter retraining")

        # Price tag slide
        if "Price Tag" in title_text or "Quick win" in title_text:
            replace_text_in_slide(slide, "Quick win: LLM swap only", "LLM upgrade + retrain")
            replace_text_in_slide(slide, "Quick win", "LLM upgrade")
            replace_text_in_slide(slide, "~$0", "~$5\u201310K")

        # "LLM Swap (immediate)" column header
        replace_text_in_slide(slide, "LLM Swap (immediate)", "LLM Upgrade (requires training)")
        replace_text_in_slide(slide, "Setup: 1-2 hours", "Training: ~2\u20134 weeks with 5K+ segments")
        replace_text_in_slide(slide, "Setup: 1–2 hours", "Training: ~2\u20134 weeks with 5K+ segments")

    # Also fix slide 9 speaker notes that mention "trivial 1-2 hour swap"
    slide9 = prs.slides[8]
    replace_in_notes(slide9, "trivial 1-2 hour swap", "architecture-compatible upgrade (requires LoRA retraining)")
    replace_in_notes(slide9, "swappable — Llama 3.1 8B is a drop-in replacement",
                     "upgradeable — Llama 3.1 8B has the same 4096 hidden size (requires adapter retraining)")

    # Fix slide 9 body text
    replace_text_in_slide(prs.slides[8],
        "LLM is swappable — Llama 3.1 8B is a drop-in replacement (same 4096 hidden size).",
        "LLM is upgradeable — Llama 3.1 8B has the same 4096 hidden size. Requires LoRA adapter retraining.")


# ─── Item 17: Arabic Risk Downgrade ──────────────────────────────────────────

def downgrade_arabic_risk(prs):
    """Item 17: Downgrade Arabic risk, shorten timeline to 2-3 months."""
    print("Item 17: Downgrading Arabic risk...")

    for i, slide in enumerate(prs.slides):
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                title_text += shape.text_frame.text + " "

        if "Arabic" in title_text and ("Roadmap" in title_text or "Pipeline" in title_text or "Replication" in title_text):
            replace_text_in_slide(slide, "RISK", "minor integration")
            replace_text_in_slide(slide, "(RISK)", "(minor integration)")
            replace_text_in_slide(slide, "3–5 months", "2\u20133 months")
            replace_text_in_slide(slide, "3-5 months", "2\u20133 months")
            replace_text_in_slide(slide, "Realistic estimate: 3–5 months", "Realistic estimate: 2\u20133 months")
            replace_text_in_slide(slide, "Realistic estimate: 3-5 months", "Realistic estimate: 2\u20133 months")

            # Downgrade risk in notes too
            replace_in_notes(slide, "RISK", "minor integration")
            replace_in_notes(slide, "3–5 months", "2\u20133 months")
            replace_in_notes(slide, "3-5 months", "2\u20133 months")
            replace_in_notes(slide, "high risk", "manageable")
            replace_in_notes(slide, "significant risk", "manageable effort")


# ─── Item 20: System + Human > Expert Animation (Slide 7) ────────────────────

def add_human_expert_animation(prs):
    """Item 20: Add 'system + human > expert' as appearing animation on slide 7."""
    print("Item 20: Adding human expert animation on slide 7...")

    slide = prs.slides[6]  # 0-indexed

    # Add text box at bottom (initially hidden via animation)
    txBox = add_textbox(slide,
        "System + human reader outperforms expert lip readers: "
        "55\u201370% vs 45\u201352% word accuracy, with near-zero hallucination risk",
        MX, Inches(5.8), CW, Inches(0.6),
        size=Pt(14), color=GOLD, bold=True, italic=True)

    # Add click-to-appear animation
    add_click_animation(slide, txBox)

    # Update speaker notes
    append_to_notes(slide,
        "\n\nAFTER VIDEO: System + human reader outperforms expert lip readers. "
        "Expert lip readers achieve ~45-52% word accuracy on unconstrained speech (Auer & Bernstein 2007). "
        "Model + context-aware human: estimated 55-70% word accuracy, 75-85% meaning capture. "
        "Key insight: the model provides candidate text (the hardest part of lip reading). "
        "The human's job becomes verification, not generation \u2014 dramatically easier. "
        "Hallucination risk drops from 20.5% to <5% with human filtering. "
        "Context is 'the great equalizer' \u2014 worth 10-20pp of disambiguation. "
        "Source: docs/evaluation/human_expert_comparison.md")


# ─── Item 21: Arabic Deep-Dive (expand to 2-3 slides) ────────────────────────

def expand_arabic_analysis(prs):
    """Item 21: Add 2 new Arabic deep-dive slides after the Arabic roadmap slide."""
    print("Item 21: Expanding Arabic analysis to multiple slides...")

    # Find the Arabic roadmap slide
    arabic_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Arabic" in shape.text_frame.text and ("Roadmap" in shape.text_frame.text or "Pipeline" in shape.text_frame.text or "Replication" in shape.text_frame.text):
                    arabic_idx = i
                    break
        if arabic_idx is not None:
            break

    if arabic_idx is None:
        print("  Warning: Could not find Arabic roadmap slide")
        return

    # We need to duplicate slides to create new ones
    # python-pptx doesn't have a native slide insert, so we duplicate the arabic slide
    # and modify the copies

    # Slide 63b: "AV-HuBERT: Why It's Not Language-Locked"
    slide_b = duplicate_slide(prs, arabic_idx)
    clear_slide_content(slide_b)
    # Add title
    add_textbox(slide_b, "AV-HuBERT: Why It's Not Language-Locked",
                MX, Inches(0.40), CW, Inches(0.75),
                size=Pt(32), color=WHITE, bold=True)
    # Accent line
    slide_b.shapes.add_shape(
        1,  # RECTANGLE
        MX, Inches(1.20), CW, Inches(0.02)
    ).fill.solid()
    slide_b.shapes[-1].fill.fore_color.rgb = TEAL
    slide_b.shapes[-1].line.fill.background()

    # Content — individual animated bullets
    add_animated_bullets(slide_b, [
        "AV-HuBERT is a self-supervised visual feature extractor",
        "Pretrained on LRS3 (English TED talks) \u2014 but not language-encoded",
        "Training loop: MFCC \u2192 K-means \u2192 pseudo-labels \u2192 masked prediction \u2192 iterate",
        'The "English-ness" is in which visual distinctions the model learned to care about',
        "Low-level features are mostly universal: lip shape, mouth opening, jaw movement",
        "Visual features are ~80% language-agnostic (mouth geometry is universal)",
        "Language specificity lives in downstream components, not the visual encoder",
    ], MX, Inches(1.50), CW, line_height=Inches(0.60),
    size=Pt(15), text_color=WHITE, bullet_color=TEAL)

    # Page number placeholder
    add_textbox(slide_b, "", MX, Inches(7.12), Inches(0.5), Inches(0.25),
                size=Pt(9), color=DGRAY)

    set_notes(slide_b,
        "AV-HuBERT self-supervised training: starts with MFCC features as initial targets, "
        "runs K-means clustering to create pseudo-labels, trains masked prediction of those labels "
        "from visual input, then iterates with better pseudo-labels. After multiple iterations, "
        "the visual encoder has learned which lip movements correspond to which sound clusters. "
        "The 'English-ness' lives in which visual distinctions the model learned to care about \u2014 "
        "e.g., it never learned to distinguish Arabic emphatics (\u0635 vs \u0633) which involve "
        "visible pharyngeal constriction. But this is an optimization target, not a hard blocker.")

    # Slide 63c: "Arabic Adaptation: What Changes"
    slide_c = duplicate_slide(prs, arabic_idx + 1)  # After slide_b
    clear_slide_content(slide_c)
    add_textbox(slide_c, "Arabic Adaptation: What Changes",
                MX, Inches(0.40), CW, Inches(0.75),
                size=Pt(32), color=WHITE, bold=True)
    # Accent line
    slide_c.shapes.add_shape(1, MX, Inches(1.20), CW, Inches(0.02)).fill.solid()
    slide_c.shapes[-1].fill.fore_color.rgb = TEAL
    slide_c.shapes[-1].line.fill.background()

    add_animated_bullets(slide_c, [
        "K-means clustering \u2014 retrain on Arabic audio features (already retrains per-dataset)",
        "LLM backbone \u2014 replace with Arabic-capable LLM (Jais, AceGPT, or multilingual Llama 3)",
        "Q-Former bridge + LoRA adapters \u2014 retrain on Arabic video-transcript pairs",
        "AV-HuBERT encoder \u2014 can reuse frozen; fine-tune later as optimization step",
        "Phase 1: Frozen AV-HuBERT + Arabic K-means + Arabic LLM + retrained Q-Former",
        "Phase 2: Fine-tune AV-HuBERT on Arabic video for language-specific distinctions",
        "Phase 3: Scale with more Arabic training data",
        "Biggest bottleneck: training data (no Arabic LRS3 equivalent at scale)",
    ], MX, Inches(1.50), CW, line_height=Inches(0.55),
    size=Pt(14), text_color=WHITE, bullet_color=TEAL)

    # Page number placeholder
    add_textbox(slide_c, "", MX, Inches(7.12), Inches(0.5), Inches(0.25),
                size=Pt(9), color=DGRAY)

    set_notes(slide_c,
        "Arabic adaptation practical bottleneck sequence:\n"
        "1. K-means: Retrain on Arabic audio features \u2014 this already retrains per-dataset in our pipeline, so Arabic clusters are essentially 'free'.\n"
        "2. LLM: Replace with Arabic-capable model. The LLM swap is the most impactful single change.\n"
        "3. Q-Former + LoRA: Retrain on Arabic video-transcript pairs. This is where the authors' actual novel contribution was.\n"
        "4. AV-HuBERT: Frozen English encoder is the starting point. Fine-tuning is Phase 2 optimization.\n\n"
        "Arabic emphatics example: \u0635 (emphatic S) vs \u0633 (plain S) have visible pharyngeal constriction "
        "that the English-pretrained encoder never learned to distinguish. Fine-tuning would teach this.\n\n"
        "Data challenge: No Arabic equivalent of LRS3. Options include custom collection from Arabic broadcast/YouTube, "
        "or cross-lingual pretraining strategies.")


def delete_slide(prs, title_fragment):
    """Delete a slide found by title fragment."""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and title_fragment in shape.text_frame.text:
                sld_id_lst = prs.slides._sldIdLst
                sld_ids = list(sld_id_lst)
                sld_id_lst.remove(sld_ids[i])
                # Also remove the slide part from the presentation
                rId = sld_ids[i].get(
                    '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
                )
                prs.part.drop_rel(rId)
                print(f"  Deleted slide {i+1} ('{title_fragment}')")
                return True
    print(f"  WARNING: Could not find slide with '{title_fragment}'")
    return False


def move_slide(prs, from_title_fragment, after_title_fragment):
    """Move a slide (found by title fragment) to after another slide (found by title fragment).

    Works by manipulating the sldIdLst XML ordering.
    """
    from_idx = None
    after_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and from_title_fragment in shape.text_frame.text:
                from_idx = i
                break
        for shape in slide.shapes:
            if shape.has_text_frame and after_title_fragment in shape.text_frame.text:
                after_idx = i
                break

    if from_idx is None:
        print(f"  WARNING: Could not find slide with '{from_title_fragment}'")
        return
    if after_idx is None:
        print(f"  WARNING: Could not find slide with '{after_title_fragment}'")
        return
    if from_idx == after_idx:
        print(f"  WARNING: Source and target are the same slide ({from_idx+1})")
        return

    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    moving = sld_ids[from_idx]
    sld_id_lst.remove(moving)

    # Recalculate after_idx since removing shifted indices
    if from_idx < after_idx:
        insert_pos = after_idx  # shifted left by 1 after removal, but we insert AFTER so +1 = after_idx
    else:
        insert_pos = after_idx + 1

    sld_id_lst.insert(insert_pos, moving)
    print(f"  Moved slide {from_idx+1} ('{from_title_fragment}') to after slide with '{after_title_fragment}' (position {insert_pos+1})")


def duplicate_slide(prs, source_idx):
    """Duplicate a slide and insert it after the source. Returns the new slide."""
    template = prs.slides[source_idx]
    slide_layout = template.slide_layout

    # Add a new blank slide
    new_slide = prs.slides.add_slide(slide_layout)

    # Move it to right after the source
    # python-pptx slide ordering uses the sldIdLst
    slides = prs.slides._sldIdLst
    # The new slide is at the end; move it to source_idx + 1
    new_sld_id = slides[-1]
    slides.remove(new_sld_id)
    # Insert after source_idx (0-indexed, but sldIdLst elements are 0-indexed too)
    slides.insert(source_idx + 1, new_sld_id)

    return new_slide


def clear_slide_content(slide):
    """Remove all shapes from a slide and set navy background."""
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        if sp.tag.endswith('}sp') or sp.tag.endswith('}pic') or sp.tag.endswith('}grpSp'):
            sp_tree.remove(sp)
    # Ensure navy background (duplicated slides may inherit white from layout)
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = NAVY


# ─── Item 22: Key Takeaways ──────────────────────────────────────────────────

def update_key_takeaways(prs):
    """Item 22: Update key takeaways slide."""
    print("Item 22: Updating key takeaways...")

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "Key Takeaway" in shape.text_frame.text:
                # Found the takeaways slide
                # Look for the bullet content shape (usually the large text box)
                for s in slide.shapes:
                    if s.has_text_frame and len(s.text_frame.paragraphs) >= 3:
                        tf = s.text_frame
                        full_text = tf.text
                        if "takeaway" in full_text.lower() or "pipeline" in full_text.lower() or "roadmap" in full_text.lower():
                            continue  # This might be the title
                        # This is likely the content shape — update its text
                        # We'll update specific numbers/phrases
                        for para in tf.paragraphs:
                            for run in para.runs:
                                # Update percentages
                                run.text = run.text.replace("40%", "61.6%")
                                run.text = run.text.replace("39.9%", "61.6%")
                                run.text = run.text.replace("IS ≥ 3.0", "IS \u2265 2.00")
                                run.text = run.text.replace("IS >= 3.0", "IS \u2265 2.00")
                                run.text = run.text.replace("3–5 months", "2\u20133 months")
                                run.text = run.text.replace("3-5 months", "2\u20133 months")
                return


# ─── Items 3, 19: NIV Threshold Global Update ────────────────────────────────

def update_niv_thresholds_global(prs):
    """Items 3, 19: Update IS thresholds across all slides."""
    print("Items 3/19: Updating NIV thresholds globally...")

    # Global text replacements across ALL slides
    global_replacements = [
        # Threshold text
        ("IS ≥ 3.0 = \"Properly Captured\"", "IS \u2265 2.00 = \"Useful Output\" (Y+P)"),
        ('IS ≥ 3.0 = "Properly Captured"', 'IS \u2265 2.00 = "Useful Output" (Y+P)'),
        ("IS ≥ 3.0 = Properly Captured", "IS \u2265 2.00 = Useful Output (Y+P)"),
        ("IS >= 3.0", "IS \u2265 2.00"),
        ("IS ≥ 3.0", "IS \u2265 2.00"),
        # Banner on slide 19
        ('IS ≥ 3.0 = "Properly Captured".  These 3 signals',
         'IS \u2265 2.00 = "Useful Output" (Y+P).  These 3 signals'),

        # Percentages — be careful with context
        # 39.9% → 61.6% (only in IS context)
        ("39.9% Properly Captured", "61.6% Useful Output"),
        ("39.9% — 3.5× more than WER's 11.4%", "61.6% — 2.4\u00d7 what WER suggests (25.5%)"),
        ("39.9%", "61.6%"),

        # Segment counts
        ("597/1,497", "922/1,497"),
        ("597 of 1,497", "922 of 1,497"),

        # 3.5x ratio
        ("3.5× more than WER", "2.4\u00d7 what WER suggests"),
        ("3.5x more than WER", "2.4\u00d7 what WER suggests"),

        # 88% agreement
        ("88% agreement", "\u03ba=0.818 agreement"),
        ("88.6% agreement", "\u03ba=0.818 agreement"),

        # 147 segments in gap
        ("147 segments", "541 segments"),  # Will be recalculated

        # 900 segments failed
        ("900 segments failed (IS < 3.0)", "575 segments below useful threshold (IS < 2.00)"),
        ("900 segments failed", "575 segments below threshold"),
        ("900 failed", "575 below threshold"),

        # 11.4% usable
        ("11.4% usable", "25.5% useful by WER"),
        ("11.4% Usable", "25.5% Useful by WER"),
        ("11.4%", "25.5%"),

        # Starting from 40%
        ("Starting from 40%", "Starting from 61.6%"),
    ]

    for slide in prs.slides:
        for old, new in global_replacements:
            replace_text_in_slide(slide, old, new)
            replace_in_notes(slide, old, new)

    # Also update slide 26 title specifically
    slide26 = prs.slides[25]
    replace_text_in_slide(slide26, "Intelligibility Score: 39.9% Properly Captured",
                         "Intelligibility Score: 61.6% Useful Output")
    replace_text_in_slide(slide26, "Intelligibility Score: 61.6% Properly Captured",
                         "Intelligibility Score: 61.6% Useful Output")

    # Slide 26 body — update the IS description
    replace_text_in_slide(slide26,
        "IS ≥ 3.0 = Properly Captured: 39.9%",
        "IS \u2265 2.00 = Useful (Y+P): 61.6% \u2014 2.4\u00d7 what WER suggests (25.5%)")
    replace_text_in_slide(slide26,
        "IS \u2265 2.00 = Properly Captured: 61.6%",
        "IS \u2265 2.00 = Useful (Y+P): 61.6%")


# ─── Item 14: IS as Surrogate Metric (Slide 36) ──────────────────────────────

def update_surrogate_metric(prs):
    """Item 14: Reframe Conservative Lower Bound as Calibrated Surrogate."""
    print("Item 14: Updating IS surrogate metric slide...")

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "Conservative" in shape.text_frame.text and "Lower Bound" in shape.text_frame.text:
                replace_text_in_shape(shape, "Conservative Lower Bound",
                                     "Calibrated Surrogate for LLM Judgment")
                replace_text_in_shape(shape, "IS: A Conservative Lower Bound",
                                     "IS: A Calibrated Surrogate")

        # Also fix in all shapes
        replace_text_in_slide(slide, "conservative lower bound", "calibrated surrogate")
        replace_text_in_slide(slide, "Conservative Lower Bound", "Calibrated Surrogate")


# ─── Item 15: WER-IS Gap with Dual Thresholds (Slide 25) ─────────────────────

def update_wer_is_gap(prs):
    """Item 15: Update WER-IS gap slide with dual thresholds."""
    print("Item 15: Updating WER-IS gap slide 25...")

    slide = prs.slides[24]  # 0-indexed

    # Replace threshold descriptions
    replace_text_in_slide(slide, "WER > 40% but IS ≥ 3.0",
                         "WER > 34% but IS \u2265 2.00")
    replace_text_in_slide(slide, "WER > 40% but IS \u2265 2.00",
                         "WER > 34% but IS \u2265 2.00")
    replace_text_in_slide(slide, "useful output that WER discards",
                         "useful meaning that WER discards (\u03ba: IS 0.818 vs WER 0.777)")

    # Update speaker notes with full dual-threshold info
    set_notes(slide,
        "WER-IS gap with NIV-calibrated thresholds.\n\n"
        "NIV Thresholds:\n"
        "Y (clearly conveyed): IS \u2265 3.80, WER \u2264 34% | IS \u03ba=0.690, WER \u03ba=0.629\n"
        "Y+P (any useful): IS \u2265 2.00, WER \u2264 77% | IS \u03ba=0.818, WER \u03ba=0.777\n\n"
        "The gap at both levels:\n"
        "Strict (Y): WER > 34% AND IS \u2265 3.80 \u2014 segments fully captured but strict WER rejects\n"
        "Forgiving (Y+P): WER > 77% AND IS \u2265 2.00 \u2014 useful meaning but even forgiving WER rejects\n\n"
        "Key point: IS beats WER at BOTH operating points. WER consistently correlates worse "
        "with the Opus judge than IS does.")


# ─── Item 6: PCA Narrative Correction (Slide 22) ─────────────────────────────

def fix_pca_narrative(prs):
    """Item 6: Fix PCA narrative — visual content AND speaker notes on slide 22."""
    print("Item 6: Fixing PCA narrative (visual + notes) on slide 22...")

    slide = prs.slides[21]  # 0-indexed

    # --- Rewrite visual content using existing shapes ---
    # Shape indices from slide inspection:
    #   0: title, 2: subtitle,
    #   4: box1 title, 5: box1 %, 6: box1 detail, 7: box1 desc,
    #   9: box2 title, 10: box2 %, 11: box2 detail, 12: box2 desc,
    #   14: box3 title, 15: box3 %, 16: box3 detail, 17: box3 desc,
    #   18: validation line
    shape_updates = {
        0: "6 Signals, 3 Principal Components",
        2: "PCA on the six IS signals reveals where the variance actually lives.",
        4: "PC1: Signal Quality",
        5: "68.4%",
        6: "Semantic + Phonetic + WER + WWER + Named Entity Accuracy",
        7: "All five content signals load equally (0.43\u20130.47). Semantic is NOT independent \u2014 it measures the same underlying quality as word accuracy.",
        9: "PC2: Output Length",
        10: "19.5%",
        11: "Length Ratio dominates (loading = 0.91)",
        12: "Independent of content quality. Catches hallucination (too long) and truncation (too short).",
        14: "PC3: Entity Swing",
        15: "5.1%",
        16: "NEA F1 loads here (below Kaiser threshold, eigenvalue 0.31)",
        17: "Minor refinement axis \u2014 names and numbers that diverge from the main quality signal.",
        18: "Together: 93% of variance in 3 components. Kaiser retains 2 (87.9%); PC3 adds nuance.",
    }

    for idx, new_text in shape_updates.items():
        shape = slide.shapes[idx]
        if shape.has_text_frame:
            # Preserve formatting of first run, replace text
            first_para = shape.text_frame.paragraphs[0]
            if first_para.runs:
                first_para.runs[0].text = new_text
                # Clear any additional runs
                for run in first_para.runs[1:]:
                    run.text = ""
            else:
                first_para.text = new_text
            # Clear any additional paragraphs
            for para in shape.text_frame.paragraphs[1:]:
                for run in para.runs:
                    run.text = ""

    # Change box 3 title color to gold (was green for "Output Sanity")
    for para in slide.shapes[14].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = GOLD

    # --- Speaker notes ---
    set_notes(slide,
        "PCA STORY (from docs/evaluation/is_pca_analysis.md):\n\n"
        "Kaiser criterion PCA on 6 standardized IS signals retains 2 principal components:\n"
        "- PC1 (68.4%): Signal Quality \u2014 all 5 content signals load equally (0.43-0.47). "
        "Semantic is NOT independent \u2014 it loads on PC1 alongside word-accuracy signals.\n"
        "- PC2 (19.5%): Output Length \u2014 Length Ratio dominates (0.91), independent of content quality.\n"
        "- PC3 (5.1%): Entity Swing \u2014 below Kaiser threshold (eigenvalue 0.31 < 1.0), minor.\n\n"
        "Together: 87.9% of variance. The old '3 dimensions' claim was from covariance clustering, "
        "not PCA. This slide now shows the correct statistical story.\n\n"
        "PCA validates the IS formula: all 5 content signals measure the same underlying quality (PC1), "
        "so the weighted sum is a reliable composite. The specific allocation (25/15/15/15/15/15) is "
        "robust \u2014 equal weighting gives r=0.999 with current weights.")


# ─── Item 7: Bad IS Sample Fix (Slide 23) ────────────────────────────────────

def fix_bad_is_sample(prs):
    """Item 7: Fix Length Ratio value on slide 23 bad sample."""
    print("Item 7: Fixing bad IS sample Length Ratio...")

    slide = prs.slides[22]  # 0-indexed

    # The bad sample shows Length Ratio = 1.00 but should be 0.25
    # Find the shape with "1.00" near "Length Ratio" context
    # There are TWO "1.00" values — one for good sample (correct) and one for bad (wrong)
    # The bad one is on the right side of the slide
    found = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            # Look for the 1.00 that's the bad sample's Length Ratio
            if text.strip() == "1.00":
                # Check position — bad sample is on the right half
                if shape.left > Inches(6):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip() == "1.00":
                                run.text = "0.25"
                                found = True
                                break
                if found:
                    break
    if not found:
        # Try alternate approach: look for "1.00" with bold and teal color on right side
        for shape in slide.shapes:
            if shape.has_text_frame and shape.left > Inches(6):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "1.00" in run.text and run.font.color.rgb == TEAL:
                            run.text = run.text.replace("1.00", "0.25")
                            found = True
        if not found:
            print("  Warning: Could not find Length Ratio 1.00 to fix on slide 23")


# ─── Item 8: Radar Chart Fix (Slide 24) ──────────────────────────────────────

def fix_radar_chart(prs):
    """Item 8: Remove 'projected/estimated' disclaimers from slide 24."""
    print("Item 8: Fixing radar chart slide 24...")

    slide = prs.slides[23]  # 0-indexed

    # Remove text about "projected" and "estimated"
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "projected" in text.lower() and "profiles" in text.lower():
                # This is the subtitle — simplify it
                replace_text_in_shape(shape,
                    "projected profiles based on literature and architecture analysis",
                    "measured profiles from LRS3 benchmark and YouTube evaluation")
            if "projected, teal" in text.lower():
                replace_text_in_shape(shape, "(projected, teal)", "(target profile, teal)")
                replace_text_in_shape(shape, "(current, coral)", "(YouTube baseline, coral)")

    # Update speaker notes
    set_notes(slide,
        "Dual radar chart: LRS3 benchmark (measured, n=170) vs YouTube real-world (measured, n=1,497).\n"
        "LRS3 values (all measured): Semantic 0.779, Phonetic 0.794, 1-WER 0.689, 1-WWER 0.662, NEA 0.683, LenRatio 0.971.\n"
        "YouTube values (all measured): Semantic 0.58, Phonetic 0.52, 1-WER 0.36, 1-WWER 0.38, NEA 0.39, LenRatio 0.72.\n"
        "Key insight: the radar shape reveals where each condition is strong and weak. "
        "Length Ratio stays high even for YouTube (model generates correct amount of text). "
        "The collapsed axes (WER, WWER, NEA) are where the domain gap hits hardest.\n"
        "Data source: commit 60ab14e, docs/_research-tools/generators/generate_dual_radar.py")

    # TODO: Regenerate and replace the radar image itself (done separately)


# ─── Item 10: Two Systems Clarity (Slide 27) ─────────────────────────────────

def clarify_two_systems(prs):
    """Item 10: Make Two Systems slide show full NIV threshold table."""
    print("Item 10: Clarifying Two Systems slide 27...")

    slide = prs.slides[26]  # 0-indexed

    # Update the IS threshold reference
    replace_text_in_slide(slide, "IS ≥ 3.0 = Captured: 39.9% (597/1,497)",
                         "IS \u2265 2.00 = Useful: 61.6% (922/1,497) | \u03ba=0.818")
    replace_text_in_slide(slide, "IS \u2265 2.00 = Captured: 61.6% (922/1,497)",
                         "IS \u2265 2.00 = Useful: 61.6% (922/1,497) | \u03ba=0.818")

    # Update the confusion matrix row header
    replace_text_in_slide(slide, "IS ≥ 3.0", "IS \u2265 2.00")

    # Update speaker notes with full comparison
    set_notes(slide,
        "Two evaluation systems comparison with NIV thresholds.\n\n"
        "NIV Threshold Table:\n"
        "Y (clearly conveyed): IS \u2265 3.80 (\u03ba=0.690) | WER \u2264 34% (\u03ba=0.629)\n"
        "Y+P (any useful):     IS \u2265 2.00 (\u03ba=0.818) | WER \u2264 77% (\u03ba=0.777)\n\n"
        "IS beats WER at BOTH operating points:\n"
        "- Y level: IS \u03ba=0.690 vs WER \u03ba=0.629 (+0.061)\n"
        "- Y+P level: IS \u03ba=0.818 vs WER \u03ba=0.777 (+0.041)\n\n"
        "Confusion matrix uses Y+P (Judge) vs IS \u2265 2.00 — the operationally meaningful comparison.\n"
        "85% Pearson correlation between IS and Opus verdicts.")


# ─── Item 12: Weight Rationale Note (Slide 22) ───────────────────────────────
# Already handled in fix_pca_narrative — combined items 6 and 12


# ─── Item 18: Pipeline ASR Separation (Slide 43) ─────────────────────────────

def separate_asr_pipeline(prs):
    """Item 18: Visually separate ASR from pipeline on slide 43."""
    print("Item 18: Separating ASR from pipeline...")

    # Find the 8-stage pipeline slide
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "8-Stage" in shape.text_frame.text:
                # Found it — add annotation near the ASR box
                # Look for the ASR shape
                for s in slide.shapes:
                    if s.has_text_frame and "ASR" in s.text_frame.text and "Whisper" in s.text_frame.text:
                        # Add annotation below/beside the ASR box
                        asr_left = s.left
                        asr_top = s.top + s.height + Inches(0.05)
                        add_textbox(slide,
                            "(Evaluation only \u2014 not in inference path)",
                            asr_left, asr_top, s.width, Inches(0.3),
                            size=Pt(9), color=CORAL, italic=True,
                            align=PP_ALIGN.CENTER)
                        break

                # Update speaker notes
                append_to_notes(slide,
                    "\n\nIMPORTANT: ASR (Whisper) is used for EVALUATION only \u2014 "
                    "it generates reference transcriptions to compare against the model's output. "
                    "It is NOT part of the actual inference pipeline. The model reads lips directly "
                    "from video without any audio processing.")
                return

    print("  Warning: Could not find 8-Stage Pipeline slide")


# ─── Item 9: Metrics Disagreement Slide ───────────────────────────────────────

def add_disagreement_slide(prs):
    """Item 9: Re-add the Part 1 metrics disagreement slide."""
    print("Item 9: Adding metrics disagreement Part 1 slide...")

    # Find slide 34 or 35 (the existing disagreement slide)
    target_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "When Metrics Disagree" in shape.text_frame.text:
                target_idx = i
                break
        if target_idx is not None:
            break

    if target_idx is None:
        print("  Warning: Could not find existing disagreement slide")
        return

    # Insert a new slide before the existing one
    new_slide = duplicate_slide(prs, target_idx)
    # Move it before the existing one (it was inserted after)
    slides = prs.slides._sldIdLst
    # Swap the last two positions
    new_id = slides[target_idx + 1]
    old_id = slides[target_idx]
    slides.remove(new_id)
    slides.insert(target_idx, new_id)

    clear_slide_content(new_slide)

    # Build the "When Metrics Disagree: What It Tells Us" slide
    add_textbox(new_slide, "When Metrics Disagree: What It Tells Us",
                MX, Inches(0.40), CW, Inches(0.75),
                size=Pt(32), color=WHITE, bold=True)

    # Accent line
    new_slide.shapes.add_shape(1, MX, Inches(1.20), CW, Inches(0.02)).fill.solid()
    new_slide.shapes[-1].fill.fore_color.rgb = TEAL
    new_slide.shapes[-1].line.fill.background()

    # Subtitle
    add_textbox(new_slide,
        "IS uses 6 signals because no single metric tells the full story. "
        "Disagreements between metrics reveal specific quality patterns:",
        MX, CT, CW, Inches(0.4),
        size=Pt(13), color=LGRAY, italic=True)

    # Four pattern cards (2x2 grid)
    patterns = [
        ("WWER \u226a WER", TEAL, "Function words wrong, content right",
         '"the team discussed quarterly" \u2192 "team discuss quarterly"\n'
         'WER 43% but WWER 15% \u2014 viewer gets the message.'),
        ("NEA high, WER high", GREEN, "Names preserved despite poor accuracy",
         '"Dr. Chen presented Q3 results" \u2192 "Dr. Chen present Q3 result"\n'
         'WER 57% but NEA F1 = 100% \u2014 critical info intact.'),
        ("Semantic high, WER high", GOLD, "Meaning preserved through paraphrasing",
         '"reduce spending" \u2192 "cut the budget"\n'
         'WER 100% but Semantic 0.87 \u2014 same meaning, different words.'),
        ("Phonetic high, Semantic low", CORAL, "Sounds right, wrong meaning",
         '"the alliance was formed" \u2192 "the lions were found"\n'
         'Phonetic 0.71 but Semantic 0.12 \u2014 dangerous deceptive error.'),
    ]

    cw = Inches(5.8)
    ch = Inches(1.8)
    gap_x = Inches(0.53)
    gap_y = Inches(0.2)
    cy = CT + Inches(0.55)

    for idx, (title, color, subtitle, body) in enumerate(patterns):
        col = idx % 2
        row = idx // 2
        x = MX + col * (cw + gap_x)
        y = cy + row * (ch + gap_y)

        # Card background
        rect = add_rounded_rect(new_slide, x, y, cw, ch, NAVY2, color)

        # Card title
        add_textbox(new_slide, f"{title}  \u2014  {subtitle}",
                    x + Inches(0.2), y + Inches(0.1), cw - Inches(0.4), Inches(0.35),
                    size=Pt(13), color=color, bold=True)

        # Card body
        add_textbox(new_slide, body,
                    x + Inches(0.2), y + Inches(0.5), cw - Inches(0.4), ch - Inches(0.6),
                    size=Pt(11), color=LGRAY)

    # Bottom text
    add_textbox(new_slide,
        "This is why IS uses 6 signals \u2014 each disagreement pattern reveals "
        "a different type of quality that a single metric would miss.",
        MX, Inches(6.3), CW, Inches(0.4),
        size=Pt(12), color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

    # Hide this slide
    hide_slide(new_slide)

    set_notes(new_slide,
        "Four key metric disagreement patterns.\n"
        "WWER<<WER: function words wrong but content preserved.\n"
        "High NEA + high WER: names survived despite poor word accuracy.\n"
        "High semantic + high WER: paraphrasing (different words, same meaning).\n"
        "High phonetic + low semantic: sounds right but wrong meaning (dangerous).\n"
        "This slide is HIDDEN \u2014 available for Q&A or deep-dive.")


# ─── Item 23: 30-Sample Slide (Slide 17) ─────────────────────────────────────

def add_30_sample_slide(prs):
    """Item 23: Replace slide 17 placeholder with 30-sample LLM judge content."""
    print("Item 23: Adding 30-sample slide on slide 17...")

    slide = prs.slides[16]  # 0-indexed

    # Remove existing content and rebuild
    clear_slide_content(slide)

    # Title
    add_textbox(slide, "LLM Judge: 30-Sample Deep Dive",
                MX, Inches(0.35), CW, Inches(0.65),
                size=Pt(32), color=WHITE, bold=True)

    # Accent line
    ln = slide.shapes.add_shape(1, MX, Inches(1.05), CW, Inches(0.02))
    ln.fill.solid()
    ln.fill.fore_color.rgb = TEAL
    ln.line.fill.background()

    # ── Left column: Summary stats ──
    lx = MX
    left_w = Inches(3.6)
    stats_y = Inches(1.25)

    add_textbox(slide, "30 Stratified Samples",
                lx, stats_y, left_w, Inches(0.30),
                size=Pt(16), color=TEAL, bold=True)

    # Summary stats as a proper table
    stats_rows = [
        ("Mean WER", "61.4%"),
        ("Mean IS", "2.67 / 5.0"),
        ("Judge Y", "7  (23.3%)"),
        ("Judge P", "12  (40.0%)"),
        ("Judge N", "11  (36.7%)"),
        ("Y + P", "19  (63.3%)"),
    ]
    stats_tbl = slide.shapes.add_table(len(stats_rows), 2,
                                        lx, stats_y + Inches(0.40),
                                        left_w, Inches(2.10)).table
    stats_tbl.columns[0].width = Inches(1.5)
    stats_tbl.columns[1].width = Inches(2.1)

    for r, (label, val) in enumerate(stats_rows):
        for c, txt in enumerate([label, val]):
            cell = stats_tbl.cell(r, c)
            cell.text = txt
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            p.font.color.rgb = WHITE if c == 1 else LGRAY
            p.font.bold = (c == 1)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY2

    add_textbox(slide, "Distribution mirrors full 1,497-segment dataset",
                lx, stats_y + Inches(2.60), left_w, Inches(0.30),
                size=Pt(11), color=MGRAY, italic=True)

    # ── Right column: Example table with ref/hyp ──
    rx = lx + left_w + Inches(0.30)
    tbl_w = CW - left_w - Inches(0.30)

    add_textbox(slide, "Representative Examples \u2014 Reference vs. Hypothesis",
                rx, stats_y, tbl_w, Inches(0.30),
                size=Pt(16), color=CORAL, bold=True)

    # Example segments: 2 Y, 2 P, 2 N showing interesting disagreements
    examples = [
        ("Y", GREEN,
         "\"...unite your empire well I can tell you what Constantine wanted...\"",
         "\"...unite your empire well I can tell you we lost he wanted...\"",
         "14.3%", "4.60"),
        ("Y", GREEN,
         "\"...talking about gruselcoin and I'm gonna break the video down...\"",
         "\"...talking about crucial coins and I'm going to break the video down...\"",
         "50.0%", "2.45"),
        ("P", GOLD,
         "\"popular with republican presidential candidates\"",
         "\"unpopular with republican presidential candidates\"",
         "20.0%", "4.54"),
        ("P", GOLD,
         "\"...syllables you may remember from our episode about syllable structure...\"",
         "\"...sentences you may remember from our episode about sentence structure...\"",
         "10.5%", "4.57"),
        ("N", CORAL,
         "\"all you have to do is unscrew\"",
         "\"all you have to do is not to\"",
         "28.6%", "3.42"),
        ("N", CORAL,
         "\"blood extraction first before going to x ray...\"",
         "\"cut your hair first before going to ashram...\"",
         "35.7%", "3.14"),
    ]

    # Build example table: Verdict | Reference / Hypothesis | WER | IS
    n_rows = len(examples) * 2  # 2 rows per example (ref + hyp)
    tbl_shape = slide.shapes.add_table(n_rows + 1, 4,
                                        rx, stats_y + Inches(0.40),
                                        tbl_w, Inches(4.70))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(0.55)
    tbl.columns[1].width = Inches(6.15)
    tbl.columns[2].width = Inches(0.70)
    tbl.columns[3].width = Inches(0.68)

    # Header row
    headers = ["", "Reference (top) / Hypothesis (bottom)", "WER", "IS"]
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = hdr
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.name = "Calibri"
        p.font.color.rgb = TEAL
        p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY2

    for i, (verdict, v_color, ref, hyp, wer, is_val) in enumerate(examples):
        ref_row = 1 + i * 2
        hyp_row = ref_row + 1

        # Merge verdict cell across 2 rows
        v_cell = tbl.cell(ref_row, 0)
        v_cell_merged = tbl.cell(hyp_row, 0)
        v_cell.merge(v_cell_merged)
        v_cell.text = verdict
        p = v_cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.name = "Calibri"
        p.font.color.rgb = v_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        v_cell.fill.solid()
        v_cell.fill.fore_color.rgb = NAVY2

        # WER cell (merged)
        w_cell = tbl.cell(ref_row, 2)
        w_cell_merged = tbl.cell(hyp_row, 2)
        w_cell.merge(w_cell_merged)
        w_cell.text = wer
        p = w_cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.name = "Calibri"
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        w_cell.fill.solid()
        w_cell.fill.fore_color.rgb = NAVY2

        # IS cell (merged)
        is_cell = tbl.cell(ref_row, 3)
        is_cell_merged = tbl.cell(hyp_row, 3)
        is_cell.merge(is_cell_merged)
        is_cell.text = is_val
        p = is_cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.name = "Calibri"
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        is_cell.fill.solid()
        is_cell.fill.fore_color.rgb = NAVY2

        # Ref text row
        ref_cell = tbl.cell(ref_row, 1)
        ref_cell.text = "REF: " + ref
        p = ref_cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.name = "Calibri"
        p.font.color.rgb = WHITE
        ref_cell.fill.solid()
        ref_cell.fill.fore_color.rgb = NAVY2

        # Hyp text row
        hyp_cell = tbl.cell(hyp_row, 1)
        hyp_cell.text = "HYP: " + hyp
        p = hyp_cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.name = "Calibri"
        p.font.color.rgb = LGRAY
        hyp_cell.fill.solid()
        hyp_cell.fill.fore_color.rgb = NAVY2

    # Bottom — insight callouts
    add_textbox(slide,
        "Y with WER 50% \u2014 LLM sees meaning IS misses  |  "
        "P flips meaning with 1 word (\"un\")  |  "
        "N at IS 3.4 \u2014 structure preserved but key verb lost",
        MX, Inches(6.35), CW, Inches(0.35),
        size=Pt(12), color=GOLD, italic=True, align=PP_ALIGN.CENTER)

    # Bottom subtitle guidance
    add_textbox(slide,
        "Each video has burned-in subtitles: reference (top) and hypothesis (bottom)",
        MX, Inches(6.70), CW, Inches(0.30),
        size=Pt(11), color=MGRAY, italic=True, align=PP_ALIGN.CENTER)

    set_notes(slide,
        "30-sample overview: stratified sample from the 1,497-segment dataset. "
        "Distribution matches full dataset closely: Y=23%, P=40%, N=37%. "
        "Mean WER 61.4% vs 64.1% full.\n\n"
        "Example highlights:\n"
        "- Y at WER 50%: 'gruselcoin' \u2192 'crucial coins' \u2014 LLM sees the meaning is preserved "
        "even though half the words differ. IS only gives 2.45.\n"
        "- P with 1-word flip: 'popular' \u2192 'unpopular' \u2014 single negation destroys meaning. "
        "LLM correctly flags as partial despite low WER.\n"
        "- N at IS 3.4: 'unscrew' \u2192 'not to' \u2014 structure is preserved but the key action verb "
        "is lost, making the instruction useless. LLM sees this, IS doesn't.\n\n"
        "These disagreements justify having BOTH metrics.")


# ─── Item 2: Page Numbering ──────────────────────────────────────────────────

def fix_page_numbers(prs):
    """Item 2: Fix sequential page numbering for visible slides."""
    print("Item 2: Fixing page numbers...")

    visible_num = 0
    for slide in prs.slides:
        if is_hidden(slide):
            continue
        visible_num += 1

        # Find the page number text box (bottom-left, 9pt, gray)
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                if len(tf.paragraphs) == 1:
                    for run in tf.paragraphs[0].runs:
                        # Check if it looks like a page number
                        if (run.font.size and run.font.size <= Pt(10) and
                            shape.top > Inches(6.5) and shape.left < Inches(2)):
                            try:
                                # Try to parse as number
                                int(run.text.strip())
                                run.text = str(visible_num)
                            except ValueError:
                                pass
                    # Also check paragraph text directly
                    p_text = tf.paragraphs[0].text.strip()
                    if (shape.top > Inches(6.5) and shape.left < Inches(2) and
                        len(p_text) <= 3):
                        try:
                            int(p_text)
                            if tf.paragraphs[0].runs:
                                tf.paragraphs[0].runs[0].text = str(visible_num)
                            else:
                                tf.paragraphs[0].text = str(visible_num)
                        except ValueError:
                            pass

    print(f"  Numbered {visible_num} visible slides")


# ─── Item 27 (Slide 27 confusion matrix) ─────────────────────────────────────
# Handled within clarify_two_systems above


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    print(f"Loading {INPUT_PATH}...")
    prs = Presentation(INPUT_PATH)
    print(f"  {len(prs.slides)} slides loaded")

    # Group 1: Text content changes
    rewrite_what_was_done(prs)          # Item 1
    expand_llm_judge_slide(prs)         # Item 4
    fill_is_motivation_slide(prs)       # Item 5
    add_lrs3_comment(prs)              # Item 11
    reframe_three_numbers(prs)          # Item 13
    clarify_llm_swap_training(prs)     # Item 16
    downgrade_arabic_risk(prs)          # Item 17
    add_human_expert_animation(prs)     # Item 20
    expand_arabic_analysis(prs)         # Item 21
    update_key_takeaways(prs)           # Item 22

    # Group 2: Number/threshold updates
    update_niv_thresholds_global(prs)   # Items 3, 19
    update_surrogate_metric(prs)        # Item 14
    update_wer_is_gap(prs)             # Item 15

    # Group 3: Narrative/visual fixes
    fix_pca_narrative(prs)             # Items 6, 12
    fix_bad_is_sample(prs)             # Item 7
    fix_radar_chart(prs)               # Item 8
    clarify_two_systems(prs)           # Item 10
    separate_asr_pipeline(prs)         # Item 18

    # Group 4: New slides
    add_disagreement_slide(prs)        # Item 9
    add_30_sample_slide(prs)           # Item 23

    # Group 5: Slide deletion
    print("Item 25: Removing 'Five Insights That Inform the Roadmap' slide...")
    delete_slide(prs, "Five Insights That Inform the Roadmap")

    # Group 6: Slide reordering and hiding
    print("Item 24: Moving 'Failure Modes: Impact' after 'LLM Upgrade: Why It Matters'...")
    move_slide(prs, "Failure Modes: Impact", "LLM Upgrade: Why It Matters")

    print("Item 26: Hiding 'LLM Upgrade: Why It Matters'...")
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "LLM Upgrade: Why It Matters" in shape.text_frame.text:
                hide_slide(slide)
                print("  Hidden.")
                break

    # Group 7: Page numbering — MUST BE LAST
    fix_page_numbers(prs)              # Item 2

    print(f"\nSaving to {OUTPUT_PATH}...")
    prs.save(OUTPUT_PATH)
    print(f"Done! {len(prs.slides)} slides saved.")


if __name__ == "__main__":
    main()
