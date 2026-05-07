#!/usr/bin/env python3
"""
audit_pptx_visual_structure.py

Deep visual + structural audit of Argos VSP PowerPoint decks.
Focus: layout, animations, videos, notes, section structure, style.

Output:
  /home/ubuntu/docs/evaluation/pptx_visual_audit.json
  /home/ubuntu/docs/evaluation/pptx_visual_audit.md
  /home/ubuntu/docs/evaluation/pptx_fix_manifest.md

This script DOES NOT modify any input PPTX.
"""

from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
from collections import Counter, defaultdict
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any, Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Constants & namespaces
# ---------------------------------------------------------------------------

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
}

DECKS = [
    "/home/ubuntu/presentation_materials_20260224/Argos_VSP_Final_84slides_Mar2026.pptx",
    "/home/ubuntu/presentation_materials_20260224/Argos_VSP_Client_v9_May2026.pptx",
    "/home/ubuntu/presentation_materials_20260224/Argos_VSP_Client_v10_May2026.pptx",
]

VIDEO_DIR = Path("/home/ubuntu/presentation_materials_20260224/06_demo_videos")
PLOT_DIR = Path("/home/ubuntu/presentation_materials_20260224/01_plots_for_slides")
BRANDING_DIR = Path("/home/ubuntu/presentation_materials_20260224/08_branding")

# Palette (per presentation/config.py, hex without leading 0x)
PALETTE = {
    "0D1B2A", "FFFFFF", "00B4D8", "E06C75", "AAAAAA", "666666",
    "333333", "4CAF50", "FFC107", "FFD54F", "FF9800", "F44336",
    "B71C1C", "152A40", "1A3550", "4F8FF7", "B066D9",
    # Common safe defaults from python-pptx
    "000000",
}

# Common fonts that are acceptable
ACCEPTABLE_FONTS = {"Calibri", "Arial", "Helvetica", "Helvetica Neue",
                    "Calibri Light", "Consolas", "Courier New", None, ""}

# Severity levels
BLOCKER = "BLOCKER"
MAJOR = "MAJOR"
MINOR = "MINOR"

# ---------------------------------------------------------------------------
# Issue model
# ---------------------------------------------------------------------------

@dataclass
class Issue:
    category: str
    severity: str
    description: str
    shape_id: Optional[str] = None
    location: Optional[str] = None  # e.g. "left=2.0in,top=3.5in"

    def to_dict(self):
        return asdict(self)


@dataclass
class SlideAudit:
    slide_num: int
    slide_title: str
    hidden: bool
    layout_name: Optional[str]
    shape_count: int
    has_video: bool
    has_logo: bool
    has_slide_number: bool
    notes_word_count: int
    animation_groups: int
    issues: list[Issue] = field(default_factory=list)

    def to_dict(self):
        d = asdict(self)
        d["issues"] = [i.to_dict() for i in self.issues]
        return d


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def emu_to_in(emu: Optional[int]) -> float:
    if emu is None:
        return 0.0
    return emu / 914400.0


def shape_bbox(shape) -> tuple[float, float, float, float]:
    """Return (left_in, top_in, right_in, bottom_in)."""
    try:
        l, t = emu_to_in(shape.left), emu_to_in(shape.top)
        w, h = emu_to_in(shape.width), emu_to_in(shape.height)
        return l, t, l + w, t + h
    except Exception:
        return 0, 0, 0, 0


def bbox_area(bb) -> float:
    return max(0, bb[2] - bb[0]) * max(0, bb[3] - bb[1])


def bbox_intersect(a, b) -> float:
    x1 = max(a[0], b[0]); y1 = max(a[1], b[1])
    x2 = min(a[2], b[2]); y2 = min(a[3], b[3])
    return max(0, x2 - x1) * max(0, y2 - y1)


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return (shape.text_frame.text or "").strip()


def is_invisible_anim_trigger(shape) -> bool:
    """Heuristic: 0×0 size or 'Anim' / 'Trigger' / 'Click' in name."""
    try:
        w, h = emu_to_in(shape.width), emu_to_in(shape.height)
    except Exception:
        return False
    if w < 0.05 or h < 0.05:
        return True
    name = (shape.name or "")
    if re.search(r"^(Anim|Trigger|Click|Hidden)", name, re.IGNORECASE):
        return True
    return False


def iter_shapes_recursive(shapes):
    """Yield (shape, depth) for all shapes including those inside groups."""
    for shape in shapes:
        yield shape, 0
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub in shape.shapes:
                    yield sub, 1
                    if sub.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for sub2 in sub.shapes:
                            yield sub2, 2
            except Exception:
                pass


def shape_visible_text(shape) -> bool:
    txt = shape_text(shape)
    return bool(txt) and not is_invisible_anim_trigger(shape)


def slide_title_text(slide) -> str:
    if slide.shapes.title is not None:
        try:
            t = (slide.shapes.title.text or "").strip()
            if t:
                return t
        except Exception:
            pass
    # fallback: first non-empty text frame
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = (shape.text_frame.text or "").strip()
            if t:
                return t.split("\n")[0][:120]
    return "(untitled)"


def slide_is_hidden(slide) -> bool:
    show = slide._element.get("show")
    return show == "0"


def slide_notes_text(slide) -> str:
    try:
        if slide.has_notes_slide:
            return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        pass
    return ""


def hex_color(rgb) -> Optional[str]:
    try:
        return str(rgb).upper()
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Animation analysis (OOXML)
# ---------------------------------------------------------------------------

def analyze_animations(slide_xml) -> dict[str, Any]:
    """Inspect <p:timing>; return summary dict.

    We use ``findall`` (with a namespace map) instead of ``xpath`` because
    python-pptx's ``CT_*`` element classes override ``xpath()`` and reject
    the ``namespaces=`` kwarg.
    """
    timing = slide_xml.find(".//p:timing", NS)
    result = {
        "has_timing": timing is not None,
        "click_groups": 0,
        "spids_referenced": [],
        "para_build_settings": [],   # list of (spid, build_type)
        "level1_par_count": 0,
        "click_step_count": 0,
        "errors": [],
    }
    if timing is None:
        return result

    try:
        # Click-step animation groups
        click_effects = timing.findall(
            ".//p:cTn[@nodeType='clickEffect']", NS)
        result["click_step_count"] = len(click_effects)

        # spTgt @spid references (anchored shape ids)
        sp_tgts = timing.findall(".//p:spTgt", NS)
        result["spids_referenced"] = [
            s.get("spid") for s in sp_tgts if s.get("spid")]

        # bldLst/bldP — per-paragraph build settings
        build_paras = slide_xml.findall(".//p:bldLst/p:bldP", NS)
        for bp in build_paras:
            spid = bp.get("spid")
            build_type = bp.get("build", "allAtOnce")
            result["para_build_settings"].append((spid, build_type))

        # Level-1 par count
        level1 = timing.findall("./p:tnLst/p:par", NS)
        result["level1_par_count"] = len(level1)

        result["click_groups"] = max(result["click_step_count"], 0)
    except Exception as e:
        result["errors"].append(f"timing parse error: {e}")

    return result


def get_slide_xml(slide) -> etree._Element:
    return slide._element


# ---------------------------------------------------------------------------
# Video analysis
# ---------------------------------------------------------------------------

def media_relationships(slide) -> list[dict[str, Any]]:
    """Return list of {rid, target_part_name, ext, size_bytes_or_None}."""
    out = []
    rels = slide.part.rels
    for rid, rel in rels.items():
        rt = rel.reltype.lower()
        if "video" in rt or "media" in rt or "movie" in rt:
            try:
                target = rel.target_part
                pname = target.partname
                # blob bytes via _blob may be heavy; just length
                blob = target.blob
                out.append({
                    "rid": rid,
                    "partname": str(pname),
                    "size": len(blob),
                    "reltype": rt,
                })
            except Exception:
                # external link
                out.append({
                    "rid": rid,
                    "partname": rel.target_ref if hasattr(rel, "target_ref") else "?",
                    "size": None,
                    "reltype": rt,
                })
    return out


def ffprobe_duration(path: Path) -> Optional[float]:
    if not path.exists():
        return None
    try:
        out = subprocess.check_output(
            ["ffprobe", "-v", "error", "-show_entries",
             "format=duration", "-of",
             "default=noprint_wrappers=1:nokey=1", str(path)],
            stderr=subprocess.DEVNULL, timeout=10).decode().strip()
        return float(out)
    except Exception:
        return None


def find_media_shapes(slide) -> list[Any]:
    out = []
    for shape, _depth in iter_shapes_recursive(slide.shapes):
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                out.append(shape)
        except Exception:
            pass
        # Some embedded videos appear as picture with movie ref — heuristic
    return out


def find_picture_shapes(slide) -> list[Any]:
    out = []
    for shape, _depth in iter_shapes_recursive(slide.shapes):
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                out.append(shape)
        except Exception:
            pass
    return out


def _has_corner_logo_picture(shapes_iter, sw_in: float, sh_in: float,
                             edge_tol: float = 1.0) -> bool:
    """Return True if any shape in `shapes_iter` is a small picture in a slide
    corner (any of the four corners). ``edge_tol`` is the max distance in
    inches from an edge that still counts as 'in the corner'.

    Loosened from the original 0.5in threshold so legitimate corner logos
    like the Argos brand mark at right_gap=0.66in are recognized.
    """
    try:
        for shape in shapes_iter:
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                l = emu_to_in(shape.left)
                t = emu_to_in(shape.top)
                w = emu_to_in(shape.width)
                h = emu_to_in(shape.height)
                r = l + w
                b = t + h
                if w <= 0.7 and h <= 0.7:
                    near_right = (sw_in - r) < edge_tol
                    near_left = l < edge_tol
                    near_top = t < edge_tol
                    near_bot = (sh_in - b) < edge_tol
                    if (near_right and (near_top or near_bot)) or \
                       (near_left and (near_top or near_bot)):
                        return True
            except Exception:
                continue
    except Exception:
        return False
    return False


def slide_has_corner_logo_anywhere(slide, sw_in: float, sh_in: float) -> bool:
    """Look for a corner logo on the slide itself, then on its layout, then
    on the slide master. Logos are commonly inherited from layout / master
    rather than added per-slide; the per-slide loop in audit_slide() can also
    miss legitimate logos when the corner threshold is too tight.
    """
    # Per-slide shapes (handles legitimate logos with slightly larger edge gap)
    if _has_corner_logo_picture(slide.shapes, sw_in, sh_in):
        return True
    try:
        layout = slide.slide_layout
    except Exception:
        layout = None
    if layout is not None:
        if _has_corner_logo_picture(layout.shapes, sw_in, sh_in):
            return True
        try:
            master = layout.slide_master
        except Exception:
            master = None
        if master is not None and _has_corner_logo_picture(
                master.shapes, sw_in, sh_in):
            return True
    return False


# ---------------------------------------------------------------------------
# Per-slide audit
# ---------------------------------------------------------------------------

def audit_slide(slide, slide_num: int, sw_in: float, sh_in: float,
                video_files_on_disk: dict[str, int]) -> SlideAudit:
    title = slide_title_text(slide)
    hidden = slide_is_hidden(slide)
    layout_name = None
    try:
        layout_name = slide.slide_layout.name
    except Exception:
        pass

    notes = slide_notes_text(slide)
    notes_wc = len(notes.split()) if notes else 0

    audit = SlideAudit(
        slide_num=slide_num,
        slide_title=title[:120],
        hidden=hidden,
        layout_name=layout_name,
        shape_count=len(list(slide.shapes)),
        has_video=False,
        has_logo=False,
        has_slide_number=False,
        notes_word_count=notes_wc,
        animation_groups=0,
    )

    issues = audit.issues

    # ------------------------------------------------------------------
    # A. Layout & Spatial
    # ------------------------------------------------------------------
    text_shapes = []
    visible_shapes = []     # for overlap/occlusion
    pic_paths = []
    media_paths = []
    logo_present = False
    slide_num_label_present = False
    slide_num_label_count = 0

    for shape, depth in iter_shapes_recursive(slide.shapes):
        try:
            l, t, r, b = shape_bbox(shape)
        except Exception:
            continue

        # Off-slide
        if l < -0.05 or t < -0.05 or r > sw_in + 0.05 or b > sh_in + 0.05:
            issues.append(Issue(
                "LAYOUT", BLOCKER,
                f"Shape '{shape.name}' extends outside slide canvas "
                f"(bbox left={l:.2f} top={t:.2f} right={r:.2f} bottom={b:.2f}; "
                f"canvas {sw_in:.2f}x{sh_in:.2f})",
                shape_id=str(getattr(shape, "shape_id", "?")),
                location=f"left={l:.2f}in,top={t:.2f}in",
            ))
        else:
            # Margin violation only when not off-slide
            if (l < 0.1 or t < 0.1 or sw_in - r < 0.1 or sh_in - b < 0.1) \
               and not is_invisible_anim_trigger(shape):
                if shape_text(shape) or shape.shape_type in (
                        MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA):
                    issues.append(Issue(
                        "LAYOUT", MINOR,
                        f"Shape '{shape.name}' within 0.1in of slide edge "
                        f"(left={l:.2f}, top={t:.2f}, right_gap={sw_in-r:.2f}, "
                        f"bot_gap={sh_in-b:.2f})",
                        shape_id=str(getattr(shape, "shape_id", "?")),
                    ))

        # Slide-number label heuristic
        # Updated: existing convention puts num in bottom-LEFT (per helpers.py
        # add_slide_num: left=MX=0.6in, top=SL_H - 0.38in ~ 7.12in)
        if shape.has_text_frame and not is_invisible_anim_trigger(shape):
            txt = shape_text(shape)
            if (re.fullmatch(r"\d{1,3}", txt) and
                emu_to_in(shape.width) <= 1.0 and
                l <= 1.5 and t >= 6.5):
                slide_num_label_present = True
                slide_num_label_count += 1

        # Logo detection (image in top-right or bottom-right corner, ~0.35in)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img_path = shape.image.filename or ""
            except Exception:
                img_path = ""
            pic_paths.append((shape, img_path))
            # logo heuristic: small (<= 0.6in), bottom-right or top-right
            sw_h = emu_to_in(shape.width)
            sh_h = emu_to_in(shape.height)
            if (sw_h <= 0.7 and sh_h <= 0.7 and
                (sw_in - r < 0.5) and
                (t < 0.5 or sh_in - b < 0.5)):
                logo_present = True

        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            media_paths.append(shape)

        # Collect for overlap analysis
        if not is_invisible_anim_trigger(shape):
            visible_shapes.append((shape, (l, t, r, b)))
            if shape.has_text_frame and shape_text(shape):
                text_shapes.append((shape, (l, t, r, b)))

        # Font-size violations.
        # Skip the slide-number label (intentionally 7-9pt per add_slide_num).
        # The renumbering loop in helpers/add_slide_num places these in the
        # bottom-left corner with a tight bbox. Match by position+size+content
        # so appendix-style labels ("A1", "A11b") and section markers
        # ("S3.2") are exempted alongside plain numerics — convention is that
        # slide-num labels are always short, ASCII, and live in the corner.
        slide_num_text = shape_text(shape) if shape.has_text_frame else ""
        is_slide_num_label = (
            shape.has_text_frame
            and slide_num_text
            and len(slide_num_text) <= 6
            # Allow numerics, appendix codes (A1, A11b), section codes (S3.2)
            and re.fullmatch(r"[A-Za-z]?\d{0,3}[A-Za-z0-9.]{0,3}", slide_num_text)
              is not None
            and emu_to_in(shape.width) <= 0.6
            and l <= 1.2 and t >= 6.8
        )
        if shape.has_text_frame and shape_text(shape) and not is_slide_num_label:
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        sz = run.font.size
                        if sz is None:
                            continue
                        pt = sz.pt
                        # Title heuristic: shape is at top of slide, large box
                        is_title_like = (
                            t < 1.2 and
                            emu_to_in(shape.width) > 6.0 and
                            len(shape_text(shape)) < 120
                        )
                        if is_title_like and pt > 44:
                            issues.append(Issue(
                                "STYLE", MINOR,
                                f"Title font {pt:.0f}pt exceeds 44pt cap "
                                f"(shape '{shape.name}')",
                                shape_id=str(getattr(shape, "shape_id", "?")),
                            ))
                        elif not is_title_like and pt < 12:
                            issues.append(Issue(
                                "STYLE", MAJOR,
                                f"Body font {pt:.0f}pt below 12pt readability "
                                f"floor (shape '{shape.name}', text \""
                                f"{shape_text(shape)[:40]}…\")",
                                shape_id=str(getattr(shape, "shape_id", "?")),
                            ))
                            break  # one report per shape
                    else:
                        continue
                    break
            except Exception:
                pass

        # Color palette violations
        try:
            if shape.has_text_frame and shape_text(shape):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            c = run.font.color.rgb
                        except Exception:
                            c = None
                        if c is None:
                            continue
                        hx = hex_color(c)
                        if hx and hx not in PALETTE:
                            issues.append(Issue(
                                "STYLE", MINOR,
                                f"Off-palette text color #{hx} "
                                f"(shape '{shape.name}')",
                                shape_id=str(getattr(shape, "shape_id", "?")),
                            ))
                            break
                    else:
                        continue
                    break
        except Exception:
            pass

        # Font family check
        try:
            if shape.has_text_frame and shape_text(shape):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        fname = run.font.name
                        if fname and fname not in ACCEPTABLE_FONTS:
                            issues.append(Issue(
                                "STYLE", MINOR,
                                f"Font '{fname}' outside acceptable set "
                                f"(shape '{shape.name}')",
                                shape_id=str(getattr(shape, "shape_id", "?")),
                            ))
                            break
                    else:
                        continue
                    break
        except Exception:
            pass

    audit.has_logo = logo_present
    audit.has_slide_number = slide_num_label_present
    audit.has_video = len(media_paths) > 0

    if not slide_num_label_present and not hidden:
        # Section dividers may legitimately omit slide number; flag MINOR
        issues.append(Issue(
            "LAYOUT", MINOR,
            "No slide-number label detected in bottom-left area "
            "(expected per add_slide_num convention).",
        ))
    if slide_num_label_count > 1:
        issues.append(Issue(
            "LAYOUT", MAJOR,
            f"{slide_num_label_count} slide-number labels detected (expected 1).",
        ))

    if not logo_present and not hidden:
        # Title slides may legitimately omit logo
        if slide_num >= 2:
            # Logos are commonly inherited from the slide layout / master
            # rather than added per-slide, AND the per-shape loop above uses
            # a tight 0.5in edge threshold that misses logos at e.g. 0.66in
            # right-gap. Re-check with a looser corner heuristic across
            # slide / layout / master before flagging.
            if not slide_has_corner_logo_anywhere(slide, sw_in, sh_in):
                issues.append(Issue(
                    "BRAND", MINOR,
                    "No corner logo (~0.35in) detected on slide, layout, or "
                    "master; expected per add_logo helper.",
                ))
            else:
                audit.has_logo = True

    # Overlap detection (visible shapes, not anim triggers).
    # Only flag pairs where BOTH shapes have visible text and where neither
    # fully contains the other (full containment is the standard
    # "text-on-card" pattern and is intentional).
    def _contains(outer, inner) -> bool:
        return (outer[0] <= inner[0] + 0.02 and
                outer[1] <= inner[1] + 0.02 and
                outer[2] >= inner[2] - 0.02 and
                outer[3] >= inner[3] - 0.02)

    overlap_pairs_seen = set()
    for i, (shA, bbA) in enumerate(visible_shapes):
        if not shape_text(shA):
            continue
        for j in range(i + 1, len(visible_shapes)):
            shB, bbB = visible_shapes[j]
            if not shape_text(shB):
                continue
            if _contains(bbA, bbB) or _contains(bbB, bbA):
                continue  # text-on-card / nested layout — expected
            inter = bbox_intersect(bbA, bbB)
            if inter <= 0:
                continue
            min_area = min(bbox_area(bbA), bbox_area(bbB))
            if min_area <= 0:
                continue
            frac = inter / min_area
            if frac > 0.05:
                key = (id(shA), id(shB))
                if key in overlap_pairs_seen:
                    continue
                overlap_pairs_seen.add(key)
                sev = MAJOR if frac > 0.30 else MINOR
                issues.append(Issue(
                    "OVERLAP", sev,
                    f"Shapes '{shA.name}' & '{shB.name}' overlap "
                    f"by {frac*100:.0f}% of smaller bbox "
                    f"(textA: \"{shape_text(shA)[:30]}\" / "
                    f"textB: \"{shape_text(shB)[:30]}\")",
                ))

    # Occlusion: text shape covered > 30% by a LATER shape
    for i, (shA, bbA) in enumerate(visible_shapes):
        if not (shape_text(shA) and shA.has_text_frame):
            continue
        areaA = bbox_area(bbA)
        if areaA <= 0:
            continue
        covered = 0.0
        for j in range(i + 1, len(visible_shapes)):
            shB, bbB = visible_shapes[j]
            covered += bbox_intersect(bbA, bbB)
        if covered / areaA > 0.30:
            issues.append(Issue(
                "OCCLUSION", MAJOR,
                f"Text shape '{shA.name}' (\""
                f"{shape_text(shA)[:40]}…\") is {(covered/areaA)*100:.0f}% "
                f"covered by later shape(s).",
                shape_id=str(getattr(shA, "shape_id", "?")),
            ))

    # ------------------------------------------------------------------
    # B. Animations
    # ------------------------------------------------------------------
    xml = get_slide_xml(slide)
    anim = analyze_animations(xml)
    audit.animation_groups = anim["click_groups"]

    if anim["click_groups"] > 8:
        issues.append(Issue(
            "ANIMATION", MAJOR,
            f"Slide has {anim['click_groups']} click-step animation groups "
            f"(>8 risks audience disengagement).",
        ))

    # Orphan animations
    spids_on_slide = set()
    for shape, _depth in iter_shapes_recursive(slide.shapes):
        try:
            spids_on_slide.add(str(shape.shape_id))
        except Exception:
            pass
    orphans = [s for s in anim["spids_referenced"]
               if s and s not in spids_on_slide]
    if orphans:
        issues.append(Issue(
            "ANIMATION", BLOCKER,
            f"Animation references {len(set(orphans))} non-existent shape id(s): "
            f"{sorted(set(orphans))[:5]}",
        ))

    # para_build with multi-paragraph shape
    for spid, btype in anim["para_build_settings"]:
        if btype in ("p", "byParagraph"):
            # Find shape, count paragraphs
            for shape, _d in iter_shapes_recursive(slide.shapes):
                if str(getattr(shape, "shape_id", "")) == str(spid):
                    if shape.has_text_frame:
                        npar = len([p for p in shape.text_frame.paragraphs
                                    if (p.text or "").strip()])
                        if npar > 1:
                            issues.append(Issue(
                                "ANIMATION", MINOR,
                                f"Shape '{shape.name}' uses para_build='{btype}' "
                                f"with {npar} paragraphs — each paragraph "
                                "animates separately.",
                                shape_id=str(spid),
                            ))
                    break

    # Missing animations on multi-card layouts (heuristic: > 4 sibling
    # rectangle/shape groups but 0 click animations)
    rect_count = 0
    for shape, _d in iter_shapes_recursive(slide.shapes):
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and \
               emu_to_in(shape.width) > 1.5 and emu_to_in(shape.height) > 1.0:
                rect_count += 1
        except Exception:
            pass
    if rect_count >= 4 and anim["click_groups"] == 0 and not hidden:
        # only suggest, MINOR
        issues.append(Issue(
            "ANIMATION", MINOR,
            f"Multi-card layout ({rect_count} large shapes) but no click "
            "animations defined — cards appear all at once.",
        ))

    if anim["errors"]:
        for e in anim["errors"]:
            issues.append(Issue("ANIMATION", BLOCKER,
                                f"Animation timing parse error: {e}"))

    # ------------------------------------------------------------------
    # C. Videos
    # ------------------------------------------------------------------
    media_rels = media_relationships(slide)
    video_meta = []

    if audit.has_video:
        # Caption / play-indicator
        play_ind_found = False
        for vshape in media_paths:
            vbb = shape_bbox(vshape)
            for shape, _d in iter_shapes_recursive(slide.shapes):
                if shape is vshape or not shape.has_text_frame:
                    continue
                txt = shape_text(shape).lower()
                if not txt:
                    continue
                if any(tok in txt for tok in ("click to play", "▶", "play",
                                              "click", "tap to play")):
                    sb = shape_bbox(shape)
                    # within ~1.5in of any edge
                    cx_v, cy_v = (vbb[0]+vbb[2])/2, (vbb[1]+vbb[3])/2
                    cx_s, cy_s = (sb[0]+sb[2])/2, (sb[1]+sb[3])/2
                    if abs(cx_v - cx_s) < 4.0 and abs(cy_v - cy_s) < 3.0:
                        play_ind_found = True
                        break
            if play_ind_found:
                break
        if not play_ind_found:
            issues.append(Issue(
                "VIDEO", MINOR,
                "Video tile lacks 'click to play' / play-indicator caption "
                "within reasonable proximity.",
            ))

    for mr in media_rels:
        # extract filename from partname (e.g. /ppt/media/media1.mp4)
        partname = mr["partname"]
        fname = Path(partname).name
        size = mr["size"]
        # Check if a same-size file exists on disk in 06_demo_videos
        match_on_disk = []
        if size is not None:
            match_on_disk = [n for n, sz in video_files_on_disk.items()
                             if sz == size]
        meta = {
            "partname": partname,
            "embedded_filename": fname,
            "embedded_size": size,
            "matched_on_disk": match_on_disk,
        }
        # ffprobe via writing the embedded blob is heavy; rely on disk match
        if match_on_disk:
            disk_path = VIDEO_DIR / match_on_disk[0]
            dur = ffprobe_duration(disk_path)
            meta["duration_s"] = dur
            if dur is not None and dur > 60:
                issues.append(Issue(
                    "VIDEO", MAJOR,
                    f"Video {match_on_disk[0]} duration "
                    f"{dur:.0f}s exceeds 60s pacing limit.",
                ))
        else:
            issues.append(Issue(
                "VIDEO", MAJOR,
                f"Embedded media ({fname}, {size or '?'} bytes) does not "
                "match any file in 06_demo_videos/ by size — likely orphaned "
                "or renamed source.",
            ))
        video_meta.append(meta)

    # Persist video meta on audit
    if video_meta:
        # Attach to first video issue or store as side-channel; we'll keep
        # it on a custom attr for the deck-level summary.
        audit_videos = video_meta  # noqa
        setattr(audit, "_videos", video_meta)

    # ------------------------------------------------------------------
    # D. Speaker notes
    # ------------------------------------------------------------------
    is_section_divider = (
        notes_wc == 0 and
        sum(1 for s, _ in iter_shapes_recursive(slide.shapes)
            if s.has_text_frame and shape_text(s)) <= 3
    )
    if notes_wc == 0 and not hidden and not is_section_divider:
        issues.append(Issue(
            "NOTES", MAJOR,
            "Body slide has no speaker notes (academic talk requires "
            "1-2 sentences min).",
        ))
    if notes_wc > 500:
        issues.append(Issue(
            "NOTES", MINOR,
            f"Speaker notes are {notes_wc} words (>500) — maintenance burden.",
        ))
    if notes:
        if re.search(r"\b(next|previous|prior|preceding|following)\s+slide\b",
                     notes, re.IGNORECASE):
            issues.append(Issue(
                "NOTES", MINOR,
                "Notes reference 'next/previous slide' — fragile if reordered.",
            ))
        # Numeric citation: any number with no adjacent doc-like reference
        nums = re.findall(r"\b\d{2,4}(?:\.\d+)?\s*%?", notes)
        if len(nums) >= 2:
            has_cite = bool(re.search(
                r"(\.md|\.csv|see\s+`|see\s+docs/|per\s+\w+\.md|file:|/home/)",
                notes, re.IGNORECASE))
            if not has_cite:
                issues.append(Issue(
                    "NOTES", MINOR,
                    f"Notes cite {len(nums)} numbers but lack any source "
                    "reference (.md/.csv path).",
                ))
        # Body-vs-notes contradiction (basic): pick first percent in notes vs body
        body_text = " ".join(
            shape_text(s) for s, _ in iter_shapes_recursive(slide.shapes)
            if s.has_text_frame)
        notes_pcts = re.findall(r"(\d{1,3})(?:\.\d+)?%", notes)
        body_pcts = re.findall(r"(\d{1,3})(?:\.\d+)?%", body_text)
        if notes_pcts and body_pcts:
            # If a body percentage is mentioned but no notes percentage is
            # within ±2pp of it, flag (very heuristic — keep MINOR)
            np_set = {int(x) for x in notes_pcts}
            for b in body_pcts:
                bi = int(b)
                if all(abs(bi - n) > 2 for n in np_set):
                    # only fire when absolute mismatch
                    if abs(bi - max(np_set, key=lambda n: -abs(n - bi))) > 2:
                        issues.append(Issue(
                            "NOTES", MINOR,
                            f"Body shows {bi}% but notes only mention "
                            f"{sorted(np_set)} — possible mismatch.",
                        ))
                        break

    return audit


# ---------------------------------------------------------------------------
# Deck-level analysis
# ---------------------------------------------------------------------------

def list_video_files() -> dict[str, int]:
    """Index video files in 06_demo_videos/.

    IMPORTANT: only top-level files and recognized live subdirs (e.g. realtalk/).
    Skip ``_archive_pre_*`` directories so archived duplicates with stale sizes
    don't shadow the current files in this size-keyed dict.
    """
    out = {}
    if not VIDEO_DIR.exists():
        return out
    # Top-level files
    for p in VIDEO_DIR.glob("*.mp4"):
        try:
            out[p.name] = p.stat().st_size
        except OSError:
            pass
    # Recognized live subdirs (non-archive)
    for sub in VIDEO_DIR.iterdir():
        if not sub.is_dir():
            continue
        if sub.name.startswith("_archive"):
            continue
        for p in sub.rglob("*.mp4"):
            # Don't overwrite a top-level entry with a same-named subdir entry
            if p.name in out:
                continue
            try:
                out[p.name] = p.stat().st_size
            except OSError:
                pass
    return out


def audit_deck(deck_path: str, video_files_on_disk: dict[str, int]) -> dict:
    prs = Presentation(deck_path)
    sw_in = emu_to_in(prs.slide_width)
    sh_in = emu_to_in(prs.slide_height)

    slide_audits: list[SlideAudit] = []
    titles_seen: list[tuple[int, str]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        sa = audit_slide(slide, idx, sw_in, sh_in, video_files_on_disk)
        slide_audits.append(sa)
        titles_seen.append((idx, sa.slide_title))

    # ------------------------------------------------------------------
    # Cross-slide checks
    # ------------------------------------------------------------------
    # Duplicate adjacent slides (>80% identical first-line text)
    for i in range(1, len(slide_audits)):
        a, b = slide_audits[i-1], slide_audits[i]
        if a.slide_title and b.slide_title and a.slide_title == b.slide_title:
            b.issues.append(Issue(
                "STRUCTURE", MINOR,
                f"Adjacent duplicate title with slide {a.slide_num} "
                f"('{a.slide_title}') — possible leftover.",
            ))

    # Hidden slides bookkeeping
    hidden_slides = [(s.slide_num, s.slide_title) for s in slide_audits
                     if s.hidden]

    # Section dividers heuristic: notes empty AND <= 3 visible text shapes AND
    # has large title-like shape
    section_dividers = []
    for sa in slide_audits:
        if sa.notes_word_count == 0 and sa.shape_count <= 6 and not sa.hidden:
            section_dividers.append((sa.slide_num, sa.slide_title))

    # Counts
    sev_counter = Counter()
    cat_counter = Counter()
    issue_total = 0
    for sa in slide_audits:
        for iss in sa.issues:
            sev_counter[iss.severity] += 1
            cat_counter[iss.category] += 1
            issue_total += 1

    # Video summary
    video_refs = []
    for sa in slide_audits:
        vmeta = getattr(sa, "_videos", None)
        if vmeta:
            for v in vmeta:
                video_refs.append({
                    "slide": sa.slide_num,
                    "title": sa.slide_title,
                    **v,
                })

    deck = {
        "deck_path": deck_path,
        "deck_name": Path(deck_path).name,
        "slide_count": len(slide_audits),
        "slide_size_in": [round(sw_in, 2), round(sh_in, 2)],
        "issue_totals": {
            "BLOCKER": sev_counter[BLOCKER],
            "MAJOR": sev_counter[MAJOR],
            "MINOR": sev_counter[MINOR],
            "TOTAL": issue_total,
        },
        "issues_by_category": dict(cat_counter),
        "hidden_slides": hidden_slides,
        "section_divider_candidates": section_dividers,
        "video_references": video_refs,
        "slides": [sa.to_dict() for sa in slide_audits],
    }
    return deck


# ---------------------------------------------------------------------------
# Markdown rendering
# ---------------------------------------------------------------------------

def render_markdown(decks: list[dict]) -> str:
    out = ["# Argos VSP — PowerPoint Visual & Structural Audit",
           "",
           "Generated by `scripts/audit_pptx_visual_structure.py`. ",
           "Read-only audit; no decks were modified.",
           "",
           "## Decks audited",
           ""]
    for d in decks:
        t = d["issue_totals"]
        out.append(f"- **{d['deck_name']}** — {d['slide_count']} slides, "
                   f"{t['TOTAL']} issues "
                   f"(BLOCKER {t['BLOCKER']}, MAJOR {t['MAJOR']}, MINOR {t['MINOR']})")
    out.append("")

    for d in decks:
        out.append(f"---\n## {d['deck_name']}")
        out.append("")
        out.append(f"- File: `{d['deck_path']}`")
        out.append(f"- Slide count: **{d['slide_count']}**")
        out.append(f"- Slide size: {d['slide_size_in'][0]}in × {d['slide_size_in'][1]}in")
        t = d["issue_totals"]
        out.append(f"- Issues: BLOCKER {t['BLOCKER']}, MAJOR {t['MAJOR']}, "
                   f"MINOR {t['MINOR']} (total {t['TOTAL']})")
        out.append("- By category:")
        for cat, n in sorted(d["issues_by_category"].items(),
                             key=lambda x: -x[1]):
            out.append(f"  - {cat}: {n}")
        out.append("")

        # Hidden slides
        if d["hidden_slides"]:
            out.append("### Hidden slides")
            out.append("")
            for n, t in d["hidden_slides"]:
                out.append(f"- Slide {n}: {t}")
            out.append("")

        # Section dividers
        if d["section_divider_candidates"]:
            out.append("### Section divider candidates")
            out.append("")
            for n, t in d["section_divider_candidates"]:
                out.append(f"- Slide {n}: {t}")
            out.append("")

        # Top-30 most problematic slides (by total issues then BLOCKER count)
        ranked = sorted(d["slides"],
                        key=lambda s: (
                            -sum(1 for i in s["issues"] if i["severity"] == BLOCKER) * 100
                            -sum(1 for i in s["issues"] if i["severity"] == MAJOR) * 10
                            -len(s["issues"])
                        ))
        ranked = [s for s in ranked if len(s["issues"]) > 0][:30]
        if ranked:
            out.append("### Top 30 most-problematic slides — fix before reusing")
            out.append("")
            out.append("| # | Slide | Title | Hidden | Issues (B/Ma/Mi) | One-line fix |")
            out.append("|---|-------|-------|--------|------------------|--------------|")
            for s in ranked:
                blockers = sum(1 for i in s["issues"] if i["severity"] == BLOCKER)
                majors = sum(1 for i in s["issues"] if i["severity"] == MAJOR)
                minors = sum(1 for i in s["issues"] if i["severity"] == MINOR)
                # Pick first BLOCKER -> first MAJOR -> first MINOR for fix hint
                top_iss = next((i for i in s["issues"] if i["severity"] == BLOCKER),
                               next((i for i in s["issues"] if i["severity"] == MAJOR),
                                    s["issues"][0]))
                fix = top_iss["description"][:160].replace("|", "\\|")
                title = (s["slide_title"] or "")[:60].replace("|", "\\|")
                out.append(f"| {ranked.index(s)+1} | {s['slide_num']} | {title} | "
                           f"{'Y' if s['hidden'] else ''} | {blockers}/{majors}/{minors} | {fix} |")
            out.append("")

        # Hidden slides worth reviving (zero issues, has notes, not a duplicate)
        revivable = [s for s in d["slides"]
                     if s["hidden"]
                     and len([i for i in s["issues"] if i["severity"] in (BLOCKER, MAJOR)]) == 0
                     and s["notes_word_count"] >= 5
                     and s["shape_count"] >= 4]
        if revivable:
            out.append("### Hidden slides worth reviving (clean layout, has notes)")
            out.append("")
            for s in revivable:
                out.append(f"- Slide {s['slide_num']}: **{s['slide_title']}** "
                           f"(notes {s['notes_word_count']} words, "
                           f"{s['shape_count']} shapes, {s['animation_groups']} anim groups)")
            out.append("")

        # Animation patterns to inherit (slides with 2-6 click groups, low issue count)
        exemplary_anim = [s for s in d["slides"]
                          if 2 <= s["animation_groups"] <= 6
                          and len([i for i in s["issues"]
                                   if i["category"] == "ANIMATION"]) == 0
                          and len(s["issues"]) <= 2]
        if exemplary_anim:
            out.append("### Animation patterns to inherit (templates)")
            out.append("")
            for s in exemplary_anim[:15]:
                out.append(f"- Slide {s['slide_num']}: **{s['slide_title']}** "
                           f"({s['animation_groups']} click steps, "
                           f"{s['shape_count']} shapes, "
                           f"{len(s['issues'])} minor cosmetic issues)")
            out.append("")

        # Video issues
        vid_issues = []
        for s in d["slides"]:
            for iss in s["issues"]:
                if iss["category"] == "VIDEO":
                    vid_issues.append((s["slide_num"], s["slide_title"], iss))
        if vid_issues or d["video_references"]:
            out.append("### Video references and issues")
            out.append("")
            if d["video_references"]:
                out.append("**Embedded videos:**")
                out.append("")
                out.append("| Slide | File | Size (bytes) | Duration (s) | Disk match |")
                out.append("|-------|------|--------------|--------------|------------|")
                for v in d["video_references"]:
                    dur = v.get("duration_s")
                    dur_s = f"{dur:.0f}" if dur is not None else "?"
                    matches = ", ".join(v.get("matched_on_disk", []) or ["—"])
                    out.append(f"| {v['slide']} | {v['embedded_filename']} | "
                               f"{v['embedded_size']} | {dur_s} | {matches} |")
                out.append("")
            if vid_issues:
                out.append("**Video issues:**")
                out.append("")
                for sn, tt, iss in vid_issues:
                    out.append(f"- Slide {sn} ({tt}): [{iss['severity']}] {iss['description']}")
                out.append("")

        # Section flow (table of contents)
        out.append("### Section flow")
        out.append("")
        out.append("| Slide | Title | Hidden | Shapes | Anim | Notes (wc) |")
        out.append("|-------|-------|--------|--------|------|-----------|")
        for s in d["slides"]:
            title = (s["slide_title"] or "")[:70].replace("|", "\\|")
            out.append(f"| {s['slide_num']} | {title} | "
                       f"{'Y' if s['hidden'] else ''} | "
                       f"{s['shape_count']} | {s['animation_groups']} | "
                       f"{s['notes_word_count']} |")
        out.append("")

    return "\n".join(out)


def render_fix_manifest(decks: list[dict]) -> str:
    out = ["# Argos VSP — PowerPoint Fix Manifest",
           "",
           "Categorized fix list for slide-writing phase. ",
           "Pulled from BLOCKER/MAJOR/MINOR issues across audited decks.",
           ""]
    for d in decks:
        out.append(f"## {d['deck_name']}")
        out.append("")
        for sev_label in (BLOCKER, MAJOR, MINOR):
            items = []
            for s in d["slides"]:
                for iss in s["issues"]:
                    if iss["severity"] == sev_label:
                        items.append((s["slide_num"], s["slide_title"], iss))
            if not items:
                out.append(f"### {sev_label}: none")
                out.append("")
                continue
            out.append(f"### {sev_label} ({len(items)})")
            out.append("")
            # Group by slide for readability
            by_slide = defaultdict(list)
            for sn, tt, iss in items:
                by_slide[(sn, tt)].append(iss)
            for (sn, tt), iss_list in sorted(by_slide.items()):
                out.append(f"**Slide {sn} — {tt}**")
                for iss in iss_list:
                    out.append(f"- [{iss['category']}] {iss['description']}")
                out.append("")
        out.append("")
    return "\n".join(out)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out-json",
                    default="/home/ubuntu/docs/evaluation/pptx_visual_audit.json")
    ap.add_argument("--out-md",
                    default="/home/ubuntu/docs/evaluation/pptx_visual_audit.md")
    ap.add_argument("--out-fix",
                    default="/home/ubuntu/docs/evaluation/pptx_fix_manifest.md")
    ap.add_argument("--decks", nargs="*", default=DECKS)
    ap.add_argument("--skip-v10", action="store_true",
                    help="Skip the v10 deck (it's identical to v9 in many builds)")
    args = ap.parse_args()

    video_files_on_disk = list_video_files()
    print(f"[audit] indexed {len(video_files_on_disk)} videos on disk", flush=True)

    decks = []
    for path in args.decks:
        if not os.path.exists(path):
            print(f"[audit] SKIP missing: {path}")
            continue
        if args.skip_v10 and "v10" in path:
            continue
        print(f"[audit] reading {path} ...", flush=True)
        try:
            d = audit_deck(path, video_files_on_disk)
        except Exception as e:
            print(f"[audit] FAILED {path}: {e}")
            continue
        t = d["issue_totals"]
        print(f"[audit]   {d['slide_count']} slides — "
              f"{t['BLOCKER']} BLOCKER, {t['MAJOR']} MAJOR, {t['MINOR']} MINOR",
              flush=True)
        decks.append(d)

    Path(args.out_json).parent.mkdir(parents=True, exist_ok=True)
    with open(args.out_json, "w") as f:
        json.dump({
            "generated_by": "scripts/audit_pptx_visual_structure.py",
            "decks": decks,
        }, f, indent=2, default=str)
    print(f"[audit] wrote {args.out_json}")

    md = render_markdown(decks)
    with open(args.out_md, "w") as f:
        f.write(md)
    print(f"[audit] wrote {args.out_md}")

    fix = render_fix_manifest(decks)
    with open(args.out_fix, "w") as f:
        f.write(fix)
    print(f"[audit] wrote {args.out_fix}")


if __name__ == "__main__":
    main()
