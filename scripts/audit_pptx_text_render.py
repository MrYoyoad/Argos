#!/usr/bin/env python3
"""Deep visual audit on a pptx — directly, without PDF render.

The bbox-based audit (`audit_pptx_visual_structure.py`) only checks shape
rectangles. It misses cases where the *rendered text* overflows the frame
(auto-fit off, text taller than frame), where text wraps and pushes
content past the slide-num zone at y=7.12, or where stacked text shapes
visually collide because their actual rendered heights exceed bbox.

This audit estimates rendered height per text shape from:
  - font size (per-run)
  - paragraph count
  - rough line-height (font_size * 1.2)
  - rough characters-per-line cap (frame_width_in / (font_size * 0.06))

For each shape, computes:
  est_rendered_h = sum_paragraphs( max(1, ceil(len(line) / cpl)) ) * line_h
And flags if est_rendered_h > frame_h * 1.05.

Also flags shapes whose bottom (top + max(frame_h, est_rendered_h))
exceeds y=7.05 (slide-num safe-zone floor).

Usage:
  python3 scripts/audit_pptx_text_render.py [pptx_path]
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Inches, Pt


PPTX = sys.argv[1] if len(sys.argv) > 1 else (
    "/home/ubuntu/presentation_materials_20260224/"
    "Argos_VSP_For_Orchard_May2026.pptx"
)

# Page geometry (from config.py): 13.33 x 7.50 widescreen
SLIDE_BOTTOM = 7.50
SLIDE_NUM_TOP = 7.12         # slide-num labels live at y >= 6.8 in bottom-left
CONTENT_BOTTOM = 7.05        # everything else must end here


def emu_to_in(emu):
    return emu / 914400.0


def estimate_rendered_height(text, font_size_pt, frame_w_in):
    """Estimate rendered text height in inches given content + font + frame width."""
    if not text or not text.strip():
        return 0.0
    line_h = (font_size_pt / 72.0) * 1.20  # 1.2 line-height factor
    # avg char width ≈ 0.55 × font height (in points)
    char_w_in = (font_size_pt * 0.55) / 72.0
    cpl = max(8, int(frame_w_in / char_w_in)) if char_w_in > 0 else 80
    h = 0.0
    for para in text.split("\n"):
        words_count = max(1, -(-len(para) // cpl))  # ceil(len/cpl)
        h += words_count * line_h
    return h


def shape_text_runs(shape):
    """Yield (text, font_size_pt, italic, color_hex_or_None) per run."""
    if not shape.has_text_frame:
        return []
    runs = []
    tf = shape.text_frame
    for para in tf.paragraphs:
        # collect runs and their sizes
        para_text = ""
        max_size = 0.0
        any_italic = False
        first_color = None
        for run in para.runs:
            para_text += run.text or ""
            sz = run.font.size
            if sz is not None:
                max_size = max(max_size, sz.pt)
            try:
                if run.font.italic:
                    any_italic = True
            except Exception:
                pass
            try:
                if first_color is None and run.font.color and run.font.color.rgb:
                    first_color = str(run.font.color.rgb)
            except Exception:
                pass
        if not para.runs:
            para_text = para.text
        if max_size == 0.0:
            # fall back to paragraph default or ambient default 18pt
            max_size = 18.0
        runs.append((para_text, max_size, any_italic, first_color))
    return runs


# Caption / footer / citation colors (per STYLE_GUIDE T1) — these may stay
# below 24pt as long as they sit in supporting roles (axis labels, footnotes,
# subtitles, credits).
_CAPTION_COLORS = {"AAAAAA", "666666", "333333"}  # LGRAY, MGRAY, DGRAY


def audit_slide(slide, slide_idx):
    issues = []
    # Walk shapes
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.left is None or shape.top is None:
            continue
        sx = emu_to_in(shape.left)
        sy = emu_to_in(shape.top)
        sw = emu_to_in(shape.width or 0)
        sh = emu_to_in(shape.height or 0)
        if sw < 0.05 or sh < 0.05:
            continue

        runs = shape_text_runs(shape)
        if not runs:
            continue
        # combined text + size for height estimate (use max font in shape)
        text = "\n".join(r[0] for r in runs if r[0])
        if not text.strip():
            continue
        max_font = max((r[1] for r in runs), default=18.0)
        any_italic = any(r[2] for r in runs)
        # use first non-None color across runs
        first_color = next((r[3] for r in runs if r[3]), None)
        is_caption_color = first_color in _CAPTION_COLORS

        est_h = estimate_rendered_height(text, max_font, sw - 0.10)  # padding

        bottom = sy + max(sh, est_h)

        # Heuristic over-estimation guard: single-paragraph single-line text
        # whose rendered chars *do* fit in a single line under our cpl
        # estimate is rarely a real overflow. The estimator over-counts
        # when many runs share a paragraph and inflate `max_font`.
        is_single_line = ("\n" not in text)
        char_w_in = (max_font * 0.55) / 72.0
        cpl = max(8, int((sw - 0.10) / char_w_in)) if char_w_in > 0 else 80
        fits_one_line = is_single_line and len(text) <= cpl
        # If it fits a single line, the only "real" overflow risk is when
        # the frame height itself is < line height (then PPT clips). Use a
        # 1-line-height check rather than the estimator's worst-case.
        line_h = (max_font / 72.0) * 1.20
        single_line_clip = fits_one_line and line_h > sh * 1.05 and (line_h - sh) > 0.15

        # Issue 1: estimated rendered height > frame * 1.05
        if (not fits_one_line) and est_h > sh * 1.05 and est_h - sh > 0.15:
            short = text.replace("\n", " ")[:70]
            issues.append({
                "slide": slide_idx,
                "kind": "TEXT_OVERFLOW_FRAME",
                "shape": shape.name,
                "frame_h": round(sh, 2),
                "est_h": round(est_h, 2),
                "font_pt": max_font,
                "lines": text.count("\n") + 1,
                "text_preview": short,
            })
        elif single_line_clip:
            short = text.replace("\n", " ")[:70]
            issues.append({
                "slide": slide_idx,
                "kind": "TEXT_OVERFLOW_FRAME",
                "shape": shape.name,
                "frame_h": round(sh, 2),
                "est_h": round(line_h, 2),
                "font_pt": max_font,
                "lines": 1,
                "text_preview": short,
            })

        # Issue 2: shape (or text) crosses content bottom — only count if
        # the actual frame (sh) is what's overflowing OR the text is multi
        # line. Single-line text inside a tall frame won't really overflow.
        text_bottom = sy + (max(sh, est_h) if not fits_one_line else max(sh, line_h))
        if text_bottom > CONTENT_BOTTOM + 0.05:
            # tolerate slide-num shapes: tiny + bottom-left
            is_slide_num = (
                sw <= 0.6 and sx <= 1.2 and sy >= 6.8
                and len(text.strip()) <= 6
            )
            if not is_slide_num:
                issues.append({
                    "slide": slide_idx,
                    "kind": "TEXT_BELOW_SAFE_ZONE",
                    "shape": shape.name,
                    "bottom": round(text_bottom, 2),
                    "limit": CONTENT_BOTTOM,
                    "font_pt": max_font,
                    "text_preview": text.replace("\n", " ")[:70],
                })

        # Issue 3: rendered text exceeds slide bottom
        if text_bottom > SLIDE_BOTTOM + 0.02:
            issues.append({
                "slide": slide_idx,
                "kind": "TEXT_OFF_CANVAS",
                "shape": shape.name,
                "bottom": round(text_bottom, 2),
                "canvas": SLIDE_BOTTOM,
                "text_preview": text.replace("\n", " ")[:70],
            })

        # Issue 4: tiny font for body text (<24pt)
        # body shapes: not in slide-num zone, not in title region (top<0.8),
        # not in footer (top>=6.0). Captions / footnotes (italic + grey,
        # or grey-coloured, or in y>=5.8 footer area) may stay at 18pt+
        # per STYLE_GUIDE T1 (footers minimum 16pt, recommended 18pt).
        is_body = sy >= 0.8 and sy <= 6.0 and not (
            sw <= 0.6 and sx <= 1.2 and sy >= 6.8
        )
        is_caption_role = (
            is_caption_color           # MGRAY / LGRAY / DGRAY text
            or any_italic              # italic = subtitle / caption / footnote
            or sy >= 5.8               # near-footer zone
            or sh <= 0.45              # tiny annotation slug (one-liner under a chart)
            or sw < 6.5                # narrow column (2-col layouts ~5-6") — one-tier-down
        )
        # Allow Pt(20) for narrow-column body (one tier under 24pt floor, per V-rules)
        # — flag only if shape is genuinely body AND font drops below 20pt.
        floor = 18 if (is_body and sw < 6.5) else 24
        if is_body and max_font < floor and len(text.strip()) > 12 and not is_caption_role:
            issues.append({
                "slide": slide_idx,
                "kind": "FONT_BELOW_24PT_BODY",
                "shape": shape.name,
                "font_pt": max_font,
                "text_preview": text.replace("\n", " ")[:70],
            })

    return issues


def main():
    p = Presentation(PPTX)
    all_issues = []
    by_kind = {}
    for i, slide in enumerate(p.slides, 1):
        if slide._element.get("show") == "0":
            # appendix hidden — still audit, but tag
            pass
        sis = audit_slide(slide, i)
        all_issues.extend(sis)
        for s in sis:
            by_kind[s["kind"]] = by_kind.get(s["kind"], 0) + 1

    print(f"=== Deep text-render audit: {Path(PPTX).name} ===")
    print(f"Slides: {len(p.slides)}  |  Issues: {len(all_issues)}")
    for k, n in sorted(by_kind.items(), key=lambda kv: -kv[1]):
        print(f"  {k:30s} {n}")
    print()

    # Print top issues by kind
    KIND_ORDER = [
        "TEXT_OFF_CANVAS",
        "TEXT_BELOW_SAFE_ZONE",
        "TEXT_OVERFLOW_FRAME",
        "FONT_BELOW_24PT_BODY",
    ]
    for kind in KIND_ORDER:
        items = [i for i in all_issues if i["kind"] == kind]
        if not items:
            continue
        print(f"--- {kind} ({len(items)}) ---")
        for it in items[:50]:
            extras = " ".join(
                f"{k}={v}" for k, v in it.items()
                if k not in ("slide", "kind", "shape", "text_preview")
            )
            print(f"  Slide {it['slide']:3d} {it['shape']:20s} {extras}")
            print(f"             text: {it['text_preview']}")
        if len(items) > 50:
            print(f"  ...and {len(items)-50} more")
        print()

    # Save JSON
    import json
    out = Path("/home/ubuntu/docs/evaluation/pptx_text_render_audit.json")
    out.write_text(json.dumps({
        "pptx": PPTX,
        "totals": {
            "slides": len(p.slides),
            "issues": len(all_issues),
            "by_kind": by_kind,
        },
        "issues": all_issues,
    }, indent=2))
    print(f"JSON: {out}")


if __name__ == "__main__":
    main()
