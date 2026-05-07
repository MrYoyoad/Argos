#!/usr/bin/env python3
"""
Extract every numeric mention from Argos VSP PPTX decks for the MBR-default
audit. Walks every slide of every input deck (recursive into grouped shapes
and tables), pulls speaker notes, and tags each number with a status:

    NEEDS_VERIFICATION  - likely shifted under MBR-default vs March top-1
    LIKELY_HOLD         - judge / kappa / r values not affected by decoder
    WRONG_FRAMING_FIX   - deprecated framing (PCA "3 dims", IS>=3.0 "captured")
    METADATA            - dates, slide numbers, page numbers, year references

Outputs CSV to /home/ubuntu/docs/evaluation/pptx_number_audit.csv.
"""
from __future__ import annotations

import csv
import os
import re
import sys
from dataclasses import dataclass
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# Inputs / outputs
# ---------------------------------------------------------------------------

DECKS = [
    ("AFTER_AMOSI_Mar2026",
     "/home/ubuntu/presentation_materials_20260224/Argos_VSP_Final_84slides_Mar2026.pptx"),
    ("Client_v9_May2026",
     "/home/ubuntu/presentation_materials_20260224/Argos_VSP_Client_v9_May2026.pptx"),
    ("Client_v10_May2026",
     "/home/ubuntu/presentation_materials_20260224/Argos_VSP_Client_v10_May2026.pptx"),
    ("Supplementary_LLM_Judge_30",
     "/home/ubuntu/presentation_materials_20260224/Supplementary_LLM_Judge_30_Examples.pptx"),
]

CSV_PATH = "/home/ubuntu/docs/evaluation/pptx_number_audit.csv"


# ---------------------------------------------------------------------------
# Number regex
#
# Matches: 12, 12.5, 12,345, 1,497, 64.1%, 0.690, r=0.925, p=0.0002,
# 1e-3, 0.85-0.95 (range -> two matches), 346/1,497, March 11, May 2 2026
# ---------------------------------------------------------------------------

NUM_RE = re.compile(
    r"""
    (?<![A-Za-z_])                  # not part of an identifier
    (?:
        \d{1,3}(?:,\d{3})+(?:\.\d+)?    # 1,497  or  1,234.56
        | \d+(?:\.\d+)?                  # 12  or  12.5  or  64.1
    )
    (?:\s?%)?                       # optional trailing %
    (?:e[-+]?\d+)?                  # optional scientific
    """,
    re.VERBOSE,
)


# ---------------------------------------------------------------------------
# Classification heuristics
# ---------------------------------------------------------------------------

# Numbers that almost certainly shift under MBR-default vs top-1 baseline.
# Check both the canonical printed form and a few common variants.
SHIFTED_VALUES = {
    # IS mean
    "2.52", "2.5",
    # Mean WER / WWER / NEA
    "64.1", "64.1%", "60.5", "60.5%", "38.9", "38.9%",
    # NIV-Y
    "23.1", "23.1%", "346",
    # NIV-Y+P
    "61.6", "61.6%", "922",
    # Salvage
    "51.1", "51.1%", "165",
    # Hallucination
    "20.5", "20.5%", "307",
    # IS tier %
    "18.4", "18.4%", "21.4", "21.4%", "21.7", "21.7%",
    "22.4", "22.4%", "16.0", "16.0%",
    # IS tier counts
    "276", "321", "325", "336", "239",
    # Three-number staircase fragments
    "40.1", "40.1%", "39.9", "39.9%",
    # Old "captured" count
    "597",
}

# Numbers that should hold under MBR-default (not decoder-dependent).
HELD_VALUES = {
    # Kappa / r values
    "0.690", "0.818", "0.925", "0.934", "0.565", "0.521", "0.629", "0.777",
    # LLM Judge blind verdicts
    "23.0", "23.0%", "41.8", "41.8%", "35.1", "35.1%", "64.9", "64.9%",
    # LLM Judge context-aware
    "15.0", "15.0%", "47.1", "47.1%", "37.9", "37.9%", "62.1", "62.1%",
    # Fine-tune
    "2.487", "2.312", "2.023", "1273", "1,273",
    "7.1", "7.1%", "12.5", "12.5%", "26.8", "26.8%",
    # 30 LLM judge examples
    "30",
    # Intra-rater
    "86.7", "86.7%",
    # Tuning
    "13", "107",
    # Cross-config
    "0.015", "16",
    # Expert heuristic delta
    "0.93", "0.93%",
    # Sample sizes
    "1497", "1,497",
}

# Phrases / token sequences that signal deprecated framing.
WRONG_FRAMING_PHRASES = [
    "3 dimensions", "three dimensions", "3 dims", "three dims",
    "3 principal", "three principal",
    "is >= 3.0", "is>=3.0", "is≥3.0", "is ≥ 3.0",
    "captured (is",
    "39.9%", "40.1%",
]

# Words that, if present near a number, suggest the number is metadata
# (slide number / year / page count / version / similar) rather than data.
METADATA_HINT_WORDS = {
    "page", "slide", "version", "ver.", "section",
}

# Months — let us catch dates like "March 11", "May 2 2026"
MONTH_NAMES = {
    "january", "february", "march", "april", "may", "june",
    "july", "august", "september", "october", "november", "december",
    "jan", "feb", "mar", "apr", "jun", "jul", "aug", "sep", "sept", "oct", "nov", "dec",
}


def normalize_number(token: str) -> str:
    """Strip whitespace, drop comma thousands separators, keep %."""
    t = token.strip().replace(" ", "")
    # collapse comma thousands but keep decimal point
    if "," in t and "." in t:
        # comma is thousands sep
        t = t.replace(",", "")
    elif "," in t:
        # ambiguous (could be euro decimal) but in our decks comma = thousands
        t = t.replace(",", "")
    return t


# Deprecated number tokens that, when they appear, indicate WRONG_FRAMING_FIX.
# These are the NUMBERS themselves (not phrases) that signal old framing.
DEPRECATED_NUMBER_TOKENS = {
    # IS>=3.0 captured framing
    "39.9", "39.9%", "40.1", "40.1%",
    # Old captured count
    "597",
}

# Markers in the surrounding text that say "we KNOW this is old, here's why we replaced it".
# When present, suppress WRONG_FRAMING_FIX (it's a historical/contrast reference, not a live claim).
HISTORICAL_MARKERS = (
    "old is", "old threshold", "superseded", "deprecated", "wrongly rejected",
    "legacy", "previously",
)


def classify(token: str, surrounding_text: str) -> str:
    raw = token.strip()
    norm = normalize_number(raw)
    norm_no_pct = norm.rstrip("%")
    surround_lower = surrounding_text.lower()
    is_historical = any(m in surround_lower for m in HISTORICAL_MARKERS)

    # Metadata: explicit slide/page/version markers, or 4-digit year
    for hint in METADATA_HINT_WORDS:
        if hint in surround_lower:
            return "METADATA"
    # Year? 19xx / 20xx
    try:
        as_int = int(norm_no_pct)
        if 1900 <= as_int <= 2099 and "." not in norm_no_pct:
            return "METADATA"
    except ValueError:
        pass
    # Month adjacency -> date
    for m in MONTH_NAMES:
        if m in surround_lower:
            return "METADATA"

    # WRONG_FRAMING_FIX: only when the TOKEN ITSELF is one of the deprecated
    # numbers AND the surrounding text doesn't already mark it as historical.
    # Also catch deprecated phrases like "3 dimensions" / "three principal components"
    # but only flag the digit "3" / "three"-adjacent numbers.
    if (norm in DEPRECATED_NUMBER_TOKENS or norm_no_pct in DEPRECATED_NUMBER_TOKENS) and not is_historical:
        return "WRONG_FRAMING_FIX"
    # PCA "3 dimensions" framing: flag any "3"-token immediately followed by "dimensions"/"principal"
    if norm_no_pct == "3" and any(p in surround_lower for p in
                                   ("3 dimensions", "3 dims", "3 principal", "three dimensions",
                                    "three dims", "three principal")):
        return "WRONG_FRAMING_FIX"

    # Shifted vs held value lookup
    if raw in SHIFTED_VALUES or norm in SHIFTED_VALUES or norm_no_pct in SHIFTED_VALUES:
        return "NEEDS_VERIFICATION"
    if raw in HELD_VALUES or norm in HELD_VALUES or norm_no_pct in HELD_VALUES:
        return "LIKELY_HOLD"

    # P-value pattern (always hold; statistical, not decoder-dependent)
    if "p=" in surround_lower or "p <" in surround_lower or "p<" in surround_lower or "p-value" in surround_lower:
        return "LIKELY_HOLD"

    # Default: if the surrounding text mentions WER / IS / NIV / hallucinat / salvage,
    # call it NEEDS_VERIFICATION; otherwise leave UNCLASSIFIED.
    for kw in ("wer", "wwer", "intelligibility", " is ", " is=", " is≥", "niv", "salvage",
               "hallucin", "captur", "nea", "named entity"):
        if kw in surround_lower:
            return "NEEDS_VERIFICATION"

    return "UNCLASSIFIED"


# ---------------------------------------------------------------------------
# PPTX walking
# ---------------------------------------------------------------------------

@dataclass
class TextChunk:
    location: str       # e.g. "title", "body", "table_cell[2,3]", "notes", "group/shape3/text"
    text: str


def _iter_table(shape, prefix: str) -> Iterable[TextChunk]:
    table = shape.table
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            txt = (cell.text or "").strip()
            if txt:
                yield TextChunk(f"{prefix}/cell[{r_idx},{c_idx}]", txt)


def _iter_shape(shape, prefix: str) -> Iterable[TextChunk]:
    """Recursively yield text chunks from a shape (handles groups + tables)."""
    try:
        st = shape.shape_type
    except Exception:
        st = None

    # Groups: recurse
    if st == MSO_SHAPE_TYPE.GROUP:
        for i, sub in enumerate(shape.shapes):
            yield from _iter_shape(sub, f"{prefix}/group[{i}]")
        return

    # Tables: walk cells
    if getattr(shape, "has_table", False):
        yield from _iter_table(shape, prefix)
        return

    # Plain text
    if shape.has_text_frame:
        tf = shape.text_frame
        for p_idx, para in enumerate(tf.paragraphs):
            txt = "".join(run.text for run in para.runs).strip()
            if not txt:
                # Some paragraphs only have a single run with no .runs registered
                txt = (para.text or "").strip()
            if txt:
                yield TextChunk(f"{prefix}/para[{p_idx}]", txt)


def _slide_title(slide) -> str:
    title_shape = None
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None
    if title_shape and title_shape.has_text_frame:
        return (title_shape.text or "").strip().replace("\n", " | ")
    # Fallback: first text shape
    for shape in slide.shapes:
        if shape.has_text_frame and (shape.text or "").strip():
            return shape.text.strip().split("\n", 1)[0]
    return ""


def _is_hidden(slide) -> bool:
    show = slide._element.get("show")
    return show == "0"


def _slide_chunks(slide) -> Iterable[TextChunk]:
    for i, shape in enumerate(slide.shapes):
        yield from _iter_shape(shape, f"shape[{i}]")
    # Speaker notes
    try:
        if slide.has_notes_slide:
            ns = slide.notes_slide
            if ns.notes_text_frame:
                for p_idx, para in enumerate(ns.notes_text_frame.paragraphs):
                    txt = (para.text or "").strip()
                    if txt:
                        yield TextChunk(f"notes/para[{p_idx}]", txt)
    except Exception:
        pass


def extract_numbers_from_chunk(chunk: TextChunk) -> list[tuple[str, str]]:
    """Return list of (number_token, surrounding_text)."""
    matches = list(NUM_RE.finditer(chunk.text))
    out = []
    for m in matches:
        token = m.group(0).strip()
        if not token:
            continue
        # Use the entire chunk text as surrounding (chunks are paragraph-level)
        out.append((token, chunk.text))
    return out


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> int:
    os.makedirs(os.path.dirname(CSV_PATH), exist_ok=True)
    rows: list[dict] = []
    for deck_label, path in DECKS:
        if not os.path.exists(path):
            print(f"[skip] {path} (missing)", file=sys.stderr)
            continue
        try:
            prs = Presentation(path)
        except Exception as e:
            print(f"[error] {path}: {e}", file=sys.stderr)
            continue
        for s_idx, slide in enumerate(prs.slides, start=1):
            title = _slide_title(slide)
            hidden = _is_hidden(slide)
            for chunk in _slide_chunks(slide):
                pairs = extract_numbers_from_chunk(chunk)
                for token, surrounding in pairs:
                    status = classify(token, surrounding)
                    rows.append({
                        "deck": deck_label,
                        "slide_num": s_idx,
                        "slide_title": title,
                        "hidden": "yes" if hidden else "no",
                        "location_in_slide": chunk.location,
                        "surrounding_text": surrounding[:400],
                        "number_mention": token,
                        "normalized": normalize_number(token),
                        "suggested_status": status,
                    })

    with open(CSV_PATH, "w", newline="", encoding="utf-8") as fh:
        fieldnames = [
            "deck", "slide_num", "slide_title", "hidden",
            "location_in_slide", "surrounding_text",
            "number_mention", "normalized", "suggested_status",
        ]
        w = csv.DictWriter(fh, fieldnames=fieldnames)
        w.writeheader()
        for r in rows:
            w.writerow(r)
    print(f"[ok] wrote {len(rows)} rows -> {CSV_PATH}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
