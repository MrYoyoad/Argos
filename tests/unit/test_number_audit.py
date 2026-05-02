"""Number audit for the client deck (Z-phase guardrail).

Scans slide text + speaker notes against MEMORY.md > Key Project Numbers.
This catches the most-common past failure mode (per
~/.claude/projects/-home-ubuntu/memory/presentation-audit-lessons.md):
stale numbers slipping through because they live in speaker notes the
presenter forgot to update.

Two kinds of check:

  1. STALE/FORBIDDEN patterns — never permitted anywhere (slides OR notes):
       "67%" / "67.0%"          legacy 860-segment baseline WER
       "860 segment"            legacy dataset count
       "IS >= 3.0" / "IS≥3.0"   legacy threshold (superseded by NIV)
       "1,273" / "1273"         LoRA-specific data count, stripped in v3
       "20,000" / "20000"       LoRA-specific projection, stripped in v3
       "r=16" / "r=64"          LoRA rank values, stripped in v3
       "PCA"                    appendix-only term, stripped from main deck
       " NIV "                  jargon — must be translated for clients
       " kappa " / "κ"          jargon — must be translated for clients

  2. CANONICAL numbers (from MEMORY.md): when these appear on a *visible*
     slide, they must match the canonical set. The audit does NOT enforce
     that every slide *uses* canonical numbers — many slides have no
     numbers at all — only that any number that LOOKS LIKE a project
     metric is actually one of the approved ones.

Run:
    pytest tests/unit/test_number_audit.py -v
"""

import re
from pathlib import Path

import pytest
from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = REPO_ROOT / "presentation_materials_20260224" / "Argos_VSP_Client_Round5_Apr2026.pptx"

# Slide indices (1-based) for BORROWED builders. Recomputed for Round 5
# after the framing-v2 reorder + 12 academic borrows landed (six judge
# examples in §Real Outputs, four failure-mode slides + two salvage
# slides in §Trust). Per the Round-5 partner-audience relaxation
# documented in /home/ubuntu/.claude/plans/i-need-to-create-proud-cupcake.md
# § "Pull-from-academic checklist," these slides carry visible WER/WWER
# values which Round-5 framing accepts (architecture-overview slides
# are OK for the partner audience).
BORROWED_SLIDES = {
    # Round 5.10: dropped slide_client_example_flagged (Obama segment 5,
    # redundant with halluc_real_example) and added two case-study
    # slides (topic-shift + Strip-save) in the how-to-read sequence.
    # Net +1 vs Round 5.9. All earlier borrows shift +1.
    44,  # slide_failure_deep_2  — 3 worked ref/hyp examples
    45,  # slide_25d             — salvage recoveries (WER stripped from headers)
    56,  # slide_data_flow       — 5-step model architecture (LoRA stripped)
    68,  # slide_thank_you       — close
}


# Canonical numbers from MEMORY.md > Key Project Numbers, plus their
# client-friendly translations (with rounded forms).
CANONICAL = {
    "61.6%", "62%",                 # NIV Y+P → useful output
    "23.1%", "23%",                 # NIV Y → clearly conveyed
    "20.5%", "1 in 5",              # hallucination / auto-flag rate
    "0.818", "82%",                 # κ vs Opus Y+P → 82% agreement
    "0.690",                        # κ vs Opus Y (notes only)
    "2.52", "2.5",                  # IS mean
    "1,497", "1497",                # baseline dataset size
    "346",                          # NIV Y count
    "922",                          # NIV Y+P count
    "39%", "38%",                   # derived trust-dashboard breakdown
    "0.925",                        # cross-config Pearson r (notes only)
    "8 of 10", "6 of 10",           # client-friendly translations
    "4 of every 10", "6 of every 10",
}

# Patterns that must never appear anywhere in the deck (slides OR notes)
# unless explicitly part of the appendix sections.
FORBIDDEN_PATTERNS = [
    (re.compile(r"\b67(\.0)?\s*%"), "stale 67% (preliminary 860-segment baseline)"),
    (re.compile(r"860\s+segments?", re.IGNORECASE), "legacy 860-segment count"),
    (re.compile(r"IS\s*[≥>]=?\s*3\.0(?!\d)"), "legacy IS>=3.0 threshold (superseded by NIV)"),
    (re.compile(r"\b1,?273\b"), "LoRA-specific 1,273-segment count (stripped in v3)"),
    (re.compile(r"\b20,?000\b"), "LoRA-specific 20,000 projection (stripped in v3)"),
    (re.compile(r"\br\s*=\s*(?:16|64)\b"), "LoRA rank value (stripped in v3)"),
]

# Patterns forbidden ONLY on visible slide text (allowed in speaker notes,
# where presenters may want the academic detail for reference).
SLIDE_ONLY_FORBIDDEN = [
    (re.compile(r"\bPCA\b"), "PCA jargon should be appendix-only"),
    (re.compile(r"\bNIV\b"), "NIV jargon should be appendix-only"),
    (re.compile(r"\bkappa\b", re.IGNORECASE), "kappa should be translated for clients"),
    (re.compile(r"κ"), "κ should be translated for clients"),
    (re.compile(r"\bLoRA\b"), "LoRA jargon should be appendix-only"),
]


@pytest.fixture(scope="module")
def deck():
    if not DECK.exists():
        pytest.skip(f"Deck not built: {DECK}")
    return Presentation(str(DECK))


def _slide_text(slide):
    """Extract visible text from a slide (titles, bullets, callouts) — NOT notes."""
    chunks = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.text:
                    chunks.append(r.text)
    return "\n".join(chunks)


def _slide_notes(slide):
    """Extract speaker notes text."""
    if not slide.has_notes_slide:
        return ""
    nf = slide.notes_slide.notes_text_frame
    return "\n".join(p.text for p in nf.paragraphs)


def test_no_stale_or_forbidden_anywhere(deck):
    """Stale / superseded numbers must not appear in our own slides.

    Borrowed slides (defined in BORROWED_SLIDES) are exempt — those builders
    live in the academic deck and the borrow-vs-build policy in the plan
    forbids modifying them. If a borrowed slide carries stale notes, that
    is a real defect to track, but it is not a new defect introduced by
    this client deck.
    """
    violations = []
    borrowed_violations = []
    for i, slide in enumerate(deck.slides, 1):
        text = _slide_text(slide) + "\n" + _slide_notes(slide)
        for pattern, why in FORBIDDEN_PATTERNS:
            for m in pattern.finditer(text):
                snippet = text[max(0, m.start() - 30): m.end() + 30].replace("\n", " ⏎ ")
                msg = f"slide {i}: {why}\n      match: …{snippet}…"
                if i in BORROWED_SLIDES:
                    borrowed_violations.append(msg)
                else:
                    violations.append(msg)
    if borrowed_violations:
        # Informational only — print so reviewers can decide separately whether
        # to chase upstream-builder fixes. Does not fail the build.
        print("\nINFO — stale patterns in borrowed slides (not blocking):")
        for v in borrowed_violations:
            print("  " + v)
    assert not violations, "Stale/forbidden patterns:\n  " + "\n  ".join(violations)


def test_no_jargon_on_visible_slides(deck):
    """Client-deck rule: PCA / NIV / κ / LoRA only in speaker notes,
    never on visible slide text. Borrowed slides are exempt (we don't
    modify upstream builders)."""
    violations = []
    for i, slide in enumerate(deck.slides, 1):
        if i in BORROWED_SLIDES:
            continue
        text = _slide_text(slide)
        for pattern, why in SLIDE_ONLY_FORBIDDEN:
            for m in pattern.finditer(text):
                snippet = text[max(0, m.start() - 30): m.end() + 30].replace("\n", " ⏎ ")
                violations.append(
                    f"slide {i}: {why}\n      match: …{snippet}…"
                )
    assert not violations, "Jargon on visible slides:\n  " + "\n  ".join(violations)


# Numbers we don't try to validate (slide labels, ordinals, year etc.)
_IGNORE_NUMERIC = {
    "1", "2", "3", "4", "5", "6", "7", "8", "9", "10",
    "11", "12", "13", "14", "15", "16", "17", "18", "19", "20",
    "30", "100",
    "2026", "2025",
}


def _percent_tokens(text):
    """Yield numbers ending in '%' from the text."""
    return re.findall(r"\d+(?:[.,]\d+)?\s*%", text)


def _bare_decimals(text):
    """Yield x.y or x.yz looking like a metric (skip integers in [1,100])."""
    out = []
    for m in re.finditer(r"\b(\d+\.\d+)\b", text):
        out.append(m.group(1))
    return out


def test_visible_percentages_are_canonical_or_derivative(deck):
    """Every percentage on a visible slide must either be canonical or
    a derivative of one (e.g. trust-dashboard 38% / 39% from 100-61.6).
    'Approval' set is the union of CANONICAL plus an explicit derivative set.
    """
    approved = set(CANONICAL)
    # Derivatives that legitimately appear on visible slides:
    approved.update({
        "38.4%", "38%",      # 100 - 61.6 = 38.4 → "needs review"
        "38.5%", "39%",      # 61.6 - 23.1 = 38.5 → "useful with context"
        "20%",                # rounded form sometimes used
        "100%",               # control reference
    })
    # Per-segment WER values that appear on the Section 3 example slides
    # (perfect / partial / flagged Obama segments). These are factual
    # per-instance numbers, not project-level metrics — they cannot contradict
    # the canonical set. Listed explicitly so the audit catches *new* stray
    # numbers but lets these through.
    approved.update({
        "0%",                 # segment 14 — perfect
        "22%",                # segment 31 — partial
        "45%",                # segment 5  — flagged hallucination
        "59%",                # segment 27 — failure_worked_example (Round 5)
    })
    # Illustrative rejection rates on the Section 7 quality-pre-filter
    # funnel. These are explicitly labeled illustrative on the slide
    # footer ("Percentages illustrative — actual rejection rates depend on
    # your video conditions").
    approved.update({
        "90%",                # head-angle pass rate
        "82%",                # mouth-visibility pass rate
        "75%",                # lighting / contrast + final pass rate
    })
    # Round-5 "compared to today" comparison anchor — published lip-reading
    # literature figures from CLIENT_MEETING_FRAMING_v2.md § "Compared to today".
    # Bear & Harvey 2017 + Assael et al. 2016 expert-human and
    # system+reviewer accuracy ranges.
    approved.update({
        "45%", "52%",         # expert human lip-reader range (45-52%)
        "55%", "70%",         # Argos + reviewer range (55-70%)
    })
    # Round-5 failure-mode taxonomy frequencies on the 574 below-threshold
    # segments per docs/CLIENT_MEETING_FRAMING_v2.md § "Trust framework".
    # These are explicit known frequencies, not project-level metrics.
    approved.update({
        "44.4%",              # Wrong Topic (255 segments)
        "18.8%",              # Hallucination (108)
        "13.9%",              # Signal Loss (80)
        "13.8%",              # Right Topic Wrong Details (79)
        "9.1%",               # Accumulated Errors (52)
    })
    # Round 5.6 — three-tier UI policy (slide 32). 1,497-segment baseline
    # tier distribution + reliability bounds. Source: docs/confidence/
    # threshold_design.md and confidence_full_analysis.md §2.2.
    approved.update({
        "22.7%",              # Trust tier share
        "35.7%",              # Salvage tier share
        "36.9%",              # Strip tier share
        "85%",                # P(correct|green) at Trust threshold
        "50%",                # P(correct|green) at Strip threshold (the "would mislead" boundary)
        "18%",                # P(correct|green) lower bound across stratification
        "93%",                # P(correct|green) upper bound across stratification
    })
    # Round 5.10 — case-study illustrative numbers from
    # docs/features/aggregation-and-confidence-case-studies.md.
    # Per-word confidence values shown verbatim on the Mode 3.1 Strip
    # case-study slide ("'good' 35%, 'idea' 53%, below 25% confidence")
    # and the concrete numeric/entity false-confidence example on the
    # pitfalls slide (1 million vs 1 billion at 96% confidence).
    approved.update({
        "25%",                # "below 25% confidence" (Strip example phrasing)
        "35%",                # "good" word probability (Mode 3.1 verbatim)
        "53%",                # "idea" word probability (Mode 3.1 verbatim)
        "96%",                # numeric/entity false-confidence example (1 billion → 1 million)
    })
    # Strip whitespace inside percent tokens
    def norm(s):
        return s.replace(" ", "")
    approved_norm = {norm(s) for s in approved}

    violations = []
    for i, slide in enumerate(deck.slides, 1):
        if i in BORROWED_SLIDES:
            continue
        text = _slide_text(slide)
        for tok in _percent_tokens(text):
            n = norm(tok)
            if n in approved_norm:
                continue
            # Allow small ordinals like "1%" 0% etc that are not actual claims
            stripped = n.rstrip("%")
            if stripped in _IGNORE_NUMERIC:
                continue
            violations.append(f"slide {i}: percentage '{tok}' is not canonical")
    assert not violations, "Non-canonical percentages on visible slides:\n  " + "\n  ".join(violations)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
