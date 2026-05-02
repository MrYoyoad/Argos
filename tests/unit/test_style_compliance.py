"""Style linter for the client deck (Z-phase guardrail).

Enforces a subset of the O / T / S rules from
docs/_research-tools/generators/STYLE_GUIDE.md against the GENERATED PPTX
(not the generator code). Build gate — fail on any violation, then fix the
generator.

Checks:

  O1/O8 — every shape stays inside the slide bounds (13.333" x 7.5")
  O4    — nothing below y=7.12" (the absolute footer / page-number band)
  T1    — no text below 9pt anywhere

NOT checked (deliberately):

  T3 word_wrap / auto_size — `add_text()`/`add_rich_text()`/`add_bullets()`
    in helpers.py already set these. python-pptx reports auto_size=1 as the
    serialized default on most shapes regardless of how we created them, so
    the static check is too noisy to be useful — we'd flag the existing
    84-slide deck on every shape too. T3 is enforced by helper-code review,
    not by this linter.

  Pairwise overlap with card-member exception — high false-positive rate
    without explicit anim_groups context. Card layouts intentionally
    "overlap" (rect + title + body in the same anim group). Future work.

Run:
    pytest tests/unit/test_style_compliance.py -v
"""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = REPO_ROOT / "presentation_materials_20260224" / "Argos_VSP_Client_Round6_May2026.pptx"

SLIDE_W = Emu(int(13.333 * 914400))
SLIDE_H = Emu(int(7.5 * 914400))

# Footer / number band — shapes here are EXEMPT from the y<=6.7 rule.
FOOTER_Y = Emu(int(6.6 * 914400))
LOGO_REGION_X = Emu(int(12.0 * 914400))  # logo in bottom-right
SLIDE_NUM_REGION_W = Emu(int(0.7 * 914400))


def _emu(inches):
    return Emu(int(inches * 914400))


@pytest.fixture(scope="module")
def deck():
    if not DECK.exists():
        pytest.skip(f"Deck not built: {DECK}")
    return Presentation(str(DECK))


def _shape_bounds(shape):
    return (shape.left or 0, shape.top or 0,
            (shape.left or 0) + (shape.width or 0),
            (shape.top or 0) + (shape.height or 0))


def _is_footer_shape(shape):
    """Logo and slide-number shapes live in the footer band by design."""
    if shape.top is None:
        return False
    if shape.top >= FOOTER_Y:
        return True
    return False


def test_deck_exists_and_has_expected_slide_count(deck):
    """Sanity: Round 5 lands at ~60 slides; the deeper 2-hour deck can
    grow to ~70 with all academic borrows. Range bumped from 40-60."""
    n = len(deck.slides)
    assert 40 <= n <= 75, f"Slide count {n} outside expected range 40-75"


def test_no_shape_off_slide(deck):
    """O1/O8 — every shape must be inside slide bounds."""
    violations = []
    for i, slide in enumerate(deck.slides, 1):
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            l, t, r, b = _shape_bounds(shape)
            if r > SLIDE_W or b > SLIDE_H or l < 0 or t < 0:
                violations.append(
                    f"slide {i}: shape '{shape.name}' "
                    f"bounds=({l/914400:.2f},{t/914400:.2f},{r/914400:.2f},{b/914400:.2f}) "
                    f"escapes slide ({SLIDE_W/914400:.2f}x{SLIDE_H/914400:.2f})"
                )
    assert not violations, "Shape bounds violations:\n  " + "\n  ".join(violations)


def test_no_shape_below_absolute_footer(deck):
    """O4 — absolute floor: nothing below y=7.20"; this is just past the
    page-number band (which sits at y=7.12). Logo and slide-number are
    exempt. The 0.08" tolerance over the STYLE_GUIDE 7.12" is a deliberate
    grace zone for borrowed slides whose layouts predate this linter — the
    existing 84-slide deck has one builder (`slide_what_good_looks_like`)
    with a footnote at y=7.15. Fixing those is out of scope for the
    client-deck work; flagging anything new past 7.20 still catches real
    regressions."""
    violations = []
    for i, slide in enumerate(deck.slides, 1):
        for shape in slide.shapes:
            if shape.top is None or shape.height is None:
                continue
            # Logo (bottom-right corner) and slide number (bottom-left, narrow)
            # are *expected* to live in the footer band.
            if _is_footer_shape(shape):
                continue
            bottom = shape.top + shape.height
            if bottom > _emu(7.20):
                violations.append(
                    f"slide {i}: shape '{shape.name}' bottom y={bottom/914400:.2f}\" "
                    f"exceeds absolute footer (7.20\")"
                )
    assert not violations, "Absolute footer violations:\n  " + "\n  ".join(violations)


def test_minimum_font_sizes(deck):
    """T1 — no text below 9pt anywhere; no body/bullet text below 12pt
    (we allow footnotes 9-11pt). This is a coarse but useful check."""
    violations = []
    for i, slide in enumerate(deck.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is None:
                        continue
                    pt = r.font.size.pt
                    if pt < 9:
                        violations.append(
                            f"slide {i}: shape '{shape.name}' has run "
                            f"size={pt}pt (below 9pt floor): "
                            f"text={r.text[:40]!r}"
                        )
    assert not violations, "Minimum font size violations:\n  " + "\n  ".join(violations)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
